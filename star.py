#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
star — Speaking Terminal Access Reader
=======================================

A fast, keyboard-driven reading and text-to-speech application designed for
students with print disabilities.
Particularly suited for graduate students in Health Sciences,
but usable enough for high school students.

Inspired by Emacspeak, Kurzweil 1000, Natural Reader, and Central Access Reader.

QUICK START
-----------
    python star.py                       # launch Qt GUI (default)
    python star.py document.pdf          # open a file in the Qt GUI
    python star.py https://example.com   # fetch and read a URL
    python star.py --tui                 # force terminal UI mode
    python star.py --plain paper.pdf     # extract text to stdout

    Qt GUI shortcuts
    ----------------
    Space          play / pause speech
    Ctrl+O         open a file
    Ctrl+H / P / T  next heading / paragraph / table
    F5             cycle theme
    F1             open README.md
    Ctrl+Q         quit

    Terminal UI shortcuts
    ---------------------
    Space        play / pause speech
    Ctrl-O       open a file
    F2           command palette
    F1           open README.md
    Ctrl-Q / q   quit

Copyright (C) 2026 Jon Pielaet
License: GNU General Public License v3 (or later)
"""

__version__ = "0.1.2"
__author__ = "Jon Pielaet"
__copyright__ = "Copyright (C) 2026 Jon Pielaet"
__license__ = "GPL-3.0-or-later"

# =============================================================================
# Standard library
# =============================================================================

import argparse
import csv
import curses
import hashlib
import json
import os
import queue
import re
import shutil
import subprocess
import sys
import tempfile
import threading
import time
import traceback
import urllib.parse
import urllib.request
import xml.etree.ElementTree as ET
import zipfile
from dataclasses import dataclass, field
from html.parser import HTMLParser
from pathlib import Path
from typing import Any, Callable, Dict, List, Optional, Tuple

# Windows curses shim
try:
    import windows_curses  # noqa: F401
except ImportError:
    pass

# =============================================================================
# Optional dependencies — all guarded; star runs with zero extras installed
# =============================================================================

# --- PDF: pdfminer.six -------------------------------------------------------
try:
    from pdfminer.high_level import extract_pages as _pdf_pages
    from pdfminer.high_level import extract_text as _pdf_text
    from pdfminer.layout import LTTextBoxHorizontal

    _PDF = "layout"
except ImportError:
    try:
        from pdfminer.high_level import extract_text as _pdf_text  # type: ignore

        _PDF = "simple"
    except ImportError:
        _PDF = ""

# --- OCR: pytesseract + PyMuPDF ---------------------------------------------
try:
    import pytesseract as _tesseract
    from PIL import Image as _PIL_Image

    _OCR = True
except ImportError:
    _OCR = False

try:
    import fitz as _fitz  # PyMuPDF

    _PYMUPDF = True
except ImportError:
    _PYMUPDF = False

# --- DOCX: python-docx -------------------------------------------------------
try:
    import docx as _docx_lib

    _DOCX = True
except ImportError:
    _DOCX = False

# --- ODT: odfpy --------------------------------------------------------------
try:
    from odf.opendocument import load as _odf_load
    from odf.text import H, ListItem, P
    from odf.text import List as OdfList

    _ODT = True
except ImportError:
    _ODT = False

# --- PowerPoint: python-pptx --------------------------------------------------
try:
    import pptx as _pptx_lib

    _PPTX = True
except ImportError:
    _PPTX = False
    _pptx_lib = None


# --- XLSX: openpyxl ----------------------------------------------------------
try:
    import openpyxl as _openpyxl

    _XLSX = True
except ImportError:
    _XLSX = False

# --- TTS: pyttsx3 ------------------------------------------------------------
try:
    import pyttsx3 as _pyttsx3

    _PYTTSX3 = True
except ImportError:
    _PYTTSX3 = False

# --- TTS: Coqui TTS (neural) -------------------------------------------------
try:
    from TTS.api import TTS as _CoquiTTS  # type: ignore[import]

    _COQUI = True
except (ImportError, Exception):  # also catches missing dependencies
    _COQUI = False
    _CoquiTTS = None  # type: ignore[assignment]

# --- Qt GUI: PyQt6 or PyQt5 --------------------------------------------------
_QT = None
try:
    from PyQt6.QtCore import Qt, QThread, QTimer, pyqtSignal
    from PyQt6.QtGui import (
        QAction,  # PyQt6 moved QAction from QtWidgets to QtGui
        QColor,
        QFont,
        QFontDatabase,
        QTextBlockFormat,
        QTextCharFormat,
        QTextCursor,
        QTextFormat,
    )
    from PyQt6.QtWidgets import (
        QApplication,
        QCheckBox,
        QComboBox,
        QDialog,
        QDialogButtonBox,
        QDockWidget,
        QDoubleSpinBox,
        QFileDialog,
        QFontDialog,
        QFormLayout,
        QHBoxLayout,
        QInputDialog,
        QLabel,
        QLineEdit,
        QListWidget,
        QListWidgetItem,
        QMainWindow,
        QMenu,
        QMessageBox,
        QPushButton,
        QSpinBox,
        QSplitter,
        QStatusBar,
        QTextBrowser,
        QTextEdit,
        QToolBar,
        QVBoxLayout,
        QWidget,
    )

    _QT = "PyQt6"
except ImportError:
    try:
        from PyQt5.QtCore import Qt, QThread, QTimer, pyqtSignal
        from PyQt5.QtGui import (
            QColor,
            QFont,
            QFontDatabase,
            QTextBlockFormat,
            QTextCharFormat,
            QTextCursor,
            QTextFormat,
        )
        from PyQt5.QtWidgets import (
            QAction,
            QApplication,
            QCheckBox,
            QComboBox,
            QDialog,
            QDialogButtonBox,
            QDockWidget,
            QDoubleSpinBox,
            QFileDialog,
            QFontDialog,
            QFormLayout,
            QHBoxLayout,
            QInputDialog,
            QLabel,
            QLineEdit,
            QListWidget,
            QListWidgetItem,
            QMainWindow,
            QMenu,
            QMessageBox,
            QPushButton,
            QSpinBox,
            QSplitter,
            QStatusBar,
            QTextBrowser,
            QTextEdit,
            QToolBar,
            QVBoxLayout,
            QWidget,
        )

        _QT = "PyQt5"
    except ImportError:
        pass

# --- Braille: louis (liblouis Python binding) --------------------------------
try:
    import louis as _louis

    _LOUIS = True
except ImportError:
    _LOUIS = False

# --- Pandoc: pypandoc (subprocess fallback also available) -------------------
try:
    import pypandoc as _pypandoc

    _PYPANDOC = True
except ImportError:
    _PYPANDOC = False

_PANDOC_BIN = shutil.which("pandoc")  # raw binary path

# --- Speech recognition: OpenAI Whisper (or faster-whisper) ------------------
# Used for dictation and recording transcription.  Both the model library and
# a microphone capture library are optional; every code path is guarded so the
# program runs identically when they are absent.
try:
    import whisper as _whisper  # openai-whisper

    _WHISPER = "openai"
except ImportError:
    try:
        from faster_whisper import WhisperModel as _FasterWhisper  # type: ignore

        _WHISPER = "faster"
    except ImportError:
        _WHISPER = ""

try:
    import numpy as _np
    import sounddevice as _sounddevice  # microphone capture

    _AUDIO_IN = True
except ImportError:
    _AUDIO_IN = False

# =============================================================================
# Metadata & paths
# =============================================================================

APP_NAME = "star"
APP_TITLE = "Speaking Terminal Access Reader"
APP_VERSION = __version__

if sys.platform == "win32":
    _CFG_ROOT = Path(os.environ.get("APPDATA", Path.home())) / APP_NAME
elif sys.platform == "darwin":
    _CFG_ROOT = Path.home() / "Library" / "Application Support" / APP_NAME
else:
    _CFG_ROOT = (
        Path(os.environ.get("XDG_CONFIG_HOME", Path.home() / ".config")) / APP_NAME
    )

SETTINGS_FILE = _CFG_ROOT / "settings.json"
CACHE_DIR = _CFG_ROOT / "cache"
THEMES_DIR = _CFG_ROOT / "themes"  # user CSS theme files live here
LOG_FILE = _CFG_ROOT / "star.log"


def _default_sans_font() -> str:
    """Return a high-quality sans-serif display font for the current platform.

    Sans-serif faces are preferred over serif for on-screen reading
    accessibility (lower visual crowding, clearer letterforms for readers
    with low vision or dyslexia).  Each value is a face that ships with the
    OS, so it resolves without the user installing anything; Qt falls back
    to its generic sans-serif if the named family is somehow missing.
    """
    if sys.platform == "darwin":
        return "Helvetica Neue"
    if sys.platform == "win32":
        return "Segoe UI"
    return "DejaVu Sans"


# =============================================================================
# Default settings
# =============================================================================

DEFAULTS: Dict[str, Any] = {
    "theme": "dark",
    "tts_backend": "auto",  # "auto"|"pyttsx3"|"espeak"|"festival"|"coqui"|"dectalk"|"none"
    "tts_rate": 265,  # words per minute — intentionally brisk
    "tts_volume": 1.0,  # 0.0 – 1.0
    "tts_voice": "",  # empty = system default (auto-resolved per platform)
    "tts_prefer_voice": "eloquence",  # substring of a preferred default voice
    "tts_auto_play": False,  # start reading on file open
    "tts_skip_code": True,  # don't read code blocks aloud
    "wrap_width": 0,  # 0 = terminal width
    "tab_width": 4,
    "show_line_numbers": False,
    "syntax_highlight": True,
    "scroll_margin": 3,
    "font_size": 0,  # 0 = terminal default; meaningful in Qt GUI only
    "ocr_lang": "eng",  # Tesseract language(s), e.g. "eng+spa"
    "braille_table": "en-ueb-g2.ctb",
    "braille_grade2": False,  # opt-in liblouis Grade 2; built-in Grade 1 is default
    "audio_export_format": "wav",  # default audio export container (no ffmpeg needed)
    "last_path": "",
    "recent_files": [],  # list of recently opened paths/URLs
    "highlight_current_word": True,
    "highlight_color": "cyan",  # color of the TTS word highlight
    "gui_width": 1000,
    "gui_height": 700,
    "footnote_mode": "inline",  # "inline" | "deferred" | "skip"
    "epub_show_chapters": True,
    "document_cache": True,
    "cache_max_size_mb": 100,
    "qt_show_toc": True,
    "qt_show_notes": False,  # Notes/annotations dock hidden until first used
    "annotations": {},  # {path: [{"char_pos", "word_idx", "anchor", "note", "tags", "cite", "ts"}]}
    "annotation_filter_presets": {},  # {name: filter-query} saved note filters
    "citations": [],  # citation library: list of CSL-ish dicts (see _citation_label)
    "whisper_model": "base",  # Whisper model size for dictation/transcription
    "transcribe_timestamps": False,  # prefix [hh:mm:ss] segment times in transcripts
    "whisper_chunk_seconds": 6,  # chunk length for live streaming dictation
    "keybindings": {},  # {default_shortcut: custom_shortcut} GUI remap overrides
    "bookmarks": {},
    "nav_history_size": 50,
    "regex_search": False,
    "qt_hidpi": True,
    "qt_font_family": _default_sans_font(),  # sans-serif for reading accessibility
    "qt_font_size": 14,
    # ── Text spacing (WCAG 1.4.12) — Qt GUI ─────────────────────────────
    # Generous, independently adjustable spacing reduces crowding effects
    # for dyslexic and low-vision readers.
    "qt_line_height": 1.5,  # line-height multiplier (1.0 = single)
    "qt_letter_spacing": 0.0,  # extra letter spacing, % of font size (0 = normal)
    "qt_word_spacing": 0.0,  # extra word spacing in px (0 = normal)
    # ── Dyslexia-friendly reading aids — Qt GUI ─────────────────────────
    "qt_dyslexia_font": False,  # prefer a bundled/installed dyslexia-friendly font
    "qt_current_line_highlight": False,  # band-highlight the line being read
    "qt_bionic_reading": False,  # embolden the leading part of each word
    # ── Karaoke word-highlight tuning — Qt GUI (highlight_speed is shared) ─
    "highlight_style": "background",  # background|underline|box|bold|color
    "highlight_lead_words": 0,  # advance the visual highlight N words (lead/lag)
    "speak_image_alts": True,
    "show_reading_level": True,
    "normalize_math": True,
    "recent_files_limit": 20,
    # Abbreviation expansion
    "expand_abbreviations": True,
    "abbrev_expansions": {},  # user overrides: {"abbrev.": "expansion"}
    # Number normalization
    "normalize_numbers": True,
    # Table reading mode
    "table_reading_mode": "structured",  # "structured" | "flat" | "skip"
    # User text highlights (persistent per-document colored annotations)
    "user_highlights": {},  # {path: [{"start": int, "end": int, "color": str}]}
    # Reading position memory
    "tts_auto_resume": True,  # restore position automatically on open
    "reading_positions": {},  # {path: {"offset": int, "pct": int, "ts": str}}
    # Speed presets
    "speed_presets": {
        "skim": 350,
        "normal": 265,
        "study": 200,
        "slow": 150,
    },
    # SSML prosody
    # Off by default: plain-text mode enables pyttsx3 word-boundary callbacks
    # which give accurate word highlighting.  SSML disables those callbacks
    # (character offsets point into the XML string, not the plain text) so the
    # timer runs blind and races ahead of speech.  Enable SSML with the
    # 'ssml-on' command or by setting use_ssml=true in the settings file.
    "use_ssml": False,  # wrap TTS text in SSML for better pausing (opt-in)
    "ssml_sentence_pause_ms": 350,  # pause after . ! ?
    "ssml_clause_pause_ms": 150,  # pause after , ; :
    # Highlight timing (fraction of speech rate the cursor advances at).
    # 1.0 = match speech WPM exactly.  The pacing guard in the highlight timer
    # (capped at _MAX_AHEAD words past the last callback-confirmed position)
    # is the true throttle, so running the timer at full speed keeps the
    # highlight tight to the audio instead of lagging behind it.
    "highlight_speed": 1.0,
}

# =============================================================================
# Settings manager
# =============================================================================


class Settings:
    """Persistent JSON settings with dot-notation access."""

    def __init__(self):
        self._data: Dict[str, Any] = dict(DEFAULTS)
        self._load()

    def _load(self) -> None:
        try:
            raw = json.loads(SETTINGS_FILE.read_text(encoding="utf-8"))
            # Accept all known keys, plus nested dicts that may have grown.
            for k, v in raw.items():
                if k in DEFAULTS:
                    # Merge nested dicts (speed_presets, reading_positions)
                    if isinstance(DEFAULTS[k], dict) and isinstance(v, dict):
                        merged = dict(DEFAULTS[k])
                        merged.update(v)
                        self._data[k] = merged
                    else:
                        self._data[k] = v
        except (FileNotFoundError, json.JSONDecodeError, OSError):
            pass
        self._migrate()

    def _migrate(self) -> None:
        """Upgrade settings files written by earlier versions.

        Older versions persisted every default into settings.json, which means
        deprecated defaults (the serif ``Georgia`` display font and the
        lagging ``0.85`` highlight speed) would otherwise be pinned forever.
        We only replace values that exactly match the old default, so a user's
        deliberate choice is never overridden.
        """
        if self._data.get("qt_font_family") == "Georgia":
            # Serif default deprecated for reading-accessibility reasons.
            self._data["qt_font_family"] = _default_sans_font()
        if self._data.get("highlight_speed") == 0.85:
            self._data["highlight_speed"] = 1.0

    def save(self) -> None:
        try:
            SETTINGS_FILE.parent.mkdir(parents=True, exist_ok=True)
            SETTINGS_FILE.write_text(
                json.dumps(self._data, indent=2, ensure_ascii=False),
                encoding="utf-8",
            )
        except OSError:
            pass

    def get(self, key: str, default: Any = None) -> Any:
        return self._data.get(key, default)

    def set(self, key: str, value: Any) -> None:
        self._data[key] = value
        self.save()

    def __getitem__(self, key: str) -> Any:
        return self._data[key]

    def __setitem__(self, key: str, value: Any) -> None:
        self.set(key, value)


# =============================================================================
# TTS system
# =============================================================================


class TTSBackend:
    """Abstract base class for text-to-speech engines."""

    name: str = "silent"

    def available(self) -> bool:
        return False

    def speak(
        self,
        text: str,
        on_word: Optional[Callable[[int, int], None]] = None,
        on_done: Optional[Callable[[], None]] = None,
    ) -> None:
        if on_done:
            on_done()

    def stop(self) -> None:
        pass

    def set_rate(self, wpm: int) -> None:
        pass

    def set_volume(self, vol: float) -> None:
        pass

    def set_voice(self, voice_id: str) -> None:
        pass

    def list_voices(self) -> List[Dict[str, str]]:
        return []

    def export_to_wav(self, text: str, wav_path: str) -> None:
        """Synthesize *text* and write the result to *wav_path* (WAV format).

        This call is **blocking**.  Raises ``RuntimeError`` if the backend
        does not support audio file export or is unavailable.
        """
        raise RuntimeError(f"Backend '{self.name}' does not support audio file export.")

    @property
    def speaking(self) -> bool:
        return False


class FestivalBackend(TTSBackend):
    """Festival speech synthesis backend (Linux / macOS).

    Drives the Festival engine via its interactive scheme interpreter so
    reading rate and voice can be controlled without additional dependencies.
    Festival binary (``festival``) must be on PATH.

    Install:
        Debian/Ubuntu:  ``sudo apt install festival``
        Fedora/RHEL:    ``sudo dnf install festival``
        macOS (Homebrew): ``brew install festival``
    """

    name = "festival"

    def __init__(self, rate: int = 265, volume: float = 1.0, voice: str = "") -> None:
        self._rate = rate
        self._volume = volume
        self._voice = voice
        self._proc: Optional[subprocess.Popen] = None
        self._speaking = False
        self._bin = shutil.which("festival")

    def available(self) -> bool:
        return self._bin is not None

    def speak(
        self,
        text: str,
        on_word: Optional[Callable[[int, int], None]] = None,
        on_done: Optional[Callable[[], None]] = None,
    ) -> None:
        self.stop()
        if not self._bin:
            if on_done:
                on_done()
            return
        self._speaking = True

        # Duration_Stretch < 1 = faster, > 1 = slower;
        # 265 wpm maps to stretch 1.0, linearly scaled.
        stretch = max(0.2, min(4.0, 265.0 / max(1, self._rate)))
        # Escape double quotes in text for the Scheme SayText call.
        escaped = text.replace("\\", "\\\\").replace('"', '\\"')

        scheme_parts: List[str] = []
        if self._voice:
            scheme_parts.append(f"({self._voice})")  # e.g. (voice_rab_diphone)
        scheme_parts.append(f"(Parameter.set 'Duration_Stretch {stretch:.3f})")
        scheme_parts.append(f'(SayText "{escaped}")')
        scheme_parts.append("(quit)")
        scheme = "\n".join(scheme_parts) + "\n"

        def _run() -> None:
            try:
                self._proc = subprocess.Popen(
                    [self._bin],
                    stdin=subprocess.PIPE,
                    stdout=subprocess.DEVNULL,
                    stderr=subprocess.DEVNULL,
                )
                if self._proc.stdin:
                    self._proc.stdin.write(scheme.encode("utf-8", errors="replace"))
                    self._proc.stdin.close()
                self._proc.wait()
            except Exception:
                pass
            finally:
                self._speaking = False
                if on_done:
                    on_done()

        threading.Thread(target=_run, daemon=True).start()

    def stop(self) -> None:
        if self._proc and self._proc.poll() is None:
            try:
                self._proc.terminate()
            except Exception:
                pass
        self._speaking = False

    def set_rate(self, wpm: int) -> None:
        self._rate = wpm

    def set_volume(self, vol: float) -> None:
        self._volume = max(0.0, min(1.0, vol))

    def set_voice(self, voice_id: str) -> None:
        """voice_id is a Festival scheme voice name, e.g. 'voice_rab_diphone'."""
        self._voice = voice_id

    def list_voices(self) -> List[Dict[str, str]]:
        """Query Festival for the list of installed voices."""
        if not self._bin:
            return []
        try:
            result = subprocess.run(
                [self._bin],
                input=b"(voice.list)\n(quit)\n",
                capture_output=True,
                timeout=10,
            )
            raw = result.stdout.decode("utf-8", errors="replace")
            # Festival prints something like: (rab_diphone en1_mbrola_3 ...)
            m = re.search(r"\(([^)]+)\)", raw)
            if m:
                return [
                    {"id": v.strip(), "name": v.strip(), "lang": ""}
                    for v in m.group(1).split()
                    if v.strip()
                ]
        except Exception:
            pass
        return []

    def export_to_wav(self, text: str, wav_path: str) -> None:
        """Synthesize *text* to *wav_path* using Festival.

        Prefers the ``text2wave`` helper (ships with most Festival
        installs); falls back to a Festival Scheme ``utt.save.wave``
        call for shorter texts when ``text2wave`` is not on PATH.
        """
        stretch = max(0.2, min(4.0, 265.0 / max(1, self._rate)))
        t2w = shutil.which("text2wave")
        if t2w:
            proc = subprocess.Popen(
                [
                    t2w,
                    "-o",
                    wav_path,
                    "-eval",
                    f"(Parameter.set 'Duration_Stretch {stretch:.3f})",
                ],
                stdin=subprocess.PIPE,
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
            )
            if proc.stdin:
                proc.stdin.write(text.encode("utf-8", errors="replace"))
                proc.stdin.close()
            proc.wait()
            return
        if not self._bin:
            raise RuntimeError("Festival is not available")
        escaped = text.replace("\\", "\\\\").replace('"', '\\"')
        voice_cmd = f"({self._voice})\n" if self._voice else ""
        scheme = (
            f"{voice_cmd}"
            f"(Parameter.set 'Duration_Stretch {stretch:.3f})\n"
            f'(utt.save.wave (utt.synth (Utterance Text "{escaped}")) "{wav_path}")\n'
            "(quit)\n"
        )
        proc = subprocess.Popen(
            [self._bin],
            stdin=subprocess.PIPE,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
        if proc.stdin:
            proc.stdin.write(scheme.encode("utf-8", errors="replace"))
            proc.stdin.close()
        proc.wait()

    @property
    def speaking(self) -> bool:
        return self._speaking


class CoquiBackend(TTSBackend):
    """Coqui TTS neural speech backend.

    Uses the ``TTS`` Python library (``pip install TTS``) to synthesize
    high-quality neural speech locally.  A suitable model is downloaded on
    first use (typically 100–400 MB) and cached under
    ``~/.local/share/tts/`` (Linux/macOS) or ``%APPDATA%\\TTS\\`` (Windows).
    A GPU is used automatically when ``torch.cuda`` is available, otherwise
    inference runs on CPU (may be slow on older machines).

    Voice selection
    ---------------
    Set ``tts_voice`` to a full Coqui model name to override the default::

        M-x tts-voice tts_models/en/vctk/vits

    Run ``star --list-voices`` after switching to the ``coqui`` backend to
    list all models available from the TTS model zoo.

    Rate control
    ------------
    Speed scaling is applied for models that accept a ``speed`` argument
    (e.g. VITS).  For other models the playback audio is resampled after
    synthesis if ``scipy`` is available; otherwise the rate setting is
    silently ignored.

    Install::

        pip install TTS        # neural TTS model + Python bindings
        pip install scipy      # optional: enables playback-speed control
    """

    name = "coqui"
    _DEFAULT_MODEL = "tts_models/en/ljspeech/tacotron2-DDC"

    def __init__(self, rate: int = 265, volume: float = 1.0, voice: str = "") -> None:
        self._rate = rate
        self._volume = volume
        self._model_name = voice or self._DEFAULT_MODEL
        self._tts_obj = None  # lazily initialized _CoquiTTS instance
        self._speaking = False
        self._stop_flag = threading.Event()
        self._play_proc: Optional[subprocess.Popen] = None

    def available(self) -> bool:
        return _COQUI

    # ── internal helpers ──────────────────────────────────────────────────────

    def _init(self) -> None:
        """Initialize (and possibly download) the TTS model.  Blocking."""
        if self._tts_obj is not None:
            return
        try:
            import torch  # type: ignore[import]

            gpu = torch.cuda.is_available()
        except ImportError:
            gpu = False
        self._tts_obj = _CoquiTTS(
            model_name=self._model_name,
            progress_bar=False,
            gpu=gpu,
        )

    @staticmethod
    def _player_cmd(wav_path: str) -> Optional[List[str]]:
        """Return a platform-appropriate command to play a WAV file."""
        if sys.platform == "win32":
            return [
                "powershell",
                "-NoProfile",
                "-Command",
                f"(New-Object System.Media.SoundPlayer '{wav_path}').PlaySync()",
            ]
        if sys.platform == "darwin":
            return ["afplay", wav_path]
        # Linux: try common players in preference order
        for player in ("aplay", "paplay", "play", "ffplay"):
            if shutil.which(player):
                if player == "ffplay":
                    return ["ffplay", "-nodisp", "-autoexit", wav_path]
                return [player, wav_path]
        return None

    # ── TTSBackend interface ────────────────────────────────────────────────

    def speak(
        self,
        text: str,
        on_word: Optional[Callable[[int, int], None]] = None,
        on_done: Optional[Callable[[], None]] = None,
    ) -> None:
        self.stop()
        if not _COQUI:
            if on_done:
                on_done()
            return
        self._speaking = True
        self._stop_flag.clear()
        rate = self._rate
        volume = self._volume
        model = self._model_name

        def _run() -> None:
            import tempfile

            tmp_path = ""
            try:
                # Lazy model initialization (downloads on first use).
                if self._stop_flag.is_set():
                    return
                self._init()
                if self._stop_flag.is_set():
                    return

                with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as tmp:
                    tmp_path = tmp.name

                # Speed: VITS and some other models accept a speed kwarg.
                try:
                    speed = 265.0 / max(1, rate)
                    self._tts_obj.tts_to_file(
                        text=text, file_path=tmp_path, speed=speed
                    )
                except TypeError:
                    # Model does not support speed — fall back to plain call.
                    self._tts_obj.tts_to_file(text=text, file_path=tmp_path)

                if self._stop_flag.is_set():
                    return

                # Optional volume scaling + pitch-neutral speed via scipy.
                if rate != 265 or volume != 1.0:
                    try:
                        _apply_wav_adjustments(tmp_path, volume)
                    except Exception:
                        pass  # adjustments are best-effort

                # Play the WAV file.
                cmd = self._player_cmd(tmp_path)
                if cmd and not self._stop_flag.is_set():
                    self._play_proc = subprocess.Popen(
                        cmd,
                        stdout=subprocess.DEVNULL,
                        stderr=subprocess.DEVNULL,
                    )
                    self._play_proc.wait()
            except Exception:
                pass
            finally:
                self._speaking = False
                self._play_proc = None
                if tmp_path:
                    try:
                        Path(tmp_path).unlink(missing_ok=True)
                    except Exception:
                        pass
                if on_done and not self._stop_flag.is_set():
                    on_done()

        threading.Thread(target=_run, daemon=True).start()

    def stop(self) -> None:
        self._stop_flag.set()
        proc = self._play_proc
        if proc and proc.poll() is None:
            try:
                proc.terminate()
            except Exception:
                pass
        self._speaking = False

    def set_rate(self, wpm: int) -> None:
        self._rate = wpm

    def set_volume(self, vol: float) -> None:
        self._volume = max(0.0, min(1.0, vol))

    def set_voice(self, voice_id: str) -> None:
        """voice_id is a full Coqui model name, e.g.
        ``tts_models/en/vctk/vits``.  Changing the model resets the
        cached TTS instance so the new model will be loaded on next speak."""
        if voice_id and voice_id != self._model_name:
            self._model_name = voice_id
            self._tts_obj = None  # force re-init

    def list_voices(self) -> List[Dict[str, str]]:
        """Return English-language models from the Coqui model zoo."""
        if not _COQUI:
            return []
        try:
            all_models = _CoquiTTS().list_models().list_tts_models()
            return [
                {
                    "id": m,
                    "name": m.split("/")[-1].replace("_", " "),
                    "lang": m.split("/")[1] if m.count("/") >= 2 else "?",
                }
                for m in all_models
            ]
        except Exception:
            # Fallback: return just the default model so list_voices never
            # crashes.
            return [
                {
                    "id": self._DEFAULT_MODEL,
                    "name": "LJSpeech Tacotron2 (default)",
                    "lang": "en",
                }
            ]

    def export_to_wav(self, text: str, wav_path: str) -> None:
        """Synthesize *text* to *wav_path* using Coqui TTS."""
        if not _COQUI:
            raise RuntimeError("Coqui TTS is not available (pip install TTS)")
        self._init()
        try:
            speed = 265.0 / max(1, self._rate)
            self._tts_obj.tts_to_file(text=text, file_path=wav_path, speed=speed)
        except TypeError:
            # Model does not accept a speed kwarg.
            self._tts_obj.tts_to_file(text=text, file_path=wav_path)
        if self._volume != 1.0:
            _apply_wav_adjustments(wav_path, self._volume)

    @property
    def speaking(self) -> bool:
        return self._speaking


def _apply_wav_adjustments(path: str, volume: float) -> None:
    """Scale WAV sample amplitudes by *volume* in-place.  Pure stdlib."""
    import struct
    import wave

    with wave.open(path, "rb") as wf:
        params = wf.getparams()
        frames = wf.readframes(params.nframes)

    sampwidth = params.sampwidth
    nchannels = params.nchannels
    if sampwidth == 2:  # 16-bit PCM (most common)
        fmt = f"<{len(frames) // 2}h"
        samples = list(struct.unpack(fmt, frames))
        samples = [max(-32768, min(32767, int(s * volume))) for s in samples]
        frames = struct.pack(fmt, *samples)
        with wave.open(path, "wb") as wf:
            wf.setparams(params)
            wf.writeframes(frames)


def _convert_audio_format(src_wav: str, dest_path: str) -> None:
    """Convert *src_wav* (WAV) to the format implied by *dest_path*'s extension.

    Supported output formats: ``.mp3``, ``.ogg``, ``.mp4``, ``.wav``.

    Conversion priority:

    1. **ffmpeg** — recommended; handles all formats cleanly.
       Install: https://ffmpeg.org/download.html
    2. **pydub** — pure-Python fallback (``pip install pydub``).
    3. **WAV copy** — if the extension is ``.wav`` no conversion is needed.

    Raises ``RuntimeError`` when the target format requires conversion but
    no suitable tool is available.
    """
    ext = Path(dest_path).suffix.lower()
    if ext == ".wav":
        shutil.copy2(src_wav, dest_path)
        return

    # --- ffmpeg -----------------------------------------------------------
    if shutil.which("ffmpeg"):
        cmd: List[str] = ["ffmpeg", "-y", "-i", src_wav]
        if ext == ".mp3":
            cmd += ["-codec:a", "libmp3lame", "-qscale:a", "2"]
        elif ext == ".ogg":
            cmd += ["-codec:a", "libvorbis", "-q:a", "4"]
        elif ext == ".mp4":
            # Audio-only MP4 (AAC inside an MPEG-4 container).
            cmd += ["-codec:a", "aac", "-b:a", "192k", "-vn"]
        # For other extensions let ffmpeg infer the codec from the name.
        cmd.append(dest_path)
        result = subprocess.run(cmd, capture_output=True)
        if result.returncode == 0:
            return
        raise RuntimeError(
            f"ffmpeg conversion failed:\n{result.stderr.decode(errors='replace')}"
        )

    # --- pydub ------------------------------------------------------------
    try:
        from pydub import AudioSegment as _AS  # type: ignore[import]

        audio = _AS.from_wav(src_wav)
        fmt_map = {".mp3": "mp3", ".ogg": "ogg", ".mp4": "mp4"}
        fmt = fmt_map.get(ext, ext.lstrip("."))
        audio.export(dest_path, format=fmt)
        return
    except ImportError:
        pass

    raise RuntimeError(
        f"Cannot convert WAV to {ext!r}.\n"
        "Install ffmpeg (https://ffmpeg.org/download.html) "
        "or run: pip install pydub"
    )


class SilentBackend(TTSBackend):
    """No-op backend when no TTS engine is available."""

    name = "silent"

    def available(self) -> bool:
        return True


class Pyttsx3Backend(TTSBackend):
    """pyttsx3 backend — uses SAPI5 (Windows), NSSpeechSynthesizer (macOS),
    or eSpeak-NG (Linux) via the pyttsx3 wrapper.

    Windows/SAPI5 note: calling engine.stop() corrupts the engine's internal
    state so that subsequent say() + runAndWait() calls are silently dropped.
    The fix is to create a *fresh* pyttsx3.init() engine inside every speech
    thread — each call is entirely self-contained.  The active engine is
    stored in self._active_engine so stop() can interrupt it at any time.
    """

    name = "pyttsx3"

    def __init__(self, rate: int = 265, volume: float = 1.0, voice: str = ""):
        self._rate = rate
        self._volume = volume
        self._voice = voice
        self._thread: Optional[threading.Thread] = None
        self._speaking = False
        self._stop_requested = False
        # Monotonically-increasing counter incremented on every speak() call.
        # Each _run closure captures the value at launch time; only the thread
        # whose generation matches the current counter is allowed to write
        # _speaking=False.  This prevents an old thread's finally-block from
        # overwriting the True set by a newer speak() call (Windows race).
        self._gen: int = 0
        # Reference to the engine currently being used by the speech thread.
        # Stored so stop() can call engine.stop() on the right object.
        self._active_engine = None
        self._available = self._probe()

    def _probe(self) -> bool:
        """Check once at startup that pyttsx3 can create an engine."""
        if not _PYTTSX3:
            return False
        try:
            eng = _pyttsx3.Engine()
            eng.stop()
            return True
        except Exception:
            return False

    def _make_engine(self):
        """Create and configure a brand-new pyttsx3 engine.

        We call ``pyttsx3.Engine()`` directly instead of ``pyttsx3.init()``.
        ``init()`` caches the engine in a WeakValueDictionary and returns the
        same object on every call as long as any thread still holds a strong
        reference to it.  After a ``stop()``-while-speaking the cached engine
        has corrupted SAPI5 internal state and silently drops every subsequent
        ``say()`` + ``runAndWait()`` call on Windows.  ``Engine()`` constructs
        a fresh COM object unconditionally and never touches the cache, so each
        speech session always gets a clean SAPI5 voice.
        """
        eng = _pyttsx3.Engine()
        eng.setProperty("rate", self._rate)
        eng.setProperty("volume", self._volume)
        if self._voice:
            try:
                eng.setProperty("voice", self._voice)
            except Exception:
                pass
        return eng

    def available(self) -> bool:
        return self._available

    def speak(
        self,
        text: str,
        on_word: Optional[Callable[[int, int], None]] = None,
        on_done: Optional[Callable[[], None]] = None,
    ) -> None:
        if not self._available:
            if on_done:
                on_done()
            return

        # Signal any running thread to stop, then clear the flag for the new one.
        self._stop_requested = True
        self._speaking = False
        if self._active_engine is not None:
            try:
                self._active_engine.stop()
            except Exception:
                pass
            self._active_engine = None

        self._stop_requested = False
        self._gen += 1
        my_gen = self._gen
        self._speaking = True

        rate = self._rate
        volume = self._volume
        voice = self._voice

        def _run() -> None:
            eng = None
            try:
                eng = self._make_engine()
                self._active_engine = eng  # expose to stop()

                if on_word is not None:
                    # Wrap in a guard so a bad callback never kills the engine.
                    def _word_cb(name: str, location: int, length: int) -> None:
                        try:
                            if not self._stop_requested:
                                on_word(location, length)
                        except Exception:
                            pass

                    eng.connect("started-word", _word_cb)

                eng.say(text)
                if not self._stop_requested:
                    eng.runAndWait()
            except Exception:
                pass
            finally:
                # Do NOT unconditionally clear self._active_engine here.
                # A newer speak() call may have already stored its own engine
                # in self._active_engine; clearing it would prevent stop()
                # from being able to interrupt that newer engine (Bug A:
                # Space stops setting _active_engine to None, so pause breaks).
                # Ownership of _active_engine is managed exclusively by
                # stop() and the top of speak().
                if self._active_engine is eng:
                    # Only clear if it's still *our* engine — no newer thread
                    # has replaced it yet.
                    self._active_engine = None
                # Stop our local engine object (no-op if already stopped).
                if eng is not None:
                    try:
                        eng.stop()
                    except Exception:
                        pass
                # Only the most-recently launched thread may clear _speaking
                # or fire on_done.  An older thread whose generation no longer
                # matches must not overwrite the True set by a newer speak()
                # call, and must not fire on_done (which would set
                # _current_word_idx=-1 and stop the newer thread's timer,
                # making replay always jump to word 0 — Bug B).
                if self._gen == my_gen:
                    self._speaking = False
                    if on_done is not None and not self._stop_requested:
                        on_done()

        self._thread = threading.Thread(target=_run, daemon=True)
        self._thread.start()

    def stop(self) -> None:
        self._stop_requested = True
        self._speaking = False
        eng = self._active_engine
        if eng is not None:
            try:
                eng.stop()
            except Exception:
                pass
            self._active_engine = None

    def set_rate(self, wpm: int) -> None:
        self._rate = wpm
        eng = self._active_engine
        if eng is not None:
            try:
                eng.setProperty("rate", wpm)
            except Exception:
                pass

    def set_volume(self, vol: float) -> None:
        self._volume = max(0.0, min(1.0, vol))
        eng = self._active_engine
        if eng is not None:
            try:
                eng.setProperty("volume", self._volume)
            except Exception:
                pass

    def set_voice(self, voice_id: str) -> None:
        self._voice = voice_id

    def list_voices(self) -> List[Dict[str, str]]:
        if not self._available:
            return []
        try:
            eng = self._make_engine()
            voices = [
                {"id": v.id, "name": v.name, "lang": getattr(v, "languages", ["?"])[0]}
                for v in eng.getProperty("voices")
            ]
            eng.stop()
            return voices
        except Exception:
            return []

    def export_to_wav(self, text: str, wav_path: str) -> None:
        """Synthesize *text* to *wav_path* using pyttsx3's ``save_to_file``."""
        if not self._available:
            raise RuntimeError("pyttsx3 is not available (pip install pyttsx3)")
        eng = self._make_engine()
        try:
            eng.save_to_file(text, wav_path)
            eng.runAndWait()
        finally:
            try:
                eng.stop()
            except Exception:
                pass

    @property
    def speaking(self) -> bool:
        return self._speaking


class ESpeakBackend(TTSBackend):
    """eSpeak-NG backend via subprocess.  Provides reliable cross-platform
    speech without the pyttsx3 dependency chain."""

    name = "espeak"

    def __init__(self, rate: int = 265, volume: int = 100, voice: str = "en-us"):
        self._rate = int(rate * 0.8)  # eSpeak rate scale ≈ 80% of wpm
        self._volume = volume
        self._voice = voice
        self._proc: Optional[subprocess.Popen] = None
        self._speaking = False
        self._bin = shutil.which("espeak-ng") or shutil.which("espeak")

    def available(self) -> bool:
        return self._bin is not None

    def speak(
        self,
        text: str,
        on_word: Optional[Callable[[int, int], None]] = None,
        on_done: Optional[Callable[[], None]] = None,
    ) -> None:
        self.stop()
        if not self._bin:
            if on_done:
                on_done()
            return
        self._speaking = True
        # Auto-detect SSML: add the -m flag so eSpeak processes <break> tags.
        already_ssml = text.lstrip().startswith("<speak>")

        # eSpeak word-callback support.
        # Wrap plain text in SSML with <mark name="N"/> tags between each
        # word so eSpeak-NG outputs "MARK N" to stdout as each word is
        # spoken.  We capture stdout and fire on_word(N, word_len) in a
        # reader thread.  Falls back gracefully if eSpeak does not emit
        # marks (the on_word callback simply never fires).
        use_marks = bool(on_word) and not already_ssml
        word_lens: List[int] = []
        if use_marks:
            tokens = re.split(r"(\s+)", text)
            ssml_parts: List[str] = ["<speak>"]
            idx = 0
            for tok in tokens:
                if tok.strip():
                    ssml_parts.append(f'<mark name="{idx}"/>{tok}')
                    word_lens.append(len(tok.strip()))
                    idx += 1
                elif tok:
                    ssml_parts.append(tok)
            ssml_parts.append("</speak>")
            text = "".join(ssml_parts)

        is_ssml = use_marks or already_ssml

        def _run() -> None:
            try:
                cmd = [
                    self._bin,
                    "-v",
                    self._voice,
                    "-s",
                    str(self._rate),
                    "-a",
                    str(self._volume),
                ]
                if is_ssml:
                    cmd.append("-m")  # SSML / markup mode
                cmd.append("--stdin")
                stdout_pipe = subprocess.PIPE if use_marks else subprocess.DEVNULL
                self._proc = subprocess.Popen(
                    cmd,
                    stdin=subprocess.PIPE,
                    stdout=stdout_pipe,
                    stderr=subprocess.DEVNULL,
                )
                # Launch mark-reader thread before writing stdin so we don't
                # miss any early marks.
                if use_marks and self._proc.stdout is not None:
                    proc_stdout = self._proc.stdout

                    def _mark_reader() -> None:
                        try:
                            for raw in proc_stdout:
                                line = raw.decode("utf-8", errors="replace").strip()
                                if line.startswith("MARK "):
                                    try:
                                        widx = int(line[5:])
                                        wlen = (
                                            word_lens[widx]
                                            if widx < len(word_lens)
                                            else 1
                                        )
                                        if on_word:
                                            on_word(widx, wlen)
                                    except (ValueError, IndexError):
                                        pass
                        except Exception:
                            pass

                    threading.Thread(target=_mark_reader, daemon=True).start()

                if self._proc.stdin:
                    self._proc.stdin.write(text.encode("utf-8", errors="replace"))
                    self._proc.stdin.close()
                self._proc.wait()
            except Exception:
                pass
            finally:
                self._speaking = False
                if on_done:
                    on_done()

        threading.Thread(target=_run, daemon=True).start()

    def stop(self) -> None:
        if self._proc and self._proc.poll() is None:
            try:
                self._proc.terminate()
            except Exception:
                pass
        self._speaking = False

    def set_rate(self, wpm: int) -> None:
        self._rate = int(wpm * 0.8)

    def set_volume(self, vol: float) -> None:
        self._volume = int(vol * 100)

    def set_voice(self, voice_id: str) -> None:
        self._voice = voice_id

    def list_voices(self) -> List[Dict[str, str]]:
        if not self._bin:
            return []
        try:
            out = subprocess.check_output(
                [self._bin, "--voices"], stderr=subprocess.DEVNULL, text=True
            )
            voices = []
            for line in out.splitlines()[1:]:
                parts = line.split()
                if len(parts) >= 4:
                    voices.append({"id": parts[2], "name": parts[3], "lang": parts[1]})
            return voices
        except Exception:
            return []

    def export_to_wav(self, text: str, wav_path: str) -> None:
        """Synthesize *text* to *wav_path* using eSpeak/eSpeak-NG (-w flag)."""
        if not self._bin:
            raise RuntimeError("espeak / espeak-ng is not available")
        cmd = [
            self._bin,
            "-v",
            self._voice,
            "-s",
            str(self._rate),
            "-a",
            str(self._volume),
            "-w",
            wav_path,
            "--stdin",
        ]
        proc = subprocess.Popen(
            cmd,
            stdin=subprocess.PIPE,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
        if proc.stdin:
            proc.stdin.write(text.encode("utf-8", errors="replace"))
            proc.stdin.close()
        proc.wait()
        if proc.returncode != 0:
            raise RuntimeError(f"espeak exited with code {proc.returncode}")

    @property
    def speaking(self) -> bool:
        return self._speaking


class DECtalkBackend(TTSBackend):
    """DECtalk backend.  Requires the DECtalk binary (dtalk or say) to be on
    PATH or pointed to via DECTALK_BIN environment variable.
    The DECtalk source code is available at github.com/dectalk/dectalk."""

    name = "dectalk"

    def __init__(self, rate: int = 265, voice: str = "Paul"):
        self._rate = rate
        self._voice = voice
        self._proc: Optional[subprocess.Popen] = None
        self._speaking = False
        self._bin = (
            os.environ.get("DECTALK_BIN")
            or shutil.which("dtalk")
            or shutil.which("dectalk")
        )

    def available(self) -> bool:
        return self._bin is not None

    def speak(
        self,
        text: str,
        on_word: Optional[Callable[[int, int], None]] = None,
        on_done: Optional[Callable[[], None]] = None,
    ) -> None:
        self.stop()
        if not self._bin:
            if on_done:
                on_done()
            return
        self._speaking = True
        # DECtalk rate: [:rate N] where N is phonemes/min; ~200 phonemes ≈ 120 wpm
        # Scale: wpm * ~1.6 to approximate phoneme rate
        dt_rate = min(600, max(75, int(self._rate * 1.6)))
        dt_text = f"[:voice {self._voice}][:rate {dt_rate}]{text}"

        def _run() -> None:
            try:
                self._proc = subprocess.Popen(
                    [self._bin],
                    stdin=subprocess.PIPE,
                    stdout=subprocess.DEVNULL,
                    stderr=subprocess.DEVNULL,
                )
                if self._proc.stdin:
                    self._proc.stdin.write(dt_text.encode("ascii", errors="replace"))
                    self._proc.stdin.close()
                self._proc.wait()
            except Exception:
                pass
            finally:
                self._speaking = False
                if on_done:
                    on_done()

        threading.Thread(target=_run, daemon=True).start()

    def stop(self) -> None:
        if self._proc and self._proc.poll() is None:
            try:
                self._proc.terminate()
            except Exception:
                pass
        self._speaking = False

    def set_rate(self, wpm: int) -> None:
        self._rate = wpm

    def set_voice(self, voice_id: str) -> None:
        self._voice = voice_id

    def export_to_wav(self, text: str, wav_path: str) -> None:
        """Synthesize *text* to *wav_path* using DECtalk's ``-w`` flag."""
        if not self._bin:
            raise RuntimeError("DECtalk is not available")
        dt_rate = min(600, max(75, int(self._rate * 1.6)))
        dt_text = f"[:voice {self._voice}][:rate {dt_rate}]{text}"
        proc = subprocess.Popen(
            [self._bin, "-w", wav_path],
            stdin=subprocess.PIPE,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
        if proc.stdin:
            proc.stdin.write(dt_text.encode("ascii", errors="replace"))
            proc.stdin.close()
        proc.wait()

    @property
    def speaking(self) -> bool:
        return self._speaking


class AppleSayBackend(TTSBackend):
    """macOS native speech via the built-in ``/usr/bin/say`` command.

    This gives Mac users Apple's high-quality system voices (including the
    Eloquence voices bundled with recent macOS releases) with **zero extra
    dependencies** — no ``pyobjc``, no Homebrew, no ``espeak``.  Because it is
    always present on macOS it is ranked above eSpeak in ``auto`` mode so the
    program never silently falls back to the robotic eSpeak voice on a Mac.

    ``say`` does not emit per-word events, so word highlighting is driven by
    the timer in :class:`TTSManager` (the same path used for Festival/Coqui).
    """

    name = "applesay"

    def __init__(self, rate: int = 265, volume: float = 1.0, voice: str = ""):
        self._rate = int(rate)
        self._volume = volume
        self._voice = voice
        self._proc: Optional[subprocess.Popen] = None
        self._speaking = False
        self._bin = shutil.which("say") if sys.platform == "darwin" else None

    def available(self) -> bool:
        return self._bin is not None

    def _cmd(self, extra: Optional[List[str]] = None) -> List[str]:
        cmd = [self._bin, "-r", str(max(50, self._rate))]
        if self._voice:
            cmd += ["-v", self._voice]
        if extra:
            cmd += extra
        return cmd

    def speak(
        self,
        text: str,
        on_word: Optional[Callable[[int, int], None]] = None,
        on_done: Optional[Callable[[], None]] = None,
    ) -> None:
        self.stop()
        if not self._bin:
            if on_done:
                on_done()
            return
        self._speaking = True

        def _run() -> None:
            try:
                # `say` reads the text to speak from stdin when no string
                # operand is given, which avoids ARG_MAX limits on long docs.
                self._proc = subprocess.Popen(
                    self._cmd(),
                    stdin=subprocess.PIPE,
                    stdout=subprocess.DEVNULL,
                    stderr=subprocess.DEVNULL,
                )
                if self._proc.stdin:
                    self._proc.stdin.write(text.encode("utf-8", errors="replace"))
                    self._proc.stdin.close()
                self._proc.wait()
            except Exception:
                pass
            finally:
                self._speaking = False
                if on_done:
                    on_done()

        threading.Thread(target=_run, daemon=True).start()

    def stop(self) -> None:
        if self._proc and self._proc.poll() is None:
            try:
                self._proc.terminate()
            except Exception:
                pass
        self._speaking = False

    def set_rate(self, wpm: int) -> None:
        self._rate = int(wpm)

    def set_volume(self, vol: float) -> None:
        self._volume = max(0.0, min(1.0, vol))

    def set_voice(self, voice_id: str) -> None:
        self._voice = voice_id

    def list_voices(self) -> List[Dict[str, str]]:
        if not self._bin:
            return []
        try:
            out = subprocess.check_output(
                [self._bin, "-v", "?"], stderr=subprocess.DEVNULL, text=True
            )
        except Exception:
            return []
        voices: List[Dict[str, str]] = []
        for line in out.splitlines():
            # Format: "Reed                en_US    # comment"
            m = re.match(r"^(.+?)\s+([a-z]{2}[-_][A-Z]{2})\s*#?", line)
            if m:
                voices.append(
                    {
                        "id": m.group(1).strip(),
                        "name": m.group(1).strip(),
                        "lang": m.group(2),
                    }
                )
        return voices

    def export_to_wav(self, text: str, wav_path: str) -> None:
        """Synthesize *text* to a WAVE file using ``say -o`` (LEI16 PCM)."""
        if not self._bin:
            raise RuntimeError("macOS 'say' command is not available")
        cmd = self._cmd(
            ["-o", wav_path, "--file-format=WAVE", "--data-format=LEI16@22050"]
        )
        proc = subprocess.Popen(
            cmd,
            stdin=subprocess.PIPE,
            stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL,
        )
        if proc.stdin:
            proc.stdin.write(text.encode("utf-8", errors="replace"))
            proc.stdin.close()
        proc.wait()
        if proc.returncode not in (0, None):
            raise RuntimeError(f"say exited with code {proc.returncode}")

    @property
    def speaking(self) -> bool:
        return self._speaking


class _SCReader:
    """Persistent single-line TTS reader for Speech Cursor mode.

    Problem being solved
    --------------------
    ``_sc_read_line`` used to call ``pyttsx3.Engine()`` on *every* line.  On
    Windows that COM initialization takes 200–500 ms, creating a window where
    ``_active_engine`` is ``None``.  If the user exits SC mode during that
    window ``Pyttsx3Backend.stop()`` cannot reach the engine via
    ``eng.stop()``; the ``_stop_requested`` flag may also lose the race,
    allowing ``runAndWait()`` to start — speech continues after the mode is
    gone.

    Solution
    --------
    One ``pyttsx3.Engine`` is built when SC mode is entered and reused for
    every line.  ``stop()`` always has a live COM object to call
    ``eng.stop()`` on, so SAPI5 is interrupted in under a frame.  If a
    mid-speech stop corrupts SAPI5 state (the known Windows issue), the next
    ``speak()`` call rebuilds the engine *inside its own background thread*
    so the curses UI never blocks.
    """

    def __init__(self, rate: int, volume: float) -> None:
        self._rate = rate
        self._volume = volume
        self._eng = None  # persistent pyttsx3 Engine
        self._thread: Optional[threading.Thread] = None
        self._stop_flag = threading.Event()  # signals _run to abort
        self._needs_rebuild = False  # engine may be corrupt

    # ── internal ──────────────────────────────────────────────────────

    def _build(self) -> "_pyttsx3.Engine":  # type: ignore[name-defined]
        eng = _pyttsx3.Engine()
        eng.setProperty("rate", self._rate)
        eng.setProperty("volume", self._volume)
        return eng

    @property
    def _busy(self) -> bool:
        return bool(self._thread and self._thread.is_alive())

    # ── public API ────────────────────────────────────────────────────

    def start(self) -> None:
        """Build the persistent engine (call once when SC mode is entered)."""
        if not _PYTTSX3:
            return
        try:
            self._eng = self._build()
        except Exception:
            self._eng = None

    def speak(self, text: str) -> None:
        """Stop current speech (if any) and read *text*.

        Always returns immediately — never blocks the UI thread.
        If a mid-speech stop corrupted the engine the rebuild happens
        inside the new speech thread, not on the caller.
        """
        if not _PYTTSX3:
            return

        if self._busy:
            # Signal old thread to abort and interrupt SAPI5.
            # Engine state may be corrupt after stop-while-busy.
            self._stop_flag.set()
            if self._eng:
                try:
                    self._eng.stop()
                except Exception:
                    pass
            self._needs_rebuild = True

        self._stop_flag.clear()
        rate = self._rate
        volume = self._volume
        needs_rebuild = self._needs_rebuild
        eng_ref = [self._eng]  # mutable cell so _run can update it
        stop_flag = self._stop_flag
        reader = self

        def _run() -> None:
            eng = eng_ref[0]
            try:
                if needs_rebuild or eng is None:
                    eng = _pyttsx3.Engine()
                    eng.setProperty("rate", rate)
                    eng.setProperty("volume", volume)
                    eng_ref[0] = eng
                    reader._eng = eng
                    reader._needs_rebuild = False
                if stop_flag.is_set():
                    return
                eng.say(text)
                if stop_flag.is_set():
                    return
                eng.runAndWait()
            except Exception:
                reader._needs_rebuild = True
            finally:
                if stop_flag.is_set():
                    reader._needs_rebuild = True  # interrupted → may be corrupt

        self._thread = threading.Thread(target=_run, daemon=True)
        self._thread.start()

    def stop(self) -> None:
        """Interrupt speech immediately.  Non-blocking — safe to call from
        the curses main loop."""
        self._stop_flag.set()
        eng = self._eng
        if eng:
            try:
                eng.stop()  # SAPI5 Skip — takes effect in < one audio frame
            except Exception:
                pass
        if self._busy:
            self._needs_rebuild = True

    def update_rate(self, rate: int) -> None:
        """Propagate a speech-rate change to the live engine."""
        self._rate = rate
        eng = self._eng
        if eng and not self._busy:
            try:
                eng.setProperty("rate", rate)
            except Exception:
                pass

    def close(self) -> None:
        """Stop speech and release the engine on SC mode exit."""
        self.stop()
        self._eng = None


class TTSManager:
    """Manages the active TTS backend and word-position tracking."""

    BACKENDS = [
        "pyttsx3",
        "applesay",
        "espeak",
        "festival",
        "coqui",
        "dectalk",
        "silent",
    ]

    def __init__(self, settings: Settings):
        self._settings = settings
        self._backend: TTSBackend = SilentBackend()
        self._word_map: List["WordPos"] = []
        self._current_word_idx: int = -1
        self._on_highlight: Optional[Callable[[int], None]] = None  # callback(word_idx)
        self._on_done: Optional[Callable[[], None]] = None
        self._timer_thread: Optional[threading.Thread] = None
        self._timer_stop = threading.Event()
        # Monotonically-increasing counter, incremented every time a new timer
        # thread is started.  Each _tick closure captures its own value so it
        # can detect that a newer timer has taken over and exit immediately.
        # This prevents multiple stale timers from calling _on_highlight
        # simultaneously, which caused the highlight to jump erratically.
        self._timer_gen: int = 0
        # Last word index confirmed by a pyttsx3 word-boundary callback.
        # -1 means no callback has fired yet for the current utterance
        # (either SSML mode where callbacks are skipped, or engine still
        # starting up).  The timer uses this to pace itself: it won't run
        # more than _MAX_AHEAD words ahead of the confirmed position.
        self._last_cb_word_idx: int = -1
        # Monotonic timestamp of the most recent pyttsx3 word callback.
        # 0.0 means no callback has fired for this utterance.  Used by the
        # timer's pacing guard: if no callback has arrived for longer than
        # _CB_TIMEOUT seconds the guard is bypassed so the highlight never
        # stalls while speech continues (SAPI5 callbacks can go silent).
        self._last_cb_time: float = 0.0
        self._select_backend(settings["tts_backend"])

    def _select_backend(self, preference: str) -> None:
        candidates: List[TTSBackend] = []
        rate = int(self._settings["tts_rate"])
        vol = float(self._settings["tts_volume"])
        voice = str(self._settings["tts_voice"])

        if preference in ("auto", "pyttsx3") and _PYTTSX3:
            candidates.append(Pyttsx3Backend(rate=rate, volume=vol, voice=voice))
        # macOS native `say` is ranked above eSpeak in auto mode so a Mac
        # always speaks with an Apple system voice (no pyobjc required) rather
        # than silently falling back to the robotic eSpeak voice.
        if preference in ("auto", "applesay"):
            candidates.append(AppleSayBackend(rate=rate, volume=vol, voice=voice))
        if preference in ("auto", "espeak"):
            candidates.append(ESpeakBackend(rate=rate, voice=voice or "en-us"))
        if preference in ("auto", "festival"):
            candidates.append(FestivalBackend(rate=rate, volume=vol, voice=voice))
        if preference in ("coqui",) and _COQUI:
            # Coqui is never selected in auto mode because model download is slow
            # and requires explicit opt-in by the user.
            candidates.append(CoquiBackend(rate=rate, volume=vol, voice=voice))
        if preference in ("auto", "dectalk"):
            candidates.append(DECtalkBackend(rate=rate))
        candidates.append(SilentBackend())

        for b in candidates:
            if b.available():
                self._backend = b
                self._backend.set_rate(rate)
                self._backend.set_volume(vol)
                self._resolve_default_voice()
                return

    def _resolve_default_voice(self) -> None:
        """Pick a sensible default voice when the user hasn't chosen one.

        When ``tts_voice`` is empty, prefer a voice whose name contains the
        ``tts_prefer_voice`` substring (default ``"eloquence"``), favoring a
        US-English variant.  This makes the bundled Eloquence voices the
        default on macOS while leaving the engine default untouched when no
        match is found.  The user's explicit voice choice always wins.
        """
        if str(self._settings.get("tts_voice", "")):
            return  # user has an explicit voice; never override it
        prefer = str(self._settings.get("tts_prefer_voice", "")).strip().lower()
        if not prefer:
            return
        try:
            voices = self._backend.list_voices()
        except Exception:
            voices = []
        if not voices:
            return
        matches = [
            v
            for v in voices
            if prefer in (v.get("name", "") + " " + v.get("id", "")).lower()
        ]
        if not matches:
            return
        # Favor a US-English variant of the preferred voice family.
        best = next(
            (m for m in matches if "us" in str(m.get("lang", "")).lower()),
            matches[0],
        )
        vid = best.get("id") or best.get("name")
        if vid:
            self._backend.set_voice(vid)

    @property
    def backend_name(self) -> str:
        return self._backend.name

    @property
    def speaking(self) -> bool:
        return self._backend.speaking

    @property
    def current_word_idx(self) -> int:
        return self._current_word_idx

    def set_word_map(self, word_map: List["WordPos"]) -> None:
        self._word_map = word_map

    def set_on_highlight(self, cb: Optional[Callable[[int], None]]) -> None:
        self._on_highlight = cb

    def set_on_done(self, cb: Optional[Callable[[], None]]) -> None:
        self._on_done = cb

    def speak(
        self,
        text: str,
        start_word_idx: int = 0,
        text_offset: int = 0,
    ) -> None:
        """Begin speaking *text*.

        Parameters
        ----------
        text:
            The string actually passed to the TTS engine.  This may be a
            *slice* of the full document plain text (everything from the
            desired start position to the end) so that the engine does not
            re-read content that has already been heard.
        start_word_idx:
            Index into the full word_map of the first word in *text*.  Used
            to seed the highlight timer at the right position.
        text_offset:
            Character offset of the first character of *text* within the
            full plain-text string.  Used to translate the byte offsets that
            pyttsx3 reports back into absolute word-map indices.
        """
        # Increment the timer generation BEFORE signalling the old timer to
        # stop.  This ensures that an old timer currently mid-loop-body will
        # see the new generation on its very next gen-check and return without
        # calling _on_highlight, preventing a stray high-word flash followed
        # by the new timer's start-word snap (the "snap back" bug).
        self._timer_gen += 1
        self._timer_stop.set()
        self._current_word_idx = max(0, start_word_idx)
        self._last_cb_word_idx = -1  # no confirmed position yet for this utterance
        self._last_cb_time = 0.0  # reset callback timestamp for this utterance

        def on_done() -> None:
            self._timer_stop.set()
            self._current_word_idx = -1
            if self._on_highlight:
                self._on_highlight(-1)
            if self._on_done:
                self._on_done()

        # pyttsx3 word callbacks supplement the timer when they fire reliably
        # (they may not on all Windows/SAPI5 configurations).  The timer is
        # always started as the primary highlight mechanism.
        if isinstance(self._backend, Pyttsx3Backend):

            def on_word_cb(location: int, length: int) -> None:
                """Translate TTS-relative location back to a word-map index.

                *location* is relative to the *text* slice passed to speak().
                Adding *text_offset* converts it to an absolute offset in the
                full plain-text string, which is what word_map stores.

                We update *_current_word_idx* here so the timer can adopt the
                accurate engine position on its next tick, but we deliberately
                do NOT call *_on_highlight* directly.  SAPI5 callbacks arrive
                asynchronously and can lag or burst; calling _on_highlight from
                the callback caused the highlight to snap backward to an older
                word while the timer had already advanced forward.
                """
                # text_offset == -1 means SSML mode: character offsets in
                # the callback point into the SSML string, not the plain
                # text.  Skip the lookup and let the timer handle highlight.
                if text_offset < 0:
                    return
                abs_loc = location + text_offset
                for i, wp in enumerate(self._word_map):
                    if wp.tts_offset <= abs_loc < wp.tts_offset + wp.tts_len + 1:
                        # Monotonic write: only advance, never retreat.
                        # Delayed or out-of-order SAPI5 callbacks for earlier
                        # words must not clobber a later confirmed position
                        # (which would make _tts_toggle save the wrong pause
                        # word and cause a backward snap on resume).
                        if i >= self._current_word_idx:
                            self._current_word_idx = i
                            self._last_cb_word_idx = i
                            self._last_cb_time = time.monotonic()
                        break

            self._backend.speak(text, on_word=on_word_cb, on_done=on_done)
        else:
            self._backend.speak(text, on_done=on_done)

        # Always start the timer — it is the reliable baseline for all backends.
        self._start_timer_highlight(start_word_idx)

    def _start_timer_highlight(self, start_idx: int) -> None:
        """Timer-based word highlight advance.  Works for every backend.

        If the word map is not yet built (async loading still running), the
        timer waits up to 10 s for it to appear before advancing.

        A monotonic *_timer_gen* counter is captured at launch.  Every loop
        iteration confirms its value still matches; if a newer timer has been
        started (via a new speak() call) the old thread exits immediately.
        This prevents multiple stale timers from racing to call _on_highlight
        with different word indices, which was the primary cause of the
        highlight jumping all over the place.
        """
        self._timer_stop.clear()
        # _timer_gen was already incremented by speak() or stop() before
        # this method was called; just capture the current value.
        my_gen = self._timer_gen
        rate = int(self._settings["tts_rate"])
        # Timer interval: run at the nominal speech rate (1.0 × wpm) so the
        # highlight tracks audio as closely as possible.  The _MAX_AHEAD guard
        # below is the true throttle for pyttsx3/SAPI5; slowing the timer
        # (< 1.0) only causes the highlight to fall behind.
        hl_speed = float(self._settings.get("highlight_speed", 1.0))
        interval = 60.0 / max(1.0, rate * max(0.1, hl_speed))
        # How many words ahead of the last callback-confirmed position the
        # timer is allowed to advance before it pauses for one tick.
        # Only active when pyttsx3 word callbacks are firing; _last_cb_word_idx
        # stays -1 in SSML mode and for non-pyttsx3 backends (guard inactive).
        #
        # 4 words of slack covers the typical SAPI5 callback delay (1-3 words
        # late) without letting the highlight race too far ahead of audio.
        _MAX_AHEAD = 4
        # If no callback has arrived within this many seconds the guard is
        # bypassed entirely: SAPI5 sometimes stops firing callbacks mid-text,
        # and without this escape the highlight would freeze while speech
        # continues.  1.5 s ≈ 6 words at 240 wpm — long enough to ride out
        # normal punctuation pauses, short enough to feel responsive.
        _CB_TIMEOUT = 1.5

        def _tick() -> None:
            # Wait for the word map to be populated (built asynchronously).
            deadline = time.monotonic() + 10.0
            while not self._timer_stop.is_set():
                if self._word_map:
                    break
                if time.monotonic() > deadline:
                    return  # gave up waiting
                time.sleep(0.05)

            # Exit immediately if a newer timer was started while we waited.
            if self._timer_gen != my_gen:
                return

            idx = max(0, start_idx)
            while not self._timer_stop.wait(interval):
                # Bail out as soon as a newer timer generation takes over.
                if self._timer_gen != my_gen:
                    return
                # Adopt the engine's position when it has run ahead of the
                # timer estimate (e.g. fast speech or SSML pauses consumed).
                # Never go backward — that would cause the highlight to jump
                # back to a word that was already spoken.
                if self._current_word_idx > idx:
                    idx = self._current_word_idx
                # Pacing guard: keep the highlight within _MAX_AHEAD words
                # of the last callback-confirmed audio position.  Only active
                # while callbacks are both firing AND recent; if SAPI5 stops
                # sending callbacks (_CB_TIMEOUT exceeded) the guard is
                # bypassed so the highlight never freezes mid-document.
                if (
                    self._last_cb_word_idx >= 0
                    and idx >= self._last_cb_word_idx + _MAX_AHEAD
                    and (time.monotonic() - self._last_cb_time) < _CB_TIMEOUT
                ):
                    continue  # hold briefly — callbacks are active but lagging
                if idx < len(self._word_map):
                    # Second gen-check immediately before the display call.
                    # Closes the narrow window between the first check above
                    # and this point where a new speak() could have bumped
                    # the generation, avoiding a stray _on_highlight flash.
                    if self._timer_gen != my_gen:
                        return
                    self._current_word_idx = idx
                    if self._on_highlight:
                        self._on_highlight(idx)
                    idx += 1
                # Don't break when we reach the end — the backend may still
                # be speaking padding/trailing punctuation.

        self._timer_thread = threading.Thread(target=_tick, daemon=True)
        self._timer_thread.start()

    def stop(self) -> None:
        # Same ordering as speak(): bump generation first so any running timer
        # exits cleanly before the stop event is processed.
        self._timer_gen += 1
        self._timer_stop.set()
        self._backend.stop()
        self._current_word_idx = -1
        self._last_cb_word_idx = -1
        self._last_cb_time = 0.0

    @property
    def last_cb_word_idx(self) -> int:
        """Last word index confirmed by a pyttsx3 word-boundary callback.
        -1 when no callback has fired for the current utterance (SSML mode
        or before the engine has produced the first word).  More accurate
        than *current_word_idx* for pause/resume because it reflects the
        actual audio position rather than the timer\'s forward estimate.
        """
        return self._last_cb_word_idx

    def set_rate(self, wpm: int) -> None:
        self._settings["tts_rate"] = wpm
        self._backend.set_rate(wpm)

    def set_volume(self, vol: float) -> None:
        self._settings["tts_volume"] = vol
        self._backend.set_volume(vol)

    def change_backend(self, name: str) -> None:
        self.stop()
        self._settings["tts_backend"] = name
        self._select_backend(name)

    def list_voices(self) -> List[Dict[str, str]]:
        return self._backend.list_voices()

    def export_audio(self, text: str, dest_path: str) -> None:
        """Synthesize *text* and save it to *dest_path*.

        The output format is inferred from the file extension:

        * ``.wav``  — written directly by the backend (no extras needed).
        * ``.mp3``  — requires **ffmpeg** or **pydub**.
        * ``.ogg``  — requires **ffmpeg** or **pydub**.
        * ``.mp4``  — requires **ffmpeg** or **pydub** (audio-only AAC).

        This method **blocks** until synthesis and conversion are complete.
        Call it from a background thread when used in a GUI to avoid
        freezing the interface.
        """
        ext = Path(dest_path).suffix.lower()
        if ext == ".wav":
            self._backend.export_to_wav(text, dest_path)
            return
        with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as tmp:
            tmp_wav = tmp.name
        try:
            self._backend.export_to_wav(text, tmp_wav)
            _convert_audio_format(tmp_wav, dest_path)
        finally:
            try:
                Path(tmp_wav).unlink(missing_ok=True)
            except Exception:
                pass


# =============================================================================
# Document data structures
# =============================================================================


@dataclass
class WordPos:
    """Maps one word in the TTS plain-text to a position in the display."""

    word: str  # the word text (stripped of punctuation)
    tts_offset: int  # char offset in the TTS plain-text string
    tts_len: int  # length in the TTS string
    disp_line: int  # rendered display line index
    disp_col: int  # starting column in that display line


@dataclass
class Document:
    path: str = ""
    title: str = ""
    markdown: str = ""  # markdown for display
    plain_text: str = ""  # clean text for TTS
    word_map: List[WordPos] = field(default_factory=list)
    metadata: Dict[str, str] = field(default_factory=dict)
    format: str = ""  # detected format
    encoding: str = "utf-8"
    # Chapter list for EPUB/DAISY navigation: [(title, href, word_idx), ...]
    chapters: List[Tuple[str, str, int]] = field(default_factory=list)


# =============================================================================
# Document loaders
# =============================================================================


def _detect_format(path: str) -> str:
    """Detect document format from extension or magic bytes."""
    p = path.lower()
    if p.startswith(("http://", "https://", "ftp://")):
        return "url"
    ext_map = {
        ".md": "markdown",
        ".markdown": "markdown",
        ".mdown": "markdown",
        ".txt": "text",
        ".text": "text",
        ".html": "html",
        ".htm": "html",
        ".xhtml": "html",
        ".pdf": "pdf",
        ".docx": "docx",
        ".doc": "doc",
        ".dot": "doc",  # legacy Word template — same binary format
        ".pptx": "pptx",
        ".ppt": "pptx",  # legacy PowerPoint — same conversion path
        ".odt": "odt",
        ".epub": "epub",
        ".csv": "csv",
        ".tsv": "tsv",
        ".xlsx": "xlsx",
        ".xls": "xlsx",
        ".tex": "latex",
        ".ltx": "latex",
        ".rst": "rst",
        ".rest": "rst",
        ".adoc": "asciidoc",
        ".asciidoc": "asciidoc",
        ".asc": "asciidoc",
        ".wiki": "mediawiki",
        ".mediawiki": "mediawiki",
        ".textile": "textile",
        ".creole": "creole",
        ".r": "r",
        ".rmd": "rmarkdown",
        ".ipynb": "notebook",
        ".xml": "xml",
        ".daisy": "daisy",
        ".opf": "daisy",
        ".ncx": "daisy",
        ".png": "image",
        ".jpg": "image",
        ".jpeg": "image",
        ".gif": "image",
        ".bmp": "image",
        ".tiff": "image",
        ".webp": "image",
        ".py": "python",
        ".js": "javascript",
        ".rs": "rust",
        ".c": "c",
        ".cpp": "c",
        ".h": "c",
        ".hpp": "c",
        ".brf": "braille",
        ".org": "orgmode",
    }
    ext = Path(path).suffix.lower()
    return ext_map.get(ext, "text")


def _strip_markdown_for_tts(
    md: str,
    skip_code: bool = True,
    table_mode: str = "structured",
) -> str:
    """Remove markdown syntax to produce clean text suitable for TTS.
    The result should sound natural when read aloud — no asterisks, slashes,
    pound signs, pipe characters, or code fences.

    table_mode is forwarded to _tables_to_narration() and controls how tables
    are rendered for speech (structured / flat / skip).
    """
    text = md

    # Remove fenced code blocks entirely if requested
    if skip_code:
        text = re.sub(r"```[\s\S]*?```", "", text)
        text = re.sub(r"~~~[\s\S]*?~~~", "", text)
        text = re.sub(r"^    .+$", "", text, flags=re.MULTILINE)  # indented code
    else:
        text = re.sub(r"```\w*\n?", "", text)
        text = re.sub(r"```", "", text)

    # Headings — keep text, drop pounds
    text = re.sub(r"^#{1,6}\s+", "", text, flags=re.MULTILINE)

    # Horizontal rules
    text = re.sub(r"^(\*{3,}|-{3,}|_{3,})\s*$", "", text, flags=re.MULTILINE)

    # Links: keep display text
    text = re.sub(r"!\[([^\]]*)\]\([^\)]*\)", r"\1", text)  # images
    text = re.sub(r"\[([^\]]*)\]\([^\)]*\)", r"\1", text)  # links

    # Inline code
    text = re.sub(r"`+(.+?)`+", r"\1", text)

    # Bold / italic
    text = re.sub(r"\*{3}(.+?)\*{3}", r"\1", text)
    text = re.sub(r"_{3}(.+?)_{3}", r"\1", text)
    text = re.sub(r"\*{2}(.+?)\*{2}", r"\1", text)
    text = re.sub(r"_{2}(.+?)_{2}", r"\1", text)
    text = re.sub(r"\*([^*\n]+?)\*", r"\1", text)
    text = re.sub(r"_([^_\n]+?)_", r"\1", text)

    # Blockquotes
    text = re.sub(r"^>\s?", "", text, flags=re.MULTILINE)

    # List markers
    text = re.sub(r"^\s*[-*+]\s+", "", text, flags=re.MULTILINE)
    text = re.sub(r"^\s*\d+[.)]\s+", "", text, flags=re.MULTILINE)

    # Table narration — must run before pipes are stripped.
    # _tables_to_narration() is defined later in the file; it operates on the
    # still-raw markdown lines and replaces table blocks with spoken prose.
    text = _tables_to_narration(text, mode=table_mode)

    # Collapse extra blank lines
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _build_word_map(plain_text: str, rendered_lines: List[str]) -> List[WordPos]:
    """Build a word map that links TTS character offsets to display positions.

    Strategy: tokenize plain text into words; for each word, scan the rendered
    lines to find a matching occurrence.  Uses a rolling search start to keep
    the match order correct even for repeated words.

    Words whose only occurrence in the display is *before* the current search
    position (e.g. column-header names repeated in structured table-row
    narration) are assigned the last confirmed forward position so the
    highlight advances linearly rather than jumping backward.
    """
    words: List[WordPos] = []
    token_re = re.compile(r"\b\w[\w'-]*")
    search_line = 0  # rolling hint: don't search lines before this
    search_col = 0  # column offset on search_line; avoids re-matching an
    # earlier occurrence of a repeated word on the same line
    last_good_line = 0  # last display line from a forward-matched word
    last_good_col = 0

    for m in token_re.finditer(plain_text):
        word = m.group()
        offset = m.start()
        word_lower = word.lower()

        found_line = last_good_line
        found_col = last_good_col
        matched = False

        # Primary forward search.  On the starting line we begin the column
        # search from search_col so we never match a word that appeared
        # earlier on the same line (prevents the highlight jumping backward
        # for common words like "the" / "a" that repeat within a line).
        for li in range(search_line, min(search_line + 80, len(rendered_lines))):
            start = search_col if li == search_line else 0
            col = rendered_lines[li].lower().find(word_lower, start)
            if col >= 0:
                found_line = li
                found_col = col
                matched = True
                break

        if not matched:
            # Extended forward scan beyond the 80-line window.
            for li in range(
                min(search_line + 80, len(rendered_lines)), len(rendered_lines)
            ):
                col = rendered_lines[li].lower().find(word_lower, 0)
                if col >= 0:
                    found_line = li
                    found_col = col
                    matched = True
                    break

        if not matched:
            # Backward-only fallback: word exists but only before the current
            # search position (e.g. a table column header repeated in row
            # narration).  Keep found_line/col at last_good_* so the highlight
            # does not regress.
            for li, rline in enumerate(rendered_lines):
                col = rline.lower().find(word_lower, 0)
                if col >= 0:
                    matched = True  # word exists — audio is fine
                    break  # found_line/found_col stay at last_good_*

        words.append(
            WordPos(
                word=word,
                tts_offset=offset,
                tts_len=len(word),
                disp_line=found_line,
                disp_col=found_col,
            )
        )
        # Only advance the search position for genuine forward matches.
        # Remove the old "-2" look-back: that was intended as a robustness
        # margin but it caused common words to cascade-match 2 lines before
        # their actual display position, making the highlight appear stuck.
        if matched and found_line >= search_line:
            search_line = found_line
            search_col = found_col + len(word)
            last_good_line = found_line
            last_good_col = found_col

    return words


# ── Loader functions ──────────────────────────────────────────────────────────


def _load_plain_text(path: str) -> str:
    try:
        return Path(path).read_text(encoding="utf-8", errors="replace")
    except OSError as e:
        return f"# Error\n\n```\n{e}\n```\n"


def _load_markdown(path: str) -> str:
    return _load_plain_text(path)


class _HTML2MD(HTMLParser):
    """Minimal HTML → Markdown converter (no third-party dependencies)."""

    _SKIP = frozenset(
        {
            "script",
            "style",
            "nav",
            "footer",
            "aside",
            "noscript",
            "svg",
            "canvas",
            "meta",
            "link",
            "base",
            "iframe",
            "template",
            "button",
            "form",
        }
    )

    def __init__(self) -> None:
        super().__init__(convert_charrefs=True)
        self._out: List[str] = []
        self._buf: List[str] = []
        self._skip = 0
        self._pre = 0
        self._bold = 0
        self._italic = 0
        self._code = 0
        self._heading = 0
        self._list: List[Tuple[str, int]] = []
        self._bq = 0
        self._trows: List[List[str]] = []
        self._trow: List[str] = []
        self._tcell: List[str] = []
        self._in_cell = False
        self._link_href = ""
        self._link_buf: List[str] = []
        self._in_link = False
        self._title_buf: List[str] = []
        self._in_title = False

    def handle_starttag(self, tag: str, attrs: list) -> None:
        a = dict(attrs)
        if tag in self._SKIP:
            self._skip += 1
            return
        if self._skip:
            return
        if tag == "pre":
            self._emit()
            self._pre += 1
            if self._pre == 1:
                lang = next(
                    (
                        c.split("-", 1)[1]
                        for c in a.get("class", "").split()
                        if c.startswith(("language-", "lang-"))
                    ),
                    "",
                )
                self._out.append(f"```{lang}")
        elif tag in ("h1", "h2", "h3", "h4", "h5", "h6"):
            self._emit()
            self._heading = int(tag[1])
        elif tag == "title":
            self._in_title = True
        elif tag in ("p", "div", "section", "article", "main", "header"):
            self._emit()
        elif tag in ("b", "strong"):
            self._bold += 1
        elif tag in ("i", "em"):
            self._italic += 1
        elif tag == "code" and not self._pre:
            self._code += 1
        elif tag == "a":
            self._in_link = True
            self._link_href = a.get("href", "")
            self._link_buf = []
        elif tag == "br":
            self._buf.append("\n")
        elif tag == "hr":
            self._emit()
            self._out += ["---", ""]
        elif tag == "ul":
            self._emit()
            self._list.append(("ul", 0))
        elif tag == "ol":
            self._emit()
            self._list.append(("ol", 0))
        elif tag == "li":
            self._emit()
        elif tag == "blockquote":
            self._emit()
            self._bq += 1
        elif tag == "table":
            self._emit()
            self._trows = []
        elif tag == "tr":
            self._trow = []
        elif tag in ("th", "td"):
            self._tcell = []
            self._in_cell = True

    def handle_endtag(self, tag: str) -> None:
        if tag in self._SKIP:
            self._skip = max(0, self._skip - 1)
            return
        if self._skip:
            return
        if tag == "pre":
            self._pre = max(0, self._pre - 1)
            if self._pre == 0:
                self._out += ["```", ""]
        elif tag in ("h1", "h2", "h3", "h4", "h5", "h6"):
            t = "".join(self._buf).strip()
            if t:
                self._out += ["#" * self._heading + " " + t, ""]
            self._buf = []
            self._heading = 0
        elif tag == "title":
            self._in_title = False
        elif tag in ("p", "section", "article", "main", "header", "div"):
            self._emit_para()
        elif tag in ("b", "strong"):
            self._bold = max(0, self._bold - 1)
        elif tag in ("i", "em"):
            self._italic = max(0, self._italic - 1)
        elif tag == "code" and not self._pre:
            self._code = max(0, self._code - 1)
        elif tag == "a":
            lt = "".join(self._link_buf).strip()
            frag = f"[{lt}]({self._link_href})" if lt and self._link_href else lt
            (self._tcell if self._in_cell else self._buf).append(frag)
            self._in_link = False
        elif tag == "li":
            t = "".join(self._buf).strip()
            if t and self._list:
                ltype, n = self._list[-1]
                pad = "  " * (len(self._list) - 1)
                if ltype == "ol":
                    n += 1
                    self._list[-1] = ("ol", n)
                    self._out.append(f"{pad}{n}. {t}")
                else:
                    self._out.append(f"{pad}* {t}")
            self._buf = []
        elif tag in ("ul", "ol"):
            if self._list:
                self._list.pop()
            if not self._list:
                self._out.append("")
        elif tag == "blockquote":
            self._emit_para()
            self._bq = max(0, self._bq - 1)
        elif tag in ("th", "td"):
            self._trow.append("".join(self._tcell).strip())
            self._tcell = []
            self._in_cell = False
            self._buf = []
        elif tag == "tr":
            if self._trow:
                self._trows.append(self._trow[:])
        elif tag == "table":
            self._flush_table()

    def handle_data(self, data: str) -> None:
        if self._skip or (not data.strip() and not self._pre):
            return
        t = data if self._pre else re.sub(r"\s+", " ", data.replace("\u00a0", " "))
        if self._pre:
            self._out.append(data.rstrip("\n"))
            return
        if self._code:
            t = f"`{t.strip()}`"
        elif self._bold and self._italic:
            t = f"***{t.strip()}***"
        elif self._bold:
            t = f"**{t.strip()}**"
        elif self._italic:
            t = f"*{t.strip()}*"
        if self._in_title:
            self._title_buf.append(t)
        elif self._in_link:
            self._link_buf.append(t)
        elif self._in_cell:
            self._tcell.append(t)
        else:
            self._buf.append(t)

    def _emit(self) -> None:
        t = "".join(self._buf).strip()
        if t:
            self._out.append(t)
        self._buf = []

    def _emit_para(self) -> None:
        t = "".join(self._buf).strip()
        if t:
            bq = "> " * self._bq
            self._out += [bq + t, ""]
        self._buf = []

    def _flush_table(self) -> None:
        rows = self._trows
        if not rows:
            return
        nc = max(len(r) for r in rows)
        hdr = (rows[0] + [""] * nc)[:nc]
        self._out.append("| " + " | ".join(hdr) + " |")
        self._out.append("|" + "|".join([" --- "] * nc) + "|")
        for row in rows[1:]:
            self._out.append("| " + " | ".join((row + [""] * nc)[:nc]) + " |")
        self._out.append("")
        self._trows = []

    def result(self) -> str:
        self._emit()
        title = "".join(self._title_buf).strip()
        lines = ([f"# {title}", ""] if title else []) + self._out
        out: List[str] = []
        prev_blank = False
        for ln in lines:
            blank = not ln.strip()
            if blank and prev_blank:
                continue
            out.append(ln)
            prev_blank = blank
        return "\n".join(out).strip()


def _load_html(path: str) -> str:
    try:
        text = Path(path).read_text(encoding="utf-8", errors="replace")
        p = _HTML2MD()
        p.feed(text)
        p.close()
        return p.result()
    except Exception as e:
        return f"# Error loading HTML\n\n```\n{e}\n```\n"


def _load_html_str(html: str) -> str:
    p = _HTML2MD()
    p.feed(html)
    p.close()
    return p.result()


def _load_epub(path: str) -> str:
    """Load EPUB via zipfile + HTML converter (no third-party package needed)."""
    import posixpath
    from urllib.parse import unquote

    _HTML_MT = frozenset({"application/xhtml+xml", "text/html"})

    def _local(elem: ET.Element) -> str:
        t = elem.tag
        return t.split("}")[-1] if "}" in t else t

    try:
        with zipfile.ZipFile(path, "r") as zf:
            try:
                container = ET.fromstring(zf.read("META-INF/container.xml"))
            except (KeyError, ET.ParseError) as e:
                return f"# EPUB Error\n\n```\n{e}\n```\n"

            opf_path = next(
                (
                    e.get("full-path")
                    for e in container.iter()
                    if _local(e) == "rootfile"
                ),
                None,
            )
            if not opf_path:
                return "# EPUB Error\n\nCould not locate OPF rootfile.\n"

            opf_root = ET.fromstring(zf.read(opf_path))
            opf_dir = (opf_path.rsplit("/", 1)[0] + "/") if "/" in opf_path else ""

            title = author = ""
            for e in opf_root.iter():
                tl = _local(e)
                if tl == "title" and not title and e.text:
                    title = e.text.strip()
                elif tl == "creator" and not author and e.text:
                    author = e.text.strip()

            manifest: Dict[str, Tuple[str, str, str]] = {}
            for e in opf_root.iter():
                if _local(e) == "item":
                    iid = e.get("id", "")
                    href = unquote(e.get("href", ""))
                    mt = e.get("media-type", "")
                    props = e.get("properties", "")
                    if iid and href:
                        manifest[iid] = (href, mt, props)

            spine = [
                e.get("idref", "")
                for e in opf_root.iter()
                if _local(e) == "itemref" and e.get("idref")
            ]

            znames = {n.lower(): n for n in zf.namelist()}

            def read_item(item_path: str) -> Optional[bytes]:
                try:
                    return zf.read(item_path)
                except KeyError:
                    canon = znames.get(item_path.lower())
                    return zf.read(canon) if canon else None

            parts: List[str] = []
            if title:
                parts += [f"# {title}", ""]
            if author:
                parts += [f"*{author}*", ""]
            if title or author:
                parts += ["---", ""]

            seen: set = set()
            for idref in spine:
                if idref not in manifest:
                    continue
                href, mt, props = manifest[idref]
                if mt not in _HTML_MT:
                    continue
                if "nav" in props.split():
                    continue
                href = href.split("#")[0]
                if not href:
                    continue
                raw = href if href.startswith("/") else opf_dir + href
                item_path = posixpath.normpath(raw).lstrip("/")
                if item_path in seen:
                    continue
                seen.add(item_path)
                data = read_item(item_path)
                if not data:
                    continue
                html_text = re.sub(
                    r"<title\b[^>]*>.*?</title>",
                    "",
                    data.decode("utf-8", "replace"),
                    flags=re.I | re.S,
                )
                ch = _load_html_str(html_text).strip()
                if ch:
                    parts += [ch, "", "---", ""]

            while parts and parts[-1] in ("", "---"):
                parts.pop()
            return (
                "\n".join(parts) + "\n"
                if parts
                else f"# {Path(path).name}\n\n*(empty EPUB)*\n"
            )
    except zipfile.BadZipFile:
        return "# EPUB Error\n\nNot a valid ZIP/EPUB file.\n"
    except Exception as e:
        return f"# EPUB Error\n\n```\n{e}\n```\n"


def _load_dtbook(path: str) -> str:
    """Load DTBook XML (DAISY digital talking book) into markdown.
    Supports DTBook 2005-3 and DAISY 3 NCX navigation."""
    _NS = {
        "dtb": "http://www.daisy.org/z3986/2005/dtbook/",
        "ncx": "http://www.daisy.org/z3986/2005/ncx/",
    }

    def _text(elem: ET.Element) -> str:
        return "".join(elem.itertext()).strip()

    def _walk(elem: ET.Element, depth: int = 0) -> List[str]:
        out: List[str] = []
        tag = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag

        if tag in (
            "dtbook",
            "book",
            "frontmatter",
            "bodymatter",
            "rearmatter",
            "pagenum",
        ):
            for child in elem:
                out.extend(_walk(child, depth))
        elif tag in ("h1", "h2", "h3", "h4", "h5", "h6"):
            level = int(tag[1])
            out.append("#" * level + " " + _text(elem))
            out.append("")
        elif tag in ("p",):
            t = _text(elem)
            if t:
                out.append(t)
                out.append("")
        elif tag in ("list",):
            for li in elem:
                li_tag = li.tag.split("}")[-1] if "}" in li.tag else li.tag
                if li_tag == "li":
                    out.append("* " + _text(li))
            out.append("")
        elif tag in ("imggroup",):
            for child in elem:
                ct = child.tag.split("}")[-1] if "}" in child.tag else child.tag
                if ct == "caption":
                    out.append(f"*{_text(child)}*")
                elif ct == "img":
                    alt = child.get("alt", "")
                    out.append(f"![{alt}]({child.get('src', '')})")
            out.append("")
        elif tag in ("sidebar", "prodnote", "annotation"):
            out.append("> " + _text(elem).replace("\n", "\n> "))
            out.append("")
        elif tag in ("level", "level1", "level2", "level3", "level4", "section", "div"):
            for child in elem:
                out.extend(_walk(child, depth + 1))
        elif tag in ("note", "footnote"):
            out.append(f"[^note]: {_text(elem)}")
        else:
            for child in elem:
                out.extend(_walk(child, depth))
        return out

    try:
        tree = ET.parse(path)
        root = tree.getroot()
        lines = _walk(root)
        # Collapse runs of blank lines
        result: List[str] = []
        blank = 0
        for ln in lines:
            if not ln.strip():
                blank += 1
                if blank <= 1:
                    result.append("")
            else:
                blank = 0
                result.append(ln)
        return "\n".join(result).strip()
    except Exception as e:
        return f"# DTBook Error\n\n```\n{e}\n```\n"


def _load_daisy_zip(path: str) -> str:
    """Load a DAISY book from a ZIP archive (Bookshare format).
    Tries DTBook XML first, then EPUB-style HTML chapters."""
    try:
        with zipfile.ZipFile(path, "r") as zf:
            names = zf.namelist()
            # Priority: DTBook XML
            for n in names:
                if n.lower().endswith((".xml", ".dtbook")) and "book" in n.lower():
                    tmp = CACHE_DIR / ("daisy_" + Path(n).name)
                    tmp.parent.mkdir(parents=True, exist_ok=True)
                    tmp.write_bytes(zf.read(n))
                    return _load_dtbook(str(tmp))
            # Fall back: any .html / .xhtml
            html_parts = []
            for n in sorted(names):
                if n.lower().endswith((".html", ".xhtml", ".htm")):
                    html_parts.append(
                        _load_html_str(zf.read(n).decode("utf-8", "replace"))
                    )
            if html_parts:
                return "\n\n---\n\n".join(html_parts)
            return "# DAISY Error\n\nNo readable content found in this ZIP.\n"
    except zipfile.BadZipFile:
        return "# DAISY Error\n\nNot a valid ZIP file.\n"
    except Exception as e:
        return f"# DAISY Error\n\n```\n{e}\n```\n"


def _load_csv_tsv(path: str, delim: str = ",") -> str:
    """Render CSV or TSV as a Markdown table."""
    try:
        rows: List[List[str]] = []
        with open(path, newline="", encoding="utf-8", errors="replace") as fh:
            reader = csv.reader(fh, delimiter=delim)
            for row in reader:
                rows.append(row)
        if not rows:
            return "*(empty file)*\n"
        nc = max(len(r) for r in rows)
        out: List[str] = [f"# {Path(path).name}", ""]
        hdr = (rows[0] + [""] * nc)[:nc]
        out.append("| " + " | ".join(h.replace("|", "\\|") for h in hdr) + " |")
        out.append("|" + "|".join([" --- "] * nc) + "|")
        for row in rows[1:]:
            cells = (row + [""] * nc)[:nc]
            out.append("| " + " | ".join(c.replace("|", "\\|") for c in cells) + " |")
        return "\n".join(out)
    except Exception as e:
        return f"# Error loading {Path(path).suffix.upper()}\n\n```\n{e}\n```\n"


def _load_xlsx(path: str) -> str:
    """Render XLSX spreadsheet as Markdown tables, one per sheet."""
    if not _XLSX:
        return (
            "# XLSX support not available\n\n"
            "Install openpyxl:  `pip install openpyxl`\n"
        )
    try:
        wb = _openpyxl.load_workbook(path, read_only=True, data_only=True)
        parts: List[str] = [f"# {Path(path).name}", ""]
        for ws in wb.worksheets:
            parts.append(f"## {ws.title}")
            parts.append("")
            rows = list(ws.iter_rows(values_only=True))
            if not rows:
                parts.append("*(empty sheet)*")
                parts.append("")
                continue
            nc = max(len(r) for r in rows)

            def _cell(v: Any) -> str:
                return str(v) if v is not None else ""

            hdr = [_cell(c) for c in (rows[0] + (None,) * nc)[:nc]]
            parts.append("| " + " | ".join(hdr) + " |")
            parts.append("|" + "|".join([" --- "] * nc) + "|")
            for row in rows[1:]:
                cells = [_cell(c) for c in (list(row) + [None] * nc)[:nc]]
                parts.append("| " + " | ".join(cells) + " |")
            parts.append("")
        return "\n".join(parts)
    except Exception as e:
        return f"# XLSX Error\n\n```\n{e}\n```\n"


def _load_docx(path: str) -> str:
    if not _DOCX:
        return (
            "# DOCX support not available\n\n"
            "Install python-docx:  `pip install python-docx`\n"
        )
    try:
        doc = _docx_lib.Document(path)
        out: List[str] = []
        for para in doc.paragraphs:
            sn = para.style.name.lower()
            txt = para.text
            if not txt.strip():
                out.append("")
                continue
            if "heading 1" in sn:
                out.append(f"# {txt}")
            elif "heading 2" in sn:
                out.append(f"## {txt}")
            elif "heading 3" in sn:
                out.append(f"### {txt}")
            elif "heading 4" in sn:
                out.append(f"#### {txt}")
            elif "heading" in sn:
                out.append(f"##### {txt}")
            elif "list" in sn:
                out.append(f"* {txt}")
            elif "code" in sn or "preformat" in sn:
                out.append(f"    {txt}")
            else:
                rich = ""
                for run in para.runs:
                    rt = run.text
                    if not rt:
                        continue
                    if run.bold and run.italic:
                        rt = f"***{rt}***"
                    elif run.bold:
                        rt = f"**{rt}**"
                    elif run.italic:
                        rt = f"*{rt}*"
                    rich += rt
                out.append(rich)
        for tbl in doc.tables:
            out.append("")
            for ri, row in enumerate(tbl.rows):
                cells = [c.text.replace("\n", " ").strip() for c in row.cells]
                out.append("| " + " | ".join(cells) + " |")
                if ri == 0:
                    out.append("|" + "|".join([" --- "] * len(cells)) + "|")
            out.append("")
        return "\n".join(out)
    except Exception as e:
        return f"# DOCX Error\n\n```\n{e}\n```\n"


def _load_doc(path: str) -> str:
    """Load a legacy binary Word (.doc / .dot) file as Markdown.

    Tries four approaches in order of preference:

    1. **python-docx** — works when the file is actually OOXML saved with a
       .doc extension (common with modern versions of Word on Windows).
    2. **antiword** subprocess — the lightweight dedicated converter for the
       true binary Word 97-2003 format.  Free binary available for Windows at
       https://www.winfield.demon.nl/  — add to PATH to enable.
    3. **LibreOffice headless** — converts .doc → .docx in a temp directory,
       then loads with python-docx.  Works if LibreOffice is installed.
    4. **Pandoc** — delegates to the existing _load_via_pandoc() helper.

    If none of these succeed a human-readable error with install instructions
    is returned.
    """
    title = Path(path).stem

    # ── 1. python-docx (file may be OOXML despite the .doc extension) ───────
    if _DOCX:
        try:
            md = _load_docx(path)
            if not md.startswith(("# DOCX Error", "# DOCX support")):
                return md
        except Exception:
            pass

    # ── 2. antiword ───────────────────────────────────────────────────
    antiword_bin = shutil.which("antiword") or shutil.which("antiword.exe")
    if antiword_bin:
        try:
            result = subprocess.run(
                [antiword_bin, "-w", "0", path],
                capture_output=True,
                timeout=30,
            )
            if result.returncode == 0:
                text = result.stdout.decode("utf-8", errors="replace").strip()
                if text:
                    return f"# {title}\n\n{text}\n"
        except Exception:
            pass

    # ── 3. LibreOffice headless (doc → docx → python-docx) ─────────────────
    lo_candidates: List[str] = ["soffice", "libreoffice"]
    if sys.platform == "win32":
        lo_candidates += [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ]
    lo_bin = next(
        (b for b in lo_candidates if shutil.which(b) or Path(b).exists()),
        None,
    )
    if lo_bin and _DOCX:
        try:
            import tempfile as _tmpmod

            with _tmpmod.TemporaryDirectory() as tmpdir:
                result = subprocess.run(
                    [
                        lo_bin,
                        "--headless",
                        "--convert-to",
                        "docx",
                        "--outdir",
                        tmpdir,
                        path,
                    ],
                    capture_output=True,
                    timeout=60,
                )
                if result.returncode == 0:
                    docx_path = Path(tmpdir) / (Path(path).stem + ".docx")
                    if docx_path.exists():
                        md = _load_docx(str(docx_path))
                        if not md.startswith(("# DOCX Error", "# DOCX support")):
                            return md
        except Exception:
            pass

    # ── 4. Pandoc ──────────────────────────────────────────────────────
    pandoc_md = _load_via_pandoc(path)
    if pandoc_md:
        return pandoc_md

    # ── Nothing worked ──────────────────────────────────────────────────
    return (
        f"# {title}\n\n"
        "**Could not load this binary Word (.doc) file.**\n\n"
        "Install one of the following to add .doc support:\n\n"
        "- **antiword** (lightest option): "
        "download the Windows binary from https://www.winfield.demon.nl/ "
        "and add it to your PATH\n"
        "- **LibreOffice** (also enables ODT/ODP/ODS conversion): "
        "https://www.libreoffice.org/\n"
        "- **Pandoc** + LibreOffice: `pip install pypandoc` then install "
        "Pandoc from https://pandoc.org/\n"
    )


def _load_pdf(path: str) -> str:
    if not _PDF and not (_OCR and _PYMUPDF):
        return (
            "# PDF support not available\n\n"
            "Install pdfminer.six:  `pip install pdfminer.six`\n"
            "For image/scanned PDFs also install:  `pip install pytesseract pymupdf`\n"
        )
    try:
        if _PDF == "layout":
            parts: List[str] = []
            for pnum, page in enumerate(_pdf_pages(path), 1):  # type: ignore[name-defined]
                page_texts = [
                    el.get_text().strip()
                    for el in page
                    if isinstance(el, LTTextBoxHorizontal)  # type: ignore[name-defined]
                    and el.get_text().strip()
                ]
                if not page_texts and _OCR and _PYMUPDF:
                    doc = _fitz.open(path)
                    pix = doc[pnum - 1].get_pixmap(matrix=_fitz.Matrix(2, 2))
                    img = _PIL_Image.frombytes(
                        "RGB", [pix.width, pix.height], pix.samples
                    )
                    doc.close()
                    page_texts = [_tesseract.image_to_string(img).strip()]
                parts.append(f"\n---\n*Page {pnum}*\n")
                parts.extend(page_texts)
            return "\n".join(parts)
        elif _PDF == "simple":
            return _strip_markdown_for_tts(_pdf_text(path) or "")  # type: ignore[name-defined]
        elif _OCR and _PYMUPDF:
            doc = _fitz.open(path)
            parts = []
            for pnum, page in enumerate(doc, 1):
                pix = page.get_pixmap(matrix=_fitz.Matrix(2, 2))
                img = _PIL_Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                parts.append(
                    f"\n---\n*Page {pnum}*\n\n"
                    + _tesseract.image_to_string(img).strip()
                )
            doc.close()
            return "\n".join(parts)
    except Exception as e:
        return f"# PDF Error\n\n```\n{e}\n```\n"
    return "# PDF Error\n\nUnknown failure.\n"


def _load_image_ocr(path: str) -> str:
    if not _OCR:
        return (
            "# OCR support not available\n\n"
            "Install pytesseract:  `pip install pytesseract`\n"
            "Also install Tesseract binary: https://github.com/tesseract-ocr/tesseract\n"
        )
    try:
        img = _PIL_Image.open(path).convert("RGB")
        text = _tesseract.image_to_string(img)
        return f"# {Path(path).name}\n\n{text.strip()}\n"
    except Exception as e:
        return f"# OCR Error\n\n```\n{e}\n```\n"


def _load_r_code(path: str) -> str:
    """Load R source with fenced code block for syntax highlighting."""
    try:
        src = Path(path).read_text(encoding="utf-8", errors="replace")
    except OSError as e:
        return f"# Error\n\n```\n{e}\n```\n"
    return f"# {Path(path).name}\n\n```r\n{src}\n```\n"


def _load_rmarkdown(path: str) -> str:
    """R Markdown: strip YAML front matter and render as markdown."""
    try:
        src = Path(path).read_text(encoding="utf-8", errors="replace")
    except OSError as e:
        return f"# Error\n\n```\n{e}\n```\n"
    src = re.sub(r"^---\s*\n.*?\n---\s*\n", "", src, flags=re.S)
    # Code chunks — wrap in fenced blocks with language tag
    src = re.sub(r"```\{r([^}]*)\}", r"```r", src)
    src = re.sub(r"```\{python([^}]*)\}", r"```python", src)
    return src


def _load_notebook(path: str) -> str:
    """Jupyter notebook: extract markdown and code cells."""
    try:
        nb = json.loads(Path(path).read_text(encoding="utf-8", errors="replace"))
        cells = nb.get("cells", [])
        parts: List[str] = [f"# {Path(path).name}", ""]
        for cell in cells:
            ct = cell.get("cell_type", "")
            src = "".join(cell.get("source", []))
            if ct == "markdown":
                parts.append(src)
                parts.append("")
            elif ct == "code":
                lang = (
                    nb.get("metadata", {})
                    .get("kernelspec", {})
                    .get("language", "python")
                )
                parts.append(f"```{lang}")
                parts.append(src)
                parts.append("```")
                parts.append("")
        return "\n".join(parts)
    except Exception as e:
        return f"# Notebook Error\n\n```\n{e}\n```\n"


def _orgmode_to_md(text: str) -> str:  # noqa: C901
    """Convert Org-mode markup to Markdown.

    Covers: headlines with TODO/tag/priority stripping, all block types
    (src, example, verbatim, quote, verse, comment), PROPERTIES and
    LOGBOOK drawers, bullet/numbered lists, tables, and the full set of
    Org inline markup (bold, italic, underline, strike-through, code,
    verbatim, links, footnotes).
    """
    lines = text.splitlines()
    out: List[str] = []
    title_prefix: List[str] = []
    in_block: Optional[str] = (
        None  # 'src'|'example'|'verbatim'|'quote'|'verse'|'comment'
    )
    i = 0

    def _inline(ln: str) -> str:
        """Apply Org inline markup to a plain line."""
        ln = re.sub(r"\*([^*\s][^*]*[^*\s]|[^*\s])\*", r"**\1**", ln)  # bold
        ln = re.sub(r"/([^/\s][^/]*[^/\s]|[^/\s])/", r"*\1*", ln)  # italic
        ln = re.sub(r"\+([^+\s][^+]*[^+\s]|[^+\s])\+", r"~~\1~~", ln)  # strike
        ln = re.sub(r"~([^~]+)~", r"`\1`", ln)  # code
        ln = re.sub(r"=([^=]+)=", r"`\1`", ln)  # verbatim
        ln = re.sub(r"\[\[([^\]]+)\]\[([^\]]+)\]\]", r"[\2](\1)", ln)  # link+desc
        ln = re.sub(r"\[\[([^\]]+)\]\]", r"[\1](\1)", ln)  # bare link
        ln = re.sub(r"\[fn:(\w+)\]", r"^[\1]^", ln)  # footnote
        return ln

    while i < len(lines):
        line = lines[i]
        stripped = line.strip()

        # ── Block begin ─────────────────────────────────────────────────
        bm = re.match(r"^\s*#\+BEGIN_(\w+)(.*)", stripped, re.I)
        if bm:
            btype = bm.group(1).lower()
            barg = bm.group(2).strip()
            in_block = btype
            if btype == "src":
                lang = barg.split()[0] if barg else ""
                out.append(f"```{lang}")
            elif btype in ("example", "verbatim"):
                out.append("```")
            # quote / verse / comment: no open fence
            i += 1
            continue

        # ── Block end ────────────────────────────────────────────────────
        if re.match(r"^\s*#\+END_", stripped, re.I):
            if in_block in ("src", "example", "verbatim"):
                out.append("```")
            in_block = None
            i += 1
            continue

        # ── Inside a block ──────────────────────────────────────────────
        if in_block == "comment":
            i += 1
            continue
        if in_block in ("src", "example", "verbatim"):
            out.append(line)
            i += 1
            continue
        if in_block in ("quote", "verse"):
            out.append("> " + _inline(stripped) if stripped else ">")
            i += 1
            continue

        # ── File-level directives ───────────────────────────────────────
        if stripped.startswith("#"):
            dm = re.match(r"#\+(\w+):\s*(.*)", stripped, re.I)
            if dm:
                key, val = dm.group(1).upper(), dm.group(2).strip()
                if key == "TITLE":
                    title_prefix[:] = [f"# {val}", ""]
                elif key == "SUBTITLE":
                    title_prefix.append(f"**{val}**")
                    title_prefix.append("")
                elif key == "AUTHOR":
                    title_prefix.append(f"*{val}*")
                elif key == "DATE":
                    title_prefix.append(f"*{val}*")
            i += 1
            continue

        # ── Drawers  :PROPERTIES: / :LOGBOOK: / arbitrary ──────────────
        if re.match(r"^:[\w-]+:\s*$", stripped) and stripped != ":END:":
            while i < len(lines) and lines[i].strip() != ":END:":
                i += 1
            i += 1  # skip :END:
            continue

        # ── Headlines ────────────────────────────────────────────────────
        hm = re.match(r"^(\*+)\s+(.*)", line)
        if hm:
            body = hm.group(2)
            body = re.sub(r"^(TODO|DONE|NEXT|WAITING|CANCELED|HOLD)\s+", "", body)
            body = re.sub(r"^COMMENT\s+", "", body)
            body = re.sub(r"\[#[A-Z]\]\s*", "", body)  # priority
            body = re.sub(r"\s+:[:\w@#%]+:\s*$", "", body)  # tags
            body = re.sub(r"\s*\[\d*/?\d*%?\]\s*", "", body)  # statistics
            out.append("#" * min(len(hm.group(1)), 6) + " " + body.strip())
            i += 1
            continue

        # ── Tables ──────────────────────────────────────────────────────
        if stripped.startswith("|"):
            # Separator row  |---+---|  → Markdown  | --- | --- |
            if re.match(r"^\|[-+]+\|?$", stripped.replace(" ", "")):
                # Count columns from previous row if available
                prev = out[-1] if out else ""
                ncols = prev.count("|") - 1 if "|" in prev else 1
                out.append("|" + "|".join([" --- "] * max(ncols, 1)) + "|")
            else:
                cells = [_inline(c.strip()) for c in stripped.strip("|").split("|")]
                out.append("| " + " | ".join(cells) + " |")
            i += 1
            continue

        # ── Lists ────────────────────────────────────────────────────────
        # Org unordered: - item  + item  (checkboxes: - [ ] item)
        lm = re.match(r"^(\s*)[-+]\s+(\[[ X-]\]\s+)?(.*)", line)
        if lm:
            checkbox = lm.group(2) or ""
            check_md = "[x] " if "X" in checkbox else "[ ] " if "[" in checkbox else ""
            out.append(lm.group(1) + "* " + check_md + _inline(lm.group(3)))
            i += 1
            continue
        # Org ordered: 1. item  1) item
        lm2 = re.match(r"^(\s*)\d+[.)]]\s+(.*)", line)
        if lm2:
            out.append(lm2.group(1) + "1. " + _inline(lm2.group(2)))
            i += 1
            continue

        # ── Plain text with inline markup ───────────────────────────────
        out.append(_inline(line))
        i += 1

    return "\n".join(title_prefix + out)


def _load_orgmode(path: str) -> str:
    """Load an Org-mode file as Markdown.

    Strategy: Pandoc → orgparse library → built-in _orgmode_to_md().
    """
    # 1. Pandoc (handles the full Org spec including exports, macros, etc.)
    md = _pandoc_convert(path, "org")
    if md:
        return md
    # 2. orgparse (pip install orgparse) — Python-native Org parser
    try:
        import orgparse as _op  # type: ignore[import]

        doc = _op.load(path)
        # Convert the orgparse tree to plain text lines then run through
        # the Markdown converter for any residual markup.
        lines_out: List[str] = []
        for node in doc.children:
            lines_out.extend(str(node).splitlines())
        return _orgmode_to_md("\n".join(lines_out))
    except Exception:
        pass
    # 3. Built-in comprehensive heuristic
    try:
        src = Path(path).read_text(encoding="utf-8", errors="replace")
        return _orgmode_to_md(src)
    except Exception as e:
        return f"# Org-mode Error\n\n```\n{e}\n```\n"


# =============================================================================
# Wiki format loaders  (RST · MediaWiki · AsciiDoc · Textile · Creole)
# =============================================================================
#
# Each loader follows the same three-tier strategy:
#   1. Pandoc with an explicit --from flag  —  highest quality when available
#   2. A dedicated Python library            —  no external binary required
#   3. A built-in regex converter            —  always works, covers ~80% of
#      real-world documents well enough for TTS
# =============================================================================


def _pandoc_convert(path: str, from_fmt: str) -> Optional[str]:
    """Run Pandoc with an explicit input format flag.  Returns Markdown or None."""
    if _PYPANDOC:
        try:
            return _pypandoc.convert_file(path, "markdown", format=from_fmt)
        except Exception:
            pass
    if _PANDOC_BIN:
        try:
            r = subprocess.run(
                [_PANDOC_BIN, "--from", from_fmt, "--to", "markdown", path],
                capture_output=True,
                text=True,
                timeout=30,
            )
            if r.returncode == 0 and r.stdout.strip():
                return r.stdout
        except Exception:
            pass
    return None


# ── reStructuredText (.rst / .rest) ──────────────────────────────────────────────


def _rst_to_md(text: str) -> str:
    """Basic reStructuredText → Markdown heuristic converter.

    Handles section headings (underline ± overline style), bold, italic,
    inline/block code, external hyperlinks, bullet / numbered lists, and
    note / warning / tip admonitions.
    """
    ADORN_CHARS = set(r"=-~^\"'`#+*@!$%&,./:<>?[]{}()")
    level_chars: List[str] = []  # adornment chars in first-encounter order
    lines = text.splitlines()
    out: List[str] = []
    i = 0

    def _heading(char: str, title: str) -> str:
        if char not in level_chars:
            level_chars.append(char)
        return "#" * min(level_chars.index(char) + 1, 6) + " " + title

    def _is_adorn(s: str) -> bool:
        return (
            bool(s)
            and len(s) >= 3
            and all(c == s[0] for c in s)
            and s[0] in ADORN_CHARS
        )

    while i < len(lines):
        raw = lines[i]
        stripped = raw.strip()

        # Overline + title + underline  (=====  Title  =====)
        if _is_adorn(stripped) and i + 2 < len(lines):
            title = lines[i + 1].strip()
            under = lines[i + 2].strip()
            if title and _is_adorn(under) and under[0] == stripped[0]:
                out.append(_heading(stripped[0], title))
                i += 3
                continue

        # Title + underline
        if stripped and i + 1 < len(lines):
            under = lines[i + 1].strip()
            if _is_adorn(under) and len(under) >= len(stripped):
                out.append(_heading(under[0], stripped))
                i += 2
                continue

        # Bare adornment line (leftover overline or separator)
        if _is_adorn(stripped):
            i += 1
            continue

        # Directives: .. code-block::, .. note::, etc.
        dm = re.match(r"\.\.\s+(\w[\w-]*)::(.*)$", stripped)
        if dm:
            directive, arg = dm.group(1).lower(), dm.group(2).strip()
            if directive in ("code", "code-block", "sourcecode"):
                out.append(f"```{arg}")
                i += 1
                while i < len(lines) and not lines[i].strip():
                    i += 1  # blank line after directive
                while i < len(lines) and (
                    lines[i].startswith("   ")
                    or lines[i].startswith("\t")
                    or not lines[i].strip()
                ):
                    body = lines[i]
                    if body.startswith("   "):
                        body = body[3:]
                    elif body.startswith("\t"):
                        body = body[1:]
                    out.append(body)
                    i += 1
                out.append("```")
                continue
            if directive in (
                "note",
                "warning",
                "tip",
                "important",
                "caution",
                "danger",
            ):
                parts = [arg] if arg else []
                i += 1
                while i < len(lines) and (
                    lines[i].startswith("   ") or not lines[i].strip()
                ):
                    if lines[i].strip():
                        parts.append(lines[i].strip())
                    i += 1
                out.append(f"> **{directive.capitalize()}:** {' '.join(parts)}")
                out.append("")
                continue
            i += 1
            continue

        # Hyperlink target  .. _name: url  (skip — simplified)
        if stripped.startswith(".. _") and ":" in stripped:
            i += 1
            continue

        # Inline markup
        raw = re.sub(r"``(.+?)``", r"`\1`", raw)  # inline code
        raw = re.sub(r"\*\*(.+?)\*\*", r"**\1**", raw)  # bold
        raw = re.sub(r"(?<!\*)\*(?!\s)(.+?)(?<!\s)\*(?!\*)", r"*\1*", raw)  # italic
        raw = re.sub(r"`([^`]+)\s+<([^>]+)>`_+", r"[\1](\2)", raw)  # hyperlink
        raw = re.sub(r"`([^`]+)`_\b", r"\1", raw)  # named ref
        out.append(raw)
        i += 1

    return "\n".join(out)


def _load_rst(path: str) -> str:
    """Load a reStructuredText file as Markdown."""
    # 1. Pandoc
    md = _pandoc_convert(path, "rst")
    if md:
        return md
    # 2. docutils (canonical Python RST library)
    try:
        from docutils.core import publish_parts  # type: ignore[import]

        src = Path(path).read_text(encoding="utf-8", errors="replace")
        html = publish_parts(src, writer_name="html")["html_body"]
        return _load_html_str(html)
    except Exception:
        pass
    # 3. Built-in heuristic converter
    try:
        src = Path(path).read_text(encoding="utf-8", errors="replace")
        return _rst_to_md(src)
    except Exception as e:
        return f"# RST Error\n\n```\n{e}\n```\n"


# ── MediaWiki (.wiki / .mediawiki) ───────────────────────────────────────────


def _mediawiki_to_md(text: str) -> str:
    """Basic MediaWiki markup → Markdown heuristic converter."""
    out_lines = []
    for line in text.splitlines():
        # Headings  == Title ==  /  === Title ===  etc.
        hm = re.match(r"^(={2,6})\s*(.+?)\s*\1\s*$", line)
        if hm:
            level = len(hm.group(1)) - 1  # == is h2 in wiki, treat as ##
            out_lines.append("#" * max(1, level) + " " + hm.group(2))
            continue

        # Template calls  {{...}}  — strip silently
        line = re.sub(r"\{\{[^}]+\}\}", "", line)
        # File / Image links  [[File:...]]  [[Image:...]]
        line = re.sub(r"\[\[(?:File|Image|Media):[^\]]*\]\]", "", line, flags=re.I)
        # Wikilinks  [[Page|display]]  or  [[Page]]
        line = re.sub(r"\[\[([^|\]]+)\|([^\]]+)\]\]", r"\2", line)
        line = re.sub(r"\[\[([^\]]+)\]\]", r"\1", line)
        # External links  [url text]
        line = re.sub(r"\[https?://\S+\s+([^\]]+)\]", r"\1", line)
        line = re.sub(r"\[https?://\S+\]", "", line)  # bare external link
        # Bold + italic  '''text'''  /  ''text''
        line = re.sub(r"'''(.+?)'''", r"**\1**", line)
        line = re.sub(r"''(.+?)''", r"*\1*", line)
        # Tables (very basic): just strip table markup
        if line.startswith("{|") or line.startswith("|}") or line.startswith("|+"):
            continue
        line = re.sub(r"^!\s*", "| ", line)  # header cell
        line = re.sub(r"^\|\|", "|", line)  # cell separator
        if re.match(r"^\|-", line):
            continue  # row separator
        out_lines.append(line)
    return "\n".join(out_lines)


def _load_mediawiki(path: str) -> str:
    """Load a MediaWiki markup file as Markdown."""
    # 1. Pandoc
    md = _pandoc_convert(path, "mediawiki")
    if md:
        return md
    # 2. mwparserfromhell (optional pure-Python MediaWiki parser)
    try:
        import mwparserfromhell as _mwp  # type: ignore[import]

        src = Path(path).read_text(encoding="utf-8", errors="replace")
        wikicode = _mwp.parse(src)
        # Strip templates and extract plain wikitext, then apply basic converter
        plain = wikicode.strip_code(normalize=True, collapse=True)
        return _mediawiki_to_md(plain)
    except Exception:
        pass
    # 3. Built-in heuristic
    try:
        src = Path(path).read_text(encoding="utf-8", errors="replace")
        return _mediawiki_to_md(src)
    except Exception as e:
        return f"# MediaWiki Error\n\n```\n{e}\n```\n"


# ── AsciiDoc (.adoc / .asciidoc / .asc) ─────────────────────────────────────


def _asciidoc_to_md(text: str) -> str:
    """Basic AsciiDoc → Markdown heuristic converter."""
    out_lines = []
    skip_next_block = False
    i = 0
    lines = text.splitlines()
    while i < len(lines):
        line = lines[i]
        # Block delimiter  ---- or ==== or ....  (listing / example blocks)
        if re.match(r"^[-=.+*_]{4,}$", line.strip()):
            skip_next_block = not skip_next_block
            if skip_next_block:
                # Check if it is a source block by looking at [source] above
                prev = out_lines[-1].strip() if out_lines else ""
                lang_m = re.match(r"\[source(?:,\s*(\w+))?", prev)
                lang = lang_m.group(1) if lang_m and lang_m.group(1) else ""
                if lang_m:
                    out_lines.pop()  # remove [source,...] line
                out_lines.append(f"```{lang}")
            else:
                out_lines.append("```")
            i += 1
            continue

        stripped = line.strip()

        # Section title  = h1  == h2  === h3 …
        hm = re.match(r"^(={1,6})\s+(.+)$", stripped)
        if hm:
            out_lines.append("#" * len(hm.group(1)) + " " + hm.group(2))
            i += 1
            continue

        # Attribute entries  :name: value  (skip)
        if re.match(r"^:[\w-]+:", stripped):
            i += 1
            continue

        # Block title  .Title
        if stripped.startswith(".") and not stripped.startswith("..."):
            out_lines.append(f"**{stripped[1:]}**")
            i += 1
            continue

        # Admonitions  NOTE: / TIP: / WARNING: / IMPORTANT: / CAUTION:
        am = re.match(r"^(NOTE|TIP|WARNING|IMPORTANT|CAUTION):\s*(.*)$", stripped)
        if am:
            out_lines.append(f"> **{am.group(1)}:** {am.group(2)}")
            i += 1
            continue

        # Inline markup
        line = re.sub(r"`(.+?)`", r"`\1`", line)  # monospace
        line = re.sub(r"\*\*(.+?)\*\*", r"**\1**", line)  # bold (unchanged)
        line = re.sub(r"\*(?!\*)(.+?)\*", r"**\1**", line)  # AsciiDoc *bold*
        line = re.sub(r"_\*(.+?)\*_|\*_(.+?)_\*", r"***\1\2***", line)  # bold italic
        line = re.sub(r"_(?!_)(.+?)_(?!_)", r"*\1*", line)  # italic
        line = re.sub(r"link:([^\[]+)\[([^\]]*)\]", r"[\2](\1)", line)  # link
        line = re.sub(
            r"https?://\S+\[([^\]]+)\]",
            lambda m: f"[{m.group(1)}]({m.group(0)[: m.group(0).index('[')]})",
            line,
        )  # URL[text]
        line = re.sub(r"<<([^,>]+),([^>]+)>>", r"\2", line)  # xref
        line = re.sub(r"<<([^>]+)>>", r"\1", line)  # bare xref
        line = re.sub(r"^\*\s+", "* ", line)  # list continuity
        line = re.sub(r"^\. ", "1. ", line)  # numbered list item
        line = re.sub(r"^\[.*?\]\s*$", "", line)  # block attribute lines
        out_lines.append(line)
        i += 1
    return "\n".join(out_lines)


def _load_asciidoc(path: str) -> str:
    """Load an AsciiDoc file as Markdown."""
    # 1. Pandoc
    md = _pandoc_convert(path, "asciidoc")
    if md:
        return md
    # 2. asciidoctor CLI  (renders to HTML, then convert)
    asciidoctor = shutil.which("asciidoctor") or shutil.which("asciidoc")
    if asciidoctor:
        try:
            r = subprocess.run(
                [asciidoctor, "-b", "html5", "-o", "-", path],
                capture_output=True,
                timeout=30,
            )
            if r.returncode == 0 and r.stdout:
                return _load_html_str(r.stdout.decode("utf-8", errors="replace"))
        except Exception:
            pass
    # 3. Built-in heuristic
    try:
        src = Path(path).read_text(encoding="utf-8", errors="replace")
        return _asciidoc_to_md(src)
    except Exception as e:
        return f"# AsciiDoc Error\n\n```\n{e}\n```\n"


# ── Textile (.textile) ───────────────────────────────────────────────────


def _textile_to_md(text: str) -> str:
    """Basic Textile → Markdown heuristic converter."""
    out_lines = []
    for line in text.splitlines():
        # Block tags  h1. h2. …  p.  bq.
        hm = re.match(r"^h([1-6])\. (.+)$", line)
        if hm:
            out_lines.append("#" * int(hm.group(1)) + " " + hm.group(2))
            continue
        if line.startswith("bq. "):
            out_lines.append("> " + line[4:])
            continue
        if line.startswith("p. "):
            line = line[3:]
        if line.startswith("pre. "):
            out_lines.append("```")
            out_lines.append(line[5:])
            out_lines.append("```")
            continue

        # Inline
        line = re.sub(r"\*\*(.+?)\*\*", r"**\1**", line)  # bold
        line = re.sub(r"\*(?!\*)(.+?)\*", r"**\1**", line)  # Textile *bold*
        line = re.sub(r"__(.+?)__", r"*\1*", line)  # italic
        line = re.sub(r"_(?!_)(.+?)_", r"*\1*", line)  # Textile _italic_
        line = re.sub(r"@(.+?)@", r"`\1`", line)  # code
        line = re.sub(r'"([^"]+)":(https?://\S+)', r"[\1](\2)", line)  # link
        line = re.sub(r"^\*{1,3} ", "* ", line)  # bullets (Textile allows */**)
        line = re.sub(r"^#{1,3} ", "1. ", line)  # numbered
        out_lines.append(line)
    return "\n".join(out_lines)


def _load_textile(path: str) -> str:
    """Load a Textile markup file as Markdown."""
    # 1. Pandoc
    md = _pandoc_convert(path, "textile")
    if md:
        return md
    # 2. textile Python library (pip install textile)
    try:
        import textile as _textile_lib  # type: ignore[import]

        src = Path(path).read_text(encoding="utf-8", errors="replace")
        html = _textile_lib.textile(src)
        return _load_html_str(html)
    except Exception:
        pass
    # 3. Built-in heuristic
    try:
        src = Path(path).read_text(encoding="utf-8", errors="replace")
        return _textile_to_md(src)
    except Exception as e:
        return f"# Textile Error\n\n```\n{e}\n```\n"


# ── Wiki Creole (.creole) ─────────────────────────────────────────────────


def _creole_to_md(text: str) -> str:
    """Basic Wiki Creole 1.0 → Markdown converter.

    Creole is a deliberately simple and regular wiki syntax; the full spec
    is small enough to implement here without external libraries.
    """
    out_lines = []
    in_nowiki_block = False
    for line in text.splitlines():
        # Nowiki blocks  {{{ / }}}
        if line.strip() == "{{{":
            in_nowiki_block = True
            out_lines.append("```")
            continue
        if line.strip() == "}}}":
            in_nowiki_block = False
            out_lines.append("```")
            continue
        if in_nowiki_block:
            out_lines.append(line)
            continue

        # Headings  == h2 ==  up to ====== h6 ======
        hm = re.match(r"^(={1,6})\s*(.+?)\s*=*\s*$", line)
        if hm:
            out_lines.append("#" * len(hm.group(1)) + " " + hm.group(2))
            continue

        # Horizontal rule  ----
        if re.match(r"^-{4,}$", line.strip()):
            out_lines.append("---")
            continue

        # Inline nowiki  {{{text}}}
        line = re.sub(r"\{\{\{(.+?)\}\}\}", r"`\1`", line)
        # Bold  **text**  → already Markdown
        # Italic  //text//
        line = re.sub(r"//(.+?)//", r"*\1*", line)
        # Links  [[url|text]]  [[url]]
        line = re.sub(r"\[\[([^|\]]+)\|([^\]]+)\]\]", r"[\2](\1)", line)
        line = re.sub(r"\[\[([^\]]+)\]\]", r"[\1](\1)", line)
        # Images  {{url|alt}}  {{url}}
        line = re.sub(r"\{\{([^|\}]+)\|([^\}]+)\}\}", r"![\2](\1)", line)
        line = re.sub(r"\{\{([^\}]+)\}\}", r"![](\1)", line)
        # Bullet lists  * item  (already Markdown-compatible)
        # Numbered lists  # item  → 1. item
        line = re.sub(r"^(#+) ", lambda m: "1. " * len(m.group(1)), line)
        # Table rows  |cell|cell|
        if line.startswith("|") and line.endswith("|"):
            cells = [c.strip() for c in line.strip("|").split("|")]
            # Header row has = in first cell convention; keep as-is
            cells = [re.sub(r"^=(.+?)=$", r"**\1**", c) for c in cells]
            line = "| " + " | ".join(cells) + " |"
        out_lines.append(line)
    return "\n".join(out_lines)


def _load_creole(path: str) -> str:
    """Load a Wiki Creole 1.0 file as Markdown."""
    # 1. Pandoc
    md = _pandoc_convert(path, "creole")
    if md:
        return md
    # 2. Built-in converter (Creole is simple and self-contained)
    try:
        src = Path(path).read_text(encoding="utf-8", errors="replace")
        return _creole_to_md(src)
    except Exception as e:
        return f"# Creole Error\n\n```\n{e}\n```\n"


# ── LaTeX (.tex / .ltx) ───────────────────────────────────────────────────


def _latex_to_md(text: str) -> str:  # noqa: C901
    """Strip LaTeX markup and convert to Markdown.

    Handles the constructs most commonly found in academic documents:
    preamble removal, sectioning, \\textbf / \\textit / \\emph,
    itemize / enumerate lists, verbatim / lstlisting / minted code
    blocks, quote / abstract environments, math stripping, special
    character normalization, and citation / cross-reference simplification.
    """
    # ─ Strip comments ───────────────────────────────────────────────────
    text = re.sub(r"%[^\n]*", "", text)

    # ─ Extract body (skip preamble) ─────────────────────────────────
    bm = re.search(r"\\begin\{document\}(.*?)\\end\{document\}", text, re.DOTALL)
    body = bm.group(1) if bm else text

    # ─ Collect title / author / date from preamble for a header block ──
    header_lines: List[str] = []
    preamble = text[: bm.start()] if bm else ""
    for cmd, fmt in (
        (r"\\title", "# {}"),
        (r"\\author", "*{} *"),
        (r"\\date", "*{} *"),
    ):
        m = re.search(cmd + r"\{([^}]+)\}", preamble + body)
        if m:
            header_lines.append(fmt.format(m.group(1).strip()))
    if header_lines:
        header_lines.append("")
        body = re.sub(r"\\maketitle\b", "", body)

    # ─ Verbatim / code environments ───────────────────────────────
    body = re.sub(
        r"\\begin\{verbatim\}(.*?)\\end\{verbatim\}",
        lambda m: "\n```\n" + m.group(1).strip() + "\n```\n",
        body,
        flags=re.DOTALL,
    )
    body = re.sub(
        r"\\begin\{lstlisting\}[^\n]*(.*?)\\end\{lstlisting\}",
        lambda m: "\n```\n" + m.group(1).strip() + "\n```\n",
        body,
        flags=re.DOTALL,
    )
    body = re.sub(
        r"\\begin\{minted\}\{(\w+)\}(.*?)\\end\{minted\}",
        lambda m: f"\n```{m.group(1)}\n" + m.group(2).strip() + "\n```\n",
        body,
        flags=re.DOTALL,
    )

    # ─ Quote / abstract environments ─────────────────────────────
    def _blockquote(m: re.Match) -> str:
        return (
            "\n" + "\n".join("> " + ln for ln in m.group(1).strip().splitlines()) + "\n"
        )

    body = re.sub(
        r"\\begin\{(quote|quotation|abstract)\}(.*?)\\end\{\1\}",
        _blockquote,
        body,
        flags=re.DOTALL,
    )

    # ─ List environments ───────────────────────────────────────────
    def _list_env(bullet: str) -> Callable[[re.Match], str]:
        def _replace(m: re.Match) -> str:
            items = re.split(r"\\item\b", m.group(1))
            result = []
            for it in items:
                it = it.strip()
                if it:
                    # \item[label] text  → **label** text
                    it = re.sub(r"^\[([^\]]+)\]\s*", r"**\1** ", it)
                    result.append(bullet + " " + it.replace("\n", " "))
            return "\n".join(result) + "\n"

        return _replace

    body = re.sub(
        r"\\begin\{itemize\}(.*?)\\end\{itemize\}",
        _list_env("*"),
        body,
        flags=re.DOTALL,
    )
    body = re.sub(
        r"\\begin\{enumerate\}(.*?)\\end\{enumerate\}",
        _list_env("1."),
        body,
        flags=re.DOTALL,
    )
    body = re.sub(
        r"\\begin\{description\}(.*?)\\end\{description\}",
        _list_env("*"),
        body,
        flags=re.DOTALL,
    )

    # ─ Math (strip display; keep inline content for TTS) ───────────
    for env_name in (
        "equation",
        "equation*",
        "align",
        "align*",
        "gather",
        "gather*",
        "multline",
        "multline*",
        "eqnarray",
        "eqnarray*",
        "math",
        "displaymath",
    ):
        body = re.sub(
            r"\\begin\{"
            + re.escape(env_name)
            + r"\}.*?\\end\{"
            + re.escape(env_name)
            + r"\}",
            "",
            body,
            flags=re.DOTALL,
        )
    body = re.sub(r"\$\$.*?\$\$", "", body, flags=re.DOTALL)
    body = re.sub(r"\\\[.*?\\\]", "", body, flags=re.DOTALL)
    body = re.sub(r"\\\((.*?)\\\)", r" \1 ", body, flags=re.DOTALL)
    body = re.sub(r"\$([^$\n]{1,120})\$", r" \1 ", body)

    # ─ Sectioning ────────────────────────────────────────────────
    for cmd, hashes in [
        ("part", "#"),
        ("chapter", "#"),
        ("section", "##"),
        ("subsection", "###"),
        ("subsubsection", "####"),
        ("paragraph", "#####"),
        ("subparagraph", "######"),
    ]:
        body = re.sub(
            r"\\" + cmd + r"\*?\{([^}]+)\}",
            lambda m, h=hashes: f"\n{h} {m.group(1)}\n",
            body,
        )

    # ─ Inline formatting ──────────────────────────────────────────
    body = re.sub(r"\\textbf\{([^}]+)\}", r"**\1**", body)
    body = re.sub(r"\\textit\{([^}]+)\}", r"*\1*", body)
    body = re.sub(r"\\emph\{([^}]+)\}", r"*\1*", body)
    body = re.sub(r"\\texttt\{([^}]+)\}", r"`\1`", body)
    body = re.sub(r"\\textsc\{([^}]+)\}", r"\1", body)
    body = re.sub(r"\\textsuperscript\{([^}]+)\}", r"^\1^", body)
    body = re.sub(r"\\textsubscript\{([^}]+)\}", r"~\1~", body)
    body = re.sub(r"\\underline\{([^}]+)\}", r"\1", body)
    body = re.sub(r"\\uline\{([^}]+)\}", r"\1", body)

    # ─ References, citations, footnotes ──────────────────────────
    body = re.sub(r"\\(?:cite|citep|citet|citealt|citealp)\{([^}]+)\}", r"[\1]", body)
    body = re.sub(r"\\(?:ref|pageref|eqref|nameref)\{[^}]+\}", "", body)
    body = re.sub(r"\\label\{[^}]+\}", "", body)
    body = re.sub(r"\\footnote\{([^}]{1,200})\}", r" (\1)", body)
    body = re.sub(r"\\footnotemark(?:\[\d+\])?", "", body)
    body = re.sub(r"\\footnotetext\{([^}]{1,200})\}", r" (\1)", body)

    # ─ URLs / hyperlinks ────────────────────────────────────────
    body = re.sub(r"\\url\{([^}]+)\}", r"\1", body)
    body = re.sub(r"\\href\{([^}]+)\}\{([^}]+)\}", r"[\2](\1)", body)

    # ─ Skip remaining environments entirely ────────────────────────
    body = re.sub(
        r"\\begin\{(figure|table|algorithm|tikzpicture|tabular)[^}]*\}.*?"
        r"\\end\{\1\}",
        "",
        body,
        flags=re.DOTALL | re.I,
    )

    # ─ Special characters and ligatures ─────────────────────────
    for latex, md_equiv in [
        ("---", "—"),
        ("--", "–"),
        (r"\\ldots", "…"),
        (r"\\dots", "…"),
        (r"\\cdots", "…"),
        ("``", "“"),
        ("''", "”"),
        (r"\\%", "%"),
        (r"\\\$", "$"),
        (r"\\&", "&"),
        (r"\\#", "#"),
        (r"\\{", "{"),
        (r"\\}", "}"),
        (r"\\textasciitilde", "~"),
        (r"\\textasciicircum", "^"),
        (r"\\textbackslash", "\\"),
        (r"\\slash", "/"),
        (r"\\,", " "),
        (r"\\;", " "),
        (r"\\!", ""),
        (r"\\quad", "  "),
    ]:
        body = re.sub(latex, md_equiv, body)

    # ─ Strip remaining block/environment tags ──────────────────────
    body = re.sub(r"\\begin\{[^}]+\}", "", body)
    body = re.sub(r"\\end\{[^}]+\}", "", body)

    # ─ Strip remaining LaTeX commands ───────────────────────────
    # Layout / spacing commands
    body = re.sub(
        r"\\(newpage|clearpage|pagebreak|noindent|par|linebreak|newline|\\)",
        "\n",
        body,
    )
    body = re.sub(
        r"\\(medskip|bigskip|smallskip|vspace\*?|hspace\*?)(?:\{[^}]*\})?",
        " ",
        body,
    )
    # Declaration-style commands
    body = re.sub(
        r"\\(centering|raggedright|raggedleft|normalfont|bfseries"
        r"|itshape|ttfamily|large|Large|LARGE|huge|Huge|small|footnotesize)",
        "",
        body,
    )
    # Any remaining command with an argument — keep the argument text
    body = re.sub(r"\\[a-zA-Z]+\*?\{([^}]{1,200})\}", r"\1", body)
    # Any remaining bare command
    body = re.sub(r"\\[a-zA-Z]+\*?\s*", " ", body)

    # ─ Clean up braces and whitespace ─────────────────────────────
    body = body.replace("{", "").replace("}", "")
    body = re.sub(r"\n{3,}", "\n\n", body)
    body = re.sub(r" {2,}", " ", body)

    return "\n".join(header_lines) + body.strip()


def _load_latex(path: str) -> str:
    """Load a LaTeX (.tex / .ltx) file as Markdown.

    Strategy: Pandoc → built-in _latex_to_md() stripper.

    Pandoc produces the best output for well-formed LaTeX (it handles
    cross-references, bibliographies, custom macros from \\newcommand, etc.).
    The built-in fallback covers the 80–90% case for typical academic papers
    and lecture notes without requiring any external tools.
    """
    # 1. Pandoc — also handles BibTeX references if they are inlined
    md = _pandoc_convert(path, "latex")
    if md:
        return md
    # 2. Built-in converter
    try:
        src = Path(path).read_text(encoding="utf-8", errors="replace")
        return _latex_to_md(src)
    except Exception as e:
        return f"# LaTeX Error\n\n```\n{e}\n```\n"


def _load_url(url: str) -> str:
    """Fetch a URL and convert to markdown."""
    try:
        req = urllib.request.Request(
            url,
            headers={"User-Agent": f"star/{APP_VERSION} (text reader; {sys.platform})"},
        )
        with urllib.request.urlopen(req, timeout=15) as resp:
            ct = resp.headers.get("Content-Type", "text/html")
            raw = resp.read()
            encoding = "utf-8"
            enc_m = re.search(r"charset=([^\s;]+)", ct)
            if enc_m:
                encoding = enc_m.group(1).strip("\"'")
            text = raw.decode(encoding, errors="replace")

        if "text/html" in ct or "xhtml" in ct:
            return _load_html_str(text)
        elif "text/plain" in ct or "markdown" in ct:
            return text
        elif "application/pdf" in ct:
            tmp = CACHE_DIR / "download.pdf"
            tmp.parent.mkdir(parents=True, exist_ok=True)
            tmp.write_bytes(raw)
            return _load_pdf(str(tmp))
        else:
            return text
    except urllib.error.URLError as e:
        return f"# URL Error\n\n```\n{e}\n```\n"
    except Exception as e:
        return f"# Fetch Error\n\n```\n{e}\n```\n"


def _load_via_pandoc(path: str) -> Optional[str]:
    """Use Pandoc binary or pypandoc to convert a file to Markdown.
    Returns None if Pandoc is not available."""
    if _PYPANDOC:
        try:
            return _pypandoc.convert_file(path, "markdown")
        except Exception:
            pass
    if _PANDOC_BIN:
        try:
            result = subprocess.run(
                [_PANDOC_BIN, "--to", "markdown", path],
                capture_output=True,
                text=True,
                timeout=30,
            )
            if result.returncode == 0:
                return result.stdout
        except Exception:
            pass
    return None


def load_document(path: str, settings: Settings) -> Document:
    """Load any supported document format and return a Document object."""
    doc = Document(path=path)
    fmt = _detect_format(path)
    doc.format = fmt

    # Check document cache before doing any parsing work
    if (
        settings.get("document_cache", True)
        and not path.startswith(("http://", "https://"))
        and fmt not in ("url",)
    ):
        fp = _settings_fingerprint(settings)
        cached = _cache_load(path, fp)
        if cached:
            doc.markdown = cached.get("markdown", "")
            doc.plain_text = cached.get("plain_text", "")
            doc.title = cached.get("title", Path(path).name)
            doc.format = cached.get("format", fmt)
            doc.metadata = cached.get("metadata", {})
            return doc

    # Dispatch table covers all registered formats; the else-branch below
    # is a last-resort Pandoc attempt for any extension we don't recognize.

    md: str = ""
    if fmt == "url":
        doc.title = path
        md = _load_url(path)
    elif fmt in ("text",):
        md = _load_plain_text(path)
    elif fmt in ("markdown", "rmarkdown"):
        md = _load_rmarkdown(path) if fmt == "rmarkdown" else _load_markdown(path)
    elif fmt == "html":
        md = _load_html(path)
    elif fmt == "epub":
        md = _load_epub(path)
    elif fmt in ("daisy", "xml") and path.lower().endswith(".xml"):
        md = _load_dtbook(path)
    elif fmt == "csv":
        md = _load_csv_tsv(path, ",")
    elif fmt == "tsv":
        md = _load_csv_tsv(path, "\t")
    elif fmt == "xlsx":
        md = _load_xlsx(path)
    elif fmt == "pptx":
        md = _load_pptx(path)
    elif fmt == "doc":
        md = _load_doc(path)
    elif fmt == "docx":
        md = _load_docx(path)
    elif fmt == "rst":
        md = _load_rst(path)
    elif fmt == "mediawiki":
        md = _load_mediawiki(path)
    elif fmt == "asciidoc":
        md = _load_asciidoc(path)
    elif fmt == "textile":
        md = _load_textile(path)
    elif fmt == "creole":
        md = _load_creole(path)
    elif fmt == "odt":
        md = _load_odt_v2(path)
    elif fmt == "pdf":
        md = _load_pdf(path)
    elif fmt == "image":
        md = _load_image_ocr(path)
    elif fmt == "r":
        md = _load_r_code(path)
    elif fmt == "notebook":
        md = _load_notebook(path)
    elif fmt == "latex":
        md = _load_latex(path)
    elif fmt == "orgmode":
        md = _load_orgmode(path)
    elif fmt in ("python", "javascript", "rust", "c"):
        lang_map = {
            "python": "python",
            "javascript": "javascript",
            "rust": "rust",
            "c": "c",
        }
        lang = lang_map.get(fmt, "")
        try:
            src = Path(path).read_text(encoding="utf-8", errors="replace")
        except OSError as e:
            src = str(e)
        md = f"# {Path(path).name}\n\n```{lang}\n{src}\n```\n"
    else:
        # Try Pandoc
        pandoc_md = _load_via_pandoc(path)
        if pandoc_md:
            md = pandoc_md
        else:
            md = _load_plain_text(path)

    # Extract title from first heading if not set
    if not doc.title:
        m = re.search(r"^#\s+(.+)$", md, re.MULTILINE)
        doc.title = m.group(1).strip() if m else Path(path).name if path else APP_TITLE

    # Apply footnote processing to markdown before stripping
    footnote_mode = str(settings.get("footnote_mode", "inline"))
    if footnote_mode != "inline":  # "inline" is already the default behavior
        md = _process_footnotes(md, mode=footnote_mode)

    doc.markdown = md
    doc.plain_text = _strip_markdown_for_tts(
        md,
        skip_code=settings["tts_skip_code"],
        table_mode=str(settings.get("table_reading_mode", "structured")),
    )

    # Extract EPUB chapter list (only for epub format)
    if fmt == "epub" and settings.get("epub_show_chapters", True):
        try:
            raw_chapters = _epub_extract_chapters(path)
            # Map hrefs to word indices via word map (built later); store titles+hrefs for now
            doc.chapters = [(t, h, 0) for t, h in raw_chapters]
        except Exception:
            doc.chapters = []

    # Cache the result
    if (
        settings.get("document_cache", True)
        and not path.startswith(("http://", "https://"))
        and len(doc.plain_text) > 1024
    ):
        fp = _settings_fingerprint(settings)
        try:
            _cache_save(
                path,
                {
                    "markdown": doc.markdown,
                    "plain_text": doc.plain_text,
                    "title": doc.title,
                    "format": doc.format,
                    "metadata": doc.metadata,
                },
                fp,
            )
        except Exception:
            pass

    return doc


# =============================================================================
# Markdown renderer (for curses TUI)
# =============================================================================

# Segment type: (text, role)
Seg = Tuple[str, str]
Line = List[Seg]

_INLINE_RE = re.compile(
    r"(`+)(.*?)\1"
    r"|\*{3}(.+?)\*{3}|_{3}(.+?)_{3}"
    r"|\*{2}(.+?)\*{2}|_{2}(.+?)_{2}"
    r"|\*([^*\n]+?)\*|_([^_\n]+?)_"
    r"|!\[([^\]]*)\]\([^)]*\)"
    r"|\[([^\]]*)\]\(([^)]*)\)",
    re.DOTALL,
)

_SYNTAX_PY_KW = set(
    "False None True and as assert async await break class continue def del "
    "elif else except finally for from global if import in is lambda nonlocal "
    "not or pass raise return try while with yield".split()
)
_SYNTAX_R_KW = set(
    "if else while for repeat in next break function NULL NA TRUE FALSE "
    "Inf NaN NA_integer_ NA_real_ NA_complex_ NA_character_".split()
)


def _parse_inline(text: str) -> List[Seg]:
    out: List[Seg] = []
    last = 0
    for m in _INLINE_RE.finditer(text):
        if m.start() > last:
            out.append((text[last : m.start()], "normal"))
        if m.group(1):
            out.append((m.group(2), "code"))
        elif m.group(3) or m.group(4):
            out.append((m.group(3) or m.group(4), "bolditalic"))
        elif m.group(5) or m.group(6):
            out.append((m.group(5) or m.group(6), "bold"))
        elif m.group(7) or m.group(8):
            out.append((m.group(7) or m.group(8), "italic"))
        elif m.group(9) is not None:
            out.append((f"[img: {m.group(9)}]", "image"))
        elif m.group(10):
            out.append((m.group(10), "link"))
        last = m.end()
    if last < len(text):
        out.append((text[last:], "normal"))
    return out or [(text, "normal")]


def _wrap_segs(segs: List[Seg], width: int) -> List[Line]:
    if width <= 0:
        return [list(segs)]
    result: List[Line] = []
    line: Line = []
    col = 0
    for text, role in segs:
        for tok in re.split(r"(\s+)", text):
            if not tok:
                continue
            is_ws = not tok.strip()
            tlen = len(tok)
            if is_ws:
                if col == 0:
                    continue
                if col + tlen <= width:
                    line.append((tok, role))
                    col += tlen
            else:
                if tlen >= width:
                    while tok:
                        chunk = tok[:width]
                        tok = tok[width:]
                        if line:
                            result.append(line)
                        line = [(chunk, role)]
                        col = len(chunk)
                elif col == 0:
                    line.append((tok, role))
                    col = tlen
                elif col + tlen <= width:
                    line.append((tok, role))
                    col += tlen
                else:
                    if line:
                        result.append(line)
                    line = [(tok, role)]
                    col = tlen
    if line:
        result.append(line)
    return result or [[]]


def render_markdown(
    md: str, width: int, tab_width: int = 4, syntax: bool = True
) -> List[Line]:
    """Convert a markdown string to a list of display lines (lists of Seg)."""
    width = max(10, width)
    lines: List[Line] = []
    src = md.splitlines()
    i = 0
    n = len(src)

    while i < n:
        ln = src[i]

        # Fenced code block
        fm = re.match(r"^(`{3,}|~{3,})\s*(\S*)", ln)
        if fm:
            fence, lang = fm.group(1), fm.group(2).lower()
            i += 1
            code_lines: List[str] = []
            while i < n and not src[i].startswith(fence[:3]):
                code_lines.append(src[i])
                i += 1
            i += 1  # closing fence
            lbl = f" {lang} " if lang else " "
            top = "┌" + lbl + "─" * max(0, width - 2 - len(lbl)) + "┐"
            lines.append([(top, "codeblock")])
            full_code = "\n".join(code_lines)
            for cl in code_lines:
                cl_exp = cl.replace("\t", " " * tab_width)
                if syntax and lang in ("python", "py"):
                    lines.append([("│ ", "codeblock")] + _lex_python_line(cl_exp))
                elif syntax and lang in ("r",):
                    lines.append([("│ ", "codeblock")] + _lex_r_line(cl_exp))
                else:
                    lines.append([("│ ", "codeblock"), (cl_exp, "codeblock")])
            bot = "└" + "─" * (width - 2) + "┘"
            lines.append([(bot, "codeblock")])
            lines.append([])
            continue

        # Setext headings
        if i + 1 < n and re.match(r"^=+\s*$", src[i + 1]) and ln.strip():
            lines.append([("# " + ln.strip(), "h1")])
            lines.append([("═" * min(len(ln) + 2, width), "h1")])
            lines.append([])
            i += 2
            continue
        if i + 1 < n and re.match(r"^-+\s*$", src[i + 1]) and ln.strip():
            lines.append([("## " + ln.strip(), "h2")])
            lines.append([("─" * min(len(ln) + 3, width), "h2")])
            lines.append([])
            i += 2
            continue

        # ATX heading
        hm = re.match(r"^(#{1,6})\s+(.*?)(?:\s+#+\s*)?$", ln)
        if hm:
            lv = min(len(hm.group(1)), 4)
            role = f"h{lv}"
            txt = hm.group(2).strip()
            prefix = {"h1": "# ", "h2": "## ", "h3": "### ", "h4": "#### "}[role]
            segs = [(prefix, role)] + [(s, role) for s, _ in _parse_inline(txt)]
            lines.append(segs)
            if lv == 1:
                lines.append([("═" * min(len(prefix + txt), width), "h1")])
            elif lv == 2:
                lines.append([("─" * min(len(prefix + txt), width), "h2")])
            lines.append([])
            i += 1
            continue

        # Horizontal rule
        if re.match(r"^(\*{3,}|-{3,}|_{3,})\s*$", ln.strip()):
            lines.append([("─" * width, "hr")])
            i += 1
            continue

        # Blockquote
        if ln.startswith(">"):
            qls: List[str] = []
            while i < n and (src[i].startswith(">") or (qls and not src[i].strip())):
                qls.append(src[i][1:].lstrip() if src[i].startswith(">") else "")
                i += 1
            text = " ".join(l for l in qls if l)
            pfx: Line = [("▌ ", "quote")]
            for wl in _wrap_segs(_parse_inline(text), width - 2):
                lines.append(pfx + wl)
            lines.append([])
            continue

        # Unordered list
        if re.match(r"^\s*[-*+]\s+\S", ln):
            while i < n and re.match(r"^\s*[-*+]\s+\S", src[i]):
                m2 = re.match(r"^\s*[-*+]\s+(.*)", src[i])
                item = m2.group(1) if m2 else ""
                rows = _wrap_segs(_parse_inline(item), width - 3)
                for k, row in enumerate(rows):
                    pfx_seg = [("  • ", "bullet")] if k == 0 else [("    ", "normal")]
                    lines.append(pfx_seg + row)
                i += 1
            lines.append([])
            continue

        # Ordered list
        if re.match(r"^\s*\d+[.)]\s+\S", ln):
            counter = 1
            while i < n and re.match(r"^\s*\d+[.)]\s+\S", src[i]):
                m2 = re.match(r"^\s*\d+[.)]\s+(.*)", src[i])
                item = m2.group(1) if m2 else ""
                pfx_str = f"  {counter}. "
                rows = _wrap_segs(_parse_inline(item), width - len(pfx_str))
                for k, row in enumerate(rows):
                    p = (
                        [(pfx_str, "ordinal")]
                        if k == 0
                        else [(" " * len(pfx_str), "normal")]
                    )
                    lines.append(p + row)
                counter += 1
                i += 1
            lines.append([])
            continue

        # Table
        if "|" in ln and i + 1 < n and re.match(r"^\|?[\s\-:|]+\|", src[i + 1]):
            tls: List[str] = []
            while i < n and "|" in src[i]:
                tls.append(src[i])
                i += 1
            for tl in tls:
                if re.match(r"^\|?[\s\-:|]+\|", tl):
                    cells = [c.strip() for c in tl.strip("|").split("|")]
                    sep = "┼".join("─" * (len(c) + 2) for c in cells)
                    lines.append([("├" + sep + "┤", "table")])
                else:
                    cells = [c.strip() for c in tl.strip("|").split("|")]
                    row_line: Line = [("│", "table")]
                    for j, cell in enumerate(cells):
                        if j:
                            row_line.append(("│", "table"))
                        row_line.append((" ", "normal"))
                        row_line.extend(_parse_inline(cell))
                        row_line.append((" ", "normal"))
                    row_line.append(("│", "table"))
                    lines.append(row_line)
            lines.append([])
            continue

        # Blank line
        if not ln.strip():
            lines.append([])
            i += 1
            continue

        # Paragraph
        pls: List[str] = [ln]
        i += 1
        while (
            i < n
            and src[i].strip()
            and not re.match(r"^(#{1,6}\s|[-*+]\s|\d+[.)]\s|>|`{3,}|~{3,})", src[i])
        ):
            pls.append(src[i])
            i += 1
        text = " ".join(l.rstrip() for l in pls)
        for wl in _wrap_segs(_parse_inline(text), width):
            lines.append(wl)
        lines.append([])

    return lines


def _lex_python_line(line: str) -> List[Seg]:
    """Tokenize a Python source line into (text, role) segments."""
    out: List[Seg] = []
    rest = line
    while rest:
        if rest.startswith("#"):
            out.append((rest, "comment"))
            break
        m = re.match(r'"[^"]*"|\'[^\']*\'', rest)
        if m:
            out.append((m.group(), "string"))
            rest = rest[m.end() :]
            continue
        m = re.match(r"[A-Za-z_]\w*", rest)
        if m:
            w = m.group()
            role = "keyword" if w in _SYNTAX_PY_KW else "code_normal"
            out.append((w, role))
            rest = rest[m.end() :]
            continue
        m = re.match(r"\d+\.?\d*", rest)
        if m:
            out.append((m.group(), "number"))
            rest = rest[m.end() :]
            continue
        out.append((rest[0], "code_normal"))
        rest = rest[1:]
    return out


def _lex_r_line(line: str) -> List[Seg]:
    """Tokenize an R source line into (text, role) segments."""
    out: List[Seg] = []
    rest = line
    while rest:
        if rest.startswith("#"):
            out.append((rest, "comment"))
            break
        m = re.match(r'"[^"]*"|\'[^\']*\'', rest)
        if m:
            out.append((m.group(), "string"))
            rest = rest[m.end() :]
            continue
        m = re.match(r"[A-Za-z_.][A-Za-z0-9_.]*", rest)
        if m:
            w = m.group()
            role = "keyword" if w in _SYNTAX_R_KW else "code_normal"
            out.append((w, role))
            rest = rest[m.end() :]
            continue
        m = re.match(r"\d+\.?\d*", rest)
        if m:
            out.append((m.group(), "number"))
            rest = rest[m.end() :]
            continue
        out.append((rest[0], "code_normal"))
        rest = rest[1:]
    return out


def lines_to_plain(rendered: List[Line]) -> str:
    """Convert rendered lines back to plain text (strip roles)."""
    return "\n".join("".join(t for t, _ in line) for line in rendered)


# =============================================================================
# Search engine
# =============================================================================


class SearchEngine:
    """Forward/backward incremental search with fuzzy matching."""

    def __init__(self) -> None:
        self.query = ""
        self._matches: List[Tuple[int, int, int]] = []  # (line, col_start, col_end)
        self._idx = -1

    @property
    def matches(self) -> List[Tuple[int, int, int]]:
        return self._matches

    @property
    def current_match(self) -> Optional[Tuple[int, int, int]]:
        if 0 <= self._idx < len(self._matches):
            return self._matches[self._idx]
        return None

    @property
    def match_count(self) -> int:
        return len(self._matches)

    @property
    def match_index(self) -> int:
        return self._idx

    def search(self, query: str, rendered: List[Line], from_line: int = 0) -> bool:
        """Run a new search.  Returns True if any matches found."""
        self.query = query
        self._matches = []
        self._idx = -1
        if not query:
            return False
        q_lower = query.lower()
        for li, segs in enumerate(rendered):
            text = "".join(t for t, _ in segs)
            tl = text.lower()
            col = 0
            while True:
                pos = tl.find(q_lower, col)
                if pos < 0:
                    break
                self._matches.append((li, pos, pos + len(query)))
                col = pos + 1
        if self._matches:
            # Start from from_line
            for i, (li, _, _) in enumerate(self._matches):
                if li >= from_line:
                    self._idx = i
                    break
            else:
                self._idx = 0
        return bool(self._matches)

    def next_match(self) -> Optional[Tuple[int, int, int]]:
        if not self._matches:
            return None
        self._idx = (self._idx + 1) % len(self._matches)
        return self._matches[self._idx]

    def prev_match(self) -> Optional[Tuple[int, int, int]]:
        if not self._matches:
            return None
        self._idx = (self._idx - 1) % len(self._matches)
        return self._matches[self._idx]

    def search_regex(
        self, pattern: str, rendered: "List[Line]", from_line: int = 0
    ) -> bool:
        """Run a regex search. Falls back to plain search if pattern is invalid."""
        try:
            rx = re.compile(pattern, re.IGNORECASE)
        except re.error:
            return self.search(pattern, rendered, from_line)
        self.query = pattern
        self._matches = []
        self._idx = -1
        for li, segs in enumerate(rendered):
            text = "".join(t for t, _ in segs)
            for m in rx.finditer(text):
                self._matches.append((li, m.start(), m.end()))
        if self._matches:
            for i, (li, _, _) in enumerate(self._matches):
                if li >= from_line:
                    self._idx = i
                    break
            else:
                self._idx = 0
        return bool(self._matches)


# =============================================================================
# Line editor (Emacs-style — shared with the other projects in this suite)
# =============================================================================


class LineEditor:
    """Single-line text editor.  Supports arrow keys, Home/End/Delete and
    the standard Ctrl shortcuts Windows and Mac users expect."""

    def __init__(self, value: str = ""):
        self.buf = list(str(value))
        self.pos = len(self.buf)
        self.hint = ""
        self.hint_full = ""
        self._kill_ring = ""

    def feed(self, ch: int) -> Optional[bool]:
        """
        Returns:
          False  — confirmed (Enter)
          None   — canceled (C-g / Esc)
          True   — still editing
        """
        if ch in (curses.KEY_ENTER, 10, 13):
            return False
        if ch in (7, 27):
            return None

        # ── Cursor movement ──────────────────────────────────────────────
        if ch in (curses.KEY_LEFT,):
            self.pos = max(0, self.pos - 1)
        elif ch in (curses.KEY_RIGHT,):
            self.pos = min(len(self.buf), self.pos + 1)
        elif ch in (curses.KEY_HOME, 1):  # Home or Ctrl+A
            self.pos = 0
        elif ch in (curses.KEY_END, 5):  # End or Ctrl+E
            self.pos = len(self.buf)
        # ── Deletion ──────────────────────────────────────────────────────
        elif ch in (curses.KEY_BACKSPACE, 127, 8):  # Backspace
            if self.pos > 0:
                self.buf.pop(self.pos - 1)
                self.pos -= 1
        elif ch in (curses.KEY_DC, 4):  # Delete key or Ctrl+D
            if self.pos < len(self.buf):
                self.buf.pop(self.pos)
        elif ch == 23:  # Ctrl+W / Ctrl+Backspace — delete word backward
            end = self.pos
            while self.pos > 0 and self.buf[self.pos - 1] == " ":
                self.pos -= 1
            while self.pos > 0 and self.buf[self.pos - 1] != " ":
                self.pos -= 1
            self._kill_ring = "".join(self.buf[self.pos : end])
            self.buf[self.pos : end] = []
        elif ch == 11:  # Ctrl+K — delete to end of line
            self._kill_ring = "".join(self.buf[self.pos :])
            self.buf = self.buf[: self.pos]
        elif ch == 21:  # Ctrl+U — delete to start of line
            self._kill_ring = "".join(self.buf[: self.pos])
            self.buf = self.buf[self.pos :]
            self.pos = 0
        elif ch == 25:  # Ctrl+Y — paste last deleted text
            for c in self._kill_ring:
                self.buf.insert(self.pos, c)
                self.pos += 1
        elif 32 <= ch <= 126 or ch > 127:
            try:
                self.buf.insert(self.pos, chr(ch))
                self.pos += 1
            except (ValueError, OverflowError):
                pass
        return True

    def accept_hint(self) -> None:
        if not self.hint or self.pos != len(self.buf):
            return
        if self.hint_full:
            text = "".join(self.buf)
            sep = max(text.rfind(" "), text.rfind("/"))
            new_val = (text[: sep + 1] if sep >= 0 else "") + self.hint_full
            self.buf = list(new_val)
        else:
            for c in self.hint:
                self.buf.append(c)
        self.pos = len(self.buf)
        self.hint = self.hint_full = ""

    @property
    def value(self) -> str:
        return "".join(self.buf)

    def set_value(self, v: str) -> None:
        self.buf = list(v)
        self.pos = len(self.buf)


# =============================================================================
# Color / theme system
# =============================================================================

# Color-pair role names → CP_* numbers
_ROLES = [
    "normal",
    "h1",
    "h2",
    "h3",
    "h4",
    "bold",
    "italic",
    "bolditalic",
    "code",
    "code_normal",
    "codeblock",
    "keyword",
    "string",
    "comment",
    "number",
    "link",
    "image",
    "quote",
    "bullet",
    "ordinal",
    "table",
    "hr",
    "current_word",  # TTS word highlight
    "search_match",  # non-current search hit
    "search_current",  # current search hit
    "status",  # status bar
    "status_hi",  # emphasized item in status bar
    "minibuf",  # minibuffer normal
    "error",  # error message
    "dim",  # hints / secondary text
    "progress",  # loading indicator
    "title_bar",  # top title bar
]
CP: Dict[str, int] = {r: i + 1 for i, r in enumerate(_ROLES)}

# (fg, bg, bold, italic, underline, dim)
_N = (7, -1, False, False, False, False)


def _t(**kw: tuple) -> Dict[str, tuple]:
    d: Dict[str, tuple] = {r: _N for r in _ROLES}
    d.update(kw)
    return d


# Dark modern theme (default) — colorblind-friendly (no red/green adjacency)
# Accent palette: cyan, blue, magenta, white.  No yellow, no green/red pairing.
THEMES: Dict[str, Dict] = {
    "dark": _t(
        normal=(7, -1, False, False, False, False),
        h1=(6, -1, True, False, False, False),  # cyan bold
        h2=(4, -1, True, False, False, False),  # blue bold
        h3=(5, -1, True, False, False, False),  # magenta bold
        h4=(7, -1, True, False, True, False),  # white bold underline
        bold=(7, -1, True, False, False, False),
        italic=(7, -1, False, True, False, False),
        bolditalic=(7, -1, True, True, False, False),
        code=(6, -1, False, False, False, False),  # cyan
        code_normal=(6, -1, False, False, False, False),
        codeblock=(6, -1, False, False, False, True),  # dim cyan
        keyword=(5, -1, True, False, False, False),  # magenta bold
        string=(6, -1, False, False, False, False),  # cyan
        comment=(4, -1, False, True, False, True),  # blue italic dim
        number=(5, -1, False, False, False, False),  # magenta
        link=(6, -1, False, False, True, False),  # cyan underline
        image=(5, -1, False, False, True, False),
        quote=(5, -1, False, True, False, False),  # magenta italic
        bullet=(6, -1, True, False, False, False),  # cyan bold
        ordinal=(6, -1, False, False, False, False),
        table=(6, -1, False, False, False, False),
        hr=(4, -1, False, False, False, True),  # blue dim
        current_word=(0, 6, True, False, False, False),  # black on cyan
        search_match=(0, 4, False, False, False, False),  # black on blue
        search_current=(0, 5, True, False, False, False),  # black on magenta
        status=(7, 4, True, False, False, False),  # white on blue
        status_hi=(6, 4, True, False, False, False),  # cyan on blue
        minibuf=(7, -1, False, False, False, False),
        error=(5, -1, True, False, False, False),  # magenta bold (no red)
        dim=(7, -1, False, False, False, True),
        progress=(6, 4, True, False, False, False),
        title_bar=(7, 4, True, False, False, False),
    ),
    "light": _t(
        normal=(0, 7, False, False, False, False),
        h1=(4, 7, True, False, False, False),  # blue bold
        h2=(5, 7, True, False, False, False),  # magenta bold
        h3=(4, 7, False, False, False, False),  # blue
        h4=(0, 7, True, False, True, False),
        bold=(0, 7, True, False, False, False),
        italic=(0, 7, False, True, False, False),
        bolditalic=(0, 7, True, True, False, False),
        code=(4, 7, False, False, False, False),
        code_normal=(4, 7, False, False, False, False),
        codeblock=(4, 7, False, False, False, True),
        keyword=(5, 7, True, False, False, False),
        string=(4, 7, False, False, False, False),
        comment=(5, 7, False, True, False, True),
        number=(5, 7, False, False, False, False),
        link=(4, 7, False, False, True, False),
        image=(5, 7, False, False, True, False),
        quote=(5, 7, False, True, False, False),
        bullet=(4, 7, True, False, False, False),
        ordinal=(4, 7, False, False, False, False),
        table=(4, 7, False, False, False, False),
        hr=(5, 7, False, False, False, True),
        current_word=(7, 4, True, False, False, False),  # white on blue
        search_match=(7, 5, False, False, False, False),  # white on magenta
        search_current=(7, 4, True, False, False, False),
        status=(7, 4, True, False, False, False),
        status_hi=(7, 5, True, False, False, False),
        minibuf=(0, 7, False, False, False, False),
        error=(5, 7, True, False, False, False),
        dim=(0, 7, False, False, False, True),
        progress=(7, 4, False, False, False, False),
        title_bar=(7, 4, True, False, False, False),
    ),
    "contrast": _t(
        # High contrast: bold white on black, cyan & magenta accents
        normal=(7, 0, False, False, False, False),
        h1=(6, 0, True, False, False, False),
        h2=(7, 0, True, False, False, False),
        h3=(5, 0, True, False, False, False),
        h4=(7, 0, True, False, True, False),
        bold=(7, 0, True, False, False, False),
        italic=(7, 0, False, True, False, False),
        bolditalic=(7, 0, True, True, False, False),
        code=(6, 0, True, False, False, False),
        code_normal=(6, 0, True, False, False, False),
        codeblock=(6, 0, False, False, False, False),
        keyword=(5, 0, True, False, False, False),
        string=(6, 0, False, False, False, False),
        comment=(7, 0, False, False, False, False),
        number=(5, 0, False, False, False, False),
        link=(6, 0, False, False, True, False),
        image=(5, 0, False, False, True, False),
        quote=(7, 0, False, False, False, False),
        bullet=(6, 0, True, False, False, False),
        ordinal=(6, 0, False, False, False, False),
        table=(7, 0, False, False, False, False),
        hr=(7, 0, False, False, False, False),
        current_word=(0, 6, True, False, False, False),
        search_match=(0, 7, False, False, False, False),
        search_current=(0, 5, True, False, False, False),
        status=(0, 7, True, False, False, False),
        status_hi=(0, 6, True, False, False, False),
        minibuf=(7, 0, True, False, False, False),
        error=(5, 0, True, False, False, False),
        dim=(7, 0, False, False, False, False),
        progress=(0, 6, False, False, False, False),
        title_bar=(0, 7, True, False, False, False),
    ),
    "phosphor": _t(
        # Classic green phosphor monochrome
        normal=(2, -1, False, False, False, False),
        h1=(2, -1, True, False, False, False),
        h2=(2, -1, True, False, True, False),
        h3=(2, -1, False, False, True, False),
        h4=(2, -1, True, False, False, False),
        bold=(2, -1, True, False, False, False),
        italic=(2, -1, False, True, False, False),
        bolditalic=(2, -1, True, True, False, False),
        code=(2, -1, True, False, False, False),
        code_normal=(2, -1, False, False, False, False),
        codeblock=(2, -1, False, False, False, True),
        keyword=(2, -1, True, False, False, False),
        string=(2, -1, False, False, False, False),
        comment=(2, -1, False, True, False, True),
        number=(2, -1, False, False, False, False),
        link=(2, -1, False, False, True, False),
        image=(2, -1, False, False, True, False),
        quote=(2, -1, False, True, False, False),
        bullet=(2, -1, True, False, False, False),
        ordinal=(2, -1, False, False, False, False),
        table=(2, -1, False, False, False, False),
        hr=(2, -1, False, False, False, True),
        current_word=(0, 2, True, False, False, False),
        search_match=(0, 2, False, False, False, False),
        search_current=(2, 0, True, False, False, False),
        status=(0, 2, True, False, False, False),
        status_hi=(0, 2, False, False, False, False),
        minibuf=(2, -1, False, False, False, False),
        error=(2, -1, True, False, False, False),
        dim=(2, -1, False, False, False, True),
        progress=(0, 2, False, False, False, False),
        title_bar=(0, 2, True, False, False, False),
    ),
}

THEME_NAMES = list(THEMES.keys())

# Roles that mark a heading line in the rendered output.
_HEADING_ROLES = frozenset({"h1", "h2", "h3", "h4"})

# Roles that mark a table line in the rendered output.
_TABLE_ROLES = frozenset({"table"})

# Sentence boundary detector.
# Matches the gap between sentences: terminal punctuation followed by
# whitespace + an uppercase letter / opening quote, or a blank line.
# Intentionally simple — some abbreviations (Dr., Mr.) will cause false
# splits, which is acceptable for a reading application.
_SENTENCE_SPLIT_RE = re.compile(
    r"(?<=[.!?\u2026])\s+(?=[A-Z\"\u201c\u2018(])"
    r"|(?<=[.!?\u2026])\s*\n"
    r"|\n{2,}"
)


def _setup_colors(theme_name: str) -> Dict[str, int]:
    """Initialize curses color pairs from a theme dict.
    Returns a mapping role → combined curses attribute integer."""
    if not curses.has_colors():
        return {r: curses.A_NORMAL for r in _ROLES}
    try:
        curses.start_color()
        curses.use_default_colors()
    except curses.error:
        pass

    theme = THEMES.get(theme_name, THEMES["dark"])
    _ATTR = {
        "bold": curses.A_BOLD,
        "italic": getattr(curses, "A_ITALIC", 0),
        "underline": curses.A_UNDERLINE,
        "dim": curses.A_DIM,
    }
    result: Dict[str, int] = {}
    for role in _ROLES:
        fg, bg, b, it, ul, dim = theme[role]
        cp = CP[role]
        try:
            curses.init_pair(cp, fg, bg)
        except curses.error:
            pass
        attr = curses.color_pair(cp)
        if b:
            attr |= curses.A_BOLD
        if it:
            attr |= _ATTR["italic"]
        if ul:
            attr |= curses.A_UNDERLINE
        if dim:
            attr |= curses.A_DIM
        result[role] = attr
    return result


# =============================================================================
# M-x command catalog
# =============================================================================

MX_COMMANDS = sorted(
    [
        "about",
        "backend",
        "book-next",
        "book-prev",
        "close",
        "contrast-up",
        "contrast-down",
        "export-braille",
        "export-markdown",
        "font-size-down",
        "font-size-up",
        "goto-line",
        "help",
        "license",
        "line-numbers",
        "next-heading",
        "next-paragraph",
        "read-next-heading",
        "read-prev-heading",
        "speech-cursor",
        "stop-speech",
        "abbrev-add",
        "abbrev-list",
        "expand-abbreviations",
        "normalize-numbers",
        "table-mode",
        "next-sentence",
        "open",
        "open-url",
        "pause",
        "play",
        "preset-add",
        "preset-list",
        "prev-heading",
        "prev-paragraph",
        "prev-sentence",
        "replay-paragraph",
        "replay-sentence",
        "save-position",
        "clear-position",
        "jump-saved",
        "speed",
        "ssml",
        "ssml-on",
        "ssml-off",
        "rate-down",
        "rate-up",
        "reload",
        "search",
        "search-backward",
        "settings",
        "speak-line",
        "speak-selection",
        "stop",
        "syntax-highlight",
        "theme",
        "tts-backend",
        "tts-voice",
        "voice-picker",
        "volume-down",
        "volume-up",
        "voice",
        "quit",
        "wrap-width",
        "abbrev-add",
        "abbrev-list",
        "bookmark-delete",
        "bookmark-goto",
        "bookmark-list",
        "bookmark-set",
        "cache-clear",
        "chapter-goto",
        "chapter-list",
        "chapter-next",
        "chapter-prev",
        "copy",
        "expand-abbreviations",
        "footnote-mode",
        "font",
        "history-back",
        "history-forward",
        "normalize-math",
        "normalize-numbers",
        "pubmed",
        "reading-level",
        "recent",
        "search-regex",
        "table-mode",
        "wiki",
        "annotate",
        "annotations-list",
        "annotations-search",
        "annotation-goto",
        "annotation-delete",
        "annotations-export",
        "shortcuts",
    ]
)


# =============================================================================
# Helper functions
# =============================================================================


def _addstr(win: "curses.window", y: int, x: int, s: str, attr: int = 0) -> None:
    h, w = win.getmaxyx()
    if y < 0 or y >= h or x < 0 or x >= w or not s:
        return
    s = s[: max(0, w - x)]
    if not s:
        return
    try:
        win.addstr(y, x, s, attr)
    except curses.error:
        pass


def _fillrow(win: "curses.window", y: int, attr: int = 0, ch: str = " ") -> None:
    h, w = win.getmaxyx()
    if y < 0 or y >= h:
        return
    try:
        win.addstr(y, 0, ch * (w - 1), attr)
    except curses.error:
        pass


# North American Braille ASCII (a.k.a. NABCC / Braille-ASCII) lookup, indexed
# by the 6-dot cell value (bit0=dot1, bit1=dot2 … bit5=dot6).  Each braille
# cell maps to exactly one ASCII byte; this is the character set BRF embossers
# expect.  This table is the canonical Unicode-braille (U+2800..U+283F) order.
_BRAILLE_ASCII = " A1B'K2L@CIF/MSP\"E3H9O6R^DJG>NTQ,*5<-U8V.%[$+X!&;:4\\0Z7(_?W]#Y)="

# Print character -> braille dot-cell value, for uncontracted (Grade 1) English.
_BRL_LETTER = {
    chr(ord("a") + i): v
    for i, v in enumerate(
        # a  b  c   d   e   f   g   h   i   j   k  l  m   n   o   p   q   r   s
        [
            1,
            3,
            9,
            25,
            17,
            11,
            27,
            19,
            10,
            26,
            5,
            7,
            13,
            29,
            21,
            15,
            31,
            23,
            14,
            # t   u   v   w   x   y   z
            30,
            37,
            39,
            58,
            45,
            61,
            53,
        ]
    )
}
_BRL_DIGIT = {
    "1": 1,
    "2": 3,
    "3": 9,
    "4": 25,
    "5": 17,
    "6": 11,
    "7": 27,
    "8": 19,
    "9": 10,
    "0": 26,
}
_BRL_PUNCT = {
    ".": 50,
    ",": 2,
    ";": 6,
    ":": 18,
    "?": 38,
    "!": 22,
    "'": 4,
    "-": 36,
    "(": 54,
    ")": 54,
    "/": 12,
    '"': 38,
}
_BRL_NUMBER_SIGN = 60  # dots 3-4-5-6
_BRL_CAPITAL_SIGN = 32  # dot 6


def _text_to_braille_grade1(text: str) -> str:
    """Translate plain text to uncontracted (Grade 1) Braille-ASCII.

    Pure-Python, dependency-free, and incapable of crashing the process — the
    reliable default for BRF export.  Handles letters, digits (with number
    sign), capital signs, common punctuation, and whitespace.  Unknown
    characters are dropped so the output stays valid Braille-ASCII.
    """
    out: List[str] = []
    in_number = False
    for ch in text:
        if ch == "\n":
            out.append("\n")
            in_number = False
            continue
        if ch == "\t" or ch == " ":
            out.append(" ")
            in_number = False
            continue
        low = ch.lower()
        if low in _BRL_LETTER:
            if ch.isupper():
                out.append(_BRAILLE_ASCII[_BRL_CAPITAL_SIGN])
            # A letter a-j immediately after a number must be separated from
            # the number by a letter sign; we end number mode on any letter.
            in_number = False
            out.append(_BRAILLE_ASCII[_BRL_LETTER[low]])
        elif ch in _BRL_DIGIT:
            if not in_number:
                out.append(_BRAILLE_ASCII[_BRL_NUMBER_SIGN])
                in_number = True
            out.append(_BRAILLE_ASCII[_BRL_DIGIT[ch]])
        elif ch in _BRL_PUNCT:
            in_number = False
            out.append(_BRAILLE_ASCII[_BRL_PUNCT[ch]])
        # else: unsupported glyph -> skip (keeps output valid Braille-ASCII)
    return "".join(out)


def _format_brf(braille: str, cells: int = 40, lines_per_page: int = 25) -> str:
    """Wrap a Braille-ASCII string into standard BRF page geometry.

    Word-wraps at *cells* columns (default 40), groups output into pages of
    *lines_per_page* lines separated by a form feed, and uses CRLF line
    endings as embossers expect.
    """
    lines: List[str] = []
    for para in braille.split("\n"):
        if not para:
            lines.append("")
            continue
        cur = ""
        for word in para.split(" "):
            if not cur:
                cur = word
            elif len(cur) + 1 + len(word) <= cells:
                cur += " " + word
            else:
                lines.append(cur)
                cur = word
            # A single word longer than a line: hard-split it.
            while len(cur) > cells:
                lines.append(cur[:cells])
                cur = cur[cells:]
        lines.append(cur)
    # Paginate.
    pages: List[str] = []
    for i in range(0, len(lines), lines_per_page):
        pages.append("\r\n".join(lines[i : i + lines_per_page]))
    return "\f".join(pages) + "\r\n"


def _export_braille(
    text: str, table: str = "en-ueb-g2.ctb", use_liblouis: bool = False
) -> str:
    """Convert text to BRF (Braille Ready Format).

    By default this uses the built-in pure-Python Grade 1 translator, which is
    always available and can never crash the host process.  This fixes the
    long-standing bug where a missing liblouis translation table caused
    liblouis to call ``exit()`` at the C level, abruptly closing the window.

    Set *use_liblouis* (``braille_grade2`` setting) to opt in to contracted
    Grade 2 translation via liblouis when it is installed and the requested
    table resolves; any failure falls back to the built-in translator.
    """
    if use_liblouis and _LOUIS:
        try:
            brl = _louis.translateString([table], text, None, 0)
            if brl:
                return _format_brf(brl)
        except Exception:
            pass  # fall through to the dependency-free translator
    return _format_brf(_text_to_braille_grade1(text))


def _format_annotations(
    items: List[Dict[str, Any]],
    ext: str,
    title: str,
    author: str,
    source: str,
) -> str:
    """Serialize document annotations to the format implied by *ext*.

    Supported extensions: ``.json``, ``.bib`` (BibTeX), ``.ris`` (RIS),
    ``.txt``, and ``.md`` (Markdown — the default for anything else).  BibTeX
    and RIS emit a single reference for the source document with the notes
    attached, which is the standard reference-manager convention.
    """
    date = time.strftime("%Y-%m-%d")
    year = time.strftime("%Y")

    if ext == ".json":
        return json.dumps(
            {
                "title": title,
                "author": author,
                "source": source,
                "exported": date,
                "annotations": items,
            },
            indent=2,
            ensure_ascii=False,
        )

    if ext == ".bib":
        key = (
            re.sub(r"\W+", "", (Path(source).stem or title or "notes"))[:40] or "notes"
        )
        notes_blob = "; ".join(
            f"[{a.get('anchor', '')}] {a.get('note', '')}".strip() for a in items
        )
        lines = [f"@misc{{{key},", f"  title = {{{title}}},"]
        if author:
            lines.append(f"  author = {{{author}}},")
        if source:
            lines.append(f"  howpublished = {{{source}}},")
        lines.append(f"  year = {{{year}}},")
        lines.append(f"  annote = {{{notes_blob}}}")
        lines.append("}")
        return "\n".join(lines) + "\n"

    if ext == ".ris":
        lines = ["TY  - GEN", f"TI  - {title}"]
        if author:
            lines.append(f"AU  - {author}")
        for a in items:
            anchor = str(a.get("anchor", "")).strip()
            note = str(a.get("note", "")).strip()
            combined = f"“{anchor}” — {note}" if anchor else note
            lines.append(f"N1  - {combined}")
        if source:
            lines.append(f"UR  - {source}")
        lines.append(f"PY  - {year}")
        lines.append("ER  - ")
        return "\r\n".join(lines) + "\r\n"

    if ext == ".txt":
        out = [f"Notes — {title}"]
        if author:
            out.append(f"Author: {author}")
        if source:
            out.append(f"Source: {source}")
        out.append(f"Exported: {date}")
        out.append("")
        for i, a in enumerate(items, 1):
            anchor = str(a.get("anchor", "")).strip()
            note = str(a.get("note", "")).strip()
            ts = str(a.get("ts", ""))
            out.append(f"{i}. {note}")
            if anchor:
                out.append(f"   context: “{anchor}”")
            if ts:
                out.append(f"   {ts}")
            out.append("")
        return "\n".join(out)

    # Default: Markdown
    md = [f"# Notes — {title}", ""]
    if author:
        md.append(f"- **Author:** {author}")
    if source:
        md.append(f"- **Source:** `{source}`")
    md.append(f"- **Exported:** {date}")
    md.append(f"- **Count:** {len(items)}")
    md.append("")
    md.append("---")
    md.append("")
    for i, a in enumerate(items, 1):
        anchor = str(a.get("anchor", "")).strip()
        note = str(a.get("note", "")).strip()
        ts = str(a.get("ts", ""))
        md.append(f"## Note {i}")
        if anchor:
            md.append("")
            md.append(f"> {anchor}")
        md.append("")
        md.append(note)
        if ts:
            md.append("")
            md.append(f"*{ts}*")
        md.append("")
    return "\n".join(md)


def _annotation_matches(a: Dict[str, Any], query: str) -> bool:
    """Return True if annotation *a* matches the free-text *query*.

    The query supports space-separated terms (all must match — AND) over the
    note body, anchor quote, and tags.  A term beginning with ``#`` is a tag
    filter (matches against the note's tags only).  An empty query matches all.
    Shared by the Qt Notes panel and the curses TUI notes list.
    """
    q = (query or "").strip().lower()
    if not q:
        return True
    note = str(a.get("note", "")).lower()
    anchor = str(a.get("anchor", "")).lower()
    tags = [str(t).lower() for t in a.get("tags", []) or []]
    for term in q.split():
        if term.startswith("#") and len(term) > 1:
            if not any(term[1:] in t for t in tags):
                return False
        else:
            if (
                term not in note
                and term not in anchor
                and not any(term in t for t in tags)
            ):
                return False
    return True


def _parse_tags(raw: str) -> List[str]:
    """Split a comma/space/`#`-separated tag string into a clean tag list."""
    parts = re.split(r"[,\s]+", (raw or "").strip())
    return [p.lstrip("#").strip() for p in parts if p.strip().lstrip("#")]


# =============================================================================
# Citation manager
# =============================================================================


def _citation_label(c: Dict[str, Any]) -> str:
    """Human-readable one-line label for a citation dict."""
    author = str(c.get("author", "")).split(" and ")[0].split(",")[0].strip()
    year = str(c.get("year", "")).strip()
    title = str(c.get("title", "")).strip()
    cid = str(c.get("id", "")).strip()
    head = " ".join(p for p in (author, f"({year})" if year else "") if p).strip()
    label = f"{head}  {title}" if head else title
    return f"[{cid}] {label}".strip() if cid else label or "(untitled)"


def _parse_bibtex(text: str) -> List[Dict[str, Any]]:
    """Parse a (subset of) BibTeX into a list of citation dicts.

    Uses brace-counting to find each entry's body so it works whether the
    closing brace is on its own line or inline, and tolerates nested braces.
    """
    items: List[Dict[str, Any]] = []
    n = len(text)
    mapping = {
        "title": "title",
        "author": "author",
        "year": "year",
        "journal": "journal",
        "booktitle": "journal",
        "doi": "doi",
        "url": "url",
        "publisher": "publisher",
    }
    for m in re.finditer(r"@(\w+)\s*\{", text):
        ctype = m.group(1).lower()
        depth, j = 1, m.end()
        while j < n and depth > 0:
            if text[j] == "{":
                depth += 1
            elif text[j] == "}":
                depth -= 1
            j += 1
        body = text[m.end() : j - 1]
        key, _sep, rest = body.partition(",")
        fields: Dict[str, Any] = {"id": key.strip(), "type": ctype}
        for fm in re.finditer(
            r"(\w+)\s*=\s*(\{(?:[^{}]|\{[^{}]*\})*\}|\"[^\"]*\"|[^,\n]+)",
            rest,
        ):
            name = fm.group(1).lower()
            val = fm.group(2).strip().strip("{}").strip('"').strip()
            if name in mapping:
                fields[mapping[name]] = val
        items.append(fields)
    return items


def _parse_ris(text: str) -> List[Dict[str, Any]]:
    """Parse RIS records into a list of citation dicts."""
    items: List[Dict[str, Any]] = []
    cur: Dict[str, Any] = {}
    authors: List[str] = []
    tag_map = {
        "TI": "title",
        "T1": "title",
        "PY": "year",
        "Y1": "year",
        "JO": "journal",
        "JF": "journal",
        "T2": "journal",
        "DO": "doi",
        "UR": "url",
        "PB": "publisher",
        "ID": "id",
    }
    for line in text.splitlines():
        mm = re.match(r"^([A-Z][A-Z0-9])  - (.*)$", line.rstrip())
        if not mm:
            continue
        tag, val = mm.group(1), mm.group(2).strip()
        if tag == "TY":
            cur = {"type": val.lower()}
            authors = []
        elif tag in ("AU", "A1"):
            authors.append(val)
        elif tag == "ER":
            if authors:
                cur["author"] = " and ".join(authors)
            if "year" in cur:
                cur["year"] = str(cur["year"])[:4]
            if not cur.get("id"):
                a0 = (authors[0].split(",")[0] if authors else "ref").strip()
                cur["id"] = re.sub(r"\W+", "", a0 + str(cur.get("year", ""))) or "ref"
            items.append(cur)
            cur = {}
        elif tag in tag_map:
            cur[tag_map[tag]] = val
    return items


def _parse_csl_json(text: str) -> List[Dict[str, Any]]:
    """Parse CSL-JSON into a list of citation dicts."""
    data = json.loads(text)
    if isinstance(data, dict):
        data = [data]
    items: List[Dict[str, Any]] = []
    for c in data:
        authors = []
        for a in c.get("author", []) or []:
            fam = a.get("family", "")
            given = a.get("given", "")
            authors.append(f"{fam}, {given}".strip(", ") if fam else given)
        issued = c.get("issued", {}) or {}
        year = ""
        if isinstance(issued, dict):
            parts = issued.get("date-parts") or [[""]]
            year = str(parts[0][0]) if parts and parts[0] else ""
        items.append(
            {
                "id": str(c.get("id", "")),
                "type": str(c.get("type", "article")),
                "title": str(c.get("title", "")),
                "author": " and ".join(authors),
                "year": year,
                "journal": str(c.get("container-title", "")),
                "doi": str(c.get("DOI", "")),
                "url": str(c.get("URL", "")),
                "publisher": str(c.get("publisher", "")),
            }
        )
    return items


def _import_citations(path: str) -> List[Dict[str, Any]]:
    """Read citations from a .bib / .ris / .json file (format by extension)."""
    text = Path(path).read_text(encoding="utf-8", errors="replace")
    ext = Path(path).suffix.lower()
    if ext == ".bib":
        return _parse_bibtex(text)
    if ext == ".ris":
        return _parse_ris(text)
    if ext in (".json", ".csl"):
        return _parse_csl_json(text)
    # Heuristic fallback by content.
    if text.lstrip().startswith("@"):
        return _parse_bibtex(text)
    if re.search(r"^TY  - ", text, re.MULTILINE):
        return _parse_ris(text)
    return _parse_csl_json(text)


def _format_citations(items: List[Dict[str, Any]], ext: str) -> str:
    """Serialize citations to BibTeX (.bib), RIS (.ris), or CSL-JSON (.json)."""
    if ext == ".ris":
        out: List[str] = []
        for c in items:
            out.append(f"TY  - {str(c.get('type', 'GEN')).upper()}")
            if c.get("title"):
                out.append(f"TI  - {c['title']}")
            for au in str(c.get("author", "")).split(" and "):
                if au.strip():
                    out.append(f"AU  - {au.strip()}")
            if c.get("year"):
                out.append(f"PY  - {c['year']}")
            if c.get("journal"):
                out.append(f"JO  - {c['journal']}")
            if c.get("doi"):
                out.append(f"DO  - {c['doi']}")
            if c.get("url"):
                out.append(f"UR  - {c['url']}")
            out.append("ER  - ")
            out.append("")
        return "\r\n".join(out)
    if ext in (".json", ".csl"):
        csl = []
        for c in items:
            authors = []
            for au in str(c.get("author", "")).split(" and "):
                au = au.strip()
                if not au:
                    continue
                if "," in au:
                    fam, given = au.split(",", 1)
                    authors.append({"family": fam.strip(), "given": given.strip()})
                else:
                    authors.append({"family": au})
            entry: Dict[str, Any] = {
                "id": c.get("id", ""),
                "type": c.get("type", "article"),
                "title": c.get("title", ""),
            }
            if authors:
                entry["author"] = authors
            if c.get("year"):
                entry["issued"] = {"date-parts": [[str(c["year"])]]}
            if c.get("journal"):
                entry["container-title"] = c["journal"]
            if c.get("doi"):
                entry["DOI"] = c["doi"]
            if c.get("url"):
                entry["URL"] = c["url"]
            csl.append(entry)
        return json.dumps(csl, indent=2, ensure_ascii=False)
    # Default: BibTeX
    out = []
    for c in items:
        cid = (
            c.get("id") or re.sub(r"\W+", "", str(c.get("title", "ref"))[:20]) or "ref"
        )
        out.append(f"@{c.get('type', 'article')}{{{cid},")
        for fld, key in (
            ("title", "title"),
            ("author", "author"),
            ("year", "year"),
            ("journal", "journal"),
            ("doi", "doi"),
            ("url", "url"),
            ("publisher", "publisher"),
        ):
            if c.get(fld):
                out.append(f"  {key} = {{{c[fld]}}},")
        out.append("}")
        out.append("")
    return "\n".join(out)


def _fetch_citation_by_doi(doi: str) -> Dict[str, Any]:
    """Look up a DOI via the Crossref REST API and return a citation dict.

    Network call (blocking) — call from a background thread in the GUI.
    Raises on network/parse errors so callers can report them.
    """
    doi = doi.strip()
    for prefix in ("https://doi.org/", "http://doi.org/", "doi:"):
        if doi.lower().startswith(prefix):
            doi = doi[len(prefix) :]
    url = "https://api.crossref.org/works/" + urllib.parse.quote(doi)
    req = urllib.request.Request(
        url, headers={"User-Agent": f"star/{APP_VERSION} (citation lookup)"}
    )
    with urllib.request.urlopen(req, timeout=15) as resp:
        data = json.loads(resp.read().decode("utf-8", errors="replace"))
    msg = data.get("message", {}) or {}
    authors = []
    for a in msg.get("author", []) or []:
        fam = a.get("family", "")
        given = a.get("given", "")
        authors.append(f"{fam}, {given}".strip(", ") if fam else given)
    year = ""
    parts = (msg.get("issued", {}) or {}).get("date-parts", [[None]])
    if parts and parts[0] and parts[0][0]:
        year = str(parts[0][0])
    title = (msg.get("title") or [""])[0]
    journal = (msg.get("container-title") or [""])[0]
    fam0 = authors[0].split(",")[0] if authors else "ref"
    cid = re.sub(r"\W+", "", fam0 + year) or "ref"
    return {
        "id": cid,
        "type": str(msg.get("type", "article")),
        "title": title,
        "author": " and ".join(authors),
        "year": year,
        "journal": journal,
        "doi": doi,
        "url": str(msg.get("URL", "")),
    }


# =============================================================================
# Speech recognition  (Whisper dictation / transcription)
# =============================================================================


def _fmt_timestamp(seconds: float) -> str:
    """Format *seconds* as ``[hh:mm:ss]`` (or ``[mm:ss]`` under an hour)."""
    s = int(max(0, seconds))
    h, rem = divmod(s, 3600)
    m, sec = divmod(rem, 60)
    return f"[{h:02d}:{m:02d}:{sec:02d}]" if h else f"[{m:02d}:{sec:02d}]"


def _transcribe_audio(
    path: str, model_name: str = "base", timestamps: bool = False
) -> str:
    """Transcribe an audio file to text using Whisper (blocking).

    Works with either ``openai-whisper`` or ``faster-whisper``.  When
    *timestamps* is True each segment is prefixed with its start time as
    ``[hh:mm:ss]`` on its own line, producing a navigable transcript.
    Raises RuntimeError with install guidance when no backend is available.
    """
    if _WHISPER == "openai":
        model = _whisper.load_model(model_name)
        result = model.transcribe(path)
        if timestamps:
            segs = result.get("segments", []) or []
            return (
                "\n".join(
                    f"{_fmt_timestamp(s.get('start', 0))} {str(s.get('text', '')).strip()}"
                    for s in segs
                ).strip()
                or str(result.get("text", "")).strip()
            )
        return str(result.get("text", "")).strip()
    if _WHISPER == "faster":
        model = _FasterWhisper(model_name)
        segments, _info = model.transcribe(path)
        if timestamps:
            return "\n".join(
                f"{_fmt_timestamp(getattr(seg, 'start', 0))} {seg.text.strip()}"
                for seg in segments
            ).strip()
        return " ".join(seg.text for seg in segments).strip()
    raise RuntimeError(
        "Speech recognition requires Whisper:\n"
        "  pip install openai-whisper   (or: pip install faster-whisper)"
    )


def _record_audio_to_wav(seconds: float, samplerate: int = 16000) -> str:
    """Record *seconds* of mono microphone audio to a temp WAV; return its path."""
    if not _AUDIO_IN:
        raise RuntimeError(
            "Microphone capture requires sounddevice + numpy:\n"
            "  pip install sounddevice numpy"
        )
    import wave

    frames = int(seconds * samplerate)
    rec = _sounddevice.rec(frames, samplerate=samplerate, channels=1, dtype="int16")
    _sounddevice.wait()
    tmp = tempfile.NamedTemporaryFile(suffix=".wav", delete=False)
    tmp.close()
    with wave.open(tmp.name, "wb") as w:
        w.setnchannels(1)
        w.setsampwidth(2)
        w.setframerate(samplerate)
        w.writeframes(rec.tobytes())
    return tmp.name


# =============================================================================
# Canonical keyboard shortcuts  (GUI/TUI parity + cheat sheet)
# =============================================================================

# The GUI bindings are the canonical set.  The TUI mirrors them where the
# terminal can express the chord; where it cannot (terminals can't see
# Ctrl+Shift+<letter>), an equivalent single-key or M-x command is listed.
_SHORTCUTS: List[Tuple[str, List[Tuple[str, str, str]]]] = [
    (
        "Playback",
        [
            ("Play / pause", "Space", "Space"),
            ("Stop", "Esc", "Esc / M-x stop"),
            ("Speed up / slow down", "Ctrl+= / Ctrl+-", "+ / -"),
            ("Play from cursor", "Ctrl+Return", "—"),
            ("Speak current line", "—", "M-x speak-line"),
        ],
    ),
    (
        "Navigation",
        [
            ("Next / prev heading", "Ctrl+H / Ctrl+Shift+H", "> / <"),
            ("Next / prev paragraph", "Ctrl+P / Ctrl+Shift+P", "] / ["),
            ("Next / prev sentence", "Alt+. / Alt+,", ". / ,"),
            ("Replay sentence", "Alt+;", ";"),
            ("Replay paragraph", "Ctrl+R", "r"),
            ("Next / prev table", "Ctrl+T / Ctrl+Shift+T", "M-x next-table"),
            ("Speech-cursor mode", "Tab", "Tab"),
        ],
    ),
    (
        "Notes / Annotations",
        [
            ("Add note at cursor", "Ctrl+Shift+A", "a / M-x annotate"),
            ("Toggle Notes panel", "Ctrl+Shift+N", "M-x annotations-list"),
            (
                "Search / filter notes",
                "(search box in panel)",
                "M-x annotations-search",
            ),
        ],
    ),
    (
        "View",
        [
            ("Cycle theme", "F5", "F5"),
            ("Toggle Contents panel", "Ctrl+\\", "—"),
            ("Reading level", "Ctrl+L", "M-x reading-level"),
            ("Choose voice", "Ctrl+Shift+V", "M-x voice-picker"),
            ("Line numbers / syntax", "—", "F6 / F7"),
        ],
    ),
    (
        "Edit & File",
        [
            ("Toggle edit mode", "Ctrl+E", "—"),
            ("Save", "Ctrl+S", "—"),
            ("Copy", "Ctrl+C", "M-x copy"),
            ("Open file", "Ctrl+O", "Ctrl+O"),
            ("Command palette", "—", "M-x / F2 / :"),
            ("Keyboard cheat sheet", "(Help menu)", "M-x shortcuts"),
            ("Help (README)", "F1", "F1"),
            ("Quit", "Ctrl+Q", "Ctrl+Q / q"),
        ],
    ),
]


def _shortcuts_text(plain: bool = False) -> str:
    """Render the canonical shortcut scheme as text.

    *plain* True produces a column layout for the curses pager; otherwise a
    Markdown table suitable for the Qt help dialog.
    """
    if plain:
        lines = ["Keyboard Shortcuts  (GUI bindings are canonical)", ""]
        for section, rows in _SHORTCUTS:
            lines.append(f"== {section} ==")
            for action, gui, tui in rows:
                lines.append(f"  {action:<26} GUI: {gui:<24} TUI: {tui}")
            lines.append("")
        return "\n".join(lines)
    md = [
        "# Keyboard Shortcuts",
        "",
        "GUI bindings are canonical; the TUI mirrors them where the terminal allows.",
        "",
    ]
    for section, rows in _SHORTCUTS:
        md.append(f"## {section}")
        md.append("")
        md.append("| Action | GUI | TUI |")
        md.append("|---|---|---|")
        for action, gui, tui in rows:
            md.append(f"| {action} | `{gui}` | `{tui}` |")
        md.append("")
    return "\n".join(md)


# =============================================================================
# SSML / prosody helpers
# =============================================================================


def _text_to_ssml(
    text: str,
    backend: str = "pyttsx3",
    sentence_ms: int = 350,
    clause_ms: int = 150,
) -> str:
    """Wrap *text* in SSML markup, inserting pauses at sentence and clause
    boundaries for more natural-sounding speech.

    The output is always wrapped in ``<speak>…</speak>`` tags:

    * eSpeak-NG accepts this directly with its ``-m`` flag.
    * SAPI5 (Windows / pyttsx3) interprets it natively.
    * DECtalk receives its own notation via ``_text_to_dectalk()`` instead.

    If *text* is already SSML (starts with ``<speak>``) it is returned
    unchanged to avoid double-wrapping.
    """
    if text.lstrip().startswith("<speak>"):
        return text  # already wrapped
    if backend == "dectalk":
        return _text_to_dectalk(text, sentence_ms, clause_ms)

    # Escape XML-reserved characters before inserting any tags.
    s = (
        text.replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace('"', "&quot;")
    )

    long_ms = sentence_ms * 2

    # Paragraph breaks → long pause
    s = re.sub(r"\n{2,}", f'<break time="{long_ms}ms"/>\n', s)
    # Sentence-ending punctuation → medium pause
    s = re.sub(
        r"([.!?\u2026])(\s+)",
        lambda m: f'{m.group(1)}<break time="{sentence_ms}ms"/>{m.group(2)}',
        s,
    )
    # Clause punctuation (comma, semicolon, colon) → short pause
    s = re.sub(
        r"([,;:])(\s+)",
        lambda m: f'{m.group(1)}<break time="{clause_ms}ms"/>{m.group(2)}',
        s,
    )
    # Em-dash / en-dash → brief pause
    s = re.sub(
        r"[\u2014\u2013]",
        f' <break time="{clause_ms}ms"/> ',
        s,
    )

    return f"<speak>{s}</speak>"


def _text_to_dectalk(
    text: str,
    sentence_ms: int = 350,
    clause_ms: int = 150,
) -> str:
    """Convert plain text to DECtalk phoneme notation for improved prosody.

    DECtalk uses ``[:pau N]`` to insert N milliseconds of silence.
    """
    s = text
    s = re.sub(
        r"([.!?\u2026])(\s+)",
        lambda m: f"{m.group(1)} [:pau {sentence_ms}] ",
        s,
    )
    s = re.sub(
        r"([,;:])(\s+)",
        lambda m: f"{m.group(1)} [:pau {clause_ms}] ",
        s,
    )
    s = re.sub(r"[\u2014\u2013]", f" [:pau {clause_ms}] ", s)
    return s


# =============================================================================
# TTS text preprocessing
# =============================================================================
# Pipeline (applied at speak-time after the plain-text slice is taken):
#   _expand_abbreviations()  ->  _normalize_numbers()  ->  _text_to_ssml()
#
# Table narration mode is baked in at document-load time inside
# _strip_markdown_for_tts() so the word map reflects the spoken structure.
# =============================================================================

# -- Abbreviation table -------------------------------------------------------
# Ordered longest-first so multi-word entries match before their components.

_ABBREV_PAIRS: List[tuple] = [
    # Multi-word Latin phrases
    ("et al.", "and others"),
    ("op. cit.", "op cit"),
    # Single-token Latin / English abbreviations
    ("e.g.,", "for example,"),  # trailing-comma variant must come first
    ("e.g.", "for example"),
    ("i.e.,", "that is,"),
    ("i.e.", "that is"),
    ("etc.", "et cetera"),
    ("cf.", "compare"),
    ("ibid.", "ibid"),
    ("n.d.", "no date"),
    ("ca.", "circa"),
    ("vs.", "versus"),
    ("approx.", "approximately"),
    # Academic / publishing
    ("Fig.", "Figure"),
    ("Figs.", "Figures"),
    ("Eq.", "Equation"),
    ("Eqs.", "Equations"),
    ("Sec.", "Section"),
    ("Chap.", "Chapter"),
    ("Ref.", "Reference"),
    ("Refs.", "References"),
    ("Vol.", "Volume"),
    ("vol.", "volume"),
    ("No.", "Number"),
    ("no.", "number"),
    ("pp.", "pages"),
    ("p.", "page"),
    ("ed.", "edition"),
    ("eds.", "editors"),
    ("Dept.", "Department"),
    ("dept.", "department"),
    ("Assoc.", "Association"),
    ("Univ.", "University"),
    ("univ.", "university"),
    # Titles / honorifics
    ("Dr.", "Doctor"),
    ("Mr.", "Mister"),
    ("Mrs.", "Missus"),
    ("Prof.", "Professor"),
    ("Jr.", "Junior"),
    ("Sr.", "Senior"),
    ("Rev.", "Reverend"),
    ("Gen.", "General"),
    ("Gov.", "Governor"),
    # Units / measurement
    ("hr.", "hour"),
    ("min.", "minutes"),
    ("sec.", "seconds"),
    ("wt.", "weight"),
    ("avg.", "average"),
    ("temp.", "temperature"),
    ("conc.", "concentration"),
    ("est.", "estimated"),
    ("max.", "maximum"),
    # Business / organizations
    ("Inc.", "Incorporated"),
    ("Corp.", "Corporation"),
    ("Ltd.", "Limited"),
]

# Compiled once: word-boundary anchor before each abbreviation token.
# We do NOT add \b after the token because abbreviations end with '.', which is
# not a word-boundary character.
_ABBREV_RE: List[tuple] = [
    (re.compile(r"\b" + re.escape(abbr)), expansion)
    for abbr, expansion in _ABBREV_PAIRS
]


def _expand_abbreviations(text: str, custom: Optional[Dict[str, str]] = None) -> str:
    """Expand common and user-defined abbreviations for natural TTS output.

    Custom expansions (from settings["abbrev_expansions"]) are applied first
    so they take precedence over the built-in list.
    """
    if custom:
        for abbr, exp in sorted(custom.items(), key=lambda x: -len(x[0])):
            text = re.sub(r"\b" + re.escape(abbr), exp, text)
    for pattern, expansion in _ABBREV_RE:
        text = pattern.sub(expansion, text)
    return text


# -- Number normalization helpers --------------------------------------------

_ONES = [
    "",
    "one",
    "two",
    "three",
    "four",
    "five",
    "six",
    "seven",
    "eight",
    "nine",
    "ten",
    "eleven",
    "twelve",
    "thirteen",
    "fourteen",
    "fifteen",
    "sixteen",
    "seventeen",
    "eighteen",
    "nineteen",
]
_TENS = [
    "",
    "",
    "twenty",
    "thirty",
    "forty",
    "fifty",
    "sixty",
    "seventy",
    "eighty",
    "ninety",
]

_MONTHS_NUM: Dict[str, str] = {
    "01": "January",
    "02": "February",
    "03": "March",
    "04": "April",
    "05": "May",
    "06": "June",
    "07": "July",
    "08": "August",
    "09": "September",
    "10": "October",
    "11": "November",
    "12": "December",
}


def _int_to_words(n: int) -> str:
    """Convert a non-negative integer to English words."""
    if n == 0:
        return "zero"
    if n < 0:
        return "negative " + _int_to_words(-n)
    if n < 20:
        return _ONES[n]
    if n < 100:
        t, o = divmod(n, 10)
        return _TENS[t] + ("-" + _ONES[o] if o else "")
    if n < 1000:
        h, r = divmod(n, 100)
        return _ONES[h] + " hundred" + (" " + _int_to_words(r) if r else "")
    for exp, label in [
        (12, "trillion"),
        (9, "billion"),
        (6, "million"),
        (3, "thousand"),
    ]:
        base = 10**exp
        if n >= base:
            q, r = divmod(n, base)
            return (
                _int_to_words(q) + " " + label + (" " + _int_to_words(r) if r else "")
            )
    return str(n)


def _year_to_words(year: int) -> str:
    """Read a 4-digit year naturally.

    1984 -> nineteen eighty-four   2024 -> twenty twenty-four
    1900 -> nineteen hundred       2000 -> two thousand
    """
    if not (100 <= year <= 2999):
        return _int_to_words(year)
    century, decade = divmod(year, 100)
    if decade == 0:
        return _int_to_words(century) + " hundred"
    if decade < 10:
        return _int_to_words(century) + " oh " + _int_to_words(decade)
    return _int_to_words(century) + " " + _int_to_words(decade)


def _ordinal_to_words(n: int) -> str:
    """Convert a positive integer to its ordinal word form.

    1 -> first   21 -> twenty-first   100 -> one hundredth
    Compound ordinals are built by recurring on the last component so that
    21 -> twenty-first (not twenty-onth), 23 -> twenty-third, etc.
    """
    specials = {
        1: "first",
        2: "second",
        3: "third",
        4: "fourth",
        5: "fifth",
        6: "sixth",
        7: "seventh",
        8: "eighth",
        9: "ninth",
        10: "tenth",
        11: "eleventh",
        12: "twelfth",
    }
    if n in specials:
        return specials[n]
    if n < 0:
        return "negative " + _ordinal_to_words(-n)
    # Teens 13-19: thirteenth, fourteenth, …, nineteenth
    if 13 <= n <= 19:
        return _int_to_words(n) + "th"
    # Compound: 21-99 with a non-zero ones digit -> apply ordinal only to ones
    if 20 <= n < 100:
        tens, ones = divmod(n, 10)
        if ones != 0:
            return _int_to_words(tens * 10) + "-" + _ordinal_to_words(ones)
        # Round tens (20, 30, …): twenty -> twentieth
        w = _int_to_words(n)
        return w[:-1] + "ieth"
    # Hundreds: apply ordinal to the remainder if non-zero
    if n < 1_000:
        hundreds, rest = divmod(n, 100)
        if rest != 0:
            return _int_to_words(hundreds * 100) + " " + _ordinal_to_words(rest)
        return _int_to_words(n) + "th"  # e.g. 'one hundredth'
    # Larger numbers: apply ordinal to the remainder, or just append 'th'
    for exp, label in [
        (12, "trillion"),
        (9, "billion"),
        (6, "million"),
        (3, "thousand"),
    ]:
        base = 10**exp
        if n >= base:
            q, r = divmod(n, base)
            if r != 0:
                return _int_to_words(q) + " " + label + " " + _ordinal_to_words(r)
            return _int_to_words(q) + " " + label + "th"
    return _int_to_words(n) + "th"


def _decimal_digits_to_words(digits: str) -> str:
    """Read decimal digits one at a time: "14" -> "one four"."""
    return " ".join(_ONES[int(d)] if d != "0" else "zero" for d in digits)


def _normalize_numbers(text: str) -> str:
    """Expand common numeric patterns to spoken English.

    Processed in specificity order:
      ISO dates     YYYY-MM-DD
      US dates      MM/DD/YYYY
      Times         HH:MM or HH:MM:SS (24-hour converted to 12-hour + AM/PM)
      Currency      dollar sign, pound sign, euro sign
      Percentages   N% or N.N%
      Ordinals      1st 2nd 3rd 4th ...
      Comma integers  1,234,567
      Decimals      N.N
      Plain integers >= 1000 (4-digit treated as years 1000-2099)
    """

    # -- ISO date  YYYY-MM-DD -------------------------------------------------
    def _iso_date(m: re.Match) -> str:
        y, mo, d = m.group(1), m.group(2), m.group(3)
        return f"{_MONTHS_NUM.get(mo, mo)} {_ordinal_to_words(int(d))}, {_year_to_words(int(y))}"

    text = re.sub(r"\b(\d{4})-(0[1-9]|1[0-2])-(0[1-9]|[12]\d|3[01])\b", _iso_date, text)

    # -- US date  MM/DD/YYYY --------------------------------------------------
    def _us_date(m: re.Match) -> str:
        mo, d, y = m.group(1), m.group(2), m.group(3)
        return f"{_MONTHS_NUM.get(mo.zfill(2), mo)} {_ordinal_to_words(int(d))}, {_year_to_words(int(y))}"

    text = re.sub(
        r"\b(0?[1-9]|1[0-2])/(0?[1-9]|[12]\d|3[01])/(\d{4})\b", _us_date, text
    )

    # -- Times  HH:MM  (optional explicit AM/PM consumed to avoid duplication) ---
    def _time_val(m: re.Match) -> str:
        h, mi = int(m.group(1)), int(m.group(2))
        explicit = (m.group(3) or "").strip().upper()  # "AM", "PM", or ""
        if h == 0 and mi == 0 and not explicit:
            return "midnight"
        if h == 12 and mi == 0 and not explicit:
            return "noon"
        if explicit == "PM" and h < 12:
            h12 = h  # 3:45 PM -> keep h=3
            period = "PM"
        elif explicit == "AM":
            h12 = h % 12 or 12
            period = "AM"
        else:
            period = "AM" if h < 12 else "PM"
            h12 = h if h <= 12 else h - 12
        if mi == 0:
            return f"{_int_to_words(h12)} {period}"
        return f"{_int_to_words(h12)} {_int_to_words(mi)} {period}"

    text = re.sub(
        r"\b([01]?\d|2[0-3]):([0-5]\d)(?::\d\d)?\s*([AaPp][Mm])?", _time_val, text
    )

    # -- Currency  $  £  € ----------------------------------------------------
    _CURR: Dict[str, tuple] = {
        "$": ("dollar", "cent"),
        "\u00a3": ("pound", "penny"),
        "\u20ac": ("euro", "cent"),
    }

    def _currency_val(m: re.Match) -> str:
        sym = m.group(1)
        whole = m.group(2).replace(",", "")
        frac = (m.group(3) or "00")[:2].ljust(2, "0")
        major, minor = _CURR.get(sym, ("dollar", "cent"))
        parts: List[str] = []
        wi = int(whole)
        if wi:
            parts.append(f"{_int_to_words(wi)} {major}{'s' if wi != 1 else ''}")
        fi = int(frac)
        if fi:
            parts.append(f"{_int_to_words(fi)} {minor}{'s' if fi != 1 else ''}")
        return " and ".join(parts) if parts else f"zero {major}s"

    text = re.sub(r"([$\u00a3\u20ac])(\d[\d,]*)\.?(\d{2})?", _currency_val, text)

    # -- Percentages  75%  3.5% -----------------------------------------------
    def _pct(m: re.Match) -> str:
        num = m.group(1)
        if "." in num:
            whole, frac = num.split(".", 1)
            return f"{_int_to_words(int(whole)) if whole else 'zero'} point {_decimal_digits_to_words(frac)} percent"
        return f"{_int_to_words(int(num))} percent"

    text = re.sub(r"\b(\d+(?:\.\d+)?)%", _pct, text)

    # -- Ordinals  1st  2nd  3rd  4th ... ------------------------------------
    text = re.sub(
        r"\b(\d+)(?:st|nd|rd|th)\b",
        lambda m: _ordinal_to_words(int(m.group(1))),
        text,
        flags=re.IGNORECASE,
    )

    # -- Large integers with comma separators  1,234,567 ----------------------
    text = re.sub(
        r"\b\d{1,3}(?:,\d{3})+\b",
        lambda m: _int_to_words(int(m.group().replace(",", ""))),
        text,
    )

    # -- Decimal numbers  3.14  0.5 -------------------------------------------
    def _decimal(m: re.Match) -> str:
        return f"{_int_to_words(int(m.group(1)))} point {_decimal_digits_to_words(m.group(2))}"

    text = re.sub(r"(?<![\d/\-])\b(\d+)\.(\d+)\b(?![-/\d])", _decimal, text)

    # -- Plain integers >= 1000 (smaller ones TTS handles reliably) -----------
    def _plain_int(m: re.Match) -> str:
        n = int(m.group())
        if 1000 <= n <= 2099:
            return _year_to_words(n)
        return _int_to_words(n)

    text = re.sub(r"(?<!\.)\b(\d{4,})\b(?!\.)", _plain_int, text)

    return text


# -- Table narration ---------------------------------------------------------


def _tables_to_narration(text: str, mode: str = "structured") -> str:
    """Convert markdown table syntax in *text* to TTS-friendly prose.

    mode="structured"  (default)
        Table with 3 columns: Name, Age, City.
        Row 1: Name is Alice, Age is 30, City is New York.
        Row 2: Name is Bob, Age is 25, City is Boston.

    mode="flat"
        Cells joined with period-space (consistent legacy behavior).

    mode="skip"
        Replace entire table with a one-line announcement.

    Must be called on raw markdown BEFORE other stripping so that pipe
    characters and separator rows are still present.
    """
    if mode not in ("structured", "flat", "skip"):
        mode = "structured"

    lines = text.split("\n")
    result: List[str] = []
    i = 0

    while i < len(lines):
        line = lines[i]
        if "|" not in line or not line.strip().startswith("|"):
            result.append(line)
            i += 1
            continue

        # Gather the full table block.
        block: List[str] = []
        while i < len(lines) and "|" in lines[i] and lines[i].strip().startswith("|"):
            block.append(lines[i])
            i += 1

        if mode == "skip":
            raw = [
                c.strip() for c in block[0].strip().strip("|").split("|") if c.strip()
            ]
            n = len(raw)
            result.append(
                f"Table with {n} column{'s' if n != 1 else ''} \u2014 skipped."
            )
            result.append("")
            continue

        # Parse cells.
        parsed: List[List[str]] = [
            [c.strip() for c in bl.strip().strip("|").split("|")] for bl in block
        ]

        # Separate header from data rows (skip separator lines).
        header: List[str] = []
        data_rows: List[List[str]] = []
        for cells in parsed:
            if bool(cells) and all(re.match(r"^[-:]+$", c) for c in cells if c):
                continue  # separator row
            clean = [c for c in cells if c]
            if not header:
                header = clean
            else:
                data_rows.append(clean)

        if mode == "flat":
            for cells in [header] + data_rows:
                clean = [c for c in cells if c]
                if clean:
                    result.append(".  ".join(clean) + ".")
            result.append("")
            continue

        # structured mode
        ncols = len(header)
        result.append(
            f"Table with {ncols} column{'s' if ncols != 1 else ''}: {', '.join(header)}."
        )
        for ri, data in enumerate(data_rows, 1):
            if header:
                parts = [
                    f"{hdr} is {data[hi]}"
                    for hi, hdr in enumerate(header)
                    if hi < len(data) and data[hi]
                ]
                if parts:
                    result.append(f"Row {ri}: {', '.join(parts)}.")
            else:
                clean = [c for c in data if c]
                if clean:
                    result.append(f"Row {ri}: {', '.join(clean)}.")

        result.append("")

    return "\n".join(result)


def _preprocess_tts_text(text: str, settings: "Settings") -> str:
    """Apply all speak-time text normalizations before SSML wrapping.

    Called in _tts_play_from_word() and _tts_speak_current_line() on the
    plain-text slice.  Each step is gated by its own settings flag so users
    can disable individual normalizations independently.

    Steps:
      1. Abbreviation expansion   (expand_abbreviations)
      2. Number normalization     (normalize_numbers)
    """
    if settings.get("expand_abbreviations", True):
        custom = settings.get("abbrev_expansions") or {}
        text = _expand_abbreviations(text, custom if custom else None)

    if settings.get("normalize_numbers", True):
        text = _normalize_numbers(text)

    if settings.get("normalize_math", True):
        text = _normalize_math_inline(text)

    return text


# =============================================================================
# New document loaders and utilities
# =============================================================================

# ---------------------------------------------------------------------------
# Feature 19 — PowerPoint (.pptx) loader
# ---------------------------------------------------------------------------


def _load_pptx(path: str) -> str:
    """Load a PowerPoint .pptx file as Markdown.

    Each slide becomes a section with its title as a heading.
    Body text, bullet points, and speaker notes are included.
    Requires: pip install python-pptx
    """
    if not _PPTX:
        return (
            "Could not load PowerPoint file: python-pptx is not installed.\n"
            "Install it with: pip install python-pptx"
        )

    prs = _pptx_lib.Presentation(path)
    sections = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        title_text = ""
        body_parts = []

        for shape in slide.shapes:
            # Images (MSO_SHAPE_TYPE.PICTURE == 13)
            if shape.shape_type == 13:
                alt = getattr(shape, "name", f"slide {slide_num} image")
                body_parts.append(f"[Image: {alt}]")
                continue

            # Tables
            if hasattr(shape, "table"):
                tbl = shape.table
                md_rows = []
                for row_idx, row in enumerate(tbl.rows):
                    cells = [c.text.strip().replace("\n", " ") for c in row.cells]
                    md_rows.append("| " + " | ".join(cells) + " |")
                    if row_idx == 0:
                        md_rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
                body_parts.append("\n".join(md_rows))
                continue

            if not shape.has_text_frame:
                continue

            is_title = (
                hasattr(shape, "placeholder_format")
                and shape.placeholder_format is not None
                and shape.placeholder_format.idx == 0
            )

            if is_title:
                title_text = shape.text_frame.text.strip()
            else:
                for para in shape.text_frame.paragraphs:
                    txt = para.text.strip()
                    if not txt:
                        continue
                    indent = "  " * para.level if para.level else ""
                    body_parts.append(f"{indent}- {txt}")

        heading = (
            f"## Slide {slide_num}: {title_text}"
            if title_text
            else f"## Slide {slide_num}"
        )
        slide_md = [heading]
        if body_parts:
            slide_md.append("")
            slide_md.extend(body_parts)

        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_md.append("")
                slide_md.append(f"> Note: {notes_text}")

        sections.append("\n".join(slide_md))

    return "\n\n".join(sections)


# ---------------------------------------------------------------------------
# Feature 20 — Math expression normalization
# ---------------------------------------------------------------------------


def _normalize_math_inline(text: str) -> str:
    """Convert LaTeX math notation to spoken English for TTS.

    Applied to plain text before it reaches the speech engine.
    Handles the most common patterns in academic/STEM writing.
    """
    # Strip display/inline delimiters (order matters: $$ before $)
    text = re.sub(r"\$\$(.+?)\$\$", r" \1 ", text, flags=re.DOTALL)
    text = re.sub(r"\$(.+?)\$", r" \1 ", text)
    text = re.sub(r"\\\[(.+?)\\\]", r" \1 ", text, flags=re.DOTALL)
    text = re.sub(r"\\\((.+?)\\\)", r" \1 ", text)

    # Statistical accents
    text = re.sub(r"\\bar\{(\w+)\}", r"\1-bar", text)
    text = re.sub(r"\\hat\{(\w+)\}", r"\1-hat", text)
    text = re.sub(r"\\tilde\{(\w+)\}", r"\1-tilde", text)
    text = re.sub(r"\\overline\{([^}]+)\}", r"\1 bar", text)

    # Fractions and roots
    text = re.sub(r"\\frac\{([^}]+)\}\{([^}]+)\}", r"\1 over \2", text)
    text = re.sub(r"\\sqrt\[([^\]]+)\]\{([^}]+)\}", r"\1-th root of \2", text)
    text = re.sub(r"\\sqrt\{([^}]+)\}", r"square root of \1", text)

    # Superscripts / powers
    text = re.sub(r"\^\{-1\}", " inverse", text)
    text = re.sub(r"\^\{2\}|\^2(?=\W)", " squared", text)
    text = re.sub(r"\^\{3\}|\^3(?=\W)", " cubed", text)
    text = re.sub(r"\^\{([^}]+)\}", r" to the \1", text)
    text = re.sub(r"\^(\w)", r" to the \1", text)

    # Subscripts
    text = re.sub(r"_\{(\w)(\w+)\}", r" sub \1 \2", text)
    text = re.sub(r"_\{(\w+)\}", r" sub \1", text)
    text = re.sub(r"_(\w)", r" sub \1", text)

    # Greek letters
    greek = {
        "alpha": "alpha",
        "beta": "beta",
        "gamma": "gamma",
        "delta": "delta",
        "epsilon": "epsilon",
        "zeta": "zeta",
        "eta": "eta",
        "theta": "theta",
        "lambda": "lambda",
        "mu": "mu",
        "nu": "nu",
        "xi": "xi",
        "pi": "pi",
        "rho": "rho",
        "sigma": "sigma",
        "tau": "tau",
        "phi": "phi",
        "chi": "chi",
        "psi": "psi",
        "omega": "omega",
        "Gamma": "Gamma",
        "Delta": "Delta",
        "Theta": "Theta",
        "Lambda": "Lambda",
        "Pi": "Pi",
        "Sigma": "Sigma",
        "Phi": "Phi",
        "Psi": "Psi",
        "Omega": "Omega",
    }
    for cmd, spoken in greek.items():
        text = text.replace(f"\\{cmd}", f" {spoken} ")

    # Trig and common functions (order: longer names first to avoid partial matches)
    trig = [
        ("arcsin", "arcsine"),
        ("arccos", "arccosine"),
        ("arctan", "arctangent"),
        ("sin", "sine"),
        ("cos", "cosine"),
        ("tan", "tangent"),
        ("cot", "cotangent"),
        ("sec", "secant"),
        ("csc", "cosecant"),
        ("ln", "natural log"),
        ("log", "log"),
        ("exp", "e to the power"),
        ("lim", "limit"),
        ("sum", "sum"),
    ]
    for cmd, spoken in trig:
        text = re.sub(rf"\\{cmd}\b", f" {spoken} ", text)

    # Operators and symbols
    ops = [
        (r"\\times\b", " times "),
        (r"×", " times "),
        (r"\\div\b", " divided by "),
        (r"÷", " divided by "),
        (r"\\pm\b", " plus or minus "),
        (r"\\neq\b", " not equal to "),
        (r"\\leq\b|\\le\b", " less than or equal to "),
        (r"≤", " less than or equal to "),
        (r"\\geq\b|\\ge\b", " greater than or equal to "),
        (r"≥", " greater than or equal to "),
        (r"\\approx\b", " approximately equal to "),
        (r"≈", " approximately equal to "),
        (r"\\infty\b", " infinity "),
        (r"∞", " infinity "),
        (r"\\rightarrow\b|\\to\b", " approaches "),
        (r"→", " approaches "),
        (r"\\leftarrow\b", " from "),
        (r"←", " from "),
        (r"\\nabla\b", " gradient of "),
        (r"\\partial\b", " partial "),
        (r"\\prod\b", " product "),
        (r"\\int\b", " integral "),
    ]
    for pattern, replacement in ops:
        text = re.sub(pattern, replacement, text)

    # Clean up residual LaTeX commands and braces
    text = re.sub(r"\\[a-zA-Z]+", " ", text)
    text = re.sub(r"[{}]", "", text)
    text = re.sub(r"  +", " ", text)
    return text.strip()


# ---------------------------------------------------------------------------
# Feature 21 — Improved ODT loader
# ---------------------------------------------------------------------------


def _load_odt_v2(path: str) -> str:
    """Load an ODT (OpenDocument Text) file to Markdown.

    Uses odfpy if available for full fidelity; otherwise falls back to
    Pandoc (if installed) or raw XML extraction.
    Supports headings, paragraphs, lists, tables, footnotes.
    """
    if _ODT:
        return _load_odt_via_odfpy(path)

    # Pandoc fallback
    try:
        import subprocess

        result = subprocess.run(
            ["pandoc", "-f", "odt", "-t", "markdown", path],
            capture_output=True,
            text=True,
            timeout=30,
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout
    except (FileNotFoundError, OSError):
        pass

    return _load_odt_raw_xml(path)


def _load_odt_via_odfpy(path: str) -> str:
    """Parse an ODT document via the odfpy library, converting to Markdown."""
    from odf.element import Element  # type: ignore  # noqa: F401
    from odf.opendocument import load as odf_load  # type: ignore

    doc = odf_load(path)
    lines: list = []
    footnotes: dict = {}
    fn_counter = [0]

    def get_text(elem) -> str:
        if not hasattr(elem, "childNodes"):
            return getattr(elem, "data", "")
        tag = elem.qname[1] if hasattr(elem, "qname") else ""
        if tag == "s":
            return " " * int(elem.getAttribute("text:c") or 1)
        if tag == "tab":
            return "\t"
        if tag == "line-break":
            return "\n"
        return "".join(get_text(c) for c in elem.childNodes)

    def walk(node, list_depth: int = 0):
        if not hasattr(node, "qname"):
            return
        tag = node.qname[1]

        if tag == "h":
            level = int(node.getAttribute("text:outline-level") or 1)
            lines.append("#" * level + " " + get_text(node).strip())
            lines.append("")

        elif tag == "p":
            txt = get_text(node).strip()
            if txt:
                if list_depth:
                    lines.append("  " * (list_depth - 1) + "- " + txt)
                else:
                    lines.append(txt)
                    lines.append("")

        elif tag == "list":
            for child in node.childNodes:
                if hasattr(child, "qname") and child.qname[1] == "list-item":
                    for sub in child.childNodes:
                        walk(sub, list_depth + 1)

        elif tag == "table":
            _odt_table_to_md(node, lines, get_text)

        elif tag == "note":
            fn_counter[0] += 1
            label = str(fn_counter[0])
            body_text = ""
            for child in node.childNodes:
                if hasattr(child, "qname") and child.qname[1] == "note-body":
                    body_text = get_text(child).strip()
            footnotes[label] = body_text
            if lines:
                lines[-1] = lines[-1] + f"[^{label}]"

        else:
            for child in node.childNodes:
                walk(child, list_depth)

    for child in doc.text.childNodes:
        walk(child)

    if footnotes:
        lines += ["", "## Footnotes", ""]
        for label, note_text in footnotes.items():
            lines.append(f"[^{label}]: {note_text}")

    return "\n".join(lines)


def _odt_table_to_md(table_node, lines: list, get_text_fn) -> None:
    """Convert an ODF table element to Markdown table rows, appended to lines."""
    rows = []
    for child in table_node.childNodes:
        if not hasattr(child, "qname"):
            continue
        if child.qname[1] == "table-row":
            cells = []
            for cell in child.childNodes:
                if hasattr(cell, "qname") and cell.qname[1] in (
                    "table-cell",
                    "covered-table-cell",
                ):
                    cells.append(get_text_fn(cell).strip().replace("\n", " "))
            rows.append(cells)

    if not rows:
        return
    col_count = max(len(r) for r in rows)
    lines.append("")
    for i, row in enumerate(rows):
        padded = row + [""] * (col_count - len(row))
        lines.append("| " + " | ".join(padded) + " |")
        if i == 0:
            lines.append("| " + " | ".join(["---"] * col_count) + " |")
    lines.append("")


def _load_odt_raw_xml(path: str) -> str:
    """Fallback ODT reader: extract text via raw ZIP + regex XML stripping."""
    import zipfile

    try:
        with zipfile.ZipFile(path, "r") as zf:
            content = zf.read("content.xml").decode("utf-8", errors="replace")
    except (KeyError, zipfile.BadZipFile) as exc:
        return f"[Could not read ODT file: {exc}]"

    text = re.sub(
        r"<text:h[^>]*>(.*?)</text:h[^>]*>", r"\n\n\1\n", content, flags=re.DOTALL
    )
    text = re.sub(r"<text:p[^>]*/>", "\n", text)
    text = re.sub(r"<text:p[^>]*>(.*?)</text:p>", r"\1\n", text, flags=re.DOTALL)
    text = re.sub(r"<[^>]+>", "", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


# ---------------------------------------------------------------------------
# Feature 22 — Markdown footnote processing
# ---------------------------------------------------------------------------


def _process_footnotes(md: str, mode: str = "inline") -> str:
    """Handle Pandoc/GitHub-style Markdown footnotes for TTS.

    mode="inline":   insert footnote text at point of reference
    mode="deferred": collect footnotes and append at end of document
    mode="skip":     remove all footnote markers and content silently

    Supported syntax:
      Definitions:  [^1]: footnote text here
      References:   word[^1]  or  word[^label]
      Multi-line:   [^label]: First line
                        continuation
    """
    definitions: dict = {}

    def _strip_definitions(text: str) -> str:
        def _replace(m):
            definitions[m.group(1)] = m.group(2).strip()
            return ""

        return re.sub(r"^\[\^([^\]]+)\]:\s*(.+)$", _replace, text, flags=re.MULTILINE)

    md = _strip_definitions(md)

    if mode == "skip":
        return re.sub(r"\[\^[^\]]+\]", "", md).strip()

    if mode == "inline":

        def _inline_sub(m):
            note = definitions.get(m.group(1), "")
            return f" (footnote: {note})" if note else ""

        return re.sub(r"\[\^([^\]]+)\]", _inline_sub, md).strip()

    if mode == "deferred":
        md = re.sub(r"\[\^[^\]]+\]", "", md).strip()
        if definitions:
            md += "\n\n## Footnotes\n\n"
            for label, text in definitions.items():
                md += f"[^{label}]: {text}\n"
        return md

    return md


# ---------------------------------------------------------------------------
# Feature 23 — EPUB chapter extraction
# ---------------------------------------------------------------------------


def _epub_extract_chapters(path: str) -> "List[Tuple[str, str]]":
    """Extract chapter/navigation data from an EPUB file.

    Reads the NCX (EPUB 2) or NAV document (EPUB 3) to return an ordered
    list of (chapter_title, href) pairs representing the book\'s structure.
    Returns an empty list if the EPUB has no navigation data.
    """
    import zipfile
    from xml.etree import ElementTree as ET

    OPF = "http://www.idpf.org/2007/opf"
    NCX = "http://www.daisy.org/z3986/2005/ncx/"
    EPUB = "http://www.idpf.org/2007/ops"
    XHTML = "http://www.w3.org/1999/xhtml"

    try:
        zf = zipfile.ZipFile(path, "r")
    except zipfile.BadZipFile:
        return []

    with zf:
        try:
            container = zf.read("META-INF/container.xml")
        except KeyError:
            return []

        cont_root = ET.fromstring(container)
        rf = cont_root.find(
            ".//{urn:oasis:names:tc:opendocument:xmlns:container}rootfile"
        )
        if rf is None:
            return []
        opf_path = rf.get("full-path", "")
        opf_dir = opf_path.rsplit("/", 1)[0] + "/" if "/" in opf_path else ""

        try:
            opf_tree = ET.fromstring(zf.read(opf_path))
        except (KeyError, ET.ParseError):
            return []

        # Build manifest: id -> (href, media-type, properties)
        manifest = {
            item.get("id", ""): (
                item.get("href", ""),
                item.get("media-type", ""),
                item.get("properties", ""),
            )
            for item in opf_tree.findall(f"{{{OPF}}}manifest/{{{OPF}}}item")
        }

        chapters: list = []

        # --- EPUB3 NAV ---
        nav_id = next(
            (k for k, (h, t, p) in manifest.items() if "nav" in p.split()), None
        )
        if nav_id:
            try:
                nav_tree = ET.fromstring(zf.read(opf_dir + manifest[nav_id][0]))
                for nav_elem in nav_tree.iter(f"{{{XHTML}}}nav"):
                    if "toc" in nav_elem.get(f"{{{EPUB}}}type", ""):
                        for li in nav_elem.iter(f"{{{XHTML}}}li"):
                            a = li.find(f"{{{XHTML}}}a")
                            if a is not None:
                                title = "".join(a.itertext()).strip()
                                href = a.get("href", "").split("#")[0]
                                if title and href:
                                    chapters.append((title, opf_dir + href))
                        break
            except (KeyError, ET.ParseError):
                pass

        # --- NCX (EPUB2 fallback) ---
        if not chapters:
            spine = opf_tree.find(f"{{{OPF}}}spine")
            toc_id = (spine.get("toc", "") if spine is not None else "") or next(
                (
                    k
                    for k, (h, t, p) in manifest.items()
                    if t == "application/x-dtbncx+xml"
                ),
                "",
            )
            if toc_id and toc_id in manifest:
                try:
                    ncx_tree = ET.fromstring(zf.read(opf_dir + manifest[toc_id][0]))
                    for np in ncx_tree.iter(f"{{{NCX}}}navPoint"):
                        label = np.find(f"{{{NCX}}}navLabel/{{{NCX}}}text")
                        content = np.find(f"{{{NCX}}}content")
                        if label is not None and content is not None:
                            title = (label.text or "").strip()
                            href = content.get("src", "").split("#")[0]
                            if title and href:
                                chapters.append((title, opf_dir + href))
                except (KeyError, ET.ParseError):
                    pass

    # Deduplicate by href, preserving first occurrence
    seen: set = set()
    unique: list = []
    for title, href in chapters:
        if href not in seen:
            seen.add(href)
            unique.append((title, href))
    return unique


# ---------------------------------------------------------------------------
# Feature 24 — Document caching
# ---------------------------------------------------------------------------

CACHE_VERSION = 1  # bump to invalidate all on-disk caches


def _cache_key(path: str, settings_fingerprint: str) -> str:
    """Return the cache filename for a given document path + settings."""
    path_hash = hashlib.md5(path.encode()).hexdigest()[:16]
    return f"{path_hash}_{settings_fingerprint}_v{CACHE_VERSION}.json"


def _cache_save(path: str, doc_data: dict, settings_fingerprint: str) -> None:
    """Save parsed document data to the local cache directory."""
    # Skip URLs and very small files (< 1 KB)
    if path.startswith(("http://", "https://")):
        return
    try:
        if Path(path).stat().st_size < 1024:
            return
        mtime = Path(path).stat().st_mtime
    except OSError:
        return

    try:
        cache_dir = Path(CACHE_DIR)
        cache_dir.mkdir(parents=True, exist_ok=True)
        payload = {
            "version": CACHE_VERSION,
            "path": path,
            "mtime": mtime,
            "settings_fingerprint": settings_fingerprint,
            "data": doc_data,
        }
        (cache_dir / _cache_key(path, settings_fingerprint)).write_text(
            json.dumps(payload, ensure_ascii=False), encoding="utf-8"
        )
    except (PermissionError, OSError):
        pass


def _cache_load(path: str, settings_fingerprint: str) -> "Optional[dict]":
    """Load cached document data if it is still valid (mtime unchanged).
    Returns None if no valid cache entry exists.
    """
    if path.startswith(("http://", "https://")):
        return None

    try:
        current_mtime = Path(path).stat().st_mtime
    except OSError:
        return None

    try:
        cache_file = Path(CACHE_DIR) / _cache_key(path, settings_fingerprint)
        if not cache_file.exists():
            return None
        payload = json.loads(cache_file.read_text(encoding="utf-8"))
        if payload.get("version") != CACHE_VERSION:
            return None
        if abs(payload.get("mtime", 0) - current_mtime) > 1e-3:
            cache_file.unlink(missing_ok=True)
            return None
        return payload.get("data")
    except (PermissionError, json.JSONDecodeError, OSError):
        return None


# ---------------------------------------------------------------------------
# Helper — settings fingerprint
# ---------------------------------------------------------------------------


def _settings_fingerprint(settings: "Settings") -> str:
    """Hash the settings values that affect document parsing output."""
    key = (
        f"{settings['tts_skip_code']}"
        f"|{settings.get('table_reading_mode', 'structured')}"
        f"|{settings.get('footnote_mode', 'inline')}"
    )
    import hashlib

    return hashlib.md5(key.encode()).hexdigest()[:8]


# =============================================================================
# Main TUI Application
# =============================================================================


class StarApp:
    """Main curses application for star — Speaking Terminal Access Reader."""

    VERSION_STRING = f"{APP_NAME} {APP_VERSION} — {APP_TITLE}"

    def __init__(
        self, stdscr: "curses.window", settings: Settings, initial_path: str = ""
    ) -> None:
        self.scr = stdscr
        self.settings = settings
        self.doc: Optional[Document] = None
        self.rendered: List[Line] = []  # rendered display lines
        self.scroll = 0  # top visible display line
        self.tts = TTSManager(settings)
        self.search = SearchEngine()
        self.search_active = False
        self.search_dir = "forward"

        # UI state
        self.theme_name: str = settings["theme"]
        self.attrs: Dict[str, int] = {}
        self.message = ""
        self.message_t = 0.0
        self.message_dur = 4.0
        self.mode = "normal"  # "normal" | "mx" | "search" | "goto"
        self.mx_ed = LineEditor()
        self.mx_completions: List[str] = []
        self.mx_comp_idx = -1
        self.mx_history: List[str] = []
        self.mx_hist_pos = -1
        # When set, _mx_update_completions draws from this list instead of
        # MX_COMMANDS.  Used by the voice picker and any other minibuffer
        # that needs its own completion source.
        self._mx_custom_completions: Optional[List[str]] = None
        self.search_ed = LineEditor()
        self.goto_ed = LineEditor()
        self.loading = False
        self.loading_msg = ""
        self._load_queue: "queue.Queue[Optional[Document]]" = queue.Queue()
        self._running = True
        self._highlight_line = -1  # display line of current TTS word
        self._highlight_col_start = -1
        self._highlight_col_end = -1
        # Navigation history
        self._nav_history: List[int] = []
        self._nav_hist_pos: int = -1

        # Speech Cursor mode state
        self._sc_line: int = 0  # display-line index of the reading cursor
        self._sc_reader: Optional[_SCReader] = None  # persistent line reader

        # Word index saved when the user pauses speech (Space).  -1 means no
        # saved position.  Used by _tts_toggle to resume from the exact word
        # where reading was paused rather than restarting from the scroll top.
        self._tts_paused_at_word: int = -1

        # Sentence map: list of word-map indices where each sentence begins.
        # Built asynchronously (same thread as the word map) after each load.
        self._sentence_starts: List[int] = [0]

        # Initialize curses
        curses.curs_set(1)
        self.scr.keypad(True)
        self.scr.timeout(150)
        os.environ.setdefault("ESCDELAY", "25")

        self._init_colors()

        # TTS word highlight callback
        self.tts.set_on_highlight(self._on_highlight)

        if initial_path:
            self._open_async(initial_path)

    def _init_colors(self) -> None:
        self.attrs = _setup_colors(self.theme_name)

    def _a(self, role: str) -> int:
        return self.attrs.get(role, curses.A_NORMAL)

    # ── Message ────────────────────────────────────────────────────────────

    def notify(self, msg: str, dur: float = 4.0, error: bool = False) -> None:
        self.message = msg
        self.message_t = time.monotonic()
        self.message_dur = dur
        if error:
            pass  # Could pipe to log

    # ── Highlight callback from TTS ────────────────────────────────────────

    def _on_highlight(self, word_idx: int) -> None:
        """Called from the TTS/timer background thread — must NOT call any
        curses functions (not thread-safe).  Only update plain attributes;
        the main draw loop reads them on the next tick and adjusts scroll."""
        if not self.settings["highlight_current_word"]:
            return
        if self.doc and 0 <= word_idx < len(self.doc.word_map):
            wp = self.doc.word_map[word_idx]
            self._highlight_line = wp.disp_line
            self._highlight_col_start = wp.disp_col
            self._highlight_col_end = wp.disp_col + wp.tts_len
        else:
            self._highlight_line = -1
            self._highlight_col_start = -1
            self._highlight_col_end = -1

    # ── Async document loading ─────────────────────────────────────────────

    def _open_async(self, path: str) -> None:
        self.loading = True
        self.loading_msg = (
            f"Loading {Path(path).name if not path.startswith('http') else path} …"
        )

        def _work() -> None:
            try:
                doc = load_document(path, self.settings)
                self._load_queue.put(doc)
            except Exception as e:
                err_doc = Document(
                    path=path,
                    title="Error",
                    format="error",
                    markdown=f"# Load Error\n\n```\n{e}\n```\n",
                )
                err_doc.plain_text = str(e)
                self._load_queue.put(err_doc)

        threading.Thread(target=_work, daemon=True).start()

    def _poll_load_queue(self) -> None:
        try:
            doc = self._load_queue.get_nowait()
        except queue.Empty:
            return
        self.loading = False
        # Persist the current document's position before replacing it.
        self._save_reading_position()
        self.doc = doc
        # Build word map in background too
        self._render_doc()
        self.scroll = 0
        self._tts_stop()  # also clears any saved pause position for old doc
        self.notify(f"Opened: {doc.title}")
        recents: List[str] = self.settings["recent_files"]
        if doc.path and doc.path not in recents:
            recents.insert(0, doc.path)
            self.settings["recent_files"] = recents[:20]
        self.settings["last_path"] = doc.path
        if self.settings["tts_auto_play"]:
            self._tts_play()

    def _render_doc(self) -> None:
        if not self.doc:
            return
        h, w = self.scr.getmaxyx()
        wrap = int(self.settings["wrap_width"]) or (w - 2)
        self.rendered = render_markdown(
            self.doc.markdown,
            wrap,
            tab_width=int(self.settings["tab_width"]),
            syntax=bool(self.settings["syntax_highlight"]),
        )

        # Build word map and sentence map asynchronously (non-blocking)
        def _build() -> None:
            flat = ["".join(t for t, _ in line) for line in self.rendered]
            self.doc.word_map = _build_word_map(self.doc.plain_text, flat)
            self.tts.set_word_map(self.doc.word_map)
            self._build_sentence_map()  # depends on word_map
            self._restore_reading_position()  # scroll to last position

        threading.Thread(target=_build, daemon=True).start()

    # ── Sentence map ──────────────────────────────────────────────────

    def _build_sentence_map(self) -> None:
        """Populate self._sentence_starts with word-map indices at which each
        sentence begins.  Runs in the background thread that also builds the
        word map, so self.doc.word_map is guaranteed to exist on entry."""
        if not self.doc or not self.doc.plain_text or not self.doc.word_map:
            self._sentence_starts = [0]
            return

        text = self.doc.plain_text
        wm = self.doc.word_map

        # Collect the character offsets where new sentences begin.
        char_starts = [0]
        for m in _SENTENCE_SPLIT_RE.finditer(text):
            char_starts.append(m.end())

        # Map each char offset to the first word at or after that offset.
        # Both char_starts and wm are ordered, so a single forward walk suffices.
        word_starts: List[int] = []
        wi = 0
        for cs in char_starts:
            while wi < len(wm) and wm[wi].tts_offset < cs:
                wi += 1
            word_starts.append(min(wi, len(wm) - 1))

        # Deduplicate while preserving order.
        seen: set = set()
        result: List[int] = []
        for ws in word_starts:
            if ws not in seen:
                seen.add(ws)
                result.append(ws)

        self._sentence_starts = result if result else [0]

    def _find_sentence_idx(self, word_idx: int) -> int:
        """Return the index in _sentence_starts of the sentence that contains
        *word_idx* (binary search; O(log n))."""
        ss = self._sentence_starts
        lo, hi, result = 0, len(ss) - 1, 0
        while lo <= hi:
            mid = (lo + hi) // 2
            if ss[mid] <= word_idx:
                result = mid
                lo = mid + 1
            else:
                hi = mid - 1
        return result

    def _current_word_for_nav(self) -> int:
        """Return the best estimate of the current reading word index.

        Priority order:
        1. Live TTS highlight (engine is actively speaking).
        2. Saved pause position (_tts_paused_at_word) set when the user
           pressed Space to pause — this is the word we stopped on, so
           replay/sentence-jump commands operate at the right place even
           while speech is not running.
        3. First word at or below the current scroll position (viewport
           fallback when no speech has started or the document was just
           opened).
        """
        if self.tts.speaking:
            # Prefer the last callback-confirmed position (actual audio
            # position) over the timer estimate (which may be ahead).
            cb = self.tts.last_cb_word_idx
            if cb >= 0:
                return cb
            idx = self.tts.current_word_idx
            if idx >= 0:
                return idx
        if self._tts_paused_at_word >= 0:
            return self._tts_paused_at_word
        if self.doc and self.doc.word_map:
            for i, wp in enumerate(self.doc.word_map):
                if wp.disp_line >= self.scroll:
                    return i
        return 0

    # ── TTS controls ──────────────────────────────────────────────────

    def _tts_play(self) -> None:
        """Start speaking from the current scroll position.
        Slices plain_text at the first word on-screen so the engine never
        re-reads content that is already above the viewport."""
        if not self.doc:
            self.notify("No document loaded.", error=True)
            return
        start_word = 0
        if self.doc.word_map:
            for i, wp in enumerate(self.doc.word_map):
                if wp.disp_line >= self.scroll:
                    start_word = i
                    break
        self._tts_play_from_word(start_word)
        self.notify(
            f"Reading at {self.settings['tts_rate']} wpm via {self.tts.backend_name}"
        )

    def _tts_play_from_word(self, word_idx: int) -> None:
        """Start or restart TTS from a specific word-map index.

        Slices ``plain_text`` so the engine only reads from *word_idx*
        onwards.  When SSML is enabled the slice is wrapped with prosody
        markup; ``text_offset=-1`` tells TTSManager to skip word-callback
        offset arithmetic (the timer provides highlight accuracy instead).
        """
        if not self.doc:
            return
        wm = self.doc.word_map
        if wm and word_idx < len(wm):
            text_offset = wm[word_idx].tts_offset
        else:
            text_offset = 0
            word_idx = 0
        text_slice = self.doc.plain_text[text_offset:]

        # Apply speak-time normalizations (abbrev expansion, number words).
        text_slice = _preprocess_tts_text(text_slice, self.settings)

        if self.settings.get("use_ssml", False):
            text_for_engine = _text_to_ssml(
                text_slice,
                backend=self.tts.backend_name,
                sentence_ms=int(self.settings.get("ssml_sentence_pause_ms", 350)),
                clause_ms=int(self.settings.get("ssml_clause_pause_ms", 150)),
            )
            # SSML shifts char offsets — use -1 sentinel, rely on timer only.
            self.tts.speak(text_for_engine, start_word_idx=word_idx, text_offset=-1)
        else:
            self.tts.speak(text_slice, start_word_idx=word_idx, text_offset=text_offset)

    def _tts_stop(self) -> None:
        """Full stop — clears both speech and any saved pause position."""
        self.tts.stop()
        self._highlight_line = -1
        self._highlight_col_start = -1
        self._highlight_col_end = -1
        self._tts_paused_at_word = -1

    def _tts_toggle(self) -> None:
        """Pause/resume speech.

        * While speaking  → pause and remember the current word index so that
          the next press resumes from exactly that word.
        * While paused    → resume from the saved word index.
        * While stopped   → start from the current scroll position (same as
          before, so opening a fresh file and pressing Space still works).
        """
        if self.tts.speaking:
            # Save the last callback-confirmed position when available — it
            # reflects actual audio position rather than the timer's forward
            # estimate.  Pausing at the timer's ahead position would cause
            # resume to skip words; pausing at the callback position may
            # repeat a word or two but is far less disorienting.
            cb = self.tts.last_cb_word_idx
            saved = cb if cb >= 0 else self.tts.current_word_idx
            self._tts_stop()  # resets _tts_paused_at_word to -1
            if saved >= 0:
                self._tts_paused_at_word = saved  # restore the paused position
        elif self._tts_paused_at_word >= 0:
            w = self._tts_paused_at_word
            self._tts_paused_at_word = -1
            self._tts_play_from_word(w)
            self.notify(
                f"Resuming at {self.settings['tts_rate']} wpm via {self.tts.backend_name}"
            )
        else:
            self._tts_play()

    def _tts_speak_current_line(self) -> None:
        if not self.rendered or self.scroll >= len(self.rendered):
            return
        line = self.rendered[self.scroll]
        text = "".join(t for t, _ in line).strip()
        if text:
            self.tts.stop()
            text = _preprocess_tts_text(text, self.settings)
            if self.settings.get("use_ssml", False):
                text = _text_to_ssml(
                    text,
                    backend=self.tts.backend_name,
                    sentence_ms=int(self.settings.get("ssml_sentence_pause_ms", 350)),
                    clause_ms=int(self.settings.get("ssml_clause_pause_ms", 150)),
                )
            self.tts._backend.speak(text)

    # ── Sentence / paragraph navigation ──────────────────────────────────

    def _sentence_jump(
        self, dest_word: int, label: str, always_play: bool = False
    ) -> None:
        """Stop TTS, scroll to dest_word's display line, and restart speech.

        Restarts if speech was already active *or* if *always_play* is True
        (used by replay commands that must always begin reading).
        """
        self._history_push()  # record position before jump
        if not self.doc or not self.doc.word_map:
            return
        dest_word = max(0, min(dest_word, len(self.doc.word_map) - 1))
        dest_line = self.doc.word_map[dest_word].disp_line
        was_speaking = self.tts.speaking
        self.tts.stop()
        self._tts_paused_at_word = -1  # navigation breaks the pause/resume chain
        self._highlight_line = self._highlight_col_start = self._highlight_col_end = -1
        total = len(self.rendered)
        self.scroll = max(0, min(dest_line, total - 1) if total else 0)
        if was_speaking or always_play:
            self._tts_play_from_word(dest_word)
        self.notify(label)

    def _skip_next_sentence(self) -> None:
        if not self.doc or not self.doc.word_map:
            return
        cur = self._current_word_for_nav()
        si = self._find_sentence_idx(cur)
        nsi = si + 1
        if nsi >= len(self._sentence_starts):
            self.notify("No next sentence")
            return
        dest = self._sentence_starts[nsi]
        # Preview the first few words of the destination sentence
        preview = " ".join(
            self.doc.word_map[i].word
            for i in range(dest, min(dest + 5, len(self.doc.word_map)))
        )
        total = len(self._sentence_starts)
        self._sentence_jump(dest, f"→ Sentence {nsi + 1}/{total}: “{preview}…”")

    def _skip_prev_sentence(self) -> None:
        if not self.doc or not self.doc.word_map:
            return
        cur = self._current_word_for_nav()
        si = self._find_sentence_idx(cur)
        # If we are well into the current sentence, jump to its start first.
        # A "well into" threshold of 3 words feels natural (like double-tap rewind).
        if cur - self._sentence_starts[si] > 3:
            psi = si
        else:
            psi = max(0, si - 1)
        dest = self._sentence_starts[psi]
        preview = " ".join(
            self.doc.word_map[i].word
            for i in range(dest, min(dest + 5, len(self.doc.word_map)))
        )
        total = len(self._sentence_starts)
        self._sentence_jump(dest, f"← Sentence {psi + 1}/{total}: “{preview}…”")

    def _replay_sentence(self) -> None:
        """Jump to the start of the current sentence and always begin reading.

        Uses *always_play=True* so that a single, authoritative
        _tts_play_from_word call is made inside _sentence_jump regardless of
        whether speech was already active.  The old pattern of checking
        ``self.tts.speaking`` after the jump caused a Windows/SAPI5 race: the
        previous speech thread\'s finally-block could set _speaking=False
        after the new thread had already set it to True, making the guard fire
        a second _tts_play_from_word call that killed the first thread.
        """
        if not self.doc or not self.doc.word_map:
            return
        cur = self._current_word_for_nav()
        si = self._find_sentence_idx(cur)
        dest = self._sentence_starts[si]
        preview = " ".join(
            self.doc.word_map[i].word
            for i in range(dest, min(dest + 5, len(self.doc.word_map)))
        )
        self._sentence_jump(
            dest, f"↺ Replaying sentence: “{preview}…”", always_play=True
        )

    def _find_current_paragraph_start(self) -> int:
        """Return the first display line of the paragraph that contains
        the current scroll position."""
        i = self.scroll
        n = len(self.rendered)
        # If we landed on a blank line, step forward to content first
        while i < n - 1 and not self.rendered[i]:
            i += 1
        # Walk backward while the previous line is content (non-blank)
        while i > 0 and self.rendered[i - 1]:
            i -= 1
        return max(0, i)

    def _replay_paragraph(self) -> None:
        """Jump to the start of the current paragraph and always begin reading."""
        if not self.doc or not self.doc.word_map:
            return
        dest_line = self._find_current_paragraph_start()
        self.tts.stop()
        self._tts_paused_at_word = -1  # navigation breaks the pause/resume chain
        self._highlight_line = self._highlight_col_start = self._highlight_col_end = -1
        total = len(self.rendered)
        self.scroll = max(0, min(dest_line, total - 1) if total else 0)
        # Find the word at dest_line
        dest_word = 0
        if self.doc.word_map:
            for i, wp in enumerate(self.doc.word_map):
                if wp.disp_line >= dest_line:
                    dest_word = i
                    break
        self._tts_play_from_word(dest_word)
        self.notify(f"↺ Replaying paragraph from line {dest_line + 1}")

    # ── Reading position memory ──────────────────────────

    def _save_reading_position(self) -> None:
        """Persist the current reading offset for the open document."""
        if not self.doc or not self.doc.path or not self.doc.word_map:
            return
        cur = self._current_word_for_nav()
        if cur < 0 or cur >= len(self.doc.word_map):
            return
        offset = self.doc.word_map[cur].tts_offset
        total_chars = len(self.doc.plain_text)
        pct = int(100 * offset / max(1, total_chars))
        positions = dict(self.settings.get("reading_positions", {}))
        positions[self.doc.path] = {
            "offset": offset,
            "pct": pct,
            "ts": time.strftime("%Y-%m-%dT%H:%M:%S"),
        }
        if len(positions) > 200:
            evict = sorted(positions, key=lambda k: positions[k].get("ts", ""))[:50]
            for k in evict:
                del positions[k]
        self.settings.set("reading_positions", positions)

    def _restore_reading_position(self, force: bool = False) -> bool:
        """Scroll to the saved position for the current document.
        Safe to call from a background thread (only writes plain attributes)."""
        if not self.doc or not self.doc.path or not self.doc.word_map:
            return False
        if not force and not self.settings.get("tts_auto_resume", True):
            return False
        saved = self.settings.get("reading_positions", {}).get(self.doc.path)
        if not saved:
            return False
        target_offset = int(saved.get("offset", 0))
        pct = int(saved.get("pct", 0))
        ts = str(saved.get("ts", ""))[:10]
        wm = self.doc.word_map
        best = len(wm) - 1
        for i, wp in enumerate(wm):
            if wp.tts_offset >= target_offset:
                best = i
                break
        dest_line = wm[best].disp_line if best < len(wm) else 0
        total = len(self.rendered)
        self.scroll = max(0, min(dest_line, total - 1) if total else 0)
        self.notify(f"Resumed at {pct}%  (saved {ts})", dur=5.0)
        return True

    def _clear_reading_position(self) -> None:
        """Delete the saved position for the current document."""
        if not self.doc or not self.doc.path:
            return
        positions = dict(self.settings.get("reading_positions", {}))
        if self.doc.path in positions:
            del positions[self.doc.path]
            self.settings.set("reading_positions", positions)
            self.notify("Reading position cleared")
        else:
            self.notify("No saved position for this document")

    # ── Speed presets ───────────────────────────────────────────────────

    def _set_speed_preset(self, name: str) -> None:
        """Apply a named speed preset or a raw wpm integer."""
        name = (name or "").strip()
        if not name:
            presets = self.settings.get("speed_presets", {})
            lines = "  ".join(f"{k}={v}wpm" for k, v in presets.items())
            self.notify(
                f"Current: {self.settings['tts_rate']} wpm  |  {lines}", dur=6.0
            )
            return
        if name.isdigit():
            wpm = max(50, min(600, int(name)))
            self.tts.set_rate(wpm)
            self.notify(f"Speed: {wpm} wpm")
            return
        presets = self.settings.get("speed_presets", {})
        if name in presets:
            wpm = int(presets[name])
            self.tts.set_rate(wpm)
            self.notify(f"Speed preset “{name}”: {wpm} wpm")
        else:
            self.notify(
                f"Unknown preset “{name}”.  Known: {', '.join(presets)}",
                error=True,
            )

    def _preset_add(self, name: str) -> None:
        """Save the current TTS rate under *name* as a new speed preset."""
        name = (name or "").strip()
        if not name:
            self.notify("Usage: preset-add <name>", error=True)
            return
        wpm = int(self.settings["tts_rate"])
        presets = dict(self.settings.get("speed_presets", {}))
        presets[name] = wpm
        self.settings.set("speed_presets", presets)
        self.notify(f"Preset “{name}” saved: {wpm} wpm")

    def _preset_list(self) -> None:
        """Show all speed presets in the status bar."""
        presets = self.settings.get("speed_presets", {})
        if not presets:
            self.notify("No speed presets defined")
            return
        parts = [
            f"{k}: {v} wpm" for k, v in sorted(presets.items(), key=lambda x: x[1])
        ]
        self.notify("Presets — " + "  |  ".join(parts), dur=7.0)

    def _cycle_speed_preset(self) -> None:
        """Cycle through speed presets in ascending WPM order (F8)."""
        presets = self.settings.get("speed_presets", {})
        if not presets:
            self.notify("No speed presets defined")
            return
        ordered = sorted(presets.items(), key=lambda x: x[1])
        cur_rate = int(self.settings["tts_rate"])
        nxt = ordered[0]
        for name, wpm in ordered:
            if wpm > cur_rate:
                nxt = (name, wpm)
                break
        self.tts.set_rate(nxt[1])
        self.notify(f"Speed: “{nxt[0]}” — {nxt[1]} wpm")

    # ── SSML toggle ─────────────────────────────────────────────────────

    # ── Bookmarks, history, search & utility commands ───────────────────
    def _bookmark_set(self, name: str = "") -> None:
        "Set a named bookmark at the current reading position."
        if not self.doc or not self.doc.word_map:
            self.notify("No document loaded.", error=True)
            return
        doc_path = self.doc.path or self.doc.title
        bookmarks = dict(self.settings.get("bookmarks", {}))
        doc_bm = dict(bookmarks.get(doc_path, {}))
        if not name:
            n = 1
            while f"mark{n}" in doc_bm:
                n += 1
            name = f"mark{n}"
        cur = self._current_word_for_nav()
        offset = self.doc.word_map[cur].tts_offset
        total_chars = len(self.doc.plain_text)
        pct = int(100 * offset / max(1, total_chars))
        doc_bm[name] = {
            "offset": offset,
            "pct": pct,
            "ts": time.strftime("%Y-%m-%dT%H:%M:%S"),
        }
        bookmarks[doc_path] = doc_bm
        self.settings.set("bookmarks", bookmarks)
        self.notify(f"Bookmark set: {name}  ({pct}%)")

    def _bookmark_goto(self, name: str) -> None:
        "Jump to a named bookmark in the current document."
        if not self.doc or not self.doc.word_map:
            self.notify("No document loaded.", error=True)
            return
        name = (name or "").strip()
        if not name:
            self.notify("Usage: bookmark-goto <name>", error=True)
            return
        doc_path = self.doc.path or self.doc.title
        doc_bm = self.settings.get("bookmarks", {}).get(doc_path, {})
        if name not in doc_bm:
            self.notify(f"Bookmark '{name}' not found.", error=True)
            return
        target_offset = int(doc_bm[name].get("offset", 0))
        wm = self.doc.word_map
        # Find the word whose tts_offset is closest to the saved offset.
        best, best_dist = 0, abs(wm[0].tts_offset - target_offset)
        for i, wp in enumerate(wm):
            dist = abs(wp.tts_offset - target_offset)
            if dist < best_dist:
                best_dist, best = dist, i
            if wp.tts_offset > target_offset + best_dist:
                break  # offsets are monotonically increasing; no closer match ahead
        self._history_push()
        self._sentence_jump(best, f"Jumping to bookmark '{name}'")

    def _bookmark_list(self) -> None:
        "List all bookmarks for the current document in the status bar."
        if not self.doc:
            self.notify("No document loaded.", error=True)
            return
        doc_path = self.doc.path or self.doc.title
        doc_bm = self.settings.get("bookmarks", {}).get(doc_path, {})
        if not doc_bm:
            self.notify("No bookmarks for this document.")
            return
        parts = [
            f"{k} ({v.get('pct', '?')}%, {str(v.get('ts', ''))[:10]})"
            for k, v in sorted(doc_bm.items())
        ]
        self.notify("Bookmarks — " + "  |  ".join(parts), dur=8.0)

    def _bookmark_delete(self, name: str) -> None:
        "Remove a named bookmark from the current document."
        if not self.doc:
            self.notify("No document loaded.", error=True)
            return
        name = (name or "").strip()
        if not name:
            self.notify("Usage: bookmark-delete <name>", error=True)
            return
        doc_path = self.doc.path or self.doc.title
        bookmarks = dict(self.settings.get("bookmarks", {}))
        doc_bm = dict(bookmarks.get(doc_path, {}))
        if name not in doc_bm:
            self.notify(f"Bookmark '{name}' not found.", error=True)
            return
        del doc_bm[name]
        bookmarks[doc_path] = doc_bm
        self.settings.set("bookmarks", bookmarks)
        self.notify(f"Bookmark '{name}' deleted.")

    # ── Feature 29 — Navigation history ─────────────────────────────────────────

    def _history_push(self, offset: int = -1) -> None:
        "Record the current TTS offset in the navigation history before a jump."
        if not self.doc or not self.doc.word_map:
            return
        if offset < 0:
            cur = self._current_word_for_nav()
            if 0 <= cur < len(self.doc.word_map):
                offset = self.doc.word_map[cur].tts_offset
            else:
                return
        # When branching off mid-history (user navigated back then jumped elsewhere),
        # discard all forward entries so the list stays consistent.
        if self._nav_hist_pos >= 0:
            self._nav_history = self._nav_history[: self._nav_hist_pos + 1]
            self._nav_hist_pos = -1
        self._nav_history.append(offset)
        max_size = int(self.settings.get("nav_history_size", 50))
        if len(self._nav_history) > max_size:
            self._nav_history = self._nav_history[-max_size:]

    def _history_back(self) -> None:
        "Navigate to the previous position in the navigation history."
        if not self._nav_history:
            self.notify("Navigation history is empty.")
            return
        total = len(self._nav_history)
        if self._nav_hist_pos == -1:
            # First back-step: jump to the most recently saved position.
            new_pos = total - 1
        elif self._nav_hist_pos > 0:
            new_pos = self._nav_hist_pos - 1
        else:
            self.notify("No earlier history.")
            return
        self._nav_hist_pos = new_pos
        self.notify(f"History: position {new_pos + 1}/{total}")
        self._jump_to_offset(self._nav_history[new_pos])

    def _history_forward(self) -> None:
        "Navigate forward after having gone back in navigation history."
        if self._nav_hist_pos < 0:
            self.notify("No forward history.")
            return
        total = len(self._nav_history)
        new_pos = self._nav_hist_pos + 1
        if new_pos >= total:
            self._nav_hist_pos = -1
            self.notify("History: at present position.")
            return
        self._nav_hist_pos = new_pos
        self.notify(f"History: position {new_pos + 1}/{total}")
        self._jump_to_offset(self._nav_history[new_pos])

    def _jump_to_offset(self, target_offset: int) -> None:
        "Scroll to the word in the current document closest to *target_offset*."
        if not self.doc or not self.doc.word_map:
            return
        wm = self.doc.word_map
        best = len(wm) - 1
        for i, wp in enumerate(wm):
            if wp.tts_offset >= target_offset:
                best = i
                break
        dest_line = wm[best].disp_line
        was_speaking = self.tts.speaking
        self.tts.stop()
        self._highlight_line = self._highlight_col_start = self._highlight_col_end = -1
        total = len(self.rendered)
        self.scroll = max(0, min(dest_line, total - 1) if total else 0)
        if was_speaking:
            self._tts_play_from_word(best)

    # ── Feature 23b — Chapter navigation ─────────────────────────────────────────

    def _chapter_next(self) -> None:
        "Jump to the next chapter in the document."
        if not self.doc or not self.doc.word_map:
            return
        chapters = getattr(self.doc, "chapters", None)
        if not chapters:
            self.notify("No chapter navigation available for this document.")
            return
        cur = self._current_word_for_nav()
        current_ch_idx = 0
        for i, (_, _, widx) in enumerate(chapters):
            if widx <= cur:
                current_ch_idx = i
        next_idx = current_ch_idx + 1
        if next_idx >= len(chapters):
            self.notify("Already at the last chapter.")
            return
        title, _, dest_word = chapters[next_idx]
        self._history_push()
        self._sentence_jump(
            dest_word, f"Chapter {next_idx + 1}/{len(chapters)}: {title}"
        )

    def _chapter_prev(self) -> None:
        "Jump to the previous chapter, or to the current chapter start if well into it."
        if not self.doc or not self.doc.word_map:
            return
        chapters = getattr(self.doc, "chapters", None)
        if not chapters:
            self.notify("No chapter navigation available for this document.")
            return
        cur = self._current_word_for_nav()
        current_ch_idx = 0
        for i, (_, _, widx) in enumerate(chapters):
            if widx <= cur:
                current_ch_idx = i
        _, _, ch_start_word = chapters[current_ch_idx]
        # Mirror the double-tap rewind idiom used by sentence navigation:
        # if the reader is more than 5 words into the chapter, replay its start;
        # otherwise go one chapter back.
        if cur - ch_start_word > 5:
            dest_idx = current_ch_idx
        elif current_ch_idx == 0:
            self.notify("Already at the first chapter.")
            return
        else:
            dest_idx = current_ch_idx - 1
        title, _, dest_word = chapters[dest_idx]
        self._history_push()
        self._sentence_jump(
            dest_word, f"Chapter {dest_idx + 1}/{len(chapters)}: {title}"
        )

    def _chapter_list(self) -> None:
        "Show all chapter titles for the current document in the status bar."
        if not self.doc:
            self.notify("No document loaded.", error=True)
            return
        chapters = getattr(self.doc, "chapters", None)
        if not chapters:
            self.notify("No chapter navigation available for this document.")
            return
        parts = [f"{i + 1}. {title}" for i, (title, _, _) in enumerate(chapters)]
        self.notify("Chapters — " + "  |  ".join(parts), dur=10.0)

    def _chapter_goto(self, name_or_num: str) -> None:
        "Jump to a chapter by 1-based number or partial case-insensitive title match."
        if not self.doc or not self.doc.word_map:
            return
        chapters = getattr(self.doc, "chapters", None)
        if not chapters:
            self.notify("No chapter navigation available for this document.")
            return
        name_or_num = (name_or_num or "").strip()
        if not name_or_num:
            self._chapter_list()
            return
        if name_or_num.isdigit():
            n = int(name_or_num) - 1  # convert to 0-based index
            if 0 <= n < len(chapters):
                title, _, dest_word = chapters[n]
                self._history_push()
                self._sentence_jump(
                    dest_word, f"Chapter {n + 1}/{len(chapters)}: {title}"
                )
            else:
                self.notify(
                    f"Chapter number out of range (1–{len(chapters)}).", error=True
                )
            return
        # Partial title match — take the first hit.
        needle = name_or_num.lower()
        for i, (title, _, dest_word) in enumerate(chapters):
            if needle in title.lower():
                self._history_push()
                self._sentence_jump(
                    dest_word, f"Chapter {i + 1}/{len(chapters)}: {title}"
                )
                return
        self.notify(f"No chapter matching '{name_or_num}'.", error=True)

    # ── Feature 30 — Regex search (StarApp wrapper) ──────────────────────────────

    def _do_search_regex(self, pattern: str) -> None:
        "Wrapper that calls search_regex and updates status."
        pattern = (pattern or "").strip()
        if not pattern or not self.rendered:
            self.notify("Usage: search-regex <pattern>", error=True)
            return
        found = self.search.search_regex(pattern, self.rendered, from_line=self.scroll)
        if found:
            m = self.search.current_match
            if m:
                self._scroll_to_line(m[0])
            self.notify(f"{self.search.match_count} match(es) for regex /{pattern}/")
        else:
            self.notify(f"No regex matches for /{pattern}/", error=True)

    # ── Feature 22b — Footnote mode toggle ───────────────────────────────────────

    def _set_footnote_mode(self, mode: str) -> None:
        "Set footnote reading mode (inline | deferred | skip) and reload."
        mode = (mode or "").strip().lower()
        if not mode:
            cur = self.settings.get("footnote_mode", "inline")
            self.notify(
                f"Footnote mode: {cur}  —  options: inline  deferred  skip",
                dur=5.0,
            )
            return
        if mode not in ("inline", "deferred", "skip"):
            self.notify(
                f"Unknown footnote mode {mode!r}.  Use: inline, deferred, or skip.",
                error=True,
            )
            return
        self.settings.set("footnote_mode", mode)
        self.notify(f"Footnote mode: {mode}  (reloading document)")
        if self.doc and self.doc.path:
            self._open_async(self.doc.path)

    # ── Feature 27 — Clipboard copy (TUI stub) ───────────────────────────────────

    def _copy_current_line(self) -> str:
        "Return the text of the current top-visible display line."
        if not self.rendered or self.scroll >= len(self.rendered):
            return ""
        return "".join(t for t, _ in self.rendered[self.scroll]).strip()

    def _copy_to_clipboard(self) -> None:
        """Copy the current top-visible line to the system clipboard.

        Uses pyperclip when available; otherwise shows the text in the
        status bar so the user can select it manually. Never raises."""
        text = self._copy_current_line()
        if not text:
            self.notify("Nothing to copy (empty line).", error=True)
            return
        try:
            import pyperclip  # type: ignore

            pyperclip.copy(text)
            truncated = text[:60] + ("…" if len(text) > 60 else "")
            self.notify(f"Copied to clipboard: {truncated}")
        except Exception:
            # pyperclip unavailable or clipboard inaccessible — surface the text.
            self.notify(f"Copy (select manually): {text}", dur=10.0)

    # ── Feature 36 — Recent files ────────────────────────────────────────────────

    def _recent_files(self) -> None:
        "Show the recent-files list and open the selected entry via minibuffer."
        recent: List[str] = self.settings.get("recent_files", [])
        if not recent:
            self.notify("No recent files.")
            return
        # Preview up to 10 entries in the status bar.
        preview_parts = [f"{i + 1}. {path}" for i, path in enumerate(recent[:10])]
        self.notify("Recent: " + "  |  ".join(preview_parts), dur=8.0)

        def _on_pick(value: str) -> None:
            value = value.strip()
            if not value:
                return
            if value.isdigit():
                n = int(value) - 1
                if 0 <= n < len(recent):
                    self._open_async(recent[n])
                else:
                    self.notify(f"No recent file #{int(value)}.", error=True)
            else:
                # User typed a path directly.
                self._open_async(value)

        self._enter_minibuffer(
            prompt=f"Open recent [1–{min(len(recent), 10)}] or path: ",
            on_commit=_on_pick,
        )

    # ── Feature 44 — Wikipedia and PubMed shortcuts ──────────────────────────────

    def _open_wikipedia(self, query: str) -> None:
        "Fetch and open the Wikipedia article for query via the URL loader."
        query = (query or "").strip()
        if not query:
            self.notify("Usage: wikipedia <query>", error=True)
            return
        # Use the standard wiki URL; _open_async → _load_url handles HTML → Markdown.
        encoded = urllib.parse.quote(query.replace(" ", "_"))
        url = f"https://en.wikipedia.org/wiki/{encoded}"
        self.notify(f"Opening Wikipedia: {query}")
        self._open_async(url)

    def _open_pubmed(self, pmid: str) -> None:
        "Fetch and open a PubMed abstract by PMID via the URL loader."
        pmid = (pmid or "").strip()
        if not pmid:
            self.notify("Usage: pubmed <PMID>", error=True)
            return
        url = (
            "https://eutils.ncbi.nlm.nih.gov/entrez/eutils/efetch.fcgi"
            f"?db=pubmed&id={urllib.parse.quote(pmid)}&rettype=abstract&retmode=text"
        )
        self.notify(f"Opening PubMed abstract: PMID {pmid}")
        self._open_async(url)

    def _cache_clear(self) -> None:
        "Delete all cached document files to free disk space."
        import shutil as _shutil

        try:
            if CACHE_DIR.exists():
                _shutil.rmtree(CACHE_DIR)
                CACHE_DIR.mkdir(parents=True, exist_ok=True)
                self.notify("Document cache cleared")
            else:
                self.notify("Cache is already empty")
        except OSError as e:
            self.notify(f"Cache clear error: {e}", error=True)

    def _ssml_toggle(self) -> None:
        val = not bool(self.settings.get("use_ssml", False))
        state = "ON" if val else "OFF"
        self.settings.set("use_ssml", val)
        self.notify(f"SSML prosody: {state}")

    # ── Voice picker ───────────────────────────────────────────────────

    def _voice_picker(self) -> None:
        """Open an interactive voice-selection minibuffer.

        Voice *names* are used as the completion source so the user never
        has to type or copy a raw Windows registry path.  Substring search
        is used (type \"zira\" to find \"Microsoft Zira Desktop\").
        Pressing Enter applies the voice and speaks a brief test phrase.
        """
        voices = self.tts.list_voices()
        if not voices:
            self.notify(
                "No voices found. Is pyttsx3 installed and the backend set to pyttsx3?",
                error=True,
            )
            return

        # Build display strings and a lookup from display name → voice dict.
        # Append the language tag for clarity; deduplicate if names collide.
        name_map: Dict[str, Dict[str, str]] = {}
        ordered: List[str] = []
        for v in voices:
            name = v.get("name", v.get("id", "Unknown"))
            lang = v.get("lang", "")
            display = f"{name}  [{lang}]" if lang else name
            key, n = display, 1
            while key in name_map:
                n += 1
                key = f"{display} ({n})"
            ordered.append(key)
            name_map[key] = v

        # Pre-fill the current voice name so the user sees their selection.
        current_id = str(self.settings.get("tts_voice", ""))
        initial = ""
        for key, v in name_map.items():
            if v.get("id") == current_id:
                initial = key
                break

        def on_select(chosen: str) -> None:
            chosen = chosen.strip()
            match = name_map.get(chosen)
            if not match:
                # Fuzzy: first case-insensitive substring hit
                low = chosen.lower()
                for key, v in name_map.items():
                    if low in key.lower():
                        match = v
                        break
            if not match:
                self.notify(f"Voice not found: {chosen!r}", error=True)
                return
            self._apply_voice(match.get("id", ""), match.get("name", chosen))

        self._enter_minibuffer(
            "Voice (Tab to browse, type to filter): ",
            initial=initial,
            on_commit=on_select,
            completions=ordered,
        )

    def _apply_voice(self, voice_id: str, voice_name: str = "") -> None:
        """Apply *voice_id* to the active backend, persist it, and speak a
        brief confirmation phrase so the user can immediately hear the change."""
        self.tts._backend.set_voice(voice_id)
        self.settings.set("tts_voice", voice_id)
        label = voice_name or voice_id or "system default"
        self.notify(f"Voice: {label}")
        # Stop any current speech then speak a one-line test so the user
        # can hear the new voice without pressing Space.
        self.tts.stop()
        self.tts._backend.speak(f"Voice changed to {label}.")

    # ── Abbreviation helpers ───────────────────────────────────────────────

    def _abbrev_add(self, arg: str) -> None:
        """Add or update a custom abbreviation expansion.
        Usage:  abbrev-add <abbrev.> <expansion words>
        Example: abbrev-add RCT randomized controlled trial
        """
        parts = arg.strip().split(None, 1)
        if len(parts) < 2:
            self.notify(
                "Usage: abbrev-add <abbreviation> <expansion>   "
                "e.g.  abbrev-add RCT randomized controlled trial",
                error=True,
            )
            return
        abbr, expansion = parts[0], parts[1].strip()
        custom = dict(self.settings.get("abbrev_expansions") or {})
        custom[abbr] = expansion
        self.settings.set("abbrev_expansions", custom)
        self.notify(f"Abbreviation saved: {abbr!r} \u2192 {expansion!r}")

    def _abbrev_list(self) -> None:
        """Show all active custom abbreviation expansions."""
        custom = self.settings.get("abbrev_expansions") or {}
        if not custom:
            self.notify("No custom abbreviations defined.  Use abbrev-add to add one.")
            return
        pairs = "  |  ".join(f"{k} \u2192 {v}" for k, v in sorted(custom.items()))
        self.notify(f"Custom abbreviations: {pairs}", dur=8.0)

    # ── Table mode helper ─────────────────────────────────────────────────

    def _set_table_mode(self, mode: str) -> None:
        """Set table reading mode.  Valid values: structured | flat | skip.
        Reloads the document so the change takes effect immediately.
        """
        mode = (mode or "").strip().lower()
        if not mode:
            cur = self.settings.get("table_reading_mode", "structured")
            self.notify(
                f"Table mode: {cur}  \u2014  options: structured  flat  skip",
                dur=5.0,
            )
            return
        if mode not in ("structured", "flat", "skip"):
            self.notify(
                f"Unknown table mode {mode!r}.  Use: structured, flat, or skip.",
                error=True,
            )
            return
        self.settings.set("table_reading_mode", mode)
        self.notify(f"Table reading mode: {mode}  (reloading document)")
        if self.doc and self.doc.path:
            self._open_async(self.doc.path)  # reload with new mode baked in

    def _compute_reading_level_tui(self) -> str:
        """Compute Flesch-Kincaid reading level for the current document."""
        if not self.doc or not self.doc.plain_text:
            return "No document loaded"
        text = self.doc.plain_text[:50000]  # cap for speed
        words = text.split()
        n_words = max(1, len(words))
        sentences = re.split(r"[.!?]+", text)
        n_sentences = max(1, len([s for s in sentences if s.strip()]))

        def _syllables(word: str) -> int:
            word = word.lower().rstrip(".,;:!?")
            if not word:
                return 1
            count = len(re.findall(r"[aeiou]+", word))
            if word.endswith("e") and count > 1:
                count -= 1
            return max(1, count)

        n_syllables = sum(_syllables(w) for w in words)
        ease = (
            206.835 - 1.015 * (n_words / n_sentences) - 84.6 * (n_syllables / n_words)
        )
        grade = 0.39 * (n_words / n_sentences) + 11.8 * (n_syllables / n_words) - 15.59
        ease = max(0.0, min(100.0, ease))
        grade = max(0.0, grade)
        if grade < 6:
            level = "elementary"
        elif grade < 9:
            level = "middle school"
        elif grade < 13:
            level = "high school"
        elif grade < 16:
            level = "college"
        else:
            level = "graduate"
        return (
            f"Reading level: Grade {grade:.1f} ({level})  "
            f"Ease: {ease:.0f}/100  "
            f"({n_words:,} words, {n_sentences:,} sentences)"
        )

    # ── Skip navigation (paragraph / heading) ────────────────────────────────────────────────────

    def _navigate_to(self, line: int) -> None:
        """Scroll to *line* and, if TTS was already playing, restart speech
        from the new position so the reader continues without interruption."""
        was_speaking = self.tts.speaking
        if was_speaking:
            self.tts.stop()
            self._highlight_line = -1
            self._highlight_col_start = -1
            self._highlight_col_end = -1
        total = len(self.rendered)
        self.scroll = max(0, min(line, total - 1) if total else 0)
        if was_speaking:
            self._tts_play()

    def _is_heading_line(self, line_idx: int) -> bool:
        """Return True if the rendered line at *line_idx* is a heading."""
        if 0 <= line_idx < len(self.rendered):
            return any(role in _HEADING_ROLES for _, role in self.rendered[line_idx])
        return False

    def _find_next_paragraph(self, from_line: int) -> int:
        """Return the first line of the paragraph that starts after the one
        containing *from_line*.  Falls back to the last line if none found."""
        n = len(self.rendered)
        i = from_line + 1
        # Walk forward through the current paragraph
        while i < n and self.rendered[i]:
            i += 1
        # Skip blank separator lines
        while i < n and not self.rendered[i]:
            i += 1
        return min(i, n - 1)

    def _find_prev_paragraph(self, from_line: int) -> int:
        """Return the first line of the paragraph that starts before the one
        containing *from_line*.  Falls back to line 0 if none found."""
        i = from_line - 1
        # Skip blank lines backward
        while i > 0 and not self.rendered[i]:
            i -= 1
        # Walk back through the previous paragraph's content
        while i > 0 and self.rendered[i]:
            i -= 1
        # Now i is on a blank line (or 0) — step forward to the first content line
        if not self.rendered[i]:  # blank
            i += 1
        return max(0, i)

    def _find_next_heading(self, from_line: int) -> Optional[int]:
        """Return the line index of the next heading after *from_line*, or
        None if there is no heading below the current position."""
        for i in range(from_line + 1, len(self.rendered)):
            if self._is_heading_line(i):
                return i
        return None

    def _find_prev_heading(self, from_line: int) -> Optional[int]:
        """Return the line index of the previous heading before *from_line*,
        or None if there is no heading above the current position."""
        for i in range(from_line - 1, -1, -1):
            if self._is_heading_line(i):
                return i
        return None

    def _skip_next_paragraph(self) -> None:
        if not self.rendered:
            return
        dest = self._find_next_paragraph(self.scroll)
        self._navigate_to(dest)
        self.notify(f"Paragraph →  line {dest + 1}")

    def _skip_prev_paragraph(self) -> None:
        if not self.rendered:
            return
        dest = self._find_prev_paragraph(self.scroll)
        self._navigate_to(dest)
        self.notify(f"Paragraph ←  line {dest + 1}")

    def _is_table_line(self, line_idx: int) -> bool:
        """Return True if the rendered line at *line_idx* contains table content."""
        if 0 <= line_idx < len(self.rendered):
            return any(role in _TABLE_ROLES for _, role in self.rendered[line_idx])
        return False

    def _find_next_table(self, from_line: int) -> Optional[int]:
        """Return the first line of the next table after *from_line*, or None."""
        n = len(self.rendered)
        i = from_line + 1
        # Skip through any table we're currently inside
        while i < n and self._is_table_line(i):
            i += 1
        for j in range(i, n):
            if self._is_table_line(j):
                return j
        return None

    def _find_prev_table(self, from_line: int) -> Optional[int]:
        """Return the first line of the previous table before *from_line*, or None."""
        i = from_line - 1
        # Skip through any table we're currently inside
        while i >= 0 and self._is_table_line(i):
            i -= 1
        # Scan backward for any table line
        while i >= 0 and not self._is_table_line(i):
            i -= 1
        if i < 0:
            return None
        # Walk back to the start of this table block
        while i > 0 and self._is_table_line(i - 1):
            i -= 1
        return i

    def _skip_next_table(self) -> None:
        dest = self._find_next_table(self.scroll)
        if dest is None:
            self.notify("No table below current position")
            return
        self._navigate_to(dest)
        self.notify(f"▼ Table — line {dest + 1}")

    def _skip_prev_table(self) -> None:
        dest = self._find_prev_table(self.scroll)
        if dest is None:
            self.notify("No table above current position")
            return
        self._navigate_to(dest)
        self.notify(f"▲ Table — line {dest + 1}")

    def _skip_next_heading(self) -> None:
        if not self.rendered:
            return
        dest = self._find_next_heading(self.scroll)
        if dest is None:
            self.notify("No heading below current position")
            return
        self._navigate_to(dest)
        # Show the heading text in the notification
        heading_text = "".join(t for t, _ in self.rendered[dest]).strip()
        self.notify(f"↓ Heading: {heading_text[:50]}")

    def _skip_prev_heading(self) -> None:
        if not self.rendered:
            return
        dest = self._find_prev_heading(self.scroll)
        if dest is None:
            self.notify("No heading above current position")
            return
        self._navigate_to(dest)
        heading_text = "".join(t for t, _ in self.rendered[dest]).strip()
        self.notify(f"↑ Heading: {heading_text[:50]}")

    # ── Read-from-heading  (’>’ / ’<’) ──────────────────────────────────────

    def _line_to_word(self, line: int) -> int:
        """Return the word-map index of the first word at or after *line*.
        Returns 0 when the word map is unavailable."""
        if self.doc and self.doc.word_map:
            for i, wp in enumerate(self.doc.word_map):
                if wp.disp_line >= line:
                    return i
        return 0

    def _read_next_heading(self) -> None:
        """Jump to the next heading and *always* begin TTS reading from it.

        Unlike ’}’ (which only resumes speech if it was already playing),
        this command starts the engine unconditionally so the user can
        move through headings while the document is stopped.
        """
        if not self.rendered:
            return
        dest_line = self._find_next_heading(self.scroll)
        if dest_line is None:
            self.notify("No heading below current position")
            return
        dest_word = self._line_to_word(dest_line)
        heading_text = "".join(t for t, _ in self.rendered[dest_line]).strip()
        self._sentence_jump(
            dest_word,
            f"⏩ Reading from: {heading_text[:60]}",
            always_play=True,
        )

    def _read_prev_heading(self) -> None:
        """Jump to the previous heading and *always* begin TTS reading from it.

        Unlike ’{’ (which only resumes speech if it was already playing),
        this command starts the engine unconditionally.
        """
        if not self.rendered:
            return
        dest_line = self._find_prev_heading(self.scroll)
        if dest_line is None:
            self.notify("No heading above current position")
            return
        dest_word = self._line_to_word(dest_line)
        heading_text = "".join(t for t, _ in self.rendered[dest_line]).strip()
        self._sentence_jump(
            dest_word,
            f"⏪ Reading from: {heading_text[:60]}",
            always_play=True,
        )

    # ── Speech Cursor (SC) mode ───────────────────────────────────────────

    def _sc_enter(self) -> None:
        """Activate Speech Cursor mode.

        Stops any running speech, positions the reading cursor at the last
        highlighted word (or the current scroll line) and switches to 'sc'
        mode where every navigation keystroke moves the cursor and reads
        just that single line.
        """
        self._tts_stop()
        if self._highlight_line >= 0:
            self._sc_line = self._highlight_line
        else:
            self._sc_line = self.scroll
        # Build the persistent engine now so the first line has no startup lag.
        if _PYTTSX3 and isinstance(self.tts._backend, Pyttsx3Backend):
            self._sc_reader = _SCReader(
                rate=int(self.settings["tts_rate"]),
                volume=float(self.settings["tts_volume"]),
            )
            self._sc_reader.start()
        else:
            self._sc_reader = None
        self.mode = "sc"
        self.notify(
            "Speech Cursor ON  ↑↓:line  ,/.:sent  [/]:para  {/}:head  t/T:table"
            "  Enter:read-on  Space:pause  Esc:exit",
            dur=7.0,
        )

    def _sc_exit(self, start_reading: bool = False) -> None:
        """Exit Speech Cursor mode.  Speech is **always** stopped first.

        If *start_reading* is True (Enter key), continuous TTS then starts
        from the cursor position so the user can resume full reading from
        wherever they browsed to.  Every other exit leaves the engine silent.
        """
        self.mode = "normal"
        # Stop the persistent SC reader first — this reaches the live SAPI5
        # engine directly without the 200–500 ms Engine() construction race.
        if self._sc_reader is not None:
            self._sc_reader.close()
            self._sc_reader = None
        self._tts_stop()  # also silence the main TTS backend
        if start_reading:
            dest_word = self._line_to_word(self._sc_line)
            self.scroll = self._sc_line
            self._tts_paused_at_word = -1
            self._tts_play_from_word(dest_word)
            self.notify(f"Reading from line {self._sc_line + 1}")

    def _sc_read_line(self, line_idx: int) -> None:
        """Stop current speech and read exactly one rendered line.

        This is the fundamental SC mode action: the cursor sits on a line,
        the engine reads that line and stops, then waits for the next
        navigation keystroke.  Blank lines are announced as \"blank\" so
        the user knows they have crossed a paragraph boundary.

        Plain text is always used regardless of the global *use_ssml* setting
        — SSML is unnecessary for a single short line and plain text keeps
        word-boundary callbacks active for accurate highlighting.
        """
        if not self.rendered or not (0 <= line_idx < len(self.rendered)):
            return
        self.tts.stop()  # stop the main TTS backend / timer
        text = "".join(t for t, _ in self.rendered[line_idx]).strip()
        if not text:
            # Announce blank lines so the user knows they've crossed a
            # paragraph boundary.
            if self._sc_reader is not None:
                self._sc_reader.speak("blank")
            else:
                self.tts._backend.speak("blank")
            return
        text = _preprocess_tts_text(text, self.settings)
        if self._sc_reader is not None:
            # Use the persistent reader: engine already warm, stop() is
            # always effective (no Engine() construction race on exit).
            self._sc_reader.speak(text)
        else:
            # Fallback for non-pyttsx3 backends (eSpeak, DECtalk).
            self.tts._backend.speak(text)

    def _sc_move(self, dest_line: int, label: str = "") -> None:
        """Move the SC cursor to *dest_line*, scroll it into view, and read
        just that single line — no continuous document reading."""
        total = len(self.rendered)
        if not total:
            return
        dest_line = max(0, min(dest_line, total - 1))
        self._sc_line = dest_line
        # Keep cursor comfortably visible
        h, _ = self.scr.getmaxyx()
        view_h = max(1, h - 4)
        margin = int(self.settings.get("scroll_margin", 3))
        if dest_line < self.scroll + margin:
            self.scroll = max(0, dest_line - margin)
        elif dest_line >= self.scroll + view_h - margin:
            self.scroll = max(0, dest_line - view_h + margin + 1)
        # Speak only this line
        self._sc_read_line(dest_line)
        if label:
            self.notify(label)

    def _handle_sc_key(self, ch: int) -> None:  # noqa: C901
        """Key handler for Speech Cursor (sc) mode."""
        # ── Exit / speech-control ─────────────────────────────────────────
        if ch == 27:  # Escape — exit + stop
            self._sc_exit(start_reading=False)
            return
        if ch == 9:  # Tab — exit SC mode and stop speech
            self._sc_exit(start_reading=False)
            return
        if ch in (curses.KEY_ENTER, 10, 13):  # Enter — exit + continuous read
            self._sc_exit(start_reading=True)
            return
        if ch == ord(" "):
            self._tts_toggle()
            return
        if ch in (0, 24):  # Ctrl+Space / Ctrl+X — full stop
            self._tts_stop()
            return

        if not self.rendered:
            return
        total = len(self.rendered)

        # ── Line navigation ───────────────────────────────────────────────
        if ch in (curses.KEY_DOWN, ord("j")):
            dest = min(self._sc_line + 1, total - 1)
            text = "".join(t for t, _ in self.rendered[dest]).strip()
            self._sc_move(dest, f"Line {dest + 1}: {text[:60]}")

        elif ch in (curses.KEY_UP, ord("k")):
            dest = max(self._sc_line - 1, 0)
            text = "".join(t for t, _ in self.rendered[dest]).strip()
            self._sc_move(dest, f"Line {dest + 1}: {text[:60]}")

        # ── Sentence navigation ───────────────────────────────────────────
        elif ch == ord("."):
            if not self.doc or not self.doc.word_map:
                return
            cur = self._line_to_word(self._sc_line)
            si = self._find_sentence_idx(cur)
            nsi = si + 1
            if nsi >= len(self._sentence_starts):
                self.notify("No next sentence")
                return
            dest_word = self._sentence_starts[nsi]
            dest_line = self.doc.word_map[dest_word].disp_line
            preview = " ".join(
                self.doc.word_map[i].word
                for i in range(dest_word, min(dest_word + 5, len(self.doc.word_map)))
            )
            total_s = len(self._sentence_starts)
            self._sc_move(
                dest_line,
                f"→ Sentence {nsi + 1}/{total_s}: “{preview}…”",
            )

        elif ch == ord(","):
            if not self.doc or not self.doc.word_map:
                return
            cur = self._line_to_word(self._sc_line)
            si = self._find_sentence_idx(cur)
            psi = si if cur - self._sentence_starts[si] > 3 else max(0, si - 1)
            dest_word = self._sentence_starts[psi]
            dest_line = self.doc.word_map[dest_word].disp_line
            preview = " ".join(
                self.doc.word_map[i].word
                for i in range(dest_word, min(dest_word + 5, len(self.doc.word_map)))
            )
            total_s = len(self._sentence_starts)
            self._sc_move(
                dest_line,
                f"← Sentence {psi + 1}/{total_s}: “{preview}…”",
            )

        # ── Paragraph navigation ──────────────────────────────────────────
        elif ch in (ord("]"), curses.KEY_NPAGE):
            dest = self._find_next_paragraph(self._sc_line)
            self._sc_move(dest, f"\u00b6 Next paragraph \u2014 line {dest + 1}")

        elif ch in (ord("["), curses.KEY_PPAGE):
            dest = self._find_prev_paragraph(self._sc_line)
            self._sc_move(dest, f"\u00b6 Prev paragraph \u2014 line {dest + 1}")

        # ── Heading navigation ────────────────────────────────────────────
        elif ch in (ord("}"), ord(">")):
            dest = self._find_next_heading(self._sc_line)
            if dest is None:
                self.notify("No heading below")
            else:
                heading_text = "".join(t for t, _ in self.rendered[dest]).strip()
                self._sc_move(dest, f"\u23e9 Heading: {heading_text[:60]}")

        elif ch in (ord("{"), ord("<")):
            dest = self._find_prev_heading(self._sc_line)
            if dest is None:
                self.notify("No heading above")
            else:
                heading_text = "".join(t for t, _ in self.rendered[dest]).strip()
                self._sc_move(dest, f"\u23ea Heading: {heading_text[:60]}")

        # ── Table navigation ──────────────────────────────────────────────
        elif ch == ord("t"):
            dest = self._find_next_table(self._sc_line)
            if dest is None:
                self.notify("No table below")
            else:
                self._sc_move(dest, f"Table at line {dest + 1}")

        elif ch == ord("T"):
            dest = self._find_prev_table(self._sc_line)
            if dest is None:
                self.notify("No table above")
            else:
                self._sc_move(dest, f"Table at line {dest + 1}")

        # ── Re-read current line ──────────────────────────────────────────
        elif ch == ord("r"):
            line_text = "".join(t for t, _ in self.rendered[self._sc_line]).strip()
            self._sc_read_line(self._sc_line)
            self.notify(f"↺ Line {self._sc_line + 1}: {line_text[:60]}")

        # ── Document boundaries ───────────────────────────────────────────
        elif ch == curses.KEY_HOME:
            self._sc_move(0, "Top of document")

        elif ch == curses.KEY_END:
            self._sc_move(total - 1, "End of document")

    def _rate_change(self, delta: int) -> None:
        new_rate = max(50, min(600, int(self.settings["tts_rate"]) + delta))
        self.tts.set_rate(new_rate)
        if self._sc_reader is not None:
            self._sc_reader.update_rate(new_rate)
        self.notify(f"Speech rate: {new_rate} wpm")

    def _volume_change(self, delta: float) -> None:
        new_vol = max(0.0, min(1.0, float(self.settings["tts_volume"]) + delta))
        self.tts.set_volume(new_vol)
        self.notify(f"Volume: {int(new_vol * 100)}%")

    # ── Scrolling ─────────────────────────────────────────────────────────

    def _scroll_by(self, delta: int) -> None:
        total = len(self.rendered)
        h, _ = self.scr.getmaxyx()
        view_h = max(1, h - 4)
        self.scroll = max(0, min(total - 1, self.scroll + delta))

    def _page_down(self) -> None:
        h, _ = self.scr.getmaxyx()
        self._scroll_by(max(1, h - 5))

    def _page_up(self) -> None:
        h, _ = self.scr.getmaxyx()
        self._scroll_by(-max(1, h - 5))

    def _goto_top(self) -> None:
        self.scroll = 0

    def _goto_bottom(self) -> None:
        self.scroll = max(0, len(self.rendered) - 1)

    def _scroll_to_line(self, display_line: int) -> None:
        h, _ = self.scr.getmaxyx()
        view_h = max(1, h - 4)
        margin = int(self.settings["scroll_margin"])
        self.scroll = max(0, min(len(self.rendered) - 1, display_line - view_h // 3))

    # ── Search ─────────────────────────────────────────────────────────────

    def _do_search(self, query: str, direction: str = "forward") -> None:
        if not query or not self.rendered:
            return
        found = self.search.search(query, self.rendered, from_line=self.scroll)
        if found:
            m = self.search.current_match
            if m:
                self._scroll_to_line(m[0])
                self.notify(f"{self.search.match_count} match(es) for '{query}'")
        else:
            self.notify(f"No matches for '{query}'", error=True)

    def _search_next(self) -> None:
        m = self.search.next_match()
        if m:
            self._scroll_to_line(m[0])
        else:
            self.notify("No search active", error=True)

    def _search_prev(self) -> None:
        m = self.search.prev_match()
        if m:
            self._scroll_to_line(m[0])
        else:
            self.notify("No search active", error=True)

    # ── File operations ────────────────────────────────────────────────────

    def _open_file_prompt(self) -> None:
        """Open a file path via minibuffer prompt."""
        last = str(self.settings.get("last_path", "")) or ""
        default = (
            (str(Path(last).parent) + os.sep) if last and Path(last).is_file() else ""
        )
        self._enter_minibuffer(
            "Find file: ", initial=default, on_commit=self._open_file_cb
        )

    def _open_file_cb(self, path: str) -> None:
        path = path.strip().rstrip("/\\")
        if not path:
            return
        path = os.path.expanduser(os.path.expandvars(path))
        if path.startswith(("http://", "https://")):
            self._open_async(path)
        elif os.path.exists(path):
            self._open_async(path)
        else:
            self.notify(f"File not found: {path}", error=True)

    def _open_url_prompt(self) -> None:
        self._enter_minibuffer(
            "Open URL: ", on_commit=lambda u: self._open_async(u.strip())
        )

    def _export_markdown(self) -> None:
        if not self.doc:
            self.notify("No document to export.", error=True)
            return
        p = Path(self.doc.path)
        default = (
            str(p.parent / (p.stem + "_export.md")) if self.doc.path else "export.md"
        )
        self._enter_minibuffer(
            "Export Markdown to: ", initial=default, on_commit=self._export_md_cb
        )

    def _export_md_cb(self, dest: str) -> None:
        dest = dest.strip()
        if not dest or not self.doc:
            return
        try:
            Path(dest).write_text(self.doc.markdown, encoding="utf-8")
            self.notify(f"Exported → {dest}")
        except OSError as e:
            self.notify(f"Export error: {e}", error=True)

    def _export_braille_cmd(self) -> None:
        if not self.doc:
            self.notify("No document to export.", error=True)
            return
        table = str(self.settings["braille_table"])
        brf = _export_braille(
            self.doc.plain_text,
            table,
            use_liblouis=bool(self.settings.get("braille_grade2", False)),
        )
        p = Path(self.doc.path) if self.doc.path else Path("export")
        dest = str(p.parent / (p.stem + ".brf"))
        try:
            Path(dest).write_text(brf, encoding="utf-8")
            self.notify(f"BRF exported → {dest}")
        except OSError as e:
            self.notify(f"BRF export error: {e}", error=True)

    def _export_audio_cmd(self, fmt: str = "") -> None:
        """Prompt for an output path and export TTS audio (M-x export-audio).

        *fmt* is the default file extension (wav, mp3, ogg, mp4).  When empty
        the ``audio_export_format`` setting is used (WAV by default — it needs
        no external tools).  Synthesis runs synchronously, so the TUI will be
        unresponsive until it finishes.  Use a shorter document or the
        espeak/pyttsx3 backend for faster results.
        """
        fmt = (fmt or str(self.settings.get("audio_export_format", "wav"))).lstrip(".")
        if not self.doc:
            self.notify("No document to export.", error=True)
            return
        p = Path(self.doc.path) if self.doc.path else Path("export")
        default = str(p.parent / (p.stem + f".{fmt}"))
        self._enter_minibuffer(
            f"Export audio ({fmt}) to: ",
            initial=default,
            on_commit=self._export_audio_cb,
        )

    def _export_audio_cb(self, dest: str) -> None:
        dest = dest.strip()
        if not dest or not self.doc:
            return
        text = _preprocess_tts_text(self.doc.plain_text, self.settings)
        self.notify("Exporting audio… please wait", dur=5.0)
        try:
            self.tts.export_audio(text, dest)
            self.notify(f"Audio exported → {dest}")
        except Exception as exc:
            self.notify(f"Audio export error: {exc}", error=True)

    def _goto_line_prompt(self) -> None:
        self._enter_minibuffer("Go to line: ", on_commit=self._goto_line_cb)

    def _goto_line_cb(self, s: str) -> None:
        try:
            ln = max(1, int(s.strip())) - 1
            self._scroll_to_line(min(ln, len(self.rendered) - 1))
        except ValueError:
            self.notify(f"Invalid line: {s}", error=True)

    # ── Theme ─────────────────────────────────────────────────────────────

    def _next_theme(self) -> None:
        idx = (
            THEME_NAMES.index(self.theme_name) if self.theme_name in THEME_NAMES else 0
        )
        self.theme_name = THEME_NAMES[(idx + 1) % len(THEME_NAMES)]
        self.settings["theme"] = self.theme_name
        self._init_colors()
        self.notify(f"Theme: {self.theme_name}")

    def _set_theme(self, name: str) -> None:
        if name in THEMES:
            self.theme_name = name
            self.settings["theme"] = name
            self._init_colors()
            self.notify(f"Theme: {name}")
        else:
            self.notify(
                f"Unknown theme '{name}'.  Available: {', '.join(THEME_NAMES)}",
                error=True,
            )

    # ── Minibuffer ─────────────────────────────────────────────────────────

    def _enter_minibuffer(
        self,
        prompt: str = "M-x: ",
        initial: str = "",
        on_commit: Optional[Callable[[str], None]] = None,
        mode: str = "mx",
        completions: Optional[List[str]] = None,
    ) -> None:
        self.mode = mode
        if mode == "mx":
            self.mx_ed = LineEditor(initial)
            self.mx_comp_idx = -1
            self.mx_hist_pos = -1
            self._mx_custom_completions = completions  # None → use MX_COMMANDS
            self._mx_update_completions()
            self._mx_prompt = prompt
            self._mx_callback = on_commit or self.execute_command
        elif mode == "search":
            self.search_ed = LineEditor(initial)
            self._search_prompt = prompt
        elif mode == "goto":
            self.goto_ed = LineEditor(initial)

    def _mx_update_completions(self) -> None:
        source = self._mx_custom_completions
        if source is not None:
            # Custom completion list (e.g. voice names): case-insensitive
            # substring match so typing "zira" finds "Microsoft Zira Desktop".
            q = self.mx_ed.value.strip().lower()
            self.mx_completions = (
                [c for c in source if q in c.lower()] if q else list(source)
            )
        else:
            # Normal M-x: prefix match against command names.
            prefix = (
                self.mx_ed.value.split()[0].lower() if self.mx_ed.value.strip() else ""
            )
            self.mx_completions = (
                [c for c in MX_COMMANDS if c.startswith(prefix)]
                if prefix
                else list(MX_COMMANDS)
            )

    def _mx_tab(self) -> None:
        matches = self.mx_completions
        if not matches:
            return
        if len(matches) == 1:
            self.mx_ed.set_value(matches[0] + " ")
            self.mx_comp_idx = 0
            return
        lcp = matches[0]
        for m in matches[1:]:
            while lcp and not m.startswith(lcp):
                lcp = lcp[:-1]
            if not lcp:
                break
        cur = self.mx_ed.value.rstrip()
        if len(lcp) > len(cur):
            self.mx_ed.set_value(lcp)
            self.mx_comp_idx = 0
        else:
            if self.mx_comp_idx < 0:
                self.mx_comp_idx = 0
            else:
                self.mx_comp_idx = (self.mx_comp_idx + 1) % len(matches)
            self.mx_ed.set_value(matches[self.mx_comp_idx])
        self._mx_update_completions()

    def _cancel_minibuffer(self) -> None:
        self.mode = "normal"
        self.search.query = ""
        self._mx_custom_completions = None

    # ── M-x command executor ──────────────────────────────────────────────

    def execute_command(self, cmd_line: str) -> None:
        cmd_line = cmd_line.strip()
        if not cmd_line:
            return
        parts = cmd_line.split()
        cmd = parts[0].lower()
        args = parts[1:]
        arg = args[0] if args else ""

        cmd_map = {
            "open": self._open_file_prompt,
            "open-url": self._open_url_prompt,
            "close": lambda: (
                setattr(self, "doc", None)
                or self._render_doc()
                or self.notify("Document closed")
            ),
            "reload": lambda: (
                self._open_async(self.doc.path)
                if self.doc
                else self.notify("No document", error=True)
            ),
            "export-markdown": self._export_markdown,
            "export-braille": self._export_braille_cmd,
            "export-audio": lambda: self._export_audio_cmd(arg or ""),
            "play": self._tts_play,
            "stop": self._tts_stop,
            "pause": self._tts_toggle,
            "speak-line": self._tts_speak_current_line,
            "search": lambda: self._enter_minibuffer(
                "Search: ",
                mode="search",
                on_commit=lambda q: self._do_search(q, "forward"),
            ),
            "search-backward": lambda: self._enter_minibuffer(
                "Search ↑: ",
                mode="search",
                on_commit=lambda q: self._do_search(q, "backward"),
            ),
            "goto-line": self._goto_line_prompt,
            "theme": lambda: self._set_theme(arg) if arg else self._next_theme(),
            "line-numbers": lambda: (
                self.settings.set(
                    "show_line_numbers", not self.settings["show_line_numbers"]
                ),
                self.notify(
                    f"Line numbers {'on' if self.settings['show_line_numbers'] else 'off'}"
                ),
            ),
            "syntax-highlight": lambda: (
                self.settings.set(
                    "syntax_highlight", not self.settings["syntax_highlight"]
                ),
                self._render_doc(),
                self.notify(
                    f"Syntax highlight {'on' if self.settings['syntax_highlight'] else 'off'}"
                ),
            ),
            "wrap-width": lambda: self._enter_minibuffer(
                "Wrap width (0=auto): ",
                on_commit=lambda v: (
                    self.settings.set(
                        "wrap_width", int(v) if v.strip().isdigit() else 0
                    ),
                    self._render_doc(),
                    self.notify(f"Wrap: {self.settings['wrap_width'] or 'auto'}"),
                ),
            ),
            "rate-up": lambda: self._rate_change(+20),
            "rate-down": lambda: self._rate_change(-20),
            "volume-up": lambda: self._volume_change(+0.1),
            "volume-down": lambda: self._volume_change(-0.1),
            "tts-backend": lambda: self._enter_minibuffer(
                "TTS backend (pyttsx3/espeak/festival/coqui/dectalk/none): ",
                on_commit=lambda v: (
                    self.tts.change_backend(v.strip()),
                    self.notify(f"TTS: {self.tts.backend_name}"),
                ),
            ),
            # tts-voice now opens the interactive picker (same as voice-picker
            # and Ctrl+T) so the user sees names rather than opaque IDs.
            "tts-voice": self._voice_picker,
            "voice-picker": self._voice_picker,
            "font-size-up": lambda: self.notify(
                "Font size is set in your terminal emulator."
            ),
            "font-size-down": lambda: self.notify(
                "Font size is set in your terminal emulator."
            ),
            "help": lambda: self._show_help(),
            "about": lambda: self.notify(
                f"{self.VERSION_STRING}  |  {__copyright__}  |  {__license__}", dur=6.0
            ),
            "license": lambda: self._show_license(),
            "quit": lambda: setattr(self, "_running", False),
            "settings": lambda: self.notify(f"Settings: {SETTINGS_FILE}"),
            # ── Skip navigation ──────────────────────────────────────────
            "next-paragraph": self._skip_next_paragraph,
            "prev-paragraph": self._skip_prev_paragraph,
            "next-heading": self._skip_next_heading,
            "prev-heading": self._skip_prev_heading,
            "read-next-heading": self._read_next_heading,
            "read-prev-heading": self._read_prev_heading,
            "speech-cursor": self._sc_enter,
            "stop-speech": self._tts_stop,
            "next-sentence": self._skip_next_sentence,
            "prev-sentence": self._skip_prev_sentence,
            "replay-sentence": self._replay_sentence,
            "replay-paragraph": self._replay_paragraph,
            # Reading position memory
            "save-position": self._save_reading_position,
            "jump-saved": lambda: self._restore_reading_position(force=True),
            "clear-position": self._clear_reading_position,
            # Speed presets
            "speed": lambda: self._set_speed_preset(arg),
            "preset-add": lambda: self._preset_add(arg),
            "preset-list": self._preset_list,
            # Bookmarks
            "bookmark-set": lambda: self._bookmark_set(arg),
            "bookmark-goto": lambda: self._bookmark_goto(arg),
            "bookmark-list": lambda: self._bookmark_list(),
            "bookmark-delete": lambda: self._bookmark_delete(arg),
            # Chapter navigation
            "chapter-next": lambda: self._chapter_next(),
            "chapter-prev": lambda: self._chapter_prev(),
            "chapter-list": lambda: self._chapter_list(),
            "chapter-goto": lambda: self._chapter_goto(arg),
            # Navigation history
            "history-back": lambda: self._history_back(),
            "history-forward": lambda: self._history_forward(),
            # Search
            "search-regex": lambda: self._enter_minibuffer(
                "Regex: ",
                mode="search",
                on_commit=lambda q: self._do_search_regex(q),
            ),
            # Utility
            "copy": lambda: self._copy_to_clipboard(),
            "recent": lambda: self._recent_files(),
            "wiki": lambda: self._open_wikipedia(arg),
            "pubmed": lambda: self._open_pubmed(arg),
            # Cache
            "cache-clear": lambda: self._cache_clear(),
            # Footnotes
            "footnote-mode": lambda: self._set_footnote_mode(arg),
            # Math normalization
            "normalize-math": lambda: (
                self.settings.set(
                    "normalize_math",
                    not self.settings.get("normalize_math", True),
                ),
                self.notify(
                    "Math normalization: "
                    + ("ON" if self.settings.get("normalize_math") else "OFF")
                ),
            ),
            # Reading level
            "reading-level": lambda: self.notify(
                self._compute_reading_level_tui(), dur=8.0
            ),
            # SSML
            "ssml": lambda: self._ssml_toggle(),
            "ssml-on": lambda: (
                self.settings.set("use_ssml", True),
                self.notify("SSML prosody: ON"),
            ),
            "ssml-off": lambda: (
                self.settings.set("use_ssml", False),
                self.notify("SSML prosody: OFF"),
            ),
            # Abbreviation expansion
            "expand-abbreviations": lambda: (
                self.settings.set(
                    "expand_abbreviations",
                    not self.settings.get("expand_abbreviations", True),
                ),
                self.notify(
                    "Abbreviation expansion: "
                    + ("ON" if self.settings.get("expand_abbreviations") else "OFF")
                ),
            ),
            "abbrev-add": lambda: self._abbrev_add(arg),
            "abbrev-list": lambda: self._abbrev_list(),
            # Number normalization
            "normalize-numbers": lambda: (
                self.settings.set(
                    "normalize_numbers",
                    not self.settings.get("normalize_numbers", True),
                ),
                self.notify(
                    "Number normalization: "
                    + ("ON" if self.settings.get("normalize_numbers") else "OFF")
                ),
            ),
            # Table reading mode
            "table-mode": lambda: self._set_table_mode(arg),
            # Annotations / notes (TUI notes panel)
            "annotate": self._annotate,
            "notes": self._notes_browser,
            "annotations-list": lambda: self._annotations_list(arg),
            "annotations-search": self._annotations_search,
            "annotation-goto": lambda: self._annotation_goto(arg),
            "annotation-delete": lambda: self._annotation_delete(arg),
            "annotations-export": self._annotations_export,
            # Keyboard cheat sheet
            "shortcuts": self._show_shortcuts,
        }

        fn = cmd_map.get(cmd)
        if fn:
            try:
                fn()
            except Exception as e:
                self.notify(f"Command error: {e}", error=True)
        else:
            self.notify(f"Unknown command '{cmd}'.  Press F1 for help.", error=True)

    # ── Help / about pager ─────────────────────────────────────────────────

    def _show_help(self) -> None:
        """Open README.md in a pager, matching the Qt GUI F1 behavior.
        Falls back to the built-in _HELP_TEXT if README.md cannot be found."""
        readme = Path(__file__).parent / "README.md"
        if readme.exists():
            try:
                tmp = load_document(str(readme), self.settings)
            except Exception:
                tmp = Document(title="star Help", markdown=_HELP_TEXT, plain_text="")
        else:
            tmp = Document(title="star Help", markdown=_HELP_TEXT, plain_text="")
        old_doc, old_rendered, old_scroll = self.doc, self.rendered, self.scroll
        self.doc = tmp
        self._render_doc()
        self.scroll = 0
        self.notify("README.md  —  q / Esc to return")
        # Pager loop
        while True:
            self.draw()
            ch = self.scr.getch()
            if ch in (ord("q"), ord("Q"), 7, 27):
                break
            elif ch in (14, curses.KEY_DOWN, ord("j")):
                self._scroll_by(1)
            elif ch in (16, curses.KEY_UP, ord("k")):
                self._scroll_by(-1)
            elif ch in (curses.KEY_NPAGE, ord(" ")):
                self._page_down()
            elif ch == curses.KEY_PPAGE:
                self._page_up()
            elif ch in (curses.KEY_HOME, 1):
                self._goto_top()
            elif ch in (curses.KEY_END, 5):
                self._goto_bottom()
        self.doc, self.rendered, self.scroll = old_doc, old_rendered, old_scroll

    def _show_license(self) -> None:
        lic_md = _LICENSE_TEXT
        tmp = Document(title="License — GPL v3", markdown=lic_md, plain_text="")
        old = (self.doc, self.rendered, self.scroll)
        self.doc = tmp
        self._render_doc()
        self.scroll = 0
        while True:
            self.draw()
            ch = self.scr.getch()
            if ch in (ord("q"), ord("Q"), 7, 27):
                break
            elif ch in (14, curses.KEY_DOWN, ord("j")):
                self._scroll_by(1)
            elif ch in (16, curses.KEY_UP, ord("k")):
                self._scroll_by(-1)
            elif ch in (curses.KEY_NPAGE, ord(" ")):
                self._page_down()
            elif ch == curses.KEY_PPAGE:
                self._page_up()
        self.doc, self.rendered, self.scroll = old

    def _show_text_pager(self, title: str, markdown: str) -> None:
        """Render *markdown* in a read-only scrollable pager (q/Esc to exit).

        Shared by the notes list and the keyboard cheat sheet; mirrors the
        navigation keys used by _show_help / _show_license.
        """
        tmp = Document(title=title, markdown=markdown, plain_text="")
        old = (self.doc, self.rendered, self.scroll)
        self.doc = tmp
        self._render_doc()
        self.scroll = 0
        while True:
            self.draw()
            ch = self.scr.getch()
            if ch in (ord("q"), ord("Q"), 7, 27):
                break
            elif ch in (14, curses.KEY_DOWN, ord("j")):
                self._scroll_by(1)
            elif ch in (16, curses.KEY_UP, ord("k")):
                self._scroll_by(-1)
            elif ch in (curses.KEY_NPAGE, ord(" ")):
                self._page_down()
            elif ch == curses.KEY_PPAGE:
                self._page_up()
            elif ch in (curses.KEY_HOME, 1):
                self._goto_top()
            elif ch in (curses.KEY_END, 5):
                self._goto_bottom()
        self.doc, self.rendered, self.scroll = old

    def _show_shortcuts(self) -> None:
        """Show the canonical keyboard cheat sheet."""
        self._show_text_pager("Keyboard Shortcuts", _shortcuts_text(plain=False))

    # ── Annotations / notes (TUI) ──────────────────────────────────────────

    def _annot_key(self) -> str:
        """Per-document key under which annotations are stored."""
        if not self.doc:
            return ""
        return self.doc.path or self.doc.title or ""

    def _load_annotations(self) -> List[Dict[str, Any]]:
        """Saved notes for the current document, sorted by position."""
        key = self._annot_key()
        if not key:
            return []
        store = self.settings.get("annotations", {}) or {}
        items = [dict(a) for a in store.get(key, [])]
        items.sort(key=lambda a: int(a.get("word_idx", a.get("char_pos", 0)) or 0))
        return items

    def _store_annotations(self, items: List[Dict[str, Any]]) -> None:
        key = self._annot_key()
        if not key:
            return
        store = dict(self.settings.get("annotations", {}) or {})
        if items:
            store[key] = items
        else:
            store.pop(key, None)
        self.settings.set("annotations", store)

    def _annotate(self) -> None:
        """Add a note at the current reading position (key 'a' / M-x annotate)."""
        if not self.doc or not self.doc.word_map:
            self.notify("No document loaded.", error=True)
            return
        self._enter_minibuffer("Note: ", on_commit=self._annotate_note_cb)

    def _annotate_note_cb(self, note: str) -> None:
        note = note.strip()
        if not note:
            return
        self._pending_note = note
        self._enter_minibuffer(
            "Tags (optional, comma-separated): ", on_commit=self._annotate_tags_cb
        )

    def _annotate_tags_cb(self, tag_str: str) -> None:
        note = getattr(self, "_pending_note", "")
        if not note:
            return
        self._pending_note = ""
        wm = self.doc.word_map if self.doc else []
        word_idx = self._current_word_for_nav()
        if word_idx < 0:
            word_idx = 0
        anchor = ""
        if wm and 0 <= word_idx < len(wm):
            dl = wm[word_idx].disp_line
            if 0 <= dl < len(self.rendered):
                anchor = "".join(t for t, _ in self.rendered[dl]).strip()[:120]
        items = self._load_annotations()
        items.append(
            {
                "char_pos": 0,  # Qt-only; TUI anchors by word_idx
                "word_idx": int(word_idx),
                "anchor": anchor,
                "note": note,
                "tags": _parse_tags(tag_str),
                "cite": "",
                "ts": time.strftime("%Y-%m-%dT%H:%M:%S"),
            }
        )
        self._store_annotations(items)
        tags = _parse_tags(tag_str)
        self.notify(
            f"Note added ({len(items)} total)"
            + (f"  tags: {', '.join(tags)}" if tags else "")
        )

    def _annotations_list(self, query: str = "") -> None:
        """Show all notes (optionally filtered) in a pager (M-x annotations-list)."""
        items = self._load_annotations()
        if not items:
            self.notify("No notes yet. Press 'a' or M-x annotate to add one.")
            return
        rows = [(i, a) for i, a in enumerate(items) if _annotation_matches(a, query)]
        if not rows:
            self.notify(f"No notes match '{query}'.")
            return
        title = self.doc.title if self.doc else "document"
        lines = [f"# Notes — {title}", ""]
        if query:
            lines.append(f"*Filter: {query} — {len(rows)}/{len(items)} shown*")
            lines.append("")
        for i, a in rows:
            first = (a.get("note", "") or "").splitlines()
            head = first[0] if first else "(empty)"
            lines.append(f"## [{i}] {head}")
            if a.get("anchor"):
                lines.append(f"> {a['anchor']}")
            meta = "  ".join(f"#{t}" for t in a.get("tags", []) or [])
            if a.get("cite"):
                meta += ("  " if meta else "") + f"@{a['cite']}"
            if meta:
                lines.append(f"`{meta}`")
            lines.append(a.get("note", ""))
            if a.get("ts"):
                lines.append(f"*{a['ts']}*")
            lines.append("")
        lines.append("---")
        lines.append(
            "M-x annotation-goto <n> · annotation-delete <n> · annotations-export"
        )
        self._show_text_pager("Notes", "\n".join(lines))

    def _annotations_search(self) -> None:
        self._enter_minibuffer(
            "Filter notes (text or #tag): ",
            on_commit=lambda q: self._annotations_list(q),
        )

    def _annotation_goto(self, arg: str) -> None:
        items = self._load_annotations()
        try:
            i = int(str(arg).strip())
        except (ValueError, TypeError):
            self.notify("Usage: annotation-goto <n>", error=True)
            return
        if not (0 <= i < len(items)):
            self.notify(f"No note #{i}", error=True)
            return
        wm = self.doc.word_map if self.doc else []
        wi = int(items[i].get("word_idx", 0) or 0)
        if wm and 0 <= wi < len(wm):
            self._scroll_to_line(wm[wi].disp_line)
            self.notify(f"Note #{i}: {items[i].get('note', '')[:50]}")
        else:
            self.notify("Note position unavailable", error=True)

    def _annotation_delete(self, arg: str) -> None:
        items = self._load_annotations()
        try:
            i = int(str(arg).strip())
        except (ValueError, TypeError):
            self.notify("Usage: annotation-delete <n>", error=True)
            return
        if not (0 <= i < len(items)):
            self.notify(f"No note #{i}", error=True)
            return
        del items[i]
        self._store_annotations(items)
        self.notify(f"Note #{i} deleted ({len(items)} left)")

    def _annotations_export(self) -> None:
        items = self._load_annotations()
        if not items:
            self.notify("No notes to export.", error=True)
            return
        p = Path(self.doc.path) if self.doc and self.doc.path else Path("notes")
        default = str(p.parent / (p.stem + "_notes.md"))
        self._enter_minibuffer(
            "Export notes to (.md/.json/.bib/.ris/.txt): ",
            initial=default,
            on_commit=self._annotations_export_cb,
        )

    def _annotations_export_cb(self, dest: str) -> None:
        dest = dest.strip()
        if not dest:
            return
        items = self._load_annotations()
        meta = getattr(self.doc, "metadata", {}) or {}
        title = self.doc.title if self.doc else "notes"
        author = meta.get("author") or meta.get("creator") or ""
        try:
            content = _format_annotations(
                items,
                Path(dest).suffix.lower(),
                title,
                author,
                self.doc.path if self.doc else "",
            )
            Path(dest).write_text(content, encoding="utf-8")
            self.notify(f"Exported {len(items)} note(s) → {dest}")
        except OSError as e:
            self.notify(f"Export error: {e}", error=True)

    # ── Saved note-filter presets ───────────────────────────────────────

    def _notes_presets(self) -> Dict[str, str]:
        return dict(self.settings.get("annotation_filter_presets", {}) or {})

    def _notes_preset_save(self, name: str, query: str) -> None:
        name = (name or "").strip()
        if not name:
            return
        presets = self._notes_presets()
        presets[name] = query
        self.settings.set("annotation_filter_presets", presets)

    # ── Inline prompt helpers (used inside the interactive notes browser) ──

    def _inline_prompt(self, prompt: str, initial: str = "") -> Optional[str]:
        """Read a line of text on the bottom row; return it, or None on Esc."""
        buf = list(initial)
        while True:
            h, w = self.scr.getmaxyx()
            _fillrow(self.scr, h - 1, self._a("minibuf"))
            shown = (prompt + "".join(buf))[: w - 1]
            _addstr(self.scr, h - 1, 0, shown, self._a("minibuf"))
            try:
                self.scr.move(h - 1, min(len(shown), w - 1))
            except curses.error:
                pass
            self.scr.refresh()
            ch = self.scr.getch()
            if ch in (10, 13, curses.KEY_ENTER):
                return "".join(buf)
            if ch in (27, 7):
                return None
            if ch in (curses.KEY_BACKSPACE, 127, 8):
                if buf:
                    buf.pop()
            elif 32 <= ch < 127:
                buf.append(chr(ch))

    def _inline_confirm(self, prompt: str) -> bool:
        """Show *prompt* on the bottom row; return True only on 'y'."""
        h, w = self.scr.getmaxyx()
        _fillrow(self.scr, h - 1, self._a("minibuf"))
        _addstr(self.scr, h - 1, 0, prompt[: w - 1], self._a("minibuf"))
        self.scr.refresh()
        return self.scr.getch() in (ord("y"), ord("Y"))

    # ── Interactive notes browser ───────────────────────────────────────

    def _notes_browser(self) -> None:
        """A dedicated, interactive notes mode for the TUI.

        Arrow keys / j,k select; Enter jumps to the note; r reads from it;
        e edits, d deletes, / filters, p cycles saved filter presets, s saves
        the current filter as a preset, x exports, q/Esc exits.
        """
        if not self.doc or not self.doc.word_map:
            self.notify("No document loaded.", error=True)
            return
        flt = ""
        sel = 0
        preset_names = list(self._notes_presets().keys())
        preset_idx = -1
        while True:
            all_items = self._load_annotations()
            rows = [
                (i, a) for i, a in enumerate(all_items) if _annotation_matches(a, flt)
            ]
            if rows:
                sel = max(0, min(sel, len(rows) - 1))
            else:
                sel = 0
            h, w = self.scr.getmaxyx()
            view_h = max(1, h - 2)
            top = 0 if sel < view_h else sel - view_h + 1
            self.scr.erase()
            hdr = f" Notes — {(self.doc.title or '')[:38]}  ({len(rows)}/{len(all_items)})"
            if flt:
                hdr += f"  filter: {flt}"
            _fillrow(self.scr, 0, self._a("title_bar"))
            _addstr(self.scr, 0, 0, hdr[: w - 1], self._a("title_bar"))
            if rows:
                for vi in range(top, min(len(rows), top + view_h)):
                    i, a = rows[vi]
                    note = (a.get("note", "") or "").splitlines()
                    head = note[0] if note else "(empty)"
                    tags = " ".join(f"#{t}" for t in a.get("tags", []) or [])
                    cite = f" @{a['cite']}" if a.get("cite") else ""
                    line = f"[{i}] {head}{('   ' + tags) if tags else ''}{cite}"
                    attr = curses.A_REVERSE if vi == sel else self._a("normal")
                    _addstr(self.scr, 1 + vi - top, 0, line[: w - 1].ljust(w - 1), attr)
            else:
                _addstr(
                    self.scr,
                    2,
                    2,
                    "No notes match.  /=filter  a=add (after exit)  q=quit",
                    self._a("dim"),
                )
            foot = (
                " ↑↓ move  ↵ jump  r read  e edit  d delete  / filter"
                "  p preset  s save  x export  q quit "
            )
            _fillrow(self.scr, h - 1, self._a("status"))
            _addstr(self.scr, h - 1, 0, foot[: w - 1], self._a("status"))
            self.scr.refresh()

            ch = self.scr.getch()
            if ch in (ord("q"), ord("Q"), 27, 7):
                return
            elif ch in (curses.KEY_DOWN, ord("j")):
                sel = min(sel + 1, max(0, len(rows) - 1))
            elif ch in (curses.KEY_UP, ord("k")):
                sel = max(sel - 1, 0)
            elif ch == curses.KEY_NPAGE:
                sel = min(sel + view_h, max(0, len(rows) - 1))
            elif ch == curses.KEY_PPAGE:
                sel = max(sel - view_h, 0)
            elif ch == curses.KEY_HOME:
                sel = 0
            elif ch == curses.KEY_END:
                sel = max(0, len(rows) - 1)
            elif ch in (10, 13, curses.KEY_ENTER) and rows:
                self._annotation_goto(str(rows[sel][0]))
                return
            elif ch == ord("r") and rows:
                self._tts_play_from_word(int(rows[sel][1].get("word_idx", 0) or 0))
            elif ch == ord("e") and rows:
                i = rows[sel][0]
                new = self._inline_prompt(f"Edit [{i}]: ", rows[sel][1].get("note", ""))
                if new is not None and new.strip():
                    items = self._load_annotations()
                    if 0 <= i < len(items):
                        items[i]["note"] = new.strip()
                        items[i]["ts"] = time.strftime("%Y-%m-%dT%H:%M:%S")
                        self._store_annotations(items)
            elif ch == ord("d") and rows:
                i = rows[sel][0]
                if self._inline_confirm(f"Delete note [{i}]? (y/n) "):
                    self._annotation_delete(str(i))
                    sel = max(0, sel - 1)
            elif ch == ord("/"):
                res = self._inline_prompt("Filter (text or #tag): ", flt)
                if res is not None:
                    flt = res.strip()
                    sel = 0
            elif ch == ord("p"):
                preset_names = list(self._notes_presets().keys())
                if preset_names:
                    preset_idx = (preset_idx + 1) % len(preset_names)
                    name = preset_names[preset_idx]
                    flt = self._notes_presets()[name]
                    sel = 0
                    self.notify(f"Preset: {name}  ({flt})")
                else:
                    self.notify("No saved presets. Press s to save one.")
            elif ch == ord("s"):
                name = self._inline_prompt("Save current filter as preset: ")
                if name and name.strip():
                    self._notes_preset_save(name, flt)
                    self.notify(f"Saved preset '{name.strip()}'")
            elif ch == ord("x"):
                return self._annotations_export()

    # ── Main run loop ──────────────────────────────────────────────────

    def run(self) -> None:
        while self._running:
            try:
                self._poll_load_queue()
                self.draw()
                ch = self.scr.getch()
                if ch == -1:
                    continue
                self._handle_key(ch)
            except KeyboardInterrupt:
                break
            except Exception:
                pass
        self.tts.stop()
        self._save_reading_position()  # remember where we stopped
        self.settings.save()

    # ── Key handling ───────────────────────────────────────────────────────

    def _handle_key(self, ch: int) -> None:
        if self.mode == "sc":
            self._handle_sc_key(ch)
        elif self.mode == "mx":
            self._handle_mx_key(ch)
        elif self.mode == "search":
            self._handle_search_key(ch)
        elif self.mode == "goto":
            self._handle_goto_key(ch)
        else:
            self._handle_normal_key(ch)

    def _handle_normal_key(self, ch: int) -> None:  # noqa: C901
        # ── Escape: stop speech and clear search ────────────────────────────
        if ch == 27:
            # Brief peek — keep ESC-x working as a silent power-user shortcut.
            self.scr.timeout(100)
            nk = self.scr.getch()
            self.scr.timeout(150)
            if nk in (ord("x"), ord("X")):
                self._enter_minibuffer("Command: ")
            else:
                self._tts_stop()
                self.search.query = ""
            return

        # ── Function keys ──────────────────────────────────────────────────
        if ch == curses.KEY_F1:
            self._show_help()
            return
        if ch == curses.KEY_F2:
            self._enter_minibuffer("Command: ")
            return
        if ch == curses.KEY_F3:
            self._search_next()
            return
        if ch == curses.KEY_F4:
            self._search_prev()
            return
        if ch == curses.KEY_F5:
            self._next_theme()
            return
        if ch == curses.KEY_F6:
            self.settings.set(
                "show_line_numbers", not self.settings["show_line_numbers"]
            )
            return
        if ch == curses.KEY_F7:
            self.settings.set("syntax_highlight", not self.settings["syntax_highlight"])
            self._render_doc()
            return
        if ch == curses.KEY_F8:
            self._cycle_speed_preset()  # cycle skim / normal / study / slow
            return
        if ch == curses.KEY_F9:
            if self.doc:
                self._open_async(self.doc.path)
            return
        if ch == curses.KEY_RESIZE:
            self._render_doc()
            return

        # ── Navigation ─────────────────────────────────────────────────────
        if ch == curses.KEY_DOWN:
            self._scroll_by(1)
        elif ch == curses.KEY_UP:
            self._scroll_by(-1)
        elif ch == curses.KEY_NPAGE:
            self._page_down()
        elif ch == curses.KEY_PPAGE:
            self._page_up()
        elif ch == curses.KEY_HOME:
            self._goto_top()
        elif ch == curses.KEY_END:
            self._goto_bottom()
        # j / k kept as silent shortcuts familiar to terminal users
        elif ch == ord("j"):
            self._scroll_by(1)
        elif ch == ord("k"):
            self._scroll_by(-1)

        # ── Speech ─────────────────────────────────────────────────────────
        elif ch == ord(" "):
            self._tts_toggle()
        elif ch in (ord("+"), ord("=")):  # speed up
            self._rate_change(+20)
        elif ch == ord("-"):  # slow down
            self._rate_change(-20)

        # ── File operations ────────────────────────────────────────────────
        elif ch == 15:  # Ctrl+O — open
            self._open_file_prompt()
        elif ch == 19:  # Ctrl+S — save/export
            self._export_markdown()
        elif ch in (17, ord("q"), ord("Q")):  # Ctrl+Q or q — quit
            self._running = False

        # ── Search ─────────────────────────────────────────────────────────
        elif ch == 6:  # Ctrl+F — find
            self._enter_minibuffer(
                "Search: ",
                mode="search",
                on_commit=lambda q: self._do_search(q, "forward"),
            )
        elif ch == ord("n"):
            self._search_next()
        elif ch == ord("N"):
            self._search_prev()

        # ── Sentence navigation ─────────────────────────────────────────
        # Alt+. / Alt+, / Alt+; match the Qt GUI sentence shortcuts.
        # Plain .  ,  ;  are kept as fallback (TUI muscle memory).
        elif ch == ord("."):
            self._skip_next_sentence()  # or Alt+.
        elif ch == ord(","):
            self._skip_prev_sentence()  # or Alt+,
        elif ch == ord(";"):
            self._replay_sentence()  # or Alt+;

        # ── Paragraph navigation ────────────────────────────────────────
        # p / P  — NVDA browse-mode convention, aligns with GUI Ctrl+P / Ctrl+Shift+P.
        # Ctrl+P (16) also added for direct GUI parity.
        # ]  [  kept as silent fallbacks.
        elif ch in (ord("p"), 16):
            self._skip_next_paragraph()  # p  Ctrl+P
        elif ch == ord("P"):
            self._skip_prev_paragraph()  # P  (Shift)
        elif ch == ord("]"):
            self._skip_next_paragraph()  # legacy
        elif ch == ord("["):
            self._skip_prev_paragraph()  # legacy
        # r / Ctrl+R  — replay paragraph; Ctrl+R (18) matches GUI Ctrl+R.
        elif ch in (ord("r"), 18):
            self._replay_paragraph()  # r  Ctrl+R

        # ── Heading navigation ──────────────────────────────────────────
        # h  — NVDA browse-mode convention, aligns with GUI Ctrl+H (next heading).
        # }  {  kept as silent fallbacks.  >  <  always-play variants.
        elif ch == ord("h"):
            self._skip_next_heading()  # h  (forward)
        elif ch == ord("}"):
            self._skip_next_heading()  # } legacy
        elif ch == ord("{"):
            self._skip_prev_heading()  # { legacy
        elif ch == ord(">"):
            self._read_next_heading()  # >  (always play)
        elif ch == ord("<"):
            self._read_prev_heading()  # <  (always play)

        # ── Table navigation ───────────────────────────────────────────
        # t / T  — NVDA browse-mode convention, aligns with GUI Ctrl+T / Ctrl+Shift+T.
        elif ch == ord("t"):
            self._skip_next_table()
        elif ch == ord("T"):
            self._skip_prev_table()

        # ── Chapter navigation ──────────────────────────────────────────────
        elif ch == curses.KEY_F10:
            self._chapter_prev()
        elif ch == curses.KEY_F11:
            self._chapter_next()
        elif ch == curses.KEY_F12:
            self._chapter_list()

        # ── Navigation history ──────────────────────────────────────────────
        elif ch == ord("H"):  # H = history back (capital to avoid conflict)
            self._history_back()
        elif ch == ord("L"):  # L = history forward
            self._history_forward()

        # ── Clipboard ────────────────────────────────────────────────────────
        elif ch == 3:  # Ctrl+C
            self._copy_to_clipboard()
        # ── Speech Cursor mode ────────────────────────────────────────────
        elif ch == 9:  # Tab — enter Speech Cursor mode
            self._sc_enter()

        # ── Voice picker (Ctrl+T) ────────────────────────────────────────
        elif ch == 20:  # Ctrl+T — T for TTS voice
            self._voice_picker()

        # ── Immediate speech stop (Ctrl+X or Ctrl+Space) ─────────────────
        elif ch in (0, 24):  # Ctrl+Space / Ctrl+X
            self._tts_stop()
            self.notify("Speech stopped")

        # ── Annotations / notes ─────────────────────────────────────────────
        elif ch == ord("a"):  # add a note at the reading position
            self._annotate()
        elif ch == ord("A"):  # interactive notes browser
            self._notes_browser()

        # ── Keyboard cheat sheet ────────────────────────────────────────────
        elif ch == ord("?"):
            self._show_shortcuts()

        # ── Command palette ────────────────────────────────────────────────
        elif ch == ord(":"):
            self._enter_minibuffer("Command: ")

    def _handle_mx_key(self, ch: int) -> None:
        if ch in (curses.KEY_ENTER, 10, 13):
            cmd = self.mx_ed.value.strip()
            if cmd:
                if not self.mx_history or self.mx_history[-1] != cmd:
                    self.mx_history.append(cmd)
                    if len(self.mx_history) > 200:
                        self.mx_history.pop(0)
            self.mode = "normal"
            self.mx_hist_pos = -1
            if cmd:
                cb = getattr(self, "_mx_callback", None)
                if cb:
                    cb(cmd)
            return
        if ch in (7, 27):
            self._cancel_minibuffer()
            return
        if ch == 9:
            self._mx_tab()
            return
        if ch == curses.KEY_UP:
            if self.mx_history:
                if self.mx_hist_pos < len(self.mx_history) - 1:
                    self.mx_hist_pos += 1
                self.mx_ed.set_value(self.mx_history[-(self.mx_hist_pos + 1)])
            return
        if ch == curses.KEY_DOWN:
            if self.mx_hist_pos > 0:
                self.mx_hist_pos -= 1
                self.mx_ed.set_value(self.mx_history[-(self.mx_hist_pos + 1)])
            elif self.mx_hist_pos == 0:
                self.mx_hist_pos = -1
                self.mx_ed.set_value("")
            return
        r = self.mx_ed.feed(ch)
        if r is True:
            self.mx_comp_idx = -1
            self._mx_update_completions()

    def _handle_search_key(self, ch: int) -> None:
        if ch in (curses.KEY_ENTER, 10, 13):
            q = self.search_ed.value.strip()
            self.mode = "normal"
            cb = getattr(self, "_mx_callback", None)
            if cb and q:
                cb(q)
            return
        if ch in (7, 27):
            self._cancel_minibuffer()
            return
        self.search_ed.feed(ch)

    def _handle_goto_key(self, ch: int) -> None:
        if ch in (curses.KEY_ENTER, 10, 13):
            s = self.goto_ed.value.strip()
            self.mode = "normal"
            self._goto_line_cb(s)
            return
        if ch in (7, 27):
            self._cancel_minibuffer()
            return
        self.goto_ed.feed(ch)

    # ── Drawing ────────────────────────────────────────────────────────────

    def draw(self) -> None:
        h, w = self.scr.getmaxyx()
        if h < 8 or w < 20:
            self.scr.erase()
            _addstr(self.scr, 0, 0, "Terminal too small (need 20×8 minimum)")
            self.scr.refresh()
            return
        self.scr.erase()
        self._draw_title(h, w)
        self._draw_document(h, w)
        self._draw_status(h, w)
        self._draw_minibuffer(h, w)
        # ── Cursor positioning for screen-reader accessibility ─────────────
        # In input modes (_draw_minibuffer already moved the cursor to the
        # insertion point).  In normal reading mode the cursor must sit on the
        # document text so that terminal screen readers (NVDA, JAWS, Orca …)
        # can follow the reading position rather than being permanently locked
        # onto the minibuffer row at the bottom of the screen.
        try:
            if self.mode not in ("mx", "search", "goto"):
                view_top = 1  # row 0 is the title bar
                view_h = max(1, h - 3 - view_top)
                if (
                    self.mode == "sc"
                    and self.scroll <= self._sc_line < self.scroll + view_h
                ):
                    # In SC mode the cursor sits on the reading-cursor line.
                    cur_row = view_top + (self._sc_line - self.scroll)
                elif (
                    self._highlight_line >= 0
                    and self.scroll <= self._highlight_line < self.scroll + view_h
                ):
                    # Cursor tracks the currently spoken word line.
                    cur_row = view_top + (self._highlight_line - self.scroll)
                else:
                    # Idle / paused: sit on the first visible document line.
                    cur_row = view_top
                self.scr.move(cur_row, 0)
        except curses.error:
            pass
        self.scr.refresh()

    def _draw_title(self, h: int, w: int) -> None:
        """Top title bar: app name | document title | TTS status | rate."""
        title = self.doc.title if self.doc else APP_TITLE
        if self.mode == "sc":
            tts_state = "▶ SC+Speaking" if self.tts.speaking else "● SC CURSOR"
        else:
            tts_state = "▶ Speaking" if self.tts.speaking else "■ Stopped"
        rate = str(self.settings["tts_rate"])
        rhs = f" {tts_state}  {rate} wpm  {self.tts.backend_name} "
        lhs = f" {APP_NAME}  │  {title} "
        gap = max(1, w - len(lhs) - len(rhs) - 1)
        bar = lhs + " " * gap + rhs
        _fillrow(self.scr, 0, self._a("title_bar"))
        _addstr(self.scr, 0, 0, bar[: w - 1], self._a("title_bar"))

    def _draw_document(self, h: int, w: int) -> None:
        """Render visible document lines into the content area (rows 1 … h-3)."""
        view_top = 1
        view_bottom = h - 3
        view_h = max(1, view_bottom - view_top)

        # Scroll to keep the current speech position visible.
        #
        # We track the *callback-confirmed* word position rather than the
        # timer's visual highlight, because the timer can race ahead of the
        # audio (engine-startup lag, SSML pauses).  When the user triggers a
        # navigation command (replay-sentence, skip, etc.) the destination is
        # also derived from the callback position, so the viewport is already
        # close to the destination — no dramatic snap-back.
        # In SSML mode (or before the first callback fires) no confirmed
        # position is available; we fall back to _highlight_line so the
        # screen still scrolls during reading.
        cb = self.tts.last_cb_word_idx
        if cb >= 0 and self.doc and cb < len(self.doc.word_map):
            _scroll_line = self.doc.word_map[cb].disp_line
        elif self._highlight_line >= 0:
            _scroll_line = self._highlight_line
        else:
            _scroll_line = -1

        if _scroll_line >= 0:
            margin = int(self.settings["scroll_margin"])
            if _scroll_line < self.scroll + margin:
                self.scroll = max(0, _scroll_line - margin)
            elif _scroll_line >= self.scroll + view_h - margin:
                self.scroll = max(0, _scroll_line - view_h + margin + 1)

        if self.loading:
            mid = view_top + view_h // 2
            _addstr(self.scr, mid, 4, self.loading_msg, self._a("progress"))
            return

        if not self.rendered:
            welcome = _WELCOME_TEXT.splitlines()
            for i, ln in enumerate(welcome[:view_h]):
                role = (
                    "h1"
                    if ln.startswith("# ")
                    else ("h2" if ln.startswith("## ") else "normal")
                )
                _addstr(self.scr, view_top + i, 2, ln, self._a(role))
            return

        total = len(self.rendered)
        cur_match = self.search.current_match

        # SC mode: remember which display-line the reading cursor is on so
        # the inner loop can highlight it.
        sc_cursor_row: int = -1
        if self.mode == "sc" and self.rendered:
            visible_sc = self.scroll <= self._sc_line < self.scroll + view_h
            if visible_sc:
                sc_cursor_row = view_top + (self._sc_line - self.scroll)

        for vi in range(view_h):
            li = self.scroll + vi
            row = view_top + vi
            if li >= total:
                break

            segs = self.rendered[li]
            col = 0
            show_ln = bool(self.settings["show_line_numbers"])
            if show_ln:
                ln_str = f"{li + 1:>4} "
                _addstr(self.scr, row, 0, ln_str, self._a("dim"))
                col = len(ln_str)

            for text, role in segs:
                if not text or col >= w - 1:
                    break
                attr = self._a(role)

                # Apply search highlighting character-by-character if needed
                if cur_match and cur_match[0] == li:
                    _, cs, ce = cur_match
                    self._draw_highlighted_text(
                        row, col, text, role, cs, ce, w, current=True
                    )
                elif any(m[0] == li for m in self.search.matches):
                    for m in self.search.matches:
                        if m[0] == li:
                            self._draw_highlighted_text(
                                row, col, text, role, m[1], m[2], w, current=False
                            )
                            break
                else:
                    # TTS current-word highlight
                    if (
                        li == self._highlight_line
                        and self.settings["highlight_current_word"]
                        and self._highlight_col_start >= 0
                    ):
                        hcs = self._highlight_col_start - col
                        hce = self._highlight_col_end - col
                        txt = text[: w - col - 1]
                        for ci, c in enumerate(txt):
                            tpos = col + ci
                            if hcs <= ci < hce:
                                _addstr(self.scr, row, tpos, c, self._a("current_word"))
                            else:
                                _addstr(self.scr, row, tpos, c, attr)
                        col += len(txt)
                    else:
                        avail = max(0, w - col - 1)
                        chunk = text[:avail]
                        _addstr(self.scr, row, col, chunk, attr)
                        col += len(chunk)

        # Scroll indicators
        if self.scroll > 0:
            _addstr(self.scr, view_top, w - 4, " ▲ ", self._a("dim"))
        if self.scroll + view_h < total:
            _addstr(self.scr, view_bottom - 1, w - 4, " ▼ ", self._a("dim"))

        # ── SC mode cursor bar ────────────────────────────────────────────────
        # Draw a full-width reverse-video bar over the SC cursor line so the
        # reading position is clearly visible even while the word highlight
        # is on a different line.
        if sc_cursor_row >= 0:
            try:
                self.scr.chgat(sc_cursor_row, 0, -1, curses.A_REVERSE)
            except curses.error:
                pass

    def _draw_highlighted_text(
        self,
        row: int,
        base_col: int,
        text: str,
        role: str,
        hl_start: int,
        hl_end: int,
        w: int,
        current: bool,
    ) -> None:
        hl_attr = self._a("search_current" if current else "search_match")
        norm_attr = self._a(role)
        col = base_col
        for i, c in enumerate(text[: w - base_col - 1]):
            tpos = base_col + i
            attr = hl_attr if hl_start <= tpos < hl_end else norm_attr
            _addstr(self.scr, row, tpos, c, attr)
            col += 1

    def _draw_status(self, h: int, w: int) -> None:
        """Status bar (second-to-last row) and hints (third-to-last row)."""
        status_row = h - 3
        hints_row = h - 2

        # Timed message
        if self.message and (time.monotonic() - self.message_t) < self.message_dur:
            _fillrow(self.scr, status_row, self._a("status"))
            _addstr(
                self.scr, status_row, 0, f" {self.message}"[: w - 1], self._a("status")
            )
        else:
            self.message = ""
            total = len(self.rendered)
            pct = int(100 * (self.scroll + 1) / max(1, total)) if total else 100
            search_info = (
                f"  [{self.search.match_index + 1}/{self.search.match_count}]"
                if self.search.match_count
                else ""
            )
            bar = (
                (
                    f" {self.doc.title[:40] if self.doc else 'No document'}  "
                    f"Line {self.scroll + 1}/{total}  {pct}%"
                    f"{search_info}"
                )
                if self.doc
                else f" {APP_TITLE}"
            )
            _fillrow(self.scr, status_row, self._a("status"))
            _addstr(self.scr, status_row, 0, bar[: w - 1], self._a("status"))

        # Hints bar
        hints = (
            "  Space:play/pause  Tab:speech-cursor  Ctrl+T:voice  Ctrl+X:stop  "
            ",/.:sent  [/]:para  {/}:head-scroll  </>:read-head  "
            ";:replay-sent  r:replay-para  "
            "Ctrl-F:search  +/-:speed  F2:commands  F1:help  Ctrl-Q:quit"
        )
        _fillrow(self.scr, hints_row, self._a("dim"))
        _addstr(self.scr, hints_row, 0, hints[: w - 1], self._a("dim"))

    def _draw_minibuffer(self, h: int, w: int) -> None:
        """Bottom minibuffer row."""
        mb_row = h - 1
        if self.mode == "mx":
            prompt = getattr(self, "_mx_prompt", "M-x: ")
            ed = self.mx_ed
            val = ed.value
            comps = self.mx_completions[:6]
            comp_str = (
                "  "
                + "  ".join(
                    f"[{c}]" if i == self.mx_comp_idx else c
                    for i, c in enumerate(comps)
                )
                if comps
                else ""
            )
            full = prompt + val + comp_str
            _fillrow(self.scr, mb_row, self._a("minibuf"))
            _addstr(self.scr, mb_row, 0, (prompt + val)[: w - 1], self._a("minibuf"))
            if comp_str and len(prompt + val) < w - 2:
                _addstr(
                    self.scr,
                    mb_row,
                    len(prompt + val),
                    comp_str[: w - len(prompt) - len(val) - 1],
                    self._a("dim"),
                )
            try:
                cx = min(len(prompt) + ed.pos, w - 1)
                self.scr.move(mb_row, cx)
            except curses.error:
                pass
        elif self.mode == "search":
            prompt = getattr(self, "_search_prompt", "Search: ")
            ed = self.search_ed
            _fillrow(self.scr, mb_row, self._a("minibuf"))
            _addstr(
                self.scr, mb_row, 0, (prompt + ed.value)[: w - 1], self._a("minibuf")
            )
            try:
                self.scr.move(mb_row, min(len(prompt) + ed.pos, w - 1))
            except curses.error:
                pass
        elif self.mode == "goto":
            prompt = "Go to line: "
            ed = self.goto_ed
            _fillrow(self.scr, mb_row, self._a("minibuf"))
            _addstr(
                self.scr, mb_row, 0, (prompt + ed.value)[: w - 1], self._a("minibuf")
            )
            try:
                self.scr.move(mb_row, min(len(prompt) + ed.pos, w - 1))
            except curses.error:
                pass
        else:
            _fillrow(self.scr, mb_row, self._a("dim"))
            if self.mode == "sc":
                idle = (
                    "  SC CURSOR  \u2191\u2193:line  ,/.:sent  [/]:para  {/}:head"
                    "  t/T:table  Enter:read-on  Space:pause  Esc:exit  Tab:normal"
                )
            else:
                idle = (
                    f"  F2:commands  Ctrl-O:open  Space:play/pause"
                    f"  Tab:speech-cursor  F1:help  Esc:stop  \u2502  {self.tts.backend_name}"
                )
            _addstr(self.scr, mb_row, 0, idle[: w - 1], self._a("dim"))


# =============================================================================
# Embedded help text
# =============================================================================

_HELP_TEXT = (
    """\
# star — Speaking Terminal Access Reader"""
    + APP_VERSION
    + """

star is a reading application with built-in text-to-speech designed for
students with print disabilities.

---

## Quick Start

Open a file:          `Ctrl+O`  or  `star document.pdf`
Start / pause:        `Space`
Stop reading:         `Esc`
Scroll up / down:     Arrow keys or Page Up / Page Down
Search:               `Ctrl+F`    then `F3` / `Shift+F3` to step through hits
Commands:             `F2`
Quit:                 `Ctrl+Q`  or  `q`

---

## Navigation

| Key | Action |
|---|---|
| `↑` / `↓` | Scroll one line |
| `Page Down` | Next page |
| `Page Up` | Previous page |
| `Home` | Beginning of document |
| `End` | End of document |

---

## Skip Navigation — Fine to Coarse

All skip keys **restart speech automatically** if reading is already in
progress — no need to press Space again after skipping.

| Key | Granularity | Action |
|---|---|---|
| `,` | Sentence | Jump to previous sentence * |
| `.` | Sentence | Jump to next sentence |
| `;` | Sentence | **Replay** current sentence from its beginning |
| `r` | Paragraph | **Replay** current paragraph from its beginning |
| `[` | Paragraph | Jump to previous paragraph |
| `]` | Paragraph | Jump to next paragraph |
| `{` (Shift+[) | Heading | Scroll to previous heading (resume if playing) |
| `}` (Shift+]) | Heading | Scroll to next heading (resume if playing) |
| `<` (Shift+,) | Heading | **Read** from previous heading (always starts TTS) |
| `>` (Shift+.) | Heading | **Read** from next heading (always starts TTS) |

---

## Speech Cursor Mode (`Tab`)

Press `Tab` to enter **Speech Cursor mode** — a dedicated navigation mode
where every movement key also starts TTS reading from the new position.
This lets you browse a document by unit (line, sentence, paragraph, heading,
table) and hear each one immediately without holding a separate play key.

| Key | Action |
|---|---|
| `↑` / `k` | Previous line — read it |
| `↓` / `j` | Next line — read it |
| `,` | Previous sentence — read it |
| `.` | Next sentence — read it |
| `[` / Page Up | Previous paragraph — read it |
| `]` / Page Down | Next paragraph — read it |
| `{` / `<` | Previous heading — read it |
| `}` / `>` | Next heading — read it |
| `t` | Next table — read it |
| `T` | Previous table — read it |
| `r` | Re-read current cursor line |
| `Space` | Pause / resume |
| `Ctrl+X` | Stop speech (cursor stays) |
| `Enter` | Exit SC mode and start continuous reading from cursor |
| `Esc` | Exit SC mode and stop speech |
| `Tab` | Exit SC mode, leave speech running |

The SC cursor is shown as a full-width reverse-video bar.  The title bar
shows `● SC CURSOR` when the mode is active.

## Instant Speech Stop

`Ctrl+X` (or `Ctrl+Space`) stops all TTS output immediately from any mode.

Note: If you are more than 3 words into the current sentence, `,` replays
the current sentence first; press it again to go to the previous one.

The status bar shows a preview of the sentence you land on (first
5 words) and its number within the document.

All commands are also available in the command palette (`F2`):
`next-sentence`, `prev-sentence`, `replay-sentence`, `replay-paragraph`,
`next-paragraph`, `prev-paragraph`, `next-heading`, `prev-heading`,
`read-next-heading`, `read-prev-heading`.

---

## Speech

| Key | Action |
|---|---|
| `Space` | Play / pause from the current position |
| `Esc` | Stop speech and clear search |
| `+` or `=` | Speed up (+20 wpm) |
| `-` | Slow down (−20 wpm) |

The word currently being spoken is highlighted in the document.

---

## File Operations

| Key | Action |
|---|---|
| `Ctrl+O` | Open a file (Tab completes the path) |
| `Ctrl+S` | Export document as a Markdown file |
| `Ctrl+Q` or `q` | Quit |
| `F9` | Reload the current document |

---

## Search

| Key | Action |
|---|---|
| `Ctrl+F` | Open search box |
| `F3` | Jump to next match |
| `F4` | Jump to previous match |
| `n` / `N` | Next / previous match (when no text box is open) |
| `Esc` | Clear search and stop speech |

---

## Chapter Navigation  (for EPUB, DAISY, and long documents)

| Key | Action |
|---|---|
| `F10` | Previous chapter |
| `F11` | Next chapter |
| `F12` | List all chapters |
| `H` | History back |
| `L` | History forward |

Also available in the command palette: `chapter-next`, `chapter-prev`, `chapter-list`, `chapter-goto`, `history-back`, `history-forward`, `bookmark-set`, `bookmark-goto`, `bookmark-list`.

---

## Command Palette  (press `F2` to open)

Type any command name and press **Enter**.  Press **Tab** to cycle through
completions.  Use `↑` / `↓` to recall previous commands.  Press **Esc** to cancel.

| Command | Description |
|---|---|
| `open` | Open a file |
| `open-url` | Open a URL |
| `export-markdown` | Export document as Markdown |
| `export-braille` | Export as BRF braille file (requires liblouis) |
| `export-audio [fmt]` | Export TTS audio as MP3/OGG/MP4/WAV (requires ffmpeg or pydub) |
| `play` | Start reading |
| `stop` | Stop reading |
| `pause` | Toggle play/pause |
| `speak-line` | Speak the current line |
| `next-paragraph` | Jump to next paragraph (restarts speech) |
| `prev-paragraph` | Jump to previous paragraph (restarts speech) |
| `next-heading` | Jump to next heading (restarts speech) |
| `prev-heading` | Jump to previous heading (restarts speech) |
| `rate-up` | Increase reading speed |
| `rate-down` | Decrease reading speed |
| `volume-up` | Increase volume |
| `volume-down` | Decrease volume |
| `tts-backend` | Switch TTS engine |
| `tts-voice` | Set TTS voice |
| `search` | Search forward |
| `search-backward` | Search backward |
| `goto-line` | Jump to line number |
| `theme [name]` | Switch color theme |
| `line-numbers` | Toggle line numbers |
| `syntax-highlight` | Toggle syntax highlighting |
| `wrap-width` | Set text wrap width |
| `reload` | Reload current document |
| `settings` | Show settings file path |
| `help` | Show this help |
| `about` | Version information |
| `license` | Show license |
| `quit` | Quit star |

---

## Themes  (`F5` to cycle, or `F2` then `theme dark`)

| Name | Description |
|---|---|
| `dark` | Modern dark — blue/cyan/magenta accents (default) |
| `light` | Light background with blue/magenta accents |
| `contrast` | High contrast — bold white on black |
| `phosphor` | Monochrome green phosphor terminal |

---

## Supported File Formats

PDF, DOCX, ODT, EPUB, HTML, DAISY/DTBook, Markdown, plain text,
CSV/TSV (rendered as tables), XLSX (rendered as tables),
LaTeX, Org-mode, Jupyter Notebook, R, R Markdown, Python, Rust, C/C++,
JavaScript, URLs (http/https).

OCR of image-based PDFs and image files (PNG, JPEG, etc.) is supported
when `pytesseract` and `pymupdf` are installed.

---

## Screen Reader & Braille

star works with NVDA, JAWS, Orca, and VoiceOver.  The terminal cursor
is always parked at the start of the minibuffer so screen readers
following the hardware cursor will track the correct position.

Braille display output via BrlTTY is automatic on Linux; on Windows,
NVDA and JAWS drive Braille displays via their built-in support.

BRF (Braille Ready Format) export: open the command palette (`F2`) and
type `export-braille`  (requires `pip install louis`)

---

Press `q` or `Esc` to close this help screen.
"""
)

_LICENSE_TEXT = (
    """\
# star — Speaking Terminal Access Reader

"""
    + __copyright__
    + """
License: """
    + __license__
    + """

---

## GNU General Public License v3

This program is free software: you can redistribute it and/or modify
it under the terms of the GNU General Public License as published by
the Free Software Foundation, either version 3 of the License, or
(at your option) any later version.

This program is distributed in the hope that it will be useful,
but WITHOUT ANY WARRANTY; without even the implied warranty of
MERCHANTABILITY or FITNESS FOR A PARTICULAR PURPOSE.  See the
GNU General Public License for more details.

Full license text: https://www.gnu.org/licenses/gpl-3.0.txt

---

Press `q` or `Esc` to close.
"""
)

_WELCOME_TEXT = (
    """\
# star — Speaking Terminal Access Reader

Version """
    + APP_VERSION
    + """   """
    + __copyright__
    + """

---

## Getting started

  **Open a file:**   Ctrl+O  (or pass a filename on the command line)
  **Open a URL:**    F2, then type  open-url
  **Read aloud:**    Space   (plays from the current position)
  **Stop:**          Esc
  **Help:**          F1
  **Commands:**      F2
  **Quit:**          Ctrl+Q  or  q

## Supported formats

  Plain Text, Markdown, LaTeX, PDF, DOCX, ODT, HTML,
  CSV, TSV, XLSX, EPUB, DAISY, DTBook, R, Python,
  Jupyter Notebook, Org-mode, URLs

## TTS backends

  pyttsx3 (built-in system voice) · eSpeak-NG · DECtalk

  Switch:  M-x tts-backend
"""
)


# =============================================================================
# CSS theme helpers (Qt GUI)
# =============================================================================

# CSS that is generated from a built-in palette dict.  Used both to write
# the seed files and as the fallback when no custom CSS overrides a palette key.
_CSS_TEMPLATE = """\
body {{
    background: {bg};
    color: {fg};
    font-family: Georgia, serif;
    margin: 14px;
    line-height: 1.6;
}}
::selection {{
    background: {sel};
}}
h1 {{ color: {h1}; margin: 8px 0 4px; }}
h2 {{ color: {h2}; margin: 6px 0 3px; }}
h3 {{ color: {h3}; margin: 4px 0 2px; }}
h4 {{ color: {h4}; margin: 4px 0 2px; }}
code {{ color: {code}; font-family: monospace; }}
pre  {{ color: {code}; font-family: monospace; white-space: pre-wrap; }}
b    {{ font-weight: bold; }}
i    {{ font-style: italic; }}
p    {{ margin: 4px 0; }}
"""


def _palette_to_css(pal: Dict[str, str]) -> str:
    """Format *_CSS_TEMPLATE* with the values from a palette dict."""
    return _CSS_TEMPLATE.format(**pal)


def _parse_css_palette(css: str) -> Dict[str, str]:
    """Extract the 8 palette keys (bg/fg/sel/h1–h4/code) from a CSS string.

    Only plain color values (hex codes, named colors, or rgb/rgba) are
    extracted; any key that cannot be found falls back to the 'dark' palette
    so the result is always a complete, valid palette dict.
    """
    # Start with dark-theme defaults so every key is always present.
    from copy import deepcopy as _dc

    _DARK = {
        "bg": "#16181d",
        "fg": "#c6ccd4",
        "sel": "#2c313a",
        "h1": "#82aaff",
        "h2": "#89ddff",
        "h3": "#c792ea",
        "h4": "#f78c6c",
        "code": "#7fdbab",
    }
    pal = _dc(_DARK)

    _val = r"([#\w][\w(),. %]*)"  # loose color value pattern

    # body { background: X; color: Y; }
    m = re.search(r"body\s*\{([^}]*)\}", css, re.DOTALL | re.IGNORECASE)
    if m:
        blk = m.group(1)
        bg = re.search(r"background(?:-color)?\s*:\s*" + _val, blk, re.I)
        fg = re.search(r"(?<!background-)color\s*:\s*" + _val, blk, re.I)
        if bg:
            pal["bg"] = bg.group(1).strip().rstrip(";")
        if fg:
            pal["fg"] = fg.group(1).strip().rstrip(";")

    # ::selection { background: Z; }
    m = re.search(r"::?selection\s*\{([^}]*)\}", css, re.DOTALL | re.IGNORECASE)
    if m:
        sel = re.search(r"background(?:-color)?\s*:\s*" + _val, m.group(1), re.I)
        if sel:
            pal["sel"] = sel.group(1).strip().rstrip(";")

    # h1 { color: X; }  …  h4 { color: X; }
    for lvl in range(1, 5):
        m = re.search(rf"h{lvl}\s*\{{([^}}]*)\}}", css, re.DOTALL | re.IGNORECASE)
        if m:
            c = re.search(r"(?<!background-)color\s*:\s*" + _val, m.group(1), re.I)
            if c:
                pal[f"h{lvl}"] = c.group(1).strip().rstrip(";")

    # code { color: X; }
    m = re.search(r"\bcode\b\s*\{([^}]*)\}", css, re.DOTALL | re.IGNORECASE)
    if m:
        c = re.search(r"(?<!background-)color\s*:\s*" + _val, m.group(1), re.I)
        if c:
            pal["code"] = c.group(1).strip().rstrip(";")

    return pal


def _load_css_themes() -> Dict[str, Dict[str, Any]]:
    """Scan *THEMES_DIR* for *.css files and return a mapping
    theme-name → palette-dict (with key '_css' holding the raw CSS text).
    """
    result: Dict[str, Dict[str, Any]] = {}
    if not THEMES_DIR.exists():
        return result
    for css_path in sorted(THEMES_DIR.glob("*.css")):
        try:
            css_text = css_path.read_text(encoding="utf-8", errors="replace")
            pal: Dict[str, Any] = _parse_css_palette(css_text)
            pal["_css"] = css_text
            result[css_path.stem] = pal
        except Exception:  # noqa: BLE001  — skip unreadable files silently
            pass
    return result


def _seed_default_css_themes() -> None:
    """Write the four built-in palettes as CSS files in *THEMES_DIR*.

    Files are only written if they do not already exist, so hand-edited
    customizations are never overwritten.  Each generated file serves as a
    ready-made starting point for user customization — copy, rename, and
    edit to create a new theme.
    """
    # Import the built-in palettes lazily to avoid a forward-reference at
    # module level (StarWindow._PALETTES is defined inside _run_qt_gui).
    _BUILT_IN: Dict[str, Dict[str, str]] = {
        "dark": {
            "bg": "#16181d",
            "fg": "#c6ccd4",
            "sel": "#2c313a",
            "h1": "#82aaff",
            "h2": "#89ddff",
            "h3": "#c792ea",
            "h4": "#f78c6c",
            "code": "#7fdbab",
        },
        "light": {
            "bg": "#fafafa",
            "fg": "#24273a",
            "sel": "#bcd0f0",
            "h1": "#1e66f5",
            "h2": "#209fb5",
            "h3": "#8839ef",
            "h4": "#e64553",
            "code": "#40a02b",
        },
        "contrast": {
            "bg": "#000000",
            "fg": "#ffffff",
            "sel": "#404040",
            "h1": "#ffff00",
            "h2": "#00ffff",
            "h3": "#ff80ff",
            "h4": "#80ff80",
            "code": "#00ff80",
        },
        "phosphor": {
            "bg": "#001200",
            "fg": "#00cc00",
            "sel": "#004400",
            "h1": "#00ff00",
            "h2": "#00ee00",
            "h3": "#00cc00",
            "h4": "#00aa00",
            "code": "#009900",
        },
    }
    try:
        THEMES_DIR.mkdir(parents=True, exist_ok=True)
    except OSError:
        return
    header = (
        "/* star CSS theme\n"
        " * Copy this file, give it a new name, and edit freely.\n"
        " * star picks up any *.css file in this directory automatically.\n"
        " * Run  View → Reload CSS Themes  (or restart) to apply changes.\n"
        " */\n\n"
    )
    for name, pal in _BUILT_IN.items():
        path = THEMES_DIR / f"{name}.css"
        if not path.exists():
            try:
                path.write_text(header + _palette_to_css(pal), encoding="utf-8")
            except OSError:
                pass


# =============================================================================
# Optional Qt GUI
# =============================================================================


def _run_qt_gui(settings: Settings, initial_path: str = "") -> None:
    """Launch the optional Qt-based GUI mode."""
    if not _QT:
        print(
            "Qt GUI requires PyQt6 or PyQt5:\n"
            "  pip install PyQt6\nor\n  pip install PyQt5",
            file=sys.stderr,
        )
        sys.exit(1)

    # QTextCursor.MoveMode enum was reorganized in PyQt6; handle both.
    try:
        _KEEP_ANCHOR = QTextCursor.MoveMode.KeepAnchor  # PyQt6
    except AttributeError:
        _KEEP_ANCHOR = QTextCursor.KeepAnchor  # PyQt5  # type: ignore[attr-defined]

    # QueuedConnection constant also changed location.
    try:
        _QUEUED = Qt.ConnectionType.QueuedConnection  # PyQt6
    except AttributeError:
        _QUEUED = Qt.QueuedConnection  # PyQt5  # type: ignore[attr-defined]

    # DockWidgetArea enum (PyQt6 uses nested enum; PyQt5 uses top-level).
    try:
        _LEFT_DOCK = Qt.DockWidgetArea.LeftDockWidgetArea  # PyQt6
        _RIGHT_DOCK = Qt.DockWidgetArea.RightDockWidgetArea  # PyQt6
    except AttributeError:
        _LEFT_DOCK = Qt.LeftDockWidgetArea  # type: ignore[attr-defined]  # PyQt5
        _RIGHT_DOCK = Qt.RightDockWidgetArea  # type: ignore[attr-defined]  # PyQt5

    # ItemDataRole.UserRole enum (PyQt6 nested; PyQt5 top-level).
    try:
        _USER_ROLE = Qt.ItemDataRole.UserRole  # PyQt6
    except AttributeError:
        _USER_ROLE = Qt.UserRole  # type: ignore[attr-defined]  # PyQt5

    # Enums used by the reading-accessibility features (spacing, reading
    # aids, highlight tuning).  Each was reorganized between PyQt5 and 6.
    try:
        _PROPORTIONAL = QTextBlockFormat.LineHeightTypes.ProportionalHeight  # PyQt6
    except AttributeError:
        _PROPORTIONAL = QTextBlockFormat.ProportionalHeight  # type: ignore[attr-defined]
    try:
        _DOC_SELECTION = QTextCursor.SelectionType.Document  # PyQt6
    except AttributeError:
        _DOC_SELECTION = QTextCursor.Document  # type: ignore[attr-defined]
    try:
        _FULL_WIDTH_SEL = QTextFormat.Property.FullWidthSelection  # PyQt6
    except AttributeError:
        _FULL_WIDTH_SEL = QTextFormat.FullWidthSelection  # type: ignore[attr-defined]
    try:
        _PCT_SPACING = QFont.SpacingType.PercentageSpacing  # PyQt6
    except AttributeError:
        _PCT_SPACING = QFont.PercentageSpacing  # type: ignore[attr-defined]
    try:
        _SINGLE_UNDERLINE = QTextCharFormat.UnderlineStyle.SingleUnderline  # PyQt6
        _WAVE_UNDERLINE = QTextCharFormat.UnderlineStyle.WaveUnderline
    except AttributeError:
        _SINGLE_UNDERLINE = QTextCharFormat.SingleUnderline  # type: ignore[attr-defined]
        _WAVE_UNDERLINE = QTextCharFormat.WaveUnderline  # type: ignore[attr-defined]

    # High DPI support — must be set before QApplication().
    # Use the module-level Qt object so PyQt6's C-extension type checks
    # receive the exact enum type they were compiled against.
    if settings.get("qt_hidpi", True):
        try:
            # PyQt6: HighDpiScaleFactorRoundingPolicy is the recommended knob.
            QApplication.setHighDpiScaleFactorRoundingPolicy(
                Qt.HighDpiScaleFactorRoundingPolicy.PassThrough  # type: ignore[attr-defined]
            )
        except (AttributeError, TypeError):
            # PyQt5 (or an older PyQt6 build): fall back to the AA_* attributes.
            try:
                QApplication.setAttribute(Qt.AA_EnableHighDpiScaling, True)  # type: ignore[attr-defined]
                QApplication.setAttribute(Qt.AA_UseHighDpiPixmaps, True)  # type: ignore[attr-defined]
            except AttributeError:
                pass  # Qt version has no HiDPI API — safe to ignore

    app = QApplication(sys.argv)
    app.setApplicationName(APP_TITLE)
    app.setApplicationVersion(APP_VERSION)

    # ── Global exception hook ─────────────────────────────────────────
    # In PyQt6, an unhandled Python exception inside a connected slot is
    # re-raised in the GUI thread and can escape app.exec(), causing the
    # process to print a traceback to a briefly-visible console window on
    # Windows and then exit.  This hook catches anything that slips through
    # and writes it to star_crash.log next to the settings file so the user
    # can read it after the window closes.
    _log_path = SETTINGS_FILE.parent / "star_crash.log"

    def _excepthook(exc_type: type, exc_value: BaseException, exc_tb: Any) -> None:
        import traceback as _tb

        msg = "".join(_tb.format_exception(exc_type, exc_value, exc_tb))
        # Write to log file first (always succeeds).
        try:
            _log_path.write_text(msg, encoding="utf-8")
        except Exception:
            pass
        # Try to show a Qt message box.
        try:
            QMessageBox.critical(
                None,
                f"{APP_NAME} — Unexpected Error",
                f"An unexpected error occurred.\n\n"
                f"{msg[:1200]}\n\n"
                f"Full details saved to:\n{_log_path}",
            )
        except Exception:
            pass  # if Qt itself is broken, silently give up

    sys.excepthook = _excepthook

    class StarWindow(QMainWindow):
        """Qt GUI window for star.

        Word-level highlight pipeline
        ─────────────────────────────
        1. TTSManager calls its _on_highlight callback with a word index.
        2. That callback emits _word_signal(idx) — a pyqtSignal — from
           whatever thread the TTS engine is running on.
        3. Because the signal is connected with QueuedConnection, Qt
           marshals the call to the GUI thread automatically.
        4. _apply_word_highlight() runs on the GUI thread, uses
           QTextEdit.setExtraSelections() to paint the highlight
           without touching the document's undo history or content,
           and calls ensureCursorVisible() to auto-scroll.
        """

        # Emitted from the TTS/timer thread; delivered to the GUI thread.
        # Carries (word_idx, session) — the session integer lets
        # _apply_word_highlight discard stale queued emissions that were
        # posted by a timer that has since been superseded by a new speak()
        # call.  Without the session stamp, a queued signal for word 0 from
        # the previous utterance scrolls the viewport back to the start of
        # the document immediately after the ToC navigation scroll.
        _word_signal = pyqtSignal(int, int)  # (word_idx, hl_session)
        # Emitted from the document-load thread to trigger _on_doc_loaded on
        # the GUI thread.  QMetaObject.invokeMethod with a plain Python method
        # name is unreliable on Windows (requires @pyqtSlot registration);
        # a proper pyqtSignal is the correct cross-thread mechanism.
        _doc_loaded_signal = pyqtSignal()
        # Emitted from the word-map build thread once the map is ready so
        # _qt_restore_reading_position can be called safely on the GUI thread.
        _restore_signal = pyqtSignal()
        # Emitted from the audio-export background thread when synthesis is
        # complete.  Carries the status-bar message to display (success or
        # error text) so the GUI thread can update safely.
        _export_audio_signal = pyqtSignal(str)
        # Emitted from the Whisper background threads so results
        # land on the GUI thread.  transcribe: (text, source_path);
        # dictate: (text, char_pos_str, anchor).
        _transcribe_signal = pyqtSignal(str, str)
        _dictate_signal = pyqtSignal(str, str, str)
        _dictate_partial_signal = pyqtSignal(str)  # live streaming dictation preview
        _doi_signal = pyqtSignal(str)  # Crossref DOI lookup result (JSON or ERROR:)

        # Per-theme CSS colors used by _md_to_html and _apply_qt_theme.
        _PALETTES: Dict[str, Dict[str, str]] = {
            "dark": {
                "bg": "#16181d",
                "fg": "#c6ccd4",
                "sel": "#2c313a",
                "h1": "#82aaff",
                "h2": "#89ddff",
                "h3": "#c792ea",
                "h4": "#f78c6c",
                "code": "#7fdbab",
            },
            "light": {
                "bg": "#fafafa",
                "fg": "#24273a",
                "sel": "#bcd0f0",
                "h1": "#1e66f5",
                "h2": "#209fb5",
                "h3": "#8839ef",
                "h4": "#e64553",
                "code": "#40a02b",
            },
            "contrast": {
                "bg": "#000000",
                "fg": "#ffffff",
                "sel": "#404040",
                "h1": "#ffff00",
                "h2": "#00ffff",
                "h3": "#ff80ff",
                "h4": "#80ff80",
                "code": "#00ff80",
            },
            "phosphor": {
                "bg": "#001200",
                "fg": "#00cc00",
                "sel": "#004400",
                "h1": "#00ff00",
                "h2": "#00ee00",
                "h3": "#00cc00",
                "h4": "#00aa00",
                "code": "#009900",
            },
        }

        def __init__(self) -> None:
            super().__init__()
            self.settings = settings
            self.doc: Optional[Document] = None
            self.tts_manager = TTSManager(settings)

            # Word index saved when the user pauses speech (Space).  -1 means no
            # saved position; used by _tts_toggle to resume from the exact word.
            self._tts_paused_at_word: int = -1

            # Maps TTS word index → absolute character offset in the Qt document.
            # Built asynchronously after each document load.
            self._qt_word_map: List[int] = []

            # QTextCharFormat applied to the currently spoken word.
            # Built from the user's highlight style/color settings so the
            # karaoke highlight can be tuned (see _rebuild_hl_fmt).
            self._hl_fmt = QTextCharFormat()
            self._rebuild_hl_fmt()

            # Session counter — incremented by every new speak() invocation.
            # The timer thread closes over the session value at speak() time;
            # _apply_word_highlight drops any delivery whose session no longer
            # matches, preventing stale queued signals from scrolling the view.
            self._hl_session: int = 0

            # CSS theme registry — must be populated before _setup_ui() because
            # _setup_ui calls _apply_qt_theme → _effective_palette → _css_themes.
            # Seed default files first (no-op if they already exist), then load.
            _seed_default_css_themes()
            self._css_themes: Dict[str, Dict[str, Any]] = _load_css_themes()

            self._setup_ui()

            # Wire the TTS highlight callback → signal (thread-safe delivery).
            self._word_signal.connect(self._apply_word_highlight, _QUEUED)
            # Seed the initial callback (session 0).  _new_hl_session() replaces
            # this at the start of every new speak() call.
            self._refresh_hl_callback()

            # Wire the doc-loaded callback → signal (thread-safe delivery).
            self._doc_loaded_signal.connect(self._on_doc_loaded, _QUEUED)
            # Wire the restore-position callback → main thread.
            self._restore_signal.connect(self._qt_restore_reading_position, _QUEUED)
            # Wire the audio-export completion signal → status bar update.
            self._export_audio_signal.connect(
                lambda msg: self.statusBar().showMessage(msg), _QUEUED
            )
            # Wire the Whisper transcription / dictation result signals.
            self._transcribe_signal.connect(self._qt_on_transcribed, _QUEUED)
            self._dictate_signal.connect(self._qt_on_dictated, _QUEUED)
            self._dictate_partial_signal.connect(self._qt_on_dictate_partial, _QUEUED)
            self._doi_signal.connect(self._qt_on_doi, _QUEUED)

            # Speech Cursor (SC) mode state for the Qt GUI.
            # Mirrors the TUI sc mode: Tab enters it, arrows navigate
            # line-by-line with TTS, Esc / Enter exits.
            self._qt_sc_mode: bool = False
            self._qt_sc_block: int = 0  # current QTextBlock index
            # Persistent SAPI5 engine for SC line reading — same role as
            # _SCReader in the TUI.  None when not in SC mode.
            self._qt_sc_reader: Optional[_SCReader] = None

            # Sentence map: list of word-map indices at which each sentence
            # begins.  Built asynchronously alongside the word map.
            self._qt_sentence_starts: List[int] = [0]

            # Edit mode state.  False = read-only (normal), True = editable.
            self._qt_edit_mode: bool = False
            self._qt_edit_dirty: bool = False  # unsaved changes in edit mode

            # Intercept keyboard events on the editor so Tab and arrow keys
            # can be processed for SC mode without fighting Qt's focus chain.
            self.editor.installEventFilter(self)

            if initial_path:
                self._open_path(initial_path)

        # ── UI construction ───────────────────────────────────────────────

        def _setup_ui(self) -> None:
            self.setWindowTitle(APP_TITLE)
            self.resize(
                int(self.settings["gui_width"]), int(self.settings["gui_height"])
            )

            # Document view
            self.editor = QTextEdit()
            self.editor.setReadOnly(True)
            # Build the base font from family/size + letter/word spacing and
            # the optional dyslexia-friendly family override.
            self.editor.setFont(self._make_editor_font())
            # Hide the blinking text cursor — this is a reader, not an editor.
            self.editor.setCursorWidth(0)
            self.setCentralWidget(self.editor)
            # Apply the current theme's colors to the editor widget.
            self._apply_qt_theme(self.settings.get("theme", "dark"))

            # AT-SPI2 accessibility metadata
            self.editor.setAccessibleName("Document View")
            self.editor.setAccessibleDescription(
                "Main reading area. Space=play/pause, Ctrl+F=search, Escape=stop."
            )

            # Registry of (label, QAction, default_shortcut) for every action
            # that carries a shortcut, so they can be remapped at runtime
            # Populated by _act and _nav below.
            self._shortcut_actions = []

            # Toolbar
            tb = self.addToolBar("Controls")
            tb.setMovable(False)

            def _act(label: str, shortcut: str, fn: Callable, tip: str = "") -> None:
                """Add one action to the toolbar with an optional tooltip."""
                a = QAction(label, self)
                if shortcut:
                    key = self._resolve_shortcut(shortcut)
                    a.setShortcut(key)
                    self._shortcut_actions.append((label, a, shortcut))
                a.setToolTip(tip or label)
                a.triggered.connect(fn)
                tb.addAction(a)

            # ── File ─────────────────────────────────────────────
            _act("Open", "Ctrl+O", self._open_dialog, "Open a file (Ctrl+O)")
            _act("URL", "", self._qt_open_url, "Open a URL")
            tb.addSeparator()
            # ── Playback ──────────────────────────────────────
            _act(
                "Play/Pause ▶⏸",
                "Space",
                self._tts_toggle,
                "Play / pause speech (Space)",
            )
            _act("Stop ■", "Escape", self._tts_stop, "Stop speech (Escape)")
            _act(
                "− Speed",
                "Ctrl+-",
                lambda: self._rate_change(-20),
                "Slow down −20 wpm (Ctrl+−)",
            )
            _act(
                "+ Speed",
                "Ctrl+=",
                lambda: self._rate_change(+20),
                "Speed up +20 wpm (Ctrl+=)",
            )
            tb.addSeparator()
            # ── Navigate ──────────────────────────────────────
            # Sentence: previous  replay  next
            _act("◀ S", "", self._qt_skip_prev_sentence, "Previous sentence (,)")
            _act("↺ S", "", self._qt_replay_sentence, "Replay sentence (;)")
            _act("S ▶", "", self._qt_skip_next_sentence, "Next sentence (.)")
            # Paragraph: previous  replay  next
            _act("◀ ¶", "", self._qt_skip_prev_paragraph, "Previous paragraph ([)")
            _act("↺ ¶", "", self._qt_replay_paragraph, "Replay paragraph (r)")
            _act("¶ ▶", "", self._qt_skip_next_paragraph, "Next paragraph (])")
            # Heading: previous  next  (always starts reading)
            _act(
                "◀ H",
                "",
                self._qt_read_prev_heading,
                "Previous heading — read aloud (<)",
            )
            _act("H ▶", "", self._qt_read_next_heading, "Next heading — read aloud (>)")
            tb.addSeparator()
            # ── Voice / cursor mode ───────────────────────────
            _act(
                "Voice…",
                "Ctrl+Shift+V",
                self._voice_picker_qt,
                "Select TTS voice (Ctrl+Shift+V)",
            )
            _act(
                "SC ○",
                "Tab",
                lambda: self._qt_sc_exit() if self._qt_sc_mode else self._qt_sc_enter(),
                "Speech cursor mode — line-by-line reading (Tab)",
            )
            tb.addSeparator()
            # ── Text ───────────────────────────────────────────
            _act(
                "Copy", "Ctrl+C", self._qt_copy, "Copy selection or paragraph (Ctrl+C)"
            )
            _act(
                "Highlight",
                "",  # Ctrl+H is now heading navigation
                self._qt_highlight,
                "Highlight selection in yellow",
            )
            _act(
                "Clear Highlights",
                "",
                self._qt_highlight_clear,
                "Remove all highlights from this document",
            )
            tb.addSeparator()
            # ── Edit ─────────────────────────────────────────
            _act(
                "Edit", "Ctrl+E", self._qt_edit_mode_toggle, "Toggle edit mode (Ctrl+E)"
            )
            _act("Save", "Ctrl+S", self._qt_save, "Save document (Ctrl+S)")
            tb.addSeparator()
            # ── View ─────────────────────────────────────────
            _act("Theme", "F5", self._next_theme, "Cycle color theme (F5)")
            _act(
                "ToC",
                "Ctrl+\\",
                self._qt_toggle_toc,
                "Toggle Contents panel (Ctrl+\\\\)",
            )
            _act(
                "Notes",
                "",  # Ctrl+Shift+N is bound as a window-level shortcut
                self._qt_toggle_annotations,
                "Toggle Notes panel (Ctrl+Shift+N)",
            )
            _act(
                "+ Note",
                "",  # Ctrl+Shift+A is bound as a window-level shortcut
                self._qt_add_annotation,
                "Add a note at the cursor (Ctrl+Shift+A)",
            )
            _act(
                "Level", "Ctrl+L", self._qt_reading_level, "Show reading level (Ctrl+L)"
            )
            _act("Font", "", self._qt_change_font_dialog, "Change display font")
            tb.addSeparator()
            # ── App ──────────────────────────────────────────
            _act("Help", "F1", self._show_about, "Open README.md (F1)")
            _act("Quit", "Ctrl+Q", self.close, "Quit star (Ctrl+Q)")

            # ── Menu bar ─────────────────────────────────────────────────────
            mb = self.menuBar()

            # File menu
            file_menu: QMenu = mb.addMenu("File")
            open_act = QAction("Open…", self)
            open_act.setShortcut("Ctrl+O")
            open_act.triggered.connect(self._open_dialog)
            file_menu.addAction(open_act)
            file_menu.addSeparator()

            export_menu: QMenu = file_menu.addMenu("Export")
            md_act = QAction("Export as Markdown…", self)
            md_act.triggered.connect(self._qt_export_markdown)
            export_menu.addAction(md_act)
            pdf_act = QAction("Export as PDF…", self)
            pdf_act.triggered.connect(self._qt_export_pdf)
            export_menu.addAction(pdf_act)
            brf_act = QAction("Export as Braille (BRF)…", self)
            brf_act.triggered.connect(self._qt_export_brf)
            export_menu.addAction(brf_act)
            audio_act = QAction("Export as Audio (MP3 / OGG / MP4)…", self)
            audio_act.triggered.connect(self._qt_export_audio)
            export_menu.addAction(audio_act)

            file_menu.addSeparator()
            quit_act = QAction("Quit", self)
            quit_act.setShortcut("Ctrl+Q")
            quit_act.triggered.connect(self.close)
            file_menu.addAction(quit_act)

            # Highlight menu
            hl_menu: QMenu = mb.addMenu("Highlight")
            _HL_COLORS = [
                ("Yellow", "#ffff00"),
                ("Green", "#90ee90"),
                ("Cyan", "#add8e6"),
                ("Pink", "#ffb6c1"),
                ("Orange", "#ffa500"),
            ]
            for _name, _color in _HL_COLORS:
                _a = QAction(f"Highlight {_name}", self)
                _a.triggered.connect(lambda _chk=False, c=_color: self._qt_highlight(c))
                hl_menu.addAction(_a)
            hl_menu.addSeparator()
            clr_act = QAction("Clear All Highlights", self)
            clr_act.triggered.connect(self._qt_highlight_clear)
            hl_menu.addAction(clr_act)

            # Notes / annotations menu
            notes_menu: QMenu = mb.addMenu("Notes")
            add_note_act = QAction("Add Note at Cursor…", self)
            add_note_act.setShortcut("Ctrl+Shift+A")
            add_note_act.triggered.connect(self._qt_add_annotation)
            notes_menu.addAction(add_note_act)
            edit_note_act = QAction("Edit Selected Note…", self)
            edit_note_act.triggered.connect(self._qt_edit_annotation)
            notes_menu.addAction(edit_note_act)
            del_note_act = QAction("Delete Selected Note", self)
            del_note_act.triggered.connect(self._qt_delete_annotation)
            notes_menu.addAction(del_note_act)
            notes_menu.addSeparator()
            toggle_notes_act = QAction("Toggle Notes Panel", self)
            toggle_notes_act.setShortcut("Ctrl+Shift+N")
            toggle_notes_act.triggered.connect(self._qt_toggle_annotations)
            notes_menu.addAction(toggle_notes_act)
            export_notes_act = QAction("Export Notes…", self)
            export_notes_act.triggered.connect(self._qt_export_annotations)
            notes_menu.addAction(export_notes_act)

            # Helper: build a menu from (label, callable) rows (None = separator).
            def _menu(title: str, rows: List[Any]) -> "QMenu":
                menu = mb.addMenu(title)
                for row in rows:
                    if row is None:
                        menu.addSeparator()
                        continue
                    label, fn = row
                    act = QAction(label, self)
                    act.triggered.connect(lambda _chk=False, f=fn: f())
                    menu.addAction(act)
                return menu

            # Speech menu — every playback command reachable without the keyboard.
            _menu(
                "Speech",
                [
                    ("Play / Pause", self._tts_toggle),
                    ("Stop", self._tts_stop),
                    ("Play from Cursor", self._qt_play_from_cursor),
                    None,
                    ("Faster (+20 wpm)", lambda: self._rate_change(+20)),
                    ("Slower (−20 wpm)", lambda: self._rate_change(-20)),
                    None,
                    ("Choose Voice…", self._voice_picker_qt),
                    (
                        "Speech Cursor Mode",
                        lambda: (
                            self._qt_sc_exit()
                            if self._qt_sc_mode
                            else self._qt_sc_enter()
                        ),
                    ),
                    None,
                    (
                        "Toggle SSML Prosody",
                        lambda: (
                            self.settings.set(
                                "use_ssml", not self.settings.get("use_ssml", False)
                            ),
                            self.statusBar().showMessage(
                                "SSML prosody: "
                                + ("ON" if self.settings.get("use_ssml") else "OFF")
                            ),
                        ),
                    ),
                ],
            )

            # Navigate menu.
            _menu(
                "Navigate",
                [
                    ("Next Sentence", self._qt_skip_next_sentence),
                    ("Previous Sentence", self._qt_skip_prev_sentence),
                    ("Replay Sentence", self._qt_replay_sentence),
                    None,
                    ("Next Paragraph", self._qt_skip_next_paragraph),
                    ("Previous Paragraph", self._qt_skip_prev_paragraph),
                    ("Replay Paragraph", self._qt_replay_paragraph),
                    None,
                    ("Next Heading", self._qt_read_next_heading),
                    ("Previous Heading", self._qt_read_prev_heading),
                    None,
                    ("Next Table", self._qt_skip_next_table),
                    ("Previous Table", self._qt_skip_prev_table),
                ],
            )

            # Edit menu.
            _menu(
                "Edit",
                [
                    ("Copy", self._qt_copy),
                    None,
                    ("Toggle Edit Mode", self._qt_edit_mode_toggle),
                    ("Save", self._qt_save),
                ],
            )

            # Citations menu.
            _menu(
                "Citations",
                [
                    ("Import…", self._qt_import_citations),
                    ("Export…", self._qt_export_citations),
                    None,
                    ("Add Citation…", self._qt_add_citation),
                    ("Add by DOI…", self._qt_add_citation_by_doi),
                    ("Insert Citation at Cursor…", self._qt_insert_citation),
                    ("Manage / Browse…", self._qt_manage_citations),
                ],
            )

            # View menu
            view_menu: QMenu = mb.addMenu("View")
            toc_act = QAction("Toggle Contents Panel", self)
            # No setShortcut here — Ctrl+\ is already on the toolbar button.
            # Duplicate window-level shortcuts make Qt treat the key as
            # ambiguous and fire neither action.
            toc_act.triggered.connect(self._qt_toggle_toc)
            view_menu.addAction(toc_act)
            notes_toggle_act = QAction("Toggle Notes Panel", self)
            notes_toggle_act.triggered.connect(self._qt_toggle_annotations)
            view_menu.addAction(notes_toggle_act)
            view_menu.addSeparator()
            theme_act = QAction("Next Theme", self)
            # No setShortcut here — F5 is already on the toolbar button.
            theme_act.triggered.connect(self._next_theme)
            view_menu.addAction(theme_act)
            pick_theme_act = QAction("Choose Theme…", self)
            pick_theme_act.triggered.connect(self._qt_pick_theme)
            view_menu.addAction(pick_theme_act)
            reload_css_act = QAction("Reload CSS Themes", self)
            reload_css_act.triggered.connect(self._qt_reload_css_themes)
            reload_css_act.setToolTip(
                f"Rescan {THEMES_DIR} for *.css files without restarting"
            )
            view_menu.addAction(reload_css_act)
            open_themes_act = QAction("Open Themes Folder", self)
            open_themes_act.triggered.connect(self._qt_open_themes_folder)
            view_menu.addAction(open_themes_act)
            view_menu.addSeparator()
            font_act = QAction("Change Font…", self)
            font_act.triggered.connect(self._qt_change_font_dialog)
            view_menu.addAction(font_act)
            level_act = QAction("Reading Level", self)
            # No setShortcut here — Ctrl+L is already on the toolbar button.
            level_act.triggered.connect(self._qt_reading_level)
            view_menu.addAction(level_act)

            # ── Reading Aids submenu (accessibility) ─────────────────────
            aids_menu: QMenu = view_menu.addMenu("Reading Aids")
            spacing_act = QAction("Text Spacing…", self)
            spacing_act.setToolTip(
                "Adjust line height, letter and word spacing (WCAG 1.4.12)"
            )
            spacing_act.triggered.connect(self._qt_text_spacing_dialog)
            aids_menu.addAction(spacing_act)
            karaoke_act = QAction("Karaoke Highlight…", self)
            karaoke_act.setToolTip(
                "Tune the spoken-word highlight style, color, speed and lead"
            )
            karaoke_act.triggered.connect(self._qt_karaoke_dialog)
            aids_menu.addAction(karaoke_act)
            aids_menu.addSeparator()
            self._dyslexia_font_act = QAction("Dyslexia-Friendly Font", self)
            self._dyslexia_font_act.setCheckable(True)
            self._dyslexia_font_act.setChecked(
                bool(self.settings.get("qt_dyslexia_font", False))
            )
            self._dyslexia_font_act.setToolTip(
                "Prefer OpenDyslexic / Atkinson Hyperlegible / Lexend if installed"
            )
            self._dyslexia_font_act.triggered.connect(self._qt_toggle_dyslexia_font)
            aids_menu.addAction(self._dyslexia_font_act)
            self._bionic_act = QAction("Bionic Reading", self)
            self._bionic_act.setCheckable(True)
            self._bionic_act.setChecked(
                bool(self.settings.get("qt_bionic_reading", False))
            )
            self._bionic_act.setToolTip("Embolden the leading part of each word")
            self._bionic_act.triggered.connect(self._qt_toggle_bionic)
            aids_menu.addAction(self._bionic_act)
            self._current_line_act = QAction("Current-Line Highlight", self)
            self._current_line_act.setCheckable(True)
            self._current_line_act.setChecked(
                bool(self.settings.get("qt_current_line_highlight", False))
            )
            self._current_line_act.setToolTip(
                "Tint the line being read with a focus band"
            )
            self._current_line_act.triggered.connect(self._qt_toggle_current_line)
            aids_menu.addAction(self._current_line_act)

            # Tools menu — transcription, dictation, and maintenance.
            _menu(
                "Tools",
                [
                    ("Transcribe Audio File…", self._qt_transcribe_file),
                    ("Dictate Note (record)…", self._qt_dictate_note),
                    (
                        "Toggle Transcript Timestamps",
                        lambda: (
                            self.settings.set(
                                "transcribe_timestamps",
                                not self.settings.get("transcribe_timestamps", False),
                            ),
                            self.statusBar().showMessage(
                                "Transcript timestamps: "
                                + (
                                    "ON"
                                    if self.settings.get("transcribe_timestamps")
                                    else "OFF"
                                )
                            ),
                        ),
                    ),
                    None,
                    ("Reading Level", self._qt_reading_level),
                    (
                        "Clear Document Cache",
                        lambda: (
                            shutil.rmtree(CACHE_DIR, ignore_errors=True),
                            self.statusBar().showMessage("Document cache cleared"),
                        ),
                    ),
                ],
            )

            # Help menu.
            _menu(
                "Help",
                [
                    ("Command Palette… (F2)", self._qt_command_palette),
                    ("Keyboard Shortcuts…", self._qt_show_shortcuts),
                    ("Customize Shortcuts…", self._qt_customize_shortcuts),
                    None,
                    ("Open README (Help)", self._show_about),
                    (
                        "About star",
                        lambda: QMessageBox.about(
                            self,
                            f"About {APP_NAME}",
                            f"<b>{APP_TITLE}</b><br>Version {APP_VERSION}<br><br>"
                            f"{__copyright__}<br>{__license__}",
                        ),
                    ),
                ],
            )

            # ── Table of Contents dock (Feature #25) ─────────────────────────
            self._toc_dock = QDockWidget("Contents", self)
            self._toc_dock.setObjectName("toc_dock")
            self._toc_dock.setAllowedAreas(
                _LEFT_DOCK | _RIGHT_DOCK  # type: ignore[operator]
            )
            self._toc_list = QListWidget()
            self._toc_list.setMinimumWidth(180)
            # Single-click / Enter: scroll to heading and stop speech so the
            # viewport and audio stay synchronized.
            self._toc_list.itemActivated.connect(self._qt_toc_navigate)
            # Double-click: stop current speech and start reading from the
            # activated heading.
            self._toc_list.itemDoubleClicked.connect(self._qt_toc_play)
            self._toc_dock.setWidget(self._toc_list)
            self.addDockWidget(_LEFT_DOCK, self._toc_dock)
            show_toc = bool(self.settings.get("qt_show_toc", True))
            self._toc_dock.setVisible(show_toc)

            # ── Annotations / Notes dock ─────────────────────────────────────
            # Mirrors the Contents dock: a list of notes anchored at character
            # positions in the document.  Single-click scrolls to the note;
            # double-click reads from there.  Notes persist per-document in
            # settings and can be exported (Markdown / JSON / BibTeX / RIS).
            self._annot_dock = QDockWidget("Notes", self)
            self._annot_dock.setObjectName("annotations_dock")
            self._annot_dock.setAllowedAreas(
                _LEFT_DOCK | _RIGHT_DOCK  # type: ignore[operator]
            )
            _annot_panel = QWidget()
            _annot_layout = QVBoxLayout(_annot_panel)
            _annot_layout.setContentsMargins(4, 4, 4, 4)
            _annot_layout.setSpacing(4)
            # Full-text search / tag filter box (type `#tag` to filter by tag).
            self._annot_filter = QLineEdit()
            self._annot_filter.setPlaceholderText("Filter notes — text or #tag…")
            self._annot_filter.setClearButtonEnabled(True)
            self._annot_filter.textChanged.connect(
                lambda _t: self._qt_build_annotations()
            )
            _annot_layout.addWidget(self._annot_filter)
            self._annot_list = QListWidget()
            self._annot_list.setMinimumWidth(200)
            self._annot_list.setWordWrap(True)
            self._annot_list.itemActivated.connect(self._qt_annotation_navigate)
            self._annot_list.itemDoubleClicked.connect(self._qt_annotation_play)
            _annot_layout.addWidget(self._annot_list)
            _btn_row = QHBoxLayout()
            for _lbl, _fn in (
                ("Add", self._qt_add_annotation),
                ("Edit", self._qt_edit_annotation),
                ("Delete", self._qt_delete_annotation),
                ("Export…", self._qt_export_annotations),
            ):
                _b = QPushButton(_lbl)
                _b.clicked.connect(lambda _chk=False, fn=_fn: fn())
                _btn_row.addWidget(_b)
            _annot_layout.addLayout(_btn_row)
            self._annot_dock.setWidget(_annot_panel)
            self.addDockWidget(_RIGHT_DOCK, self._annot_dock)
            self._annot_dock.setVisible(bool(self.settings.get("qt_show_notes", False)))

            # ── Window-level shortcuts ─────────────────────────────────────
            # Aligned with conventional screen-reader key bindings so users
            # familiar with NVDA/JAWS/VoiceOver have the same muscle memory.
            # Modifier rules:
            #   Ctrl+<letter>          → forward navigation
            #   Ctrl+Shift+<letter>    → backward navigation
            #   Alt+<punctuation>      → sentence navigation (avoids
            #                            collisions with text-editing chords)
            def _nav(key: str, fn: Callable, label: str = "") -> None:
                na = QAction(self)
                resolved = self._resolve_shortcut(key)
                na.setShortcut(resolved)
                na.triggered.connect(fn)
                self.addAction(na)  # window-level — fires regardless of focus
                self._shortcut_actions.append((label or key, na, key))

            # Play from cursor position
            _nav("Ctrl+Return", self._qt_play_from_cursor, "Play from cursor")

            # Heading navigation (H = Heading, screen-reader convention)
            _nav("Ctrl+H", self._qt_skip_next_heading, "Next heading")
            _nav("Ctrl+Shift+H", self._qt_skip_prev_heading, "Previous heading")

            # Paragraph navigation (P = Paragraph)
            _nav("Ctrl+P", self._qt_skip_next_paragraph, "Next paragraph")
            _nav("Ctrl+Shift+P", self._qt_skip_prev_paragraph, "Previous paragraph")
            _nav("Ctrl+R", self._qt_replay_paragraph, "Replay paragraph")

            # Table navigation (T = Table)
            _nav("Ctrl+T", self._qt_skip_next_table, "Next table")
            _nav("Ctrl+Shift+T", self._qt_skip_prev_table, "Previous table")

            # Sentence navigation (Alt+punctuation avoids text-editing conflicts)
            _nav("Alt+.", self._qt_skip_next_sentence, "Next sentence")
            _nav("Alt+,", self._qt_skip_prev_sentence, "Previous sentence")
            _nav("Alt+;", self._qt_replay_sentence, "Replay sentence")

            # Document editing
            _nav("Ctrl+E", self._qt_edit_mode_toggle, "Toggle edit mode")
            _nav("Ctrl+S", self._qt_save, "Save")

            # Annotations / notes (A = annotation, N = notes panel)
            _nav("Ctrl+Shift+A", self._qt_add_annotation, "Add note")
            _nav("Ctrl+Shift+N", self._qt_toggle_annotations, "Toggle notes panel")

            # Command palette (F2 — mirrors the TUI command palette key)
            _nav("F2", self._qt_command_palette, "Command palette")

            # Show the welcome screen immediately so the window is never
            # blank at launch.  _on_doc_loaded replaces it when a file loads.
            self.editor.setHtml(self._welcome_html())
            self._apply_block_spacing()
            self.statusBar().showMessage(APP_TITLE)

        def _welcome_html(self) -> str:
            """Return the splash-screen HTML, styled with the current theme."""
            theme_name = str(self.settings.get("theme", "dark"))
            pal = self._effective_palette(theme_name)
            custom_css = str(pal.get("_css", ""))
            if custom_css:
                # Inject the theme CSS plus some extras for the welcome page
                # (kbd highlight, centered layout) that the theme may omit.
                style = (
                    custom_css + f"kbd{{background:{pal.get('sel', '#313244')};"
                    f"color:{pal['fg']};padding:1px 5px;"
                    "border-radius:3px;font-family:monospace;font-size:0.9em}"
                    "td{padding:2px 12px 2px 0}"
                )
            else:
                style = (
                    f"body{{background:{pal['bg']};color:{pal['fg']};"
                    "     font-family:Georgia,serif;margin:24px;line-height:1.6}}"
                    f"h1{{color:{pal['h1']};margin-bottom:4px}}"
                    f"h2{{color:{pal['h2']};margin-top:20px;margin-bottom:6px}}"
                    f"kbd{{background:{pal.get('sel', '#313244')};"
                    f"color:{pal['fg']};padding:1px 5px;"
                    "border-radius:3px;font-family:monospace;font-size:0.9em}"
                    "td{padding:2px 12px 2px 0}"
                )
            return (
                "<html><head><style>" + style + "</style></head><body>"
                f"<h1>star &#8212; Speaking Terminal Access Reader</h1>"
                f"<p>Version {APP_VERSION} &nbsp;&middot;&nbsp; {APP_TITLE}</p>"
                "<h2>Getting started</h2>"
                "<table>"
                "<tr><td><kbd>Ctrl+O</kbd></td>"
                "    <td>Open a file</td></tr>"
                "<tr><td><kbd>Space</kbd></td>"
                "    <td>Play / pause speech</td></tr>"
                "<tr><td><kbd>Esc</kbd></td>"
                "    <td>Stop speech</td></tr>"
                "<tr><td><kbd>+</kbd> / <kbd>-</kbd></td>"
                "    <td>Speed up / slow down</td></tr>"
                "<tr><td><kbd>F1</kbd></td>"
                "    <td>Help</td></tr>"
                "<tr><td><kbd>Ctrl+Q</kbd></td>"
                "    <td>Quit</td></tr>"
                "</table>"
                "<h2>Supported formats</h2>"
                "<p>PDF &nbsp; DOCX &nbsp; ODT &nbsp; EPUB &nbsp; HTML &nbsp;"
                "Markdown &nbsp; plain text &nbsp; CSV / TSV / XLSX &nbsp;"
                "DAISY / DTBook &nbsp; URLs</p>"
                "</body></html>"
            )

        # ── Document loading ─────────────────────────────────────────────────────

        def _qt_open_url(self) -> None:
            """Prompt for a URL and open it as a document."""
            url, ok = QInputDialog.getText(self, "Open URL", "Enter a web address:")
            if ok and url.strip():
                self._open_path(url.strip())

        def _open_dialog(self) -> None:
            path, _ = QFileDialog.getOpenFileName(
                self,
                "Open Document",
                "",
                "All Supported "
                "(*.pdf *.doc *.dot *.docx *.ppt *.pptx *.odt "
                "*.epub *.html *.htm *.md *.txt "
                "*.rst *.rest *.adoc *.asciidoc *.asc "
                "*.wiki *.mediawiki *.textile *.creole "
                "*.tex *.ltx *.org "
                "*.csv *.tsv *.xlsx *.r *.rmd)"
                ";;Word / PowerPoint (*.doc *.dot *.docx *.ppt *.pptx)"
                ";;Wiki / Markup "
                "(*.rst *.rest *.adoc *.asciidoc *.asc "
                "*.wiki *.mediawiki *.textile *.creole)"
                ";;LaTeX (*.tex *.ltx)"
                ";;Org-mode (*.org)"
                ";;PDF (*.pdf)"
                ";;All Files (*)",
            )
            if path:
                self._open_path(path)

        def _open_path(self, path: str) -> None:
            # Save where we were in the *current* document before replacing it.
            self._qt_save_reading_position()
            self.statusBar().showMessage(f"Loading {Path(path).name} …")
            QApplication.processEvents()
            self._pending_doc: Optional[Document] = None

            def _work() -> None:
                try:
                    self._pending_doc = load_document(path, self.settings)
                except Exception as _exc:  # noqa: BLE001
                    # Never let the background thread die silently and leave
                    # the UI frozen.  Create a minimal error document instead
                    # so _on_doc_loaded always has something to display.
                    _err = Document(
                        path=path,
                        title=f"Error — {Path(path).name}",
                        markdown=(
                            f"# Could not open {Path(path).name}\n\n"
                            f"```\n{_exc}\n```\n\n"
                            "Check that the file exists and is not locked."
                        ),
                        plain_text=str(_exc),
                        format="error",
                    )
                    self._pending_doc = _err
                # Signal the GUI thread to call _on_doc_loaded.
                # Using a pyqtSignal guarantees safe cross-thread delivery;
                # QMetaObject.invokeMethod with a plain string requires the
                # method to be a registered @pyqtSlot and fails silently on
                # Windows when that registration is missing.
                self._doc_loaded_signal.emit()

            threading.Thread(target=_work, daemon=True).start()

        def _on_doc_loaded(self) -> None:
            # Wrap the entire slot body so that any exception is caught here
            # rather than propagating through PyQt6’s event loop and crashing
            # app.exec() — the symptom on Windows is a brief console window
            # that closes too quickly to read the traceback.
            try:
                self._on_doc_loaded_impl()
            except Exception as _exc:  # noqa: BLE001
                import traceback as _tb

                detail = _tb.format_exc()
                self.statusBar().showMessage(f"Error displaying document: {_exc}")
                try:
                    QMessageBox.critical(
                        self,
                        "Document Display Error",
                        f"Could not display the document.\n\n{detail[:1000]}",
                    )
                except Exception:
                    pass

        def _on_doc_loaded_impl(self) -> None:
            """Inner implementation of _on_doc_loaded, called inside a
            try/except wrapper to prevent slot exceptions from escaping
            into the Qt event loop."""
            doc = getattr(self, "_pending_doc", None)
            if not doc:
                return
            self.doc = doc
            self.editor.setHtml(self._md_to_html(doc.markdown or ""))
            self._apply_block_spacing()  # line-height (reset by setHtml)
            self.editor.setExtraSelections([])  # clear leftover TTS highlights
            # Build the ToC panel from the new document's headings.
            self._qt_build_toc()
            # Populate the Notes panel from saved annotations for this document.
            self._qt_build_annotations()
            # Restore any saved user highlights for this document.
            self._qt_apply_user_highlights()

            # Read Qt plain text NOW (main thread required) then hand off.
            qt_plain = self.editor.document().toPlainText()

            def _build() -> None:
                try:
                    plain = doc.plain_text or ""
                    # TTSManager word map (line-based, used for timer highlighting)
                    flat = qt_plain.splitlines()
                    doc.word_map = _build_word_map(plain, flat)
                    self.tts_manager.set_word_map(doc.word_map)
                    # Qt char-offset map (used by _apply_word_highlight)
                    self._build_qt_word_map(plain, qt_plain)
                    # Sentence map — same algorithm as StarApp._build_sentence_map
                    wm = doc.word_map
                    if wm and plain:
                        char_starts = [0]
                        for _m in _SENTENCE_SPLIT_RE.finditer(plain):
                            char_starts.append(_m.end())
                        _wi = 0
                        word_starts: List[int] = []
                        for cs in char_starts:
                            while _wi < len(wm) and wm[_wi].tts_offset < cs:
                                _wi += 1
                            word_starts.append(min(_wi, len(wm) - 1))
                        seen: set = set()
                        result: List[int] = []
                        for ws in word_starts:
                            if ws not in seen:
                                seen.add(ws)
                                result.append(ws)
                        self._qt_sentence_starts = result if result else [0]
                    # Signal the main thread to restore the reading position
                    # now that the word map and sentence map are both ready.
                    self._restore_signal.emit()
                except Exception:
                    pass  # word map is best-effort; TTS works without it

            threading.Thread(target=_build, daemon=True).start()
            self.statusBar().showMessage(f"Opened: {doc.title}")

        # ── Word-position mapping ─────────────────────────────────────────

        def _build_qt_word_map(self, plain_text: str, qt_text: str) -> None:
            """Populate self._qt_word_map: a list where index i is the
            absolute character offset of the i-th TTS word inside the Qt
            document text.

            Uses a rolling forward search so that repeated words match in
            document order.  Runs in a background thread (no Qt calls).

            Words whose only occurrence in the document text is *before*
            search_from (e.g. column-header names repeated in structured
            table-row narration) are assigned last_good_pos so the highlight
            advances linearly rather than jumping backward to the header row.
            """
            result: List[int] = []
            qt_lower = qt_text.lower()
            token_re = re.compile(r"\b\w[\w'-]*")
            search_from = 0
            last_good_pos = 0  # last position from a forward-matched word

            for m in token_re.finditer(plain_text):
                word = m.group().lower()
                pos = qt_lower.find(word, search_from)
                if pos >= 0:
                    result.append(pos)
                    search_from = pos + 1  # advance past this occurrence
                    last_good_pos = pos
                else:
                    # Forward search failed.  Check if the word exists at all.
                    global_pos = qt_lower.find(word, 0)
                    if global_pos >= 0:
                        # Only a backward match exists (e.g. a table column
                        # header repeated in row narration).  Use last_good_pos
                        # so the highlight doesn't jump back to the header.
                        result.append(last_good_pos)
                    else:
                        result.append(0)  # not found anywhere
                    # Don't update search_from; we haven't moved forward.

            self._qt_word_map = result

        # ── HTML rendering ────────────────────────────────────────────────

        def _effective_palette(self, theme_name: str) -> Dict[str, Any]:
            """Return the palette dict for *theme_name*.

            Lookup order:
              1. Custom CSS themes loaded from THEMES_DIR (have key '_css').
              2. Built-in _PALETTES.
              3. Fallback: dark built-in palette.

            Uses getattr so this method is safe to call even before
            _css_themes is assigned (e.g. during very early init).
            """
            css_themes: Dict[str, Any] = getattr(self, "_css_themes", {})
            if theme_name in css_themes:
                return css_themes[theme_name]
            return self._PALETTES.get(theme_name, self._PALETTES["dark"])

        @property
        def _all_theme_names(self) -> List[str]:
            """Built-in theme names followed by any custom CSS theme names."""
            extra = [n for n in sorted(self._css_themes) if n not in THEME_NAMES]
            return THEME_NAMES + extra

        def _md_to_html(self, md: str) -> str:
            """Convert internal Markdown to styled HTML for QTextEdit.

            When the active theme came from a CSS file the raw CSS is
            injected verbatim so every selector the user wrote is honored.
            For built-in palettes the CSS is generated from the palette dict.
            """
            theme_name = self.settings.get("theme", "dark")
            pal = self._effective_palette(theme_name)
            custom_css: str = str(pal.get("_css", ""))

            out: List[str] = []
            for line in (md or "").splitlines():
                if line.startswith("# "):
                    out.append(f"<h1>{line[2:]}</h1>")
                elif line.startswith("## "):
                    out.append(f"<h2>{line[3:]}</h2>")
                elif line.startswith("### "):
                    out.append(f"<h3>{line[4:]}</h3>")
                elif line.startswith("#### "):
                    out.append(f"<h4>{line[5:]}</h4>")
                elif not line.strip():
                    out.append("<br>")
                else:
                    line = re.sub(r"\*\*(.+?)\*\*", r"<b>\1</b>", line)
                    line = re.sub(r"\*(.+?)\*", r"<i>\1</i>", line)
                    line = re.sub(r"`(.+?)`", r"<code>\1</code>", line)
                    out.append(f"<p>{line}</p>")
            body = "\n".join(out)
            # Bionic-reading: embolden the leading part of each word so the
            # eye is pulled forward through the text (a dyslexia reading aid).
            if self.settings.get("qt_bionic_reading", False):
                body = self._bionic_html(body)
            fam = self._effective_font_family()
            if custom_css:
                style = custom_css
            else:
                style = (
                    f"body{{background:{pal['bg']};color:{pal['fg']};"
                    f"      font-family:{fam},sans-serif;margin:14px;}}"
                    f"h1{{color:{pal['h1']}}}"
                    f"h2{{color:{pal['h2']}}}"
                    f"h3{{color:{pal['h3']}}}"
                    f"h4{{color:{pal['h4']}}}"
                    f"code{{color:{pal['code']};font-family:monospace;}}"
                )
            return (
                "<html><head><style>" + style + "</style></head>"
                f"<body>{body}</body></html>"
            )

        # ── Reading accessibility: fonts, spacing & reading aids ──────────

        # Candidate dyslexia-friendly families, in order of preference.
        # OpenDyslexic is purpose-built; the others are widely available
        # fallbacks with the high legibility traits dyslexic readers benefit
        # from (open apertures, distinct letterforms, generous spacing).
        _DYSLEXIA_FONTS = (
            "OpenDyslexic",
            "OpenDyslexic3",
            "OpenDyslexicAlta",
            "Atkinson Hyperlegible",
            "Lexend",
            "Lexend Deca",
            "Comic Sans MS",
            "Comic Neue",
        )

        def _find_dyslexia_font(self) -> str:
            """Return the first installed dyslexia-friendly family, or "".

            QFontDatabase.families() is a static method in PyQt6 but an
            instance method in PyQt5, so both call styles are attempted.
            """
            try:
                fams = QFontDatabase.families()  # PyQt6 (static)
            except TypeError:
                fams = QFontDatabase().families()  # PyQt5 (instance)
            except Exception:
                return ""
            available = {str(f).lower() for f in fams}
            for cand in self._DYSLEXIA_FONTS:
                if cand.lower() in available:
                    return cand
            return ""

        def _effective_font_family(self) -> str:
            """The display family, honoring the dyslexia-font preference.

            When 'qt_dyslexia_font' is on and a dyslexia-friendly family is
            installed it wins; otherwise the user's chosen family is used.
            """
            if self.settings.get("qt_dyslexia_font", False):
                fam = self._find_dyslexia_font()
                if fam:
                    return fam
            return str(self.settings.get("qt_font_family", "Georgia"))

        def _make_editor_font(self) -> "QFont":
            """Construct the editor's base QFont from family, size, and the
            letter/word-spacing accessibility settings.

            Letter and word spacing are applied through QFont (Qt's rich-text
            CSS subset does not support letter-spacing/word-spacing), while
            line height is applied separately via _apply_block_spacing.
            """
            fam = self._effective_font_family()
            size = int(self.settings.get("font_size", 0)) or int(
                self.settings.get("qt_font_size", 14)
            )
            f = QFont(fam, max(6, size))
            ls = float(self.settings.get("qt_letter_spacing", 0.0))
            # PercentageSpacing: 100 == normal; we store *extra* percent.
            f.setLetterSpacing(_PCT_SPACING, 100.0 + ls)
            ws = float(self.settings.get("qt_word_spacing", 0.0))
            f.setWordSpacing(ws)
            return f

        def _apply_block_spacing(self) -> None:
            """Apply the line-height multiplier to every block in the document.

            Run after each setHtml() because block formats are per-document
            and are reset whenever the HTML is replaced.
            """
            try:
                mult = float(self.settings.get("qt_line_height", 1.5))
            except (TypeError, ValueError):
                mult = 1.5
            pct = max(100.0, mult * 100.0)
            cur = QTextCursor(self.editor.document())
            cur.select(_DOC_SELECTION)
            bf = QTextBlockFormat()
            # setLineHeight's second arg is an int height-type.  PyQt5 enums
            # are plain ints; PyQt6 enums expose the int via .value.
            ht = getattr(_PROPORTIONAL, "value", _PROPORTIONAL)
            bf.setLineHeight(pct, int(ht))
            cur.mergeBlockFormat(bf)

        def _apply_text_spacing(self) -> None:
            """Re-apply font (letter/word spacing) and block (line-height)
            spacing to the live document, then refresh."""
            self.editor.setFont(self._make_editor_font())
            self._apply_block_spacing()

        def _bionic_word(self, m: "re.Match") -> str:
            """Embolden the leading ~40% of a single word for bionic reading."""
            w = m.group(0)
            if len(w) <= 1:
                return f"<b>{w}</b>"
            n = max(1, round(len(w) * 0.4))
            return f"<b>{w[:n]}</b>{w[n:]}"

        def _bionic_html(self, html: str) -> str:
            """Apply bionic-reading emphasis to the text runs of an HTML body.

            Splits on tags and HTML entities so markup and entities are left
            intact; text inside <code> spans is skipped (code stays verbatim).
            """
            parts = re.split(r"(<[^>]+>|&[a-zA-Z]+;|&#\d+;)", html)
            in_code = False
            out: List[str] = []
            word_re = re.compile(r"[^\W\d_]{2,}", re.UNICODE)
            for p in parts:
                if not p:
                    continue
                if p.startswith("<"):
                    low = p.lower()
                    if low.startswith("<code"):
                        in_code = True
                    elif low.startswith("</code"):
                        in_code = False
                    out.append(p)
                elif p.startswith("&"):
                    out.append(p)  # HTML entity — leave untouched
                elif in_code:
                    out.append(p)
                else:
                    out.append(word_re.sub(self._bionic_word, p))
            return "".join(out)

        def _rebuild_hl_fmt(self) -> None:
            """Rebuild the spoken-word highlight format from the user's
            'highlight_style' and 'highlight_color' settings.

            Styles: background (filled), underline, box (wavy underline),
            bold (colored bold text), color (colored text).
            """
            style = str(self.settings.get("highlight_style", "background"))
            color = QColor(str(self.settings.get("highlight_color", "cyan")))
            if not color.isValid():
                color = QColor("#06b6d4")  # cyan-500 fallback
            fmt = QTextCharFormat()
            if style == "underline":
                fmt.setFontUnderline(True)
                fmt.setUnderlineColor(color)
                fmt.setUnderlineStyle(_SINGLE_UNDERLINE)
                fmt.setFontWeight(700)
            elif style == "box":
                fmt.setUnderlineColor(color)
                fmt.setUnderlineStyle(_WAVE_UNDERLINE)
                fmt.setFontWeight(700)
            elif style == "bold":
                fmt.setForeground(color)
                fmt.setFontWeight(900)
            elif style == "color":
                fmt.setForeground(color)
                fmt.setFontWeight(700)
            else:  # "background" (default)
                fmt.setBackground(color)
                # Pick a readable foreground based on the fill's luminance.
                lum = 0.299 * color.red() + 0.587 * color.green() + 0.114 * color.blue()
                fmt.setForeground(QColor("#000000" if lum > 140 else "#ffffff"))
                fmt.setFontWeight(700)
            self._hl_fmt = fmt

        # ── TTS controls ──────────────────────────────────────────────────

        def _refresh_hl_callback(self) -> None:
            """Re-register the TTS highlight callback, capturing the current
            session number.

            Every new speech invocation must call this *after* incrementing
            self._hl_session so that the closed-over ``_s`` value matches the
            session that _apply_word_highlight expects.  Any emission still in
            the Qt event queue from a prior session carries the old ``_s`` and
            is therefore silently dropped by _apply_word_highlight.
            """
            _s = self._hl_session
            self.tts_manager.set_on_highlight(
                lambda idx, __s=_s: self._word_signal.emit(idx, __s)
            )

        def _tts_play(self) -> None:
            """Start speech from the beginning (word 0)."""
            if not self.doc:
                return
            self._hl_session += 1  # new session — must come first
            self._refresh_hl_callback()  # re-wire before any stop/speak
            self.editor.setExtraSelections([])  # clear stale highlight
            self.tts_manager.stop()
            self._tts_paused_at_word = -1
            wm = getattr(self.doc, "word_map", [])
            text_offset = wm[0].tts_offset if wm else 0
            text_slice = self.doc.plain_text[text_offset:]
            self.tts_manager.speak(
                text_slice, start_word_idx=0, text_offset=text_offset
            )
            self.statusBar().showMessage(
                f"Reading at {self.settings['tts_rate']} wpm …"
            )

        def _tts_play_from_word(self, word_idx: int) -> None:
            """Resume or start speech from a specific word-map index."""
            if not self.doc:
                return
            self._hl_session += 1  # new session — must come first
            self._refresh_hl_callback()  # re-wire before any stop/speak
            self.editor.setExtraSelections([])
            self.tts_manager.stop()
            self._tts_paused_at_word = -1
            wm = getattr(self.doc, "word_map", [])
            if wm and word_idx < len(wm):
                text_offset = wm[word_idx].tts_offset
            else:
                text_offset = 0
                word_idx = 0
            text_slice = self.doc.plain_text[text_offset:]
            self.tts_manager.speak(
                text_slice, start_word_idx=word_idx, text_offset=text_offset
            )
            self.statusBar().showMessage(
                f"Resuming at {self.settings['tts_rate']} wpm …"
            )

        def _qt_play_from_cursor(self) -> None:
            """Start speech from the current text-cursor or selection position.

            If text is selected (e.g. the user click-dragged or used
            Shift+arrow), speech begins at the *start* of the selection so
            the user hears the passage they highlighted.  If there is no
            selection, speech begins at the plain cursor position (the
            word whose character offset is closest to or just after the
            cursor).

            Shortcut: Ctrl+Return.
            """
            if not self.doc:
                self.statusBar().showMessage("No document loaded")
                return
            cursor = self.editor.textCursor()
            # Use selection start so clicking and highlighting a passage
            # then pressing Ctrl+Return reads from the top of the highlight.
            char_pos = (
                cursor.selectionStart() if cursor.hasSelection() else cursor.position()
            )
            word_idx = self._qt_char_to_word(char_pos)
            self._tts_play_from_word(word_idx)

        def _tts_toggle(self) -> None:
            """Pause/resume speech (Space bar).

            * While speaking → pause, saving the current word so the next
              press resumes from exactly that word.
            * While paused   → resume from the saved word index.
            * While stopped  → start from the beginning.
            """
            if self.tts_manager.speaking:
                saved = self.tts_manager.current_word_idx
                self._tts_stop()  # clears _tts_paused_at_word
                if saved >= 0:
                    self._tts_paused_at_word = saved
            elif self._tts_paused_at_word >= 0:
                w = self._tts_paused_at_word
                self._tts_paused_at_word = -1
                self._tts_play_from_word(w)
            else:
                self._tts_play()

        def _tts_stop(self) -> None:
            """Full stop — saves position, clears speech."""
            self._qt_save_reading_position()
            self.tts_manager.stop()
            self.editor.setExtraSelections([])  # clear highlight immediately
            self._tts_paused_at_word = -1
            self.statusBar().showMessage("Stopped.")

        # ── Word highlight (called on GUI thread via _word_signal) ─────────

        def _apply_word_highlight(self, word_idx: int, session: int) -> None:
            """Paint the currently spoken word using QTextEdit.setExtraSelections().

            setExtraSelections() is non-destructive: it does not modify the
            document, does not touch the undo stack, and is cleared simply
            by passing an empty list.  User highlights are merged so they
            are not erased on each TTS word advance.

            The *session* parameter is compared against self._hl_session to
            reject stale deliveries.  Qt's QueuedConnection buffers signal
            emissions; when _tts_play_from_word() starts a new utterance, a
            signal from the *old* timer that was already in the queue can
            fire AFTER the viewport has scrolled to the ToC target and
            immediately scroll it back.  Dropping any delivery whose session
            doesn't match the current one eliminates that race entirely.
            """
            # Discard stale deliveries from superseded speech sessions.
            if session != self._hl_session:
                return

            # Build base list from persistent user highlights so they are
            # never erased by TTS word advances.
            selections = self._get_user_highlight_selections()

            if word_idx < 0 or not self._qt_word_map:
                self.editor.setExtraSelections(selections)
                return

            # Lead/lag tuning: shift the *visual* highlight ahead of (or
            # behind) the spoken word so the cursor can be made to anticipate
            # the audio for users who track best slightly ahead.
            try:
                lead = int(self.settings.get("highlight_lead_words", 0))
            except (TypeError, ValueError):
                lead = 0
            vis_idx = word_idx + lead
            vis_idx = max(0, min(vis_idx, len(self._qt_word_map) - 1))

            char_pos = self._qt_word_map[vis_idx]

            # Word length from the word map built in _build_qt_word_map.
            word_len = 1
            if self.doc and vis_idx < len(self.doc.word_map):
                word_len = max(1, self.doc.word_map[vis_idx].tts_len)

            # Clamp to actual document length.
            doc_obj = self.editor.document()
            doc_len = doc_obj.characterCount()
            if char_pos >= doc_len:
                return
            word_len = min(word_len, doc_len - char_pos - 1)
            if word_len <= 0:
                return

            # Build the selection cursor that spans exactly this word.
            cursor = QTextCursor(doc_obj)
            cursor.setPosition(char_pos)
            cursor.setPosition(char_pos + word_len, _KEEP_ANCHOR)

            # Optional current-line band: a full-width tint behind the line
            # being read, drawn *under* the word highlight (reading aid #9).
            if self.settings.get("qt_current_line_highlight", False):
                pal = self._effective_palette(self.settings.get("theme", "dark"))
                band = QColor(str(pal.get("sel", "#2c313a")))
                line_fmt = QTextCharFormat()
                line_fmt.setBackground(band)
                # Property id is an int; PyQt6 enums expose it via .value.
                _fw = getattr(_FULL_WIDTH_SEL, "value", _FULL_WIDTH_SEL)
                line_fmt.setProperty(int(_fw), True)
                line_cur = QTextCursor(doc_obj)
                line_cur.setPosition(char_pos)
                line_sel = QTextEdit.ExtraSelection()
                line_sel.format = line_fmt
                line_sel.cursor = line_cur
                selections = selections + [line_sel]

            # Wrap format + cursor in an ExtraSelection and apply it,
            # prepending the persistent user highlights (and line band).
            sel = QTextEdit.ExtraSelection()
            sel.format = self._hl_fmt
            sel.cursor = cursor
            self.editor.setExtraSelections(selections + [sel])

            # Scroll so the highlighted word is always visible.
            # setTextCursor + ensureCursorVisible is the reliable approach;
            # the cursor width is 0 so no blinking caret is visible.
            self.editor.setTextCursor(cursor)
            self.editor.ensureCursorVisible()

            # Status bar: word text + reading progress.
            if self.doc and word_idx < len(self.doc.word_map):
                word_text = self.doc.word_map[word_idx].word
                pct = int(100 * word_idx / max(1, len(self.doc.word_map)))
                self.statusBar().showMessage(
                    f"▶  “{word_text}”  —  {pct}%  —  {self.settings['tts_rate']} wpm"
                )

        # ── Speech Cursor mode (Qt) ─────────────────────────────────────

        def eventFilter(self, obj: Any, event: Any) -> bool:
            """Intercept keyboard events on the editor for SC mode.

            Tab    — enter / exit Speech Cursor mode
            While in SC mode:
              ↑ / ↓  — previous / next block, read it
              Enter  — exit SC mode and start continuous reading
              Esc    — exit SC mode, stop speech
            """
            if obj is not self.editor:
                return super().eventFilter(obj, event)
            try:
                from PyQt6.QtCore import QEvent  # noqa: F401

                kp = QEvent.Type.KeyPress
                Key = Qt.Key
            except (ImportError, AttributeError):
                from PyQt5.QtCore import QEvent  # type: ignore

                kp = QEvent.KeyPress  # type: ignore
                Key = Qt  # type: ignore

            if event.type() != kp:
                return super().eventFilter(obj, event)

            key = event.key()

            try:
                k_tab = Key.Key_Tab
                k_up = Key.Key_Up
                k_dn = Key.Key_Down
                k_esc = Key.Key_Escape
                k_ret = Key.Key_Return
                k_ent = Key.Key_Enter
            except AttributeError:  # PyQt5 uses Qt.Key_Tab etc. directly
                k_tab = Qt.Key_Tab  # type: ignore
                k_up = Qt.Key_Up  # type: ignore
                k_dn = Qt.Key_Down  # type: ignore
                k_esc = Qt.Key_Escape  # type: ignore
                k_ret = Qt.Key_Return  # type: ignore
                k_ent = Qt.Key_Enter  # type: ignore

            if key == k_tab:
                if self._qt_sc_mode:
                    self._qt_sc_exit()
                else:
                    self._qt_sc_enter()
                return True  # consume; don't let Tab move focus

            if not self._qt_sc_mode:
                return super().eventFilter(obj, event)

            if key in (k_esc,):
                self._qt_sc_exit()
                return True
            if key in (k_ret, k_ent):
                self._qt_sc_exit(start_reading=True)
                return True
            if key == k_up:
                self._qt_sc_move(-1)
                return True
            if key == k_dn:
                self._qt_sc_move(+1)
                return True

            return super().eventFilter(obj, event)

        def _qt_sc_enter(self) -> None:
            """Enter Qt Speech Cursor mode.

            Stops any running TTS, positions the reading cursor at the
            visible text-cursor position, and builds the persistent
            _SCReader engine so the first line has no startup lag.
            """
            self.tts_manager.stop()
            self.editor.setExtraSelections([])
            # Position SC cursor at the block under the text cursor.
            tc = self.editor.textCursor()
            self._qt_sc_block = tc.blockNumber()
            # Build the persistent reader so the first keystroke has no
            # SAPI5 engine-creation delay and stop() is always effective.
            if _PYTTSX3 and isinstance(self.tts_manager._backend, Pyttsx3Backend):
                self._qt_sc_reader = _SCReader(
                    rate=int(self.settings["tts_rate"]),
                    volume=float(self.settings["tts_volume"]),
                )
                self._qt_sc_reader.start()
            else:
                self._qt_sc_reader = None
            self._qt_sc_mode = True
            self._qt_sc_highlight()
            self.setWindowTitle(f"● SC CURSOR — {APP_TITLE}")
            self.statusBar().showMessage(
                "Speech Cursor ON  ↑↓:line  Enter:read-on  Esc/Tab:exit"
            )

        def _qt_sc_exit(self, start_reading: bool = False) -> None:
            """Exit Qt Speech Cursor mode.  Speech is always stopped first."""
            self._qt_sc_mode = False
            # Stop the persistent reader before stopping the main manager
            # so the SAPI5 engine is always reachable.
            if self._qt_sc_reader is not None:
                self._qt_sc_reader.close()
                self._qt_sc_reader = None
            self.tts_manager.stop()
            self.editor.setExtraSelections([])
            self.setWindowTitle(APP_TITLE)
            if start_reading:
                # Start continuous reading from the SC cursor position.
                doc_obj = self.editor.document()
                block = doc_obj.findBlockByNumber(self._qt_sc_block)
                if block.isValid():
                    char_pos = block.position()
                    wm = getattr(self.doc, "word_map", [])
                    if wm:
                        qwm = self._qt_word_map
                        wi = 0
                        for i, off in enumerate(qwm):
                            if off >= char_pos:
                                wi = i
                                break
                        self._tts_play_from_word(wi)
                        self.statusBar().showMessage(
                            f"Reading from line {self._qt_sc_block + 1}"
                        )
                        return
                self._tts_play()
            else:
                self.statusBar().showMessage("Speech Cursor OFF")

        def _qt_sc_move(self, delta: int) -> None:
            """Move the SC cursor by *delta* blocks and read the new block."""
            doc_obj = self.editor.document()
            n_blocks = doc_obj.blockCount()
            if n_blocks == 0:
                return
            self._qt_sc_block = max(0, min(self._qt_sc_block + delta, n_blocks - 1))
            self._qt_sc_highlight()
            self._qt_sc_read_block()

        def _qt_sc_highlight(self) -> None:
            """Highlight the current SC cursor block with a bar selection."""
            doc_obj = self.editor.document()
            block = doc_obj.findBlockByNumber(self._qt_sc_block)
            if not block.isValid():
                return
            tc = QTextCursor(block)
            self.editor.setTextCursor(tc)
            self.editor.ensureCursorVisible()
            tc2 = QTextCursor(block)
            tc2.select(
                QTextCursor.SelectionType.BlockUnderCursor
                if hasattr(QTextCursor, "SelectionType")
                else QTextCursor.BlockUnderCursor  # type: ignore
            )
            fmt = QTextCharFormat()
            fmt.setBackground(QColor("#313244"))
            sel = QTextEdit.ExtraSelection()
            sel.cursor = tc2
            sel.format = fmt
            self.editor.setExtraSelections([sel])

        def _qt_sc_read_block(self) -> None:
            """Speak the plain text of the current SC cursor block.

            Uses the persistent _SCReader when the pyttsx3 backend is
            active so that stop() is always effective and there is no
            200–500 ms SAPI5 engine-construction window between navigation
            keystrokes.  Falls back to direct backend.speak() for eSpeak
            and other subprocess-based backends where termination is
            trivially reliable.
            """
            doc_obj = self.editor.document()
            block = doc_obj.findBlockByNumber(self._qt_sc_block)
            text = block.text().strip() if block.isValid() else ""
            # Stop the main TTS manager (clears timer, highlight, etc.) then
            # drive speech through the SC reader or backend directly.
            self.tts_manager.stop()
            if not text:
                if self._qt_sc_reader is not None:
                    self._qt_sc_reader.speak("blank")
                else:
                    self.tts_manager._backend.speak("blank")
                self.statusBar().showMessage(
                    f"SC  —  line {self._qt_sc_block + 1}: (blank)"
                )
                return
            text = _preprocess_tts_text(text, self.settings)
            if self._qt_sc_reader is not None:
                self._qt_sc_reader.speak(text)
            else:
                self.tts_manager._backend.speak(text)
            self.statusBar().showMessage(
                f"SC  —  line {self._qt_sc_block + 1}: {text[:80]}"
            )

        # ── Navigation (sentence · paragraph · heading) ────────────────────────

        # ─ helpers ─────────────────────────────────────────────────────────

        def _qt_current_word_for_nav(self) -> int:
            """Best estimate of the current reading position (word index).

            Priority:
            1. Callback-confirmed audio word (most accurate when speaking).
            2. Timer estimate (when speaking but no callback yet).
            3. QTextEdit text-cursor position (when TTS is idle).
               The editor cursor is moved by _apply_word_highlight during TTS
               and by _qt_navigate_to_word during manual navigation, so it
               always reflects the last reading or navigation point.  This
               ensures _qt_save_reading_position records a useful position
               even after TTS has been stopped.
            """
            if self.tts_manager.speaking:
                cb = self.tts_manager.last_cb_word_idx
                if cb >= 0:
                    return cb
                idx = self.tts_manager.current_word_idx
                if idx >= 0:
                    return idx
            # TTS idle: derive position from the editor’s text cursor.
            qwm = self._qt_word_map
            if qwm:
                char_pos = self.editor.textCursor().position()
                for i, off in enumerate(qwm):
                    if off >= char_pos:
                        return i
                return len(qwm) - 1  # cursor is past the last mapped word
            return 0

        def _qt_find_sentence_idx(self, word_idx: int) -> int:
            """Binary-search _qt_sentence_starts for the sentence containing
            *word_idx* (same algorithm as StarApp._find_sentence_idx)."""
            ss = self._qt_sentence_starts
            lo, hi, result = 0, len(ss) - 1, 0
            while lo <= hi:
                mid = (lo + hi) // 2
                if ss[mid] <= word_idx:
                    result = mid
                    lo = mid + 1
                else:
                    hi = mid - 1
            return result

        def _qt_navigate_to_word(
            self, word_idx: int, always_play: bool = False
        ) -> None:
            """Stop TTS, scroll the editor to *word_idx*, and restart speech
            if it was already running (or if *always_play* is True)."""
            was_speaking = self.tts_manager.speaking
            self.tts_manager.stop()
            self.editor.setExtraSelections([])
            wm = getattr(self.doc, "word_map", []) if self.doc else []
            if not wm or word_idx >= len(wm):
                return
            qwm = self._qt_word_map
            if word_idx < len(qwm):
                cursor = QTextCursor(self.editor.document())
                cursor.setPosition(qwm[word_idx])
                self.editor.setTextCursor(cursor)
                self.editor.ensureCursorVisible()
            if was_speaking or always_play:
                self._tts_play_from_word(word_idx)

        def _qt_block_to_word(self, block_num: int) -> int:
            """Return the word-map index of the first word inside QTextBlock
            *block_num*, searching forward through _qt_word_map."""
            block = self.editor.document().findBlockByNumber(block_num)
            if not block.isValid():
                return 0
            char_pos = block.position()
            for i, off in enumerate(self._qt_word_map):
                if off >= char_pos:
                    return i
            return 0

        def _qt_current_block(self) -> int:
            """Block number that best represents the current reading position.
            Uses the word map when speaking, otherwise the text cursor."""
            cur_word = self._qt_current_word_for_nav()
            qwm = self._qt_word_map
            if cur_word < len(qwm):
                char_pos = qwm[cur_word]
                cursor = QTextCursor(self.editor.document())
                cursor.setPosition(char_pos)
                return cursor.blockNumber()
            return self.editor.textCursor().blockNumber()

        def _qt_is_heading_block(self, block_num: int) -> bool:
            """Return True if the QTextBlock at *block_num* is a heading."""
            block = self.editor.document().findBlockByNumber(block_num)
            if not block.isValid():
                return False
            try:
                return block.blockFormat().headingLevel() > 0
            except AttributeError:
                return False  # Qt < 5.12 — heading level API unavailable

        # ─ sentence ──────────────────────────────────────────────────────────

        def _qt_skip_next_sentence(self) -> None:
            """Jump to the next sentence; restart speech if it was playing."""
            if not self.doc or not self.doc.word_map:
                return
            cur = self._qt_current_word_for_nav()
            si = self._qt_find_sentence_idx(cur)
            nsi = si + 1
            if nsi >= len(self._qt_sentence_starts):
                self.statusBar().showMessage("No next sentence")
                return
            dest = self._qt_sentence_starts[nsi]
            preview = " ".join(
                self.doc.word_map[i].word
                for i in range(dest, min(dest + 5, len(self.doc.word_map)))
            )
            total = len(self._qt_sentence_starts)
            self._qt_navigate_to_word(dest)
            self.statusBar().showMessage(f"→ Sentence {nsi + 1}/{total}: “{preview}…”")

        def _qt_skip_prev_sentence(self) -> None:
            """Jump to the previous sentence (or replay the current one if
            more than 3 words in); restart speech if it was playing."""
            if not self.doc or not self.doc.word_map:
                return
            cur = self._qt_current_word_for_nav()
            si = self._qt_find_sentence_idx(cur)
            psi = si if cur - self._qt_sentence_starts[si] > 3 else max(0, si - 1)
            dest = self._qt_sentence_starts[psi]
            preview = " ".join(
                self.doc.word_map[i].word
                for i in range(dest, min(dest + 5, len(self.doc.word_map)))
            )
            total = len(self._qt_sentence_starts)
            self._qt_navigate_to_word(dest)
            self.statusBar().showMessage(f"← Sentence {psi + 1}/{total}: “{preview}…”")

        def _qt_replay_sentence(self) -> None:
            """Jump to the start of the current sentence and *always* begin
            reading, matching the TUI\'s ’;’ key behavior."""
            if not self.doc or not self.doc.word_map:
                return
            cur = self._qt_current_word_for_nav()
            si = self._qt_find_sentence_idx(cur)
            dest = self._qt_sentence_starts[si]
            preview = " ".join(
                self.doc.word_map[i].word
                for i in range(dest, min(dest + 5, len(self.doc.word_map)))
            )
            self._qt_navigate_to_word(dest, always_play=True)
            self.statusBar().showMessage(f"↺ Replaying: “{preview}…”")

        # ─ paragraph ─────────────────────────────────────────────────────────

        def _qt_skip_next_paragraph(self) -> None:
            """Jump to the next paragraph; restart speech if it was playing."""
            doc_obj = self.editor.document()
            n = doc_obj.blockCount()
            cur_block = self._qt_current_block()
            i = cur_block + 1
            # Skip through any remaining content of the current paragraph
            while i < n and doc_obj.findBlockByNumber(i).text().strip():
                i += 1
            # Skip blank separator blocks
            while i < n and not doc_obj.findBlockByNumber(i).text().strip():
                i += 1
            if i >= n:
                self.statusBar().showMessage("No next paragraph")
                return
            self._qt_navigate_to_word(self._qt_block_to_word(i))
            self.statusBar().showMessage(f"¶  Next paragraph — block {i + 1}")

        def _qt_skip_prev_paragraph(self) -> None:
            """Jump to the previous paragraph; restart speech if it was playing."""
            doc_obj = self.editor.document()
            cur_block = self._qt_current_block()
            i = cur_block - 1
            # Skip blank lines backward
            while i > 0 and not doc_obj.findBlockByNumber(i).text().strip():
                i -= 1
            # Walk back through the previous paragraph’s content
            while i > 0 and doc_obj.findBlockByNumber(i - 1).text().strip():
                i -= 1
            i = max(0, i)
            self._qt_navigate_to_word(self._qt_block_to_word(i))
            self.statusBar().showMessage(f"¶  Prev paragraph — block {i + 1}")

        def _qt_replay_paragraph(self) -> None:
            """Jump to the start of the current paragraph and *always* begin
            reading, matching the TUI\'s ’r’ key behavior."""
            doc_obj = self.editor.document()
            cur_block = self._qt_current_block()
            # Walk back to the first block of this paragraph
            i = cur_block
            while i > 0 and doc_obj.findBlockByNumber(i - 1).text().strip():
                i -= 1
            # Step forward past any leading blank lines
            n = doc_obj.blockCount()
            while i < n - 1 and not doc_obj.findBlockByNumber(i).text().strip():
                i += 1
            self._qt_navigate_to_word(self._qt_block_to_word(i), always_play=True)
            self.statusBar().showMessage(f"↺ Replaying paragraph from block {i + 1}")

        # ─ heading ───────────────────────────────────────────────────────────

        def _qt_skip_next_heading(self) -> None:
            """Scroll to next heading; restart speech if it was playing."""
            doc_obj = self.editor.document()
            n = doc_obj.blockCount()
            start = self._qt_current_block() + 1
            for i in range(start, n):
                if self._qt_is_heading_block(i):
                    heading_text = doc_obj.findBlockByNumber(i).text().strip()
                    self._qt_navigate_to_word(self._qt_block_to_word(i))
                    self.statusBar().showMessage(f"↓ Heading: {heading_text[:60]}")
                    return
            self.statusBar().showMessage("No heading below current position")

        def _qt_skip_prev_heading(self) -> None:
            """Scroll to previous heading; restart speech if it was playing."""
            doc_obj = self.editor.document()
            start = self._qt_current_block() - 1
            for i in range(start, -1, -1):
                if self._qt_is_heading_block(i):
                    heading_text = doc_obj.findBlockByNumber(i).text().strip()
                    self._qt_navigate_to_word(self._qt_block_to_word(i))
                    self.statusBar().showMessage(f"↑ Heading: {heading_text[:60]}")
                    return
            self.statusBar().showMessage("No heading above current position")

        def _qt_read_next_heading(self) -> None:
            """Jump to next heading and *always* begin reading (TUI ’>’)."""
            doc_obj = self.editor.document()
            n = doc_obj.blockCount()
            start = self._qt_current_block() + 1
            for i in range(start, n):
                if self._qt_is_heading_block(i):
                    heading_text = doc_obj.findBlockByNumber(i).text().strip()
                    self._qt_navigate_to_word(
                        self._qt_block_to_word(i), always_play=True
                    )
                    self.statusBar().showMessage(
                        f"⏩ Reading from: {heading_text[:60]}"
                    )
                    return
            self.statusBar().showMessage("No heading below current position")

        def _qt_read_prev_heading(self) -> None:
            """Jump to previous heading and *always* begin reading (TUI '<')."""
            doc_obj = self.editor.document()
            start = self._qt_current_block() - 1
            for i in range(start, -1, -1):
                if self._qt_is_heading_block(i):
                    heading_text = doc_obj.findBlockByNumber(i).text().strip()
                    self._qt_navigate_to_word(
                        self._qt_block_to_word(i), always_play=True
                    )
                    self.statusBar().showMessage(
                        f"⏪ Reading from: {heading_text[:60]}"
                    )
                    return
            self.statusBar().showMessage("No heading above current position")

        # ─ table ──────────────────────────────────────────────────────────────

        def _qt_is_table_block(self, block_num: int) -> bool:
            """Return True when block *block_num* is a markdown table line.

            In the current renderer, table rows appear as plain paragraphs
            whose text starts with the pipe character '|'.
            """
            block = self.editor.document().findBlockByNumber(block_num)
            if not block.isValid():
                return False
            return block.text().lstrip().startswith("|")

        def _qt_skip_next_table(self) -> None:
            """Jump to the first row of the next table (Ctrl+T)."""
            doc_obj = self.editor.document()
            n = doc_obj.blockCount()
            cur = self._qt_current_block()
            # Skip out of any table we're currently inside.
            i = cur + 1
            while (
                i < n and self._qt_is_table_block(i - 1) and self._qt_is_table_block(i)
            ):
                i += 1
            # Skip non-table blocks to find the next table start.
            while i < n and not self._qt_is_table_block(i):
                i += 1
            if i >= n:
                self.statusBar().showMessage("No table below current position")
                return
            self._qt_navigate_to_word(self._qt_block_to_word(i))
            preview = doc_obj.findBlockByNumber(i).text()[:60]
            self.statusBar().showMessage(f"▼ Table: {preview}")

        def _qt_skip_prev_table(self) -> None:
            """Jump to the first row of the previous table (Ctrl+Shift+T)."""
            doc_obj = self.editor.document()
            cur = self._qt_current_block()
            # Skip back through the current table (if any).
            i = cur - 1
            while i > 0 and self._qt_is_table_block(i):
                i -= 1
            # Skip non-table blocks backward to find the end of the prev table.
            while i > 0 and not self._qt_is_table_block(i):
                i -= 1
            if i < 0 or not self._qt_is_table_block(i):
                self.statusBar().showMessage("No table above current position")
                return
            # Walk to the first row of that table.
            while i > 0 and self._qt_is_table_block(i - 1):
                i -= 1
            self._qt_navigate_to_word(self._qt_block_to_word(i))
            preview = doc_obj.findBlockByNumber(i).text()[:60]
            self.statusBar().showMessage(f"▲ Table: {preview}")

        # ─ document editing ───────────────────────────────────────────────

        def _qt_edit_mode_toggle(self) -> None:
            """Toggle between read mode and edit mode (Ctrl+E).

            In *read mode* the editor displays rendered HTML and is
            read-only.  In *edit mode* the raw Markdown source is shown
            as plain text so the user can make changes with any standard
            text-editing shortcut (Ctrl+Z/Y, Ctrl+X/C/V, arrow keys,
            Delete, Home, End, …).  Ctrl+S saves; Ctrl+E exits without
            saving.
            """
            if not self._qt_edit_mode:
                self._qt_enter_edit_mode()
            else:
                self._qt_exit_edit_mode(save=False)

        def _qt_enter_edit_mode(self) -> None:
            """Switch the editor to editable Markdown source view."""
            if not self.doc:
                self.statusBar().showMessage("No document to edit")
                return
            self.tts_manager.stop()
            self.editor.setReadOnly(False)
            self.editor.setCursorWidth(2)  # visible text cursor
            # Show raw Markdown so the user edits the source, not HTML.
            self.editor.setPlainText(self.doc.markdown or "")
            self.editor.document().setModified(False)
            self._qt_edit_mode = True
            self._qt_edit_dirty = False
            self.editor.document().contentsChanged.connect(
                self._qt_on_edit_contents_changed
            )
            self.statusBar().showMessage(
                "✏  EDIT MODE — Markdown source  ·  "
                "Ctrl+S: save  ·  Ctrl+E: discard & exit"
            )

        def _qt_on_edit_contents_changed(self) -> None:
            """Mark the document dirty when the user types in edit mode."""
            if not self._qt_edit_dirty:
                self._qt_edit_dirty = True
                title = self.doc.title if self.doc else "document"
                self.statusBar().showMessage(
                    f"✏  EDIT MODE  ·  {title}  [modified — Ctrl+S to save]"
                )

        def _qt_exit_edit_mode(self, save: bool = False) -> None:
            """Leave edit mode, optionally saving first."""
            if not self._qt_edit_mode:
                return
            if save:
                self._qt_save()
            # Disconnect the dirty listener.
            try:
                self.editor.document().contentsChanged.disconnect(
                    self._qt_on_edit_contents_changed
                )
            except (RuntimeError, TypeError):
                pass
            self._qt_edit_mode = False
            self._qt_edit_dirty = False
            # Re-render the (possibly updated) Markdown.
            md = self.doc.markdown if self.doc else ""
            self.editor.setReadOnly(True)
            self.editor.setCursorWidth(0)
            self.editor.setHtml(self._md_to_html(md))
            self._apply_block_spacing()
            self._qt_apply_user_highlights()
            self.statusBar().showMessage("Read mode")

        def _qt_save(self) -> None:
            """Save the current document.

            In *edit mode*: the edited Markdown is written back to the
            original file (for .md / .markdown / .txt / .rst / .org);
            for any other format a Save-As dialog is shown.  The document
            is then re-rendered and the TTS word maps are rebuilt.

            In *read mode*: falls through to the Markdown export dialog
            (same as File → Export → Export as Markdown…).
            """
            if not self._qt_edit_mode:
                # Not editing — offer markdown export.
                self._qt_export_markdown()
                return

            if not self.doc:
                return

            # Capture the edited source from the plain-text editor.
            new_md = self.editor.toPlainText()

            # --- persist to disk -------------------------------------------
            orig = Path(self.doc.path) if self.doc.path else None
            text_exts = {
                ".md",
                ".markdown",
                ".txt",
                ".rst",
                ".org",
                ".adoc",
                ".asc",
                ".asciidoc",
            }
            if orig and orig.suffix.lower() in text_exts:
                try:
                    orig.write_text(new_md, encoding="utf-8")
                    saved_path = str(orig)
                except OSError as exc:
                    self.statusBar().showMessage(f"Save error: {exc}")
                    return
            else:
                # Binary or non-text format — prompt for a .md path.
                stem = orig.stem if orig else "document"
                parent = str(orig.parent) if orig else ""
                dest, _ = QFileDialog.getSaveFileName(
                    self,
                    "Save As Markdown",
                    str(Path(parent) / (stem + ".md")),
                    "Markdown (*.md *.markdown)",
                )
                if not dest:
                    return
                try:
                    Path(dest).write_text(new_md, encoding="utf-8")
                    saved_path = dest
                except OSError as exc:
                    self.statusBar().showMessage(f"Save error: {exc}")
                    return

            # --- update in-memory document ---------------------------------
            self.doc.markdown = new_md
            self.doc.plain_text = _strip_markdown_for_tts(
                new_md,
                skip_code=bool(self.settings.get("tts_skip_code", True)),
                table_mode=str(self.settings.get("table_reading_mode", "structured")),
            )
            self._qt_edit_dirty = False

            # --- exit edit mode and re-render ------------------------------
            try:
                self.editor.document().contentsChanged.disconnect(
                    self._qt_on_edit_contents_changed
                )
            except (RuntimeError, TypeError):
                pass
            self._qt_edit_mode = False
            self.editor.setReadOnly(True)
            self.editor.setCursorWidth(0)
            self.editor.setHtml(self._md_to_html(new_md))
            self._apply_block_spacing()
            self._qt_apply_user_highlights()
            self._qt_build_toc()
            self._qt_build_annotations()

            # Rebuild word maps asynchronously (same flow as _on_doc_loaded_impl)
            qt_plain = self.editor.document().toPlainText()
            doc_ref = self.doc

            def _rebuild() -> None:
                try:
                    flat = qt_plain.splitlines()
                    doc_ref.word_map = _build_word_map(doc_ref.plain_text, flat)
                    self.tts_manager.set_word_map(doc_ref.word_map)
                    self._build_qt_word_map(doc_ref.plain_text, qt_plain)
                except Exception:
                    pass

            import threading as _threading

            _threading.Thread(target=_rebuild, daemon=True).start()

            self.statusBar().showMessage(f"Saved → {saved_path}")

        # ── Reading position memory ──────────────────────────────────────

        def _qt_save_reading_position(self) -> None:
            """Persist the current reading offset for the open document.
            Identical logic to StarApp._save_reading_position."""
            if not self.doc or not self.doc.path or not self.doc.word_map:
                return
            cur = self._qt_current_word_for_nav()
            wm = self.doc.word_map
            if cur < 0 or cur >= len(wm):
                return
            offset = wm[cur].tts_offset
            total_chars = len(self.doc.plain_text or "")
            pct = int(100 * offset / max(1, total_chars))
            positions = dict(self.settings.get("reading_positions", {}))
            positions[self.doc.path] = {
                "offset": offset,
                "pct": pct,
                "ts": time.strftime("%Y-%m-%dT%H:%M:%S"),
            }
            if len(positions) > 200:
                evict = sorted(positions, key=lambda k: positions[k].get("ts", ""))[:50]
                for k in evict:
                    del positions[k]
            self.settings.set("reading_positions", positions)

        def _qt_restore_reading_position(self) -> None:
            """Scroll to the saved position for the current document and
            optionally resume TTS.  Called on the GUI thread via
            _restore_signal after the word map has been built."""
            if not self.doc or not self.doc.path or not self.doc.word_map:
                return
            if not self.settings.get("tts_auto_resume", True):
                return
            saved = self.settings.get("reading_positions", {}).get(self.doc.path)
            if not saved:
                return
            target_offset = int(saved.get("offset", 0))
            pct = int(saved.get("pct", 0))
            ts = str(saved.get("ts", ""))[:10]
            # Find the word-map entry whose tts_offset is at or beyond the
            # saved offset.
            wm = self.doc.word_map
            best = len(wm) - 1
            for i, wp in enumerate(wm):
                if wp.tts_offset >= target_offset:
                    best = i
                    break
            # Scroll the editor to that position and move the text cursor
            # there.  Moving the cursor is essential: _qt_current_word_for_nav
            # reads the cursor position when TTS is idle, so without this
            # a subsequent _qt_save_reading_position call (e.g. from
            # closeEvent before the user ever starts TTS) would see position 0
            # and overwrite the just-restored offset with the start of the doc.
            qwm = self._qt_word_map
            if best < len(qwm):
                cursor = QTextCursor(self.editor.document())
                cursor.setPosition(qwm[best])
                self.editor.setTextCursor(cursor)
                self.editor.ensureCursorVisible()
            self.statusBar().showMessage(f"Resumed at {pct}%  (saved {ts})", 5000)

        # ── Display options ─────────────────────────────────────────────────────

        def _rate_change(self, delta: int) -> None:
            new_rate = max(50, min(600, int(self.settings["tts_rate"]) + delta))
            self.tts_manager.set_rate(new_rate)
            if self._qt_sc_reader is not None:
                self._qt_sc_reader.update_rate(new_rate)
            self.statusBar().showMessage(f"Rate: {new_rate} wpm")

        def _voice_picker_qt(self) -> None:
            """Open a native voice-selection dialog.

            Lists all voices available from the active TTS backend by their
            friendly name (not the raw Windows registry ID).  The user picks
            one from a scrollable list; pressing OK applies the voice and
            speaks a short confirmation so the change is immediately audible.
            Ctrl+T opens this dialog from anywhere in the GUI.
            """
            voices = self.tts_manager.list_voices()
            if not voices:
                self.statusBar().showMessage(
                    "No TTS voices found. Is pyttsx3 installed "
                    "and the backend set to pyttsx3?"
                )
                return

            # Build display strings (name + language tag).
            items: List[str] = []
            for v in voices:
                name = v.get("name", v.get("id", "Unknown"))
                lang = v.get("lang", "")
                items.append(f"{name}  [{lang}]" if lang else name)

            # Highlight the current voice.
            current_id = str(self.settings.get("tts_voice", ""))
            current_idx = 0
            for i, v in enumerate(voices):
                if v.get("id") == current_id:
                    current_idx = i
                    break

            chosen, ok = QInputDialog.getItem(
                self,
                "Select TTS Voice",
                "Choose a voice  (Ctrl+T to reopen):",
                items,
                current_idx,
                False,  # not editable — must pick from the list
            )
            if not ok:
                return

            idx = items.index(chosen)
            v = voices[idx]
            voice_id = v.get("id", "")
            voice_name = v.get("name", voice_id)
            self.tts_manager._backend.set_voice(voice_id)
            self.settings.set("tts_voice", voice_id)
            self.statusBar().showMessage(f"Voice: {voice_name}")
            # Stop any running speech and speak a brief test phrase so the
            # user hears the new voice without having to press Play.
            self.tts_manager.stop()
            self.editor.setExtraSelections([])
            self.tts_manager._backend.speak(f"Voice changed to {voice_name}.")

        def _apply_qt_theme(self, theme_name: str) -> None:
            """Apply *theme_name* to the editor widget and re-render.

            Uses _effective_palette so CSS-file themes are honored for
            both the QTextEdit stylesheet and the HTML body styling.
            """
            pal = self._effective_palette(theme_name)
            font_size = int(self.settings.get("font_size", 0)) or 14
            sheet = (
                "QTextEdit, QTextBrowser {"
                f"  background-color: {pal['bg']};"
                f"  color: {pal['fg']};"
                f"  font-size: {font_size}pt;"
                f"  selection-background-color: {pal['sel']};"
                "}"
            )
            self.editor.setStyleSheet(sheet)
            # Re-render so the HTML body styles match the new theme.
            if self.doc is not None:
                self.editor.setHtml(self._md_to_html(self.doc.markdown))
            else:
                self.editor.setHtml(self._welcome_html())
            # Line height is a per-document block format reset by setHtml.
            self._apply_block_spacing()

        def _next_theme(self) -> None:
            """Cycle to the next theme (built-in + custom CSS)."""
            all_themes = self._all_theme_names
            current = str(self.settings.get("theme", "dark"))
            idx = all_themes.index(current) if current in all_themes else 0
            new_theme = all_themes[(idx + 1) % len(all_themes)]
            self.settings["theme"] = new_theme
            self._apply_qt_theme(new_theme)
            self.statusBar().showMessage(f"Theme: {new_theme}")

        def _qt_pick_theme(self) -> None:
            """Open a dialog listing all available themes for direct selection."""
            all_themes = self._all_theme_names
            current = str(self.settings.get("theme", "dark"))
            current_idx = all_themes.index(current) if current in all_themes else 0
            chosen, ok = QInputDialog.getItem(
                self,
                "Choose Theme",
                "Select a color theme:",
                all_themes,
                current_idx,
                False,
            )
            if ok and chosen:
                self.settings["theme"] = chosen
                self._apply_qt_theme(chosen)
                self.statusBar().showMessage(f"Theme: {chosen}")

        def _qt_reload_css_themes(self) -> None:
            """Re-scan THEMES_DIR and reload all *.css theme files.

            Useful after dropping a new .css file in the themes folder
            without restarting the application.  The current theme is
            re-applied immediately if it was a CSS theme that changed.
            """
            self._css_themes = _load_css_themes()
            n = len(self._css_themes)
            current = str(self.settings.get("theme", "dark"))
            # Re-apply in case the active theme's CSS was modified on disk.
            self._apply_qt_theme(current)
            self.statusBar().showMessage(
                f"CSS themes reloaded — {n} file(s) found in {THEMES_DIR}"
            )

        def _qt_open_themes_folder(self) -> None:
            """Open THEMES_DIR in the system file manager."""
            try:
                THEMES_DIR.mkdir(parents=True, exist_ok=True)
            except OSError:
                pass
            import subprocess as _sp

            try:
                if sys.platform == "win32":
                    _sp.Popen(["explorer", str(THEMES_DIR)])
                elif sys.platform == "darwin":
                    _sp.Popen(["open", str(THEMES_DIR)])
                else:
                    _sp.Popen(["xdg-open", str(THEMES_DIR)])
                self.statusBar().showMessage(f"Opened: {THEMES_DIR}")
            except OSError as exc:
                self.statusBar().showMessage(f"Could not open folder: {exc}")

        def _show_about(self) -> None:
            """Open README.md as a document (F1 help).

            README.md lives in the same directory as star.py and is the
            canonical reference for every feature, keyboard shortcut, and
            setting.  Opening it in the main window gives the user full TTS,
            navigation, highlighting, and search — far more useful than a
            limited modal dialog.
            """
            readme = Path(__file__).parent / "README.md"
            if readme.exists():
                self._open_path(str(readme))
            else:
                self.statusBar().showMessage(
                    f"{APP_TITLE} v{APP_VERSION} — README.md not found", 6000
                )

        def _set_font(self, family: str = "", size: int = 0) -> None:
            """Change the display font family and/or size.

            Persists both *qt_font_family* and *qt_font_size*, then calls
            _apply_qt_theme so the QTextEdit stylesheet (which embeds the
            font-size) is refreshed in the same step.  Without this, a theme
            cycle after a size change would silently revert to the old size.
            """
            if family:
                self.settings["qt_font_family"] = family
            if size > 0:
                self.settings["qt_font_size"] = size
                # Keep font_size in sync so _apply_qt_theme always reads the
                # correct value regardless of which key it checks.
                self.settings["font_size"] = size
            fam = self._effective_font_family()
            fsz = int(self.settings.get("qt_font_size", 14))
            # setFont sets the base font for the widget (with letter/word
            # spacing); _apply_qt_theme then re-renders the HTML with the
            # matching font-size in the stylesheet and re-applies line height.
            self.editor.setFont(self._make_editor_font())
            self._apply_qt_theme(str(self.settings.get("theme", "dark")))
            self.statusBar().showMessage(f"Font: {fam} {fsz}pt")

        def _qt_copy(self) -> None:
            "Copy selected text or current paragraph to the clipboard."
            cursor = self.editor.textCursor()
            if cursor.hasSelection():
                text = cursor.selectedText()
            else:
                try:
                    cursor.select(QTextCursor.SelectionType.BlockUnderCursor)
                except AttributeError:
                    cursor.select(QTextCursor.BlockUnderCursor)  # PyQt5
                text = cursor.selectedText()
            if text:
                QApplication.clipboard().setText(text)
                self.statusBar().showMessage(f"Copied {len(text)} chars")

        def _qt_reading_level(self) -> None:
            "Show Flesch-Kincaid reading level for the current document."
            if not self.doc or not self.doc.plain_text:
                self.statusBar().showMessage("No document loaded")
                return
            text = self.doc.plain_text[:50000]
            words = text.split()
            n_words = max(1, len(words))
            sentences = re.split(r"[.!?]+", text)
            n_sentences = max(1, len([s for s in sentences if s.strip()]))

            def _syl(w):
                w = w.lower().rstrip(".,;:!?")
                c = len(re.findall(r"[aeiou]+", w))
                if w.endswith("e") and c > 1:
                    c -= 1
                return max(1, c)

            n_syllables = sum(_syl(w) for w in words)
            ease = (
                206.835
                - 1.015 * (n_words / n_sentences)
                - 84.6 * (n_syllables / n_words)
            )
            grade = (
                0.39 * (n_words / n_sentences) + 11.8 * (n_syllables / n_words) - 15.59
            )
            ease = max(0.0, min(100.0, ease))
            grade = max(0.0, grade)
            level = (
                "elementary"
                if grade < 6
                else "middle school"
                if grade < 9
                else "high school"
                if grade < 13
                else "college"
                if grade < 16
                else "graduate"
            )
            self.statusBar().showMessage(
                f"Reading level: Grade {grade:.1f} ({level})  Ease: {ease:.0f}/100",
                8000,
            )

        # ── Export methods ───────────────────────────────────────────────

        def _qt_export_markdown(self) -> None:
            """Save the current document as a Markdown file."""
            if not self.doc:
                self.statusBar().showMessage("No document loaded")
                return
            p = Path(self.doc.path) if self.doc.path else Path("export")
            default = str(p.parent / (p.stem + ".md"))
            dest, _ = QFileDialog.getSaveFileName(
                self, "Export as Markdown", default, "Markdown (*.md *.markdown)"
            )
            if not dest:
                return
            try:
                Path(dest).write_text(self.doc.markdown, encoding="utf-8")
                self.statusBar().showMessage(f"Exported Markdown → {dest}")
            except OSError as e:
                self.statusBar().showMessage(f"Export error: {e}")

        def _qt_export_pdf(self) -> None:
            """Save the current document as a PDF (via Qt's built-in printer).

            User highlights are rendered into the PDF because we temporarily
            apply them as document char-format before printing, then reload
            the HTML to revert.
            """
            if not self.doc:
                self.statusBar().showMessage("No document loaded")
                return
            p = Path(self.doc.path) if self.doc.path else Path("export")
            default = str(p.parent / (p.stem + ".pdf"))
            dest, _ = QFileDialog.getSaveFileName(
                self, "Export as PDF", default, "PDF Files (*.pdf)"
            )
            if not dest:
                return
            try:
                try:
                    from PyQt6.QtPrintSupport import QPrinter  # type: ignore

                    _pdf_format = QPrinter.OutputFormat.PdfFormat
                    _hi_res = QPrinter.PrinterMode.HighResolution
                except ImportError:
                    from PyQt5.QtPrintSupport import QPrinter  # type: ignore

                    _pdf_format = QPrinter.PdfFormat
                    _hi_res = QPrinter.HighResolution
            except ImportError:
                self.statusBar().showMessage(
                    "PDF export requires PyQt6.QtPrintSupport or PyQt5.QtPrintSupport"
                )
                return

            printer = QPrinter(_hi_res)
            printer.setOutputFormat(_pdf_format)
            printer.setOutputFileName(dest)

            # Apply user highlights to the document temporarily so they
            # are baked into the PDF output.
            doc_obj = self.editor.document()
            path_key = (self.doc.path or "__no_path__") if self.doc else "__no_path__"
            highlights = self.settings._data.get("user_highlights", {}).get(
                path_key, []
            )
            for hl in highlights:
                cur = QTextCursor(doc_obj)
                cur.setPosition(hl.get("start", 0))
                cur.setPosition(hl.get("end", 0), _KEEP_ANCHOR)
                fmt = QTextCharFormat()
                fmt.setBackground(QColor(hl.get("color", "#ffff00")))
                cur.mergeCharFormat(fmt)

            doc_obj.print_(printer)
            self.statusBar().showMessage(f"Exported PDF → {dest}")

            # Revert: reload the original HTML (erases inline format changes).
            self.editor.setHtml(self._md_to_html(self.doc.markdown or ""))
            self._qt_apply_user_highlights()

        def _qt_export_brf(self) -> None:
            """Save the current document as a Braille-Ready File (.brf)."""
            if not self.doc:
                self.statusBar().showMessage("No document loaded")
                return
            p = Path(self.doc.path) if self.doc.path else Path("export")
            default = str(p.parent / (p.stem + ".brf"))
            dest, _ = QFileDialog.getSaveFileName(
                self, "Export as Braille", default, "Braille (*.brf)"
            )
            if not dest:
                return
            table = str(self.settings.get("braille_table", "en-ueb-g2.ctb"))
            brf = _export_braille(
                self.doc.plain_text,
                table,
                use_liblouis=bool(self.settings.get("braille_grade2", False)),
            )
            try:
                Path(dest).write_text(brf, encoding="utf-8")
                self.statusBar().showMessage(f"Exported BRF → {dest}")
            except OSError as e:
                self.statusBar().showMessage(f"Export error: {e}")

        def _qt_export_audio(self) -> None:
            """Export the full document as a TTS audio file.

            Synthesis runs in a background thread so the GUI stays responsive.
            **WAV is the default** because it needs no external tools; MP3,
            OGG, and MP4 output additionally require **ffmpeg** (recommended)
            or **pydub** (``pip install pydub``).
            """
            if not self.doc:
                self.statusBar().showMessage("No document loaded")
                return
            p = Path(self.doc.path) if self.doc.path else Path("export")
            fmt = str(self.settings.get("audio_export_format", "wav")).lstrip(".")
            default = str(p.parent / (p.stem + f".{fmt}"))
            dest, _ = QFileDialog.getSaveFileName(
                self,
                "Export as Audio",
                default,
                "Audio Files (*.wav *.mp3 *.ogg *.mp4);;All Files (*)",
            )
            if not dest:
                return
            text = _preprocess_tts_text(self.doc.plain_text, self.settings)
            fmt = Path(dest).suffix.upper().lstrip(".") or "MP3"
            self.statusBar().showMessage(
                f"Exporting {fmt} audio \u2026 this may take a while"
            )

            def _do_export() -> None:
                try:
                    self.tts_manager.export_audio(text, dest)
                    self._export_audio_signal.emit(f"Audio exported \u2192 {dest}")
                except Exception as exc:
                    self._export_audio_signal.emit(f"Audio export error: {exc}")

            threading.Thread(target=_do_export, daemon=True).start()

        # ── User highlights ────────────────────────────────────────────────

        def _qt_highlight(self, color: str = "#ffff00") -> None:
            """Highlight the current selection in the given color and persist it.

            If no text is selected, shows a hint in the status bar.
            Default color is yellow; the Highlight menu exposes five presets.
            Shortcut: Ctrl+H (applies yellow highlight).
            """
            cursor = self.editor.textCursor()
            if not cursor.hasSelection():
                self.statusBar().showMessage(
                    "Select text first, then press Ctrl+H or choose Highlight menu"
                )
                return
            if not self.doc:
                return
            start = cursor.selectionStart()
            end = cursor.selectionEnd()
            path_key = self.doc.path or "__no_path__"
            hl_store: dict = self.settings._data.setdefault("user_highlights", {})
            hl_list: list = hl_store.setdefault(path_key, [])
            hl_list.append({"start": start, "end": end, "color": color})
            self.settings.save()
            self._qt_apply_user_highlights()
            self.statusBar().showMessage(f"Highlight added ({color})")

        def _qt_highlight_clear(self) -> None:
            """Remove all user highlights for the current document."""
            if not self.doc:
                return
            path_key = self.doc.path or "__no_path__"
            hl_store: dict = self.settings._data.get("user_highlights", {})
            if path_key in hl_store:
                hl_store[path_key] = []
            self.settings.save()
            self.editor.setExtraSelections([])
            self.statusBar().showMessage("Highlights cleared")

        def _qt_apply_user_highlights(self) -> None:
            """Re-apply all persisted user highlights for the current document
            as extra selections (non-destructive, does not modify the document)."""
            self.editor.setExtraSelections(self._get_user_highlight_selections())

        def _get_user_highlight_selections(self) -> list:
            """Return a list of ExtraSelection objects for all saved user highlights.
            Called both by _qt_apply_user_highlights and _apply_word_highlight
            so TTS word highlights are always merged on top."""
            if not self.doc:
                return []
            path_key = self.doc.path or "__no_path__"
            highlights = self.settings._data.get("user_highlights", {}).get(
                path_key, []
            )
            doc_obj = self.editor.document()
            doc_len = doc_obj.characterCount()
            selections = []
            for hl in highlights:
                start = hl.get("start", 0)
                end = hl.get("end", 0)
                if start >= end or end > doc_len:
                    continue
                fmt = QTextCharFormat()
                fmt.setBackground(QColor(hl.get("color", "#ffff00")))
                cur = QTextCursor(doc_obj)
                cur.setPosition(start)
                cur.setPosition(end, _KEEP_ANCHOR)
                sel = QTextEdit.ExtraSelection()
                sel.format = fmt
                sel.cursor = cur
                selections.append(sel)
            return selections

        # ── Table of Contents panel ───────────────────────────────────────────

        def _qt_build_toc(self) -> None:
            """Populate the Contents dock from the current document headings.

            Each item stores two roles:
              _USER_ROLE     – the raw heading title (used for display).
              _USER_ROLE + 1 – the pre-computed character offset of that
                               heading inside the rendered QTextDocument.

            The character offset is found with a *rolling forward search* so
            that repeated section titles (e.g. "Introduction" in multiple
            chapters) each resolve to their own occurrence rather than always
            the first.  -1 is stored when the title cannot be located.
            """
            self._toc_list.clear()
            if not self.doc:
                return
            search_from = 0  # advance past each match to avoid re-matching
            for line in (self.doc.markdown or "").splitlines():
                m = re.match(r"^(#{1,6})\s+(.*)", line)
                if not m:
                    continue
                level = len(m.group(1))
                title = m.group(2).strip()
                if not title:
                    continue
                indent = "\u2002" * (level - 1) * 2  # en-spaces for visual indent
                item = QListWidgetItem(indent + title)
                item.setData(_USER_ROLE, title)
                # Locate this heading in the rendered document using a
                # rolling cursor so identical heading titles in different
                # chapters resolve to their own paragraph, not always the
                # first occurrence (which is the bug that made ToC always
                # start speech from word 0).
                c = self.editor.document().find(title, search_from)
                if not c.isNull():
                    char_pos: int = c.selectionStart()
                    search_from = c.selectionEnd() + 1
                else:
                    char_pos = -1
                item.setData(_USER_ROLE + 1, char_pos)
                self._toc_list.addItem(item)

        def _qt_char_to_word(self, char_pos: int) -> int:
            """Return the word-map index of the first TTS word at or after
            *char_pos* (a character offset in the rendered QTextDocument).

            Handles three edge cases robustly:
              • *char_pos* is before the first word  → returns 0.
              • *char_pos* is past the last word     → returns the last index.
              • _qt_word_map is still empty (async build in progress)
                → falls back to a proportional estimate using doc.word_map
                  length so at least the correct region of the document is
                  reached rather than always word 0.
            """
            wm = self._qt_word_map
            if wm:
                for i, off in enumerate(wm):
                    if off >= char_pos:
                        return i
                # char_pos is past every mapped word — start from the last.
                return len(wm) - 1

            # Fallback: word map not yet built by the background thread.
            # Estimate position proportionally so the user lands in roughly
            # the right place instead of always restarting from word 0.
            if self.doc and self.doc.word_map:
                doc_len = self.editor.document().characterCount()
                if doc_len > 1:
                    pct = char_pos / doc_len
                    return max(
                        0,
                        min(
                            int(pct * len(self.doc.word_map)),
                            len(self.doc.word_map) - 1,
                        ),
                    )
            return 0

        def _qt_toc_navigate(self, item: QListWidgetItem) -> None:
            """Single-click / Enter: scroll the viewport to the heading.

            Intentionally does *not* stop or redirect speech — the user
            may be browsing the ToC while the book reads.  To jump speech
            to a heading, double-click instead.
            """
            char_pos_data = item.data(_USER_ROLE + 1)
            char_pos: int = char_pos_data if char_pos_data is not None else -1
            if char_pos < 0:
                # Heading wasn't found at build time — try a live search.
                title = item.data(_USER_ROLE) or ""
                c = self.editor.document().find(title)
                if c.isNull():
                    return
                char_pos = c.selectionStart()
            cursor = QTextCursor(self.editor.document())
            cursor.setPosition(char_pos)
            self.editor.setTextCursor(cursor)
            self.editor.ensureCursorVisible()
            title = item.data(_USER_ROLE) or ""
            self.statusBar().showMessage(
                f"Navigated to: {title}  \u00b7  double-click to start reading here"
            )

        def _qt_toc_play(self, item: QListWidgetItem) -> None:
            """Double-click: stop current speech and start reading from the heading.

            Word-index resolution uses a direct search of *doc.plain_text* —
            the same text the speech engine will speak — rather than the
            Qt-char-offset → _qt_word_map indirection.  The indirection fails
            silently in several common situations:

              • _qt_word_map is still being built asynchronously.
              • Char positions become stale after a theme re-render.
              • document().find() matches the title in body text that
                precedes the actual heading, giving char_pos ≈ 0.

            Duplicate heading titles are handled by counting the number of
            items before this one in the ToC list that carry the same title
            and picking the corresponding occurrence in plain_text.
            """
            if not self.doc:
                return
            title = item.data(_USER_ROLE) or ""
            if not title:
                return

            # ── Scroll the viewport to the heading ───────────────────────
            char_pos_data = item.data(_USER_ROLE + 1)
            char_pos: int = char_pos_data if char_pos_data is not None else -1
            if char_pos < 0:
                c = self.editor.document().find(title)
                if not c.isNull():
                    char_pos = c.selectionStart()
            if char_pos >= 0:
                scroll_cur = QTextCursor(self.editor.document())
                scroll_cur.setPosition(char_pos)
                self.editor.setTextCursor(scroll_cur)
                self.editor.ensureCursorVisible()

            # ── Determine which occurrence of this title we want ─────────
            # Some documents have identical headings in multiple chapters
            # (e.g. “Introduction” repeated).  Count prior ToC rows with the
            # same title to pick the correct occurrence in plain_text.
            row = self._toc_list.row(item)
            occurrence = sum(
                1
                for r in range(row)
                if (
                    (self._toc_list.item(r) or QListWidgetItem()).data(_USER_ROLE)
                    == title
                )
            )

            # ── Search doc.plain_text for the heading ────────────────────
            plain = self.doc.plain_text or ""
            plain_lower = plain.lower()
            title_lower = title.lower()
            search_pos = 0
            tts_pos = -1
            for _ in range(occurrence + 1):
                idx = plain_lower.find(title_lower, search_pos)
                if idx < 0:
                    break
                tts_pos = idx
                search_pos = idx + max(1, len(title_lower))

            # ── Map plain-text position → word index ────────────────────
            wm = getattr(self.doc, "word_map", [])
            if tts_pos >= 0 and wm:
                word_idx = len(wm) - 1  # default: last word
                for i, wp in enumerate(wm):
                    if wp.tts_offset >= tts_pos:
                        word_idx = i
                        break
            else:
                # word_map not ready yet (rare) — fall back to Qt mapping
                word_idx = self._qt_char_to_word(char_pos) if char_pos >= 0 else 0

            # ── Start speech from the heading ────────────────────────────
            self._tts_play_from_word(word_idx)
            self.statusBar().showMessage(f"▶  Reading from: {title}")

        def _qt_toggle_toc(self) -> None:
            """Toggle the visibility of the Contents dock panel."""
            visible = not self._toc_dock.isVisible()
            self._toc_dock.setVisible(visible)
            self.settings["qt_show_toc"] = visible

        # ── Annotations / Notes panel ────────────────────────────────────────

        def _annot_key(self) -> str:
            """Per-document key under which annotations are stored."""
            if not self.doc:
                return ""
            return self.doc.path or self.doc.title or ""

        def _qt_load_annotations(self) -> List[Dict[str, Any]]:
            """Return a mutable copy of the saved annotations for this document,
            sorted by document position."""
            key = self._annot_key()
            if not key:
                return []
            store = self.settings.get("annotations", {}) or {}
            items = [dict(a) for a in store.get(key, [])]
            items.sort(key=lambda a: int(a.get("char_pos", 0)))
            return items

        def _qt_store_annotations(self, items: List[Dict[str, Any]]) -> None:
            """Persist *items* as the annotation list for this document."""
            key = self._annot_key()
            if not key:
                return
            store = dict(self.settings.get("annotations", {}) or {})
            if items:
                store[key] = items
            else:
                store.pop(key, None)
            self.settings.set("annotations", store)

        def _qt_build_annotations(self) -> None:
            """Populate the Notes dock list from saved annotations.

            Each row stores the annotation's char position (_USER_ROLE + 1)
            and its index in the saved list (_USER_ROLE) so edit/delete can
            target the right entry even when the visible list is filtered.
            The filter box performs full-text search over the note, anchor,
            and tags; a `#tag` term filters by tag.
            """
            self._annot_list.clear()
            items = self._qt_load_annotations()
            query = self._annot_filter.text() if hasattr(self, "_annot_filter") else ""
            doc_len = max(1, self.editor.document().characterCount())
            shown = 0
            for idx, a in enumerate(items):
                if not _annotation_matches(a, query):
                    continue
                shown += 1
                note = str(a.get("note", "")).strip()
                anchor = str(a.get("anchor", "")).strip()
                tags = a.get("tags", []) or []
                cite = str(a.get("cite", "")).strip()
                pct = int(100 * int(a.get("char_pos", 0)) / doc_len)
                first = note.splitlines()[0] if note else "(empty note)"
                label = f"{pct:>3}%  {first}"
                meta = ""
                if tags:
                    meta += "  ".join(f"#{t}" for t in tags)
                if cite:
                    meta += ("  " if meta else "") + f"@{cite}"
                if anchor:
                    label += f"\n        “{anchor[:48]}”"
                if meta:
                    label += f"\n        {meta}"
                item = QListWidgetItem(label)
                item.setData(_USER_ROLE, idx)
                item.setData(_USER_ROLE + 1, int(a.get("char_pos", 0)))
                item.setToolTip(note or anchor)
                self._annot_list.addItem(item)
            if query and hasattr(self, "_annot_dock"):
                self._annot_dock.setWindowTitle(f"Notes ({shown}/{len(items)})")
            elif hasattr(self, "_annot_dock"):
                self._annot_dock.setWindowTitle(
                    f"Notes ({len(items)})" if items else "Notes"
                )

        def _qt_current_anchor(self) -> Tuple[int, str]:
            """Return (char_pos, anchor_text) for a new note.

            Uses the selection if there is one (so the user can annotate a
            specific passage), otherwise the cursor's paragraph as context.
            """
            cursor = self.editor.textCursor()
            if cursor.hasSelection():
                char_pos = cursor.selectionStart()
                anchor = cursor.selectedText()
            else:
                char_pos = cursor.position()
                anchor = cursor.block().text()
            # selectedText() uses U+2029 for line breaks; normalize whitespace.
            anchor = " ".join(anchor.split())[:120]
            return char_pos, anchor

        def _qt_add_annotation(self) -> None:
            """Prompt for note text and attach it at the current position."""
            if not self.doc:
                self.statusBar().showMessage("Open a document before adding notes")
                return
            char_pos, anchor = self._qt_current_anchor()
            text, ok = QInputDialog.getMultiLineText(
                self,
                "Add Note",
                f"Note for: “{anchor[:60]}”" if anchor else "Note:",
                "",
            )
            if not ok or not text.strip():
                return
            # Optional tags (comma/space separated; leading # optional).
            tag_str, _ok2 = QInputDialog.getText(
                self, "Tags (optional)", "Tags, comma-separated:"
            )
            tags = _parse_tags(tag_str)
            items = self._qt_load_annotations()
            items.append(
                {
                    "char_pos": int(char_pos),
                    "word_idx": self._qt_char_to_word(int(char_pos)),
                    "anchor": anchor,
                    "note": text.strip(),
                    "tags": tags,
                    "cite": "",
                    "ts": time.strftime("%Y-%m-%dT%H:%M:%S"),
                }
            )
            self._qt_store_annotations(items)
            self._qt_build_annotations()
            if not self._annot_dock.isVisible():
                self._annot_dock.setVisible(True)
                self.settings["qt_show_notes"] = True
            self.statusBar().showMessage(
                f"Note added{(' with tags: ' + ', '.join(tags)) if tags else ''}"
            )

        def _qt_selected_annotation_index(self) -> int:
            """Index (into the saved list) of the selected note, or -1."""
            item = self._annot_list.currentItem()
            if item is None:
                return -1
            data = item.data(_USER_ROLE)
            return int(data) if data is not None else -1

        def _qt_edit_annotation(self) -> None:
            """Edit the text of the selected note."""
            idx = self._qt_selected_annotation_index()
            items = self._qt_load_annotations()
            if idx < 0 or idx >= len(items):
                self.statusBar().showMessage("Select a note to edit")
                return
            text, ok = QInputDialog.getMultiLineText(
                self, "Edit Note", "Note:", str(items[idx].get("note", ""))
            )
            if not ok:
                return
            if not text.strip():
                # Empty text deletes the note.
                del items[idx]
            else:
                items[idx]["note"] = text.strip()
                items[idx]["ts"] = time.strftime("%Y-%m-%dT%H:%M:%S")
            self._qt_store_annotations(items)
            self._qt_build_annotations()
            self.statusBar().showMessage("Note updated")

        def _qt_delete_annotation(self) -> None:
            """Delete the selected note."""
            idx = self._qt_selected_annotation_index()
            items = self._qt_load_annotations()
            if idx < 0 or idx >= len(items):
                self.statusBar().showMessage("Select a note to delete")
                return
            del items[idx]
            self._qt_store_annotations(items)
            self._qt_build_annotations()
            self.statusBar().showMessage("Note deleted")

        def _qt_annotation_charpos(self, item: QListWidgetItem) -> int:
            """Resolve a Qt character position for a notes-list *item*.

            Prefers the stored Qt char position; falls back to the note's
            word index (how TUI-created notes are anchored) mapped through
            the Qt word map so notes made in either UI navigate correctly.
            """
            data = item.data(_USER_ROLE + 1)
            char_pos = int(data) if data is not None else -1
            if char_pos and char_pos > 0:
                return char_pos
            idx = item.data(_USER_ROLE)
            items = self._qt_load_annotations()
            if idx is not None and 0 <= int(idx) < len(items):
                wi = int(items[int(idx)].get("word_idx", 0) or 0)
                if self._qt_word_map and 0 <= wi < len(self._qt_word_map):
                    return self._qt_word_map[wi]
            return char_pos

        def _qt_annotation_navigate(self, item: QListWidgetItem) -> None:
            """Single-click / Enter: scroll the viewport to the note anchor."""
            char_pos = self._qt_annotation_charpos(item)
            if char_pos < 0:
                return
            doc_len = self.editor.document().characterCount()
            char_pos = max(0, min(char_pos, doc_len - 1))
            cursor = QTextCursor(self.editor.document())
            cursor.setPosition(char_pos)
            self.editor.setTextCursor(cursor)
            self.editor.ensureCursorVisible()

        def _qt_annotation_play(self, item: QListWidgetItem) -> None:
            """Double-click: start reading from the note anchor."""
            if not self.doc:
                return
            char_pos = self._qt_annotation_charpos(item)
            if char_pos < 0:
                return
            self._tts_play_from_word(self._qt_char_to_word(char_pos))

        def _qt_toggle_annotations(self) -> None:
            """Toggle the visibility of the Notes dock panel."""
            visible = not self._annot_dock.isVisible()
            self._annot_dock.setVisible(visible)
            self.settings["qt_show_notes"] = visible

        def _qt_export_annotations(self) -> None:
            """Export the current document's notes to Markdown / JSON / BibTeX / RIS.

            The format is chosen by the file extension in the save dialog.
            BibTeX and RIS emit a single reference for the source document with
            the notes attached (the standard reference-manager convention).
            """
            if not self.doc:
                self.statusBar().showMessage("No document loaded")
                return
            items = self._qt_load_annotations()
            if not items:
                self.statusBar().showMessage("No notes to export")
                return
            p = Path(self.doc.path) if self.doc.path else Path("notes")
            default = str(p.parent / (p.stem + "_notes.md"))
            dest, _flt = QFileDialog.getSaveFileName(
                self,
                "Export Notes",
                default,
                "Markdown (*.md);;JSON (*.json);;BibTeX (*.bib);;RIS (*.ris);;Text (*.txt)",
            )
            if not dest:
                return
            ext = Path(dest).suffix.lower()
            meta = getattr(self.doc, "metadata", {}) or {}
            title = self.doc.title or Path(self.doc.path or "document").stem
            author = (
                meta.get("author") or meta.get("creator") or meta.get("Author") or ""
            )
            try:
                content = _format_annotations(
                    items, ext, title, author, self.doc.path or ""
                )
                Path(dest).write_text(content, encoding="utf-8")
                self.statusBar().showMessage(f"Exported {len(items)} note(s) → {dest}")
            except OSError as exc:
                self.statusBar().showMessage(f"Export error: {exc}")

        # ── Citation manager ──────────────────────────────

        def _qt_load_citations(self) -> List[Dict[str, Any]]:
            return [dict(c) for c in (self.settings.get("citations", []) or [])]

        def _qt_store_citations(self, items: List[Dict[str, Any]]) -> None:
            self.settings.set("citations", items)

        def _qt_import_citations(self) -> None:
            """Import citations from a BibTeX / RIS / CSL-JSON file."""
            src, _flt = QFileDialog.getOpenFileName(
                self,
                "Import Citations",
                "",
                "Citations (*.bib *.ris *.json *.csl);;All Files (*)",
            )
            if not src:
                return
            try:
                imported = _import_citations(src)
            except Exception as exc:  # noqa: BLE001
                self.statusBar().showMessage(f"Import error: {exc}")
                return
            existing = self._qt_load_citations()
            by_id = {c.get("id"): c for c in existing if c.get("id")}
            added = 0
            for c in imported:
                cid = c.get("id")
                if cid and cid in by_id:
                    by_id[cid].update(c)  # refresh existing
                else:
                    existing.append(c)
                    if cid:
                        by_id[cid] = c
                    added += 1
            self._qt_store_citations(existing)
            self.statusBar().showMessage(
                f"Imported {len(imported)} citation(s) ({added} new); "
                f"library now {len(existing)}"
            )

        def _qt_export_citations(self) -> None:
            """Export the citation library to BibTeX / RIS / CSL-JSON."""
            items = self._qt_load_citations()
            if not items:
                self.statusBar().showMessage("Citation library is empty")
                return
            dest, _flt = QFileDialog.getSaveFileName(
                self,
                "Export Citations",
                "citations.bib",
                "BibTeX (*.bib);;RIS (*.ris);;CSL-JSON (*.json)",
            )
            if not dest:
                return
            try:
                Path(dest).write_text(
                    _format_citations(items, Path(dest).suffix.lower()),
                    encoding="utf-8",
                )
                self.statusBar().showMessage(
                    f"Exported {len(items)} citation(s) → {dest}"
                )
            except OSError as exc:
                self.statusBar().showMessage(f"Export error: {exc}")

        def _qt_add_citation(self) -> None:
            """Add a citation to the library by hand (a few quick prompts)."""
            fields = [
                ("id", "Citation key (e.g. smith2020):"),
                ("title", "Title:"),
                ("author", "Author(s) (Last, First and Last, First):"),
                ("year", "Year:"),
                ("journal", "Journal / container (optional):"),
                ("doi", "DOI (optional):"),
            ]
            cite: Dict[str, Any] = {"type": "article"}
            for key, prompt in fields:
                val, ok = QInputDialog.getText(self, "Add Citation", prompt)
                if not ok:
                    return
                cite[key] = val.strip()
            if not (cite.get("title") or cite.get("id")):
                self.statusBar().showMessage("Citation needs at least a title or key")
                return
            if not cite.get("id"):
                cite["id"] = (
                    re.sub(r"\W+", "", cite.get("author", "ref").split(",")[0])
                    + str(cite.get("year", ""))
                ) or "ref"
            items = self._qt_load_citations()
            items.append(cite)
            self._qt_store_citations(items)
            self.statusBar().showMessage(f"Added citation [{cite['id']}]")

        def _qt_add_citation_by_doi(self) -> None:
            """Fetch a citation from a DOI via Crossref and add it (background)."""
            doi, ok = QInputDialog.getText(
                self, "Add Citation by DOI", "DOI (e.g. 10.1038/nature12373):"
            )
            if not ok or not doi.strip():
                return
            self.statusBar().showMessage(f"Looking up {doi.strip()} …")

            def _work() -> None:
                try:
                    cite = _fetch_citation_by_doi(doi)
                    self._doi_signal.emit(json.dumps(cite))
                except Exception as exc:  # noqa: BLE001
                    self._doi_signal.emit(f"ERROR: {exc}")

            threading.Thread(target=_work, daemon=True).start()

        def _qt_on_doi(self, payload: str) -> None:
            """Main-thread handler for a Crossref DOI lookup result."""
            if payload.startswith("ERROR: "):
                self.statusBar().showMessage(f"DOI lookup failed: {payload[7:]}")
                return
            try:
                cite = json.loads(payload)
            except ValueError:
                self.statusBar().showMessage("DOI lookup returned bad data")
                return
            items = self._qt_load_citations()
            if any(c.get("id") == cite.get("id") for c in items):
                self.statusBar().showMessage(f"[{cite.get('id')}] already in library")
                return
            items.append(cite)
            self._qt_store_citations(items)
            self.statusBar().showMessage(
                f"Added [{cite.get('id')}] {str(cite.get('title', ''))[:50]}"
            )

        def _qt_insert_citation(self) -> None:
            """Insert a Pandoc-style `[@key]` marker at the cursor (or copy it).

            In edit mode the marker is inserted inline; otherwise it is copied
            to the clipboard so it can be pasted after entering edit mode.
            """
            items = self._qt_load_citations()
            if not items:
                self.statusBar().showMessage("No citations — import or add one first")
                return
            labels = [_citation_label(c) for c in items]
            choice, ok = QInputDialog.getItem(
                self, "Insert Citation", "Insert which reference?", labels, 0, False
            )
            if not ok or choice not in labels:
                return
            marker = f"[@{items[labels.index(choice)].get('id', '')}]"
            if self._qt_edit_mode:
                self.editor.textCursor().insertText(marker)
                self.statusBar().showMessage(f"Inserted {marker}")
            else:
                QApplication.clipboard().setText(marker)
                self.statusBar().showMessage(
                    f"Copied {marker} — enter edit mode (Ctrl+E) to insert inline"
                )

        def _qt_manage_citations(self) -> None:
            """Browse the citation library; copy a key or delete an entry."""
            items = self._qt_load_citations()
            if not items:
                self.statusBar().showMessage(
                    "No citations yet — use Citations → Import or Add"
                )
                return
            labels = [_citation_label(c) for c in items]
            choice, ok = QInputDialog.getItem(
                self, "Citations", "Library (Cancel to close):", labels, 0, False
            )
            if not ok or choice not in labels:
                return
            c = items[labels.index(choice)]
            action, ok2 = QInputDialog.getItem(
                self,
                _citation_label(c),
                "Action:",
                ["Copy key to clipboard", "Link to selected note", "Delete"],
                0,
                False,
            )
            if not ok2:
                return
            if action.startswith("Copy"):
                QApplication.clipboard().setText(str(c.get("id", "")))
                self.statusBar().showMessage(f"Copied [{c.get('id', '')}]")
            elif action.startswith("Link"):
                self._qt_link_citation(c.get("id", ""))
            elif action == "Delete":
                items.remove(c)
                self._qt_store_citations(items)
                self.statusBar().showMessage("Citation deleted")

        def _qt_link_citation(self, cite_id: str) -> None:
            """Attach a citation key to the currently selected note."""
            idx = self._qt_selected_annotation_index()
            items = self._qt_load_annotations()
            if idx < 0 or idx >= len(items):
                self.statusBar().showMessage(
                    "Select a note first, then link a citation"
                )
                return
            items[idx]["cite"] = cite_id
            self._qt_store_annotations(items)
            self._qt_build_annotations()
            self.statusBar().showMessage(f"Linked note to [{cite_id}]")

        # ── Speech recognition (Whisper) ───────────────────

        def _qt_transcribe_file(self) -> None:
            """Transcribe an audio file with Whisper and open it as a document."""
            if not _WHISPER:
                QMessageBox.information(
                    self,
                    "Speech recognition unavailable",
                    "Transcription requires Whisper:\n\n"
                    "    pip install openai-whisper\n"
                    "or  pip install faster-whisper",
                )
                return
            src, _flt = QFileDialog.getOpenFileName(
                self,
                "Transcribe Audio",
                "",
                "Audio (*.wav *.mp3 *.m4a *.ogg *.flac *.aac *.mp4);;All Files (*)",
            )
            if not src:
                return
            model = str(self.settings.get("whisper_model", "base"))
            ts = bool(self.settings.get("transcribe_timestamps", False))
            self.statusBar().showMessage(
                f"Transcribing with Whisper ({model})… this may take a while"
            )

            def _work() -> None:
                try:
                    text = _transcribe_audio(src, model, timestamps=ts)
                    self._transcribe_signal.emit(text, src)
                except Exception as exc:  # noqa: BLE001
                    self._transcribe_signal.emit("", f"ERROR: {exc}")

            threading.Thread(target=_work, daemon=True).start()

        def _qt_on_transcribed(self, text: str, src: str) -> None:
            """Main-thread handler for a completed transcription."""
            if src.startswith("ERROR: "):
                self.statusBar().showMessage(src[7:])
                return
            if not text:
                self.statusBar().showMessage("Transcription produced no text")
                return
            name = Path(src).stem if src else "transcription"
            md = f"# Transcription — {name}\n\n{text}\n"
            self._pending_doc = Document(
                path="",
                title=f"Transcription — {name}",
                markdown=md,
                plain_text=text,
                format="transcription",
            )
            self._on_doc_loaded()
            self.statusBar().showMessage(f"Transcribed {name} ({len(text)} chars)")

        def _qt_dictate_note(self) -> None:
            """Record a short voice memo and add it as a note (Whisper)."""
            if not self.doc:
                self.statusBar().showMessage("Open a document before dictating a note")
                return
            if not _WHISPER or not _AUDIO_IN:
                QMessageBox.information(
                    self,
                    "Dictation unavailable",
                    "Voice dictation requires Whisper and a microphone library:\n\n"
                    "    pip install openai-whisper sounddevice numpy",
                )
                return
            secs, ok = QInputDialog.getInt(
                self, "Dictate Note", "Record for how many seconds?", 8, 2, 300
            )
            if not ok:
                return
            char_pos, anchor = self._qt_current_anchor()
            model = str(self.settings.get("whisper_model", "base"))
            chunk = max(2, int(self.settings.get("whisper_chunk_seconds", 6)))
            self.statusBar().showMessage(f"Recording {secs}s… speak now")

            def _work() -> None:
                # Live streaming: record in short chunks, transcribing each as
                # it arrives so the user sees text accumulate instead of waiting
                # for one long blocking pass.  Chunks are concatenated into the
                # final note.
                try:
                    remaining = secs
                    parts: List[str] = []
                    while remaining > 0:
                        seg = min(chunk, remaining)
                        remaining -= seg
                        wav = _record_audio_to_wav(seg)
                        try:
                            piece = _transcribe_audio(wav, model)
                        finally:
                            Path(wav).unlink(missing_ok=True)
                        if piece:
                            parts.append(piece)
                            self._dictate_partial_signal.emit(" ".join(parts))
                    self._dictate_signal.emit(
                        " ".join(parts), str(int(char_pos)), anchor
                    )
                except Exception as exc:  # noqa: BLE001
                    self._dictate_signal.emit("", "ERROR", str(exc))

            threading.Thread(target=_work, daemon=True).start()

        def _qt_on_dictate_partial(self, text: str) -> None:
            """Live status update as dictation chunks are transcribed."""
            preview = text[-80:]
            self.statusBar().showMessage(f"🎙 …{preview}")

        def _qt_on_dictated(self, text: str, char_pos_s: str, anchor: str) -> None:
            """Main-thread handler for a completed dictation → save as a note."""
            if char_pos_s == "ERROR":
                self.statusBar().showMessage(f"Dictation error: {anchor}")
                return
            if not text:
                self.statusBar().showMessage("Dictation produced no text")
                return
            items = self._qt_load_annotations()
            items.append(
                {
                    "char_pos": int(char_pos_s or 0),
                    "word_idx": self._qt_char_to_word(int(char_pos_s or 0)),
                    "anchor": anchor,
                    "note": text,
                    "tags": ["dictated"],
                    "cite": "",
                    "ts": time.strftime("%Y-%m-%dT%H:%M:%S"),
                }
            )
            self._qt_store_annotations(items)
            self._qt_build_annotations()
            if not self._annot_dock.isVisible():
                self._annot_dock.setVisible(True)
                self.settings["qt_show_notes"] = True
            self.statusBar().showMessage(f"Dictated note added ({len(text)} chars)")

        # ── Keyboard cheat sheet ────────────────────────

        def _qt_show_shortcuts(self) -> None:
            """Show the canonical keyboard shortcut cheat sheet in a dialog."""
            dlg = QDialog(self)
            dlg.setWindowTitle("Keyboard Shortcuts")
            dlg.resize(640, 600)
            layout = QVBoxLayout(dlg)
            view = QTextBrowser()
            view.setHtml(self._md_to_html(_shortcuts_text(plain=False)))
            layout.addWidget(view)
            try:
                _ok_btn = QDialogButtonBox.StandardButton.Ok  # PyQt6
            except AttributeError:
                _ok_btn = QDialogButtonBox.Ok  # type: ignore[attr-defined]  # PyQt5
            buttons = QDialogButtonBox(_ok_btn)
            buttons.accepted.connect(dlg.accept)
            layout.addWidget(buttons)
            dlg.exec() if _QT == "PyQt6" else dlg.exec_()

        # ── Remappable keybindings ─────────────────

        def _resolve_shortcut(self, default: str) -> str:
            """Return the user's override for *default* shortcut, or *default*."""
            overrides = self.settings.get("keybindings", {}) or {}
            return str(overrides.get(default, default))

        def _qt_customize_shortcuts(self) -> None:
            """Let the user remap any window/toolbar shortcut; persists overrides."""
            actions = getattr(self, "_shortcut_actions", [])
            if not actions:
                self.statusBar().showMessage("No remappable shortcuts")
                return
            labels = [
                f"{label}   —   {act.shortcut().toString() or '(none)'}"
                for label, act, _default in actions
            ]
            labels.append("→ Reset all to defaults")
            choice, ok = QInputDialog.getItem(
                self, "Customize Shortcuts", "Action to rebind:", labels, 0, False
            )
            if not ok:
                return
            if choice == "→ Reset all to defaults":
                self.settings.set("keybindings", {})
                for _label, act, default in actions:
                    act.setShortcut(default)
                self.statusBar().showMessage("Shortcuts reset to defaults")
                return
            idx = labels.index(choice)
            label, act, default = actions[idx]
            new_key, ok2 = QInputDialog.getText(
                self,
                "Rebind Shortcut",
                f"New shortcut for “{label}”\n"
                f"(e.g. Ctrl+Shift+K; blank = restore default {default}):",
                text=act.shortcut().toString(),
            )
            if not ok2:
                return
            overrides = dict(self.settings.get("keybindings", {}) or {})
            new_key = new_key.strip()
            if not new_key or new_key == default:
                overrides.pop(default, None)
                act.setShortcut(default)
            else:
                overrides[default] = new_key
                act.setShortcut(new_key)
            self.settings.set("keybindings", overrides)
            self.statusBar().showMessage(
                f"“{label}” → {act.shortcut().toString() or default}"
            )

        # ── Command palette ───────────────────────

        def _qt_command_registry(self) -> List[Tuple[str, Callable]]:
            """All commands the palette can search and run."""
            return [
                ("Open File…", self._open_dialog),
                ("Open URL…", self._qt_open_url),
                ("Play / Pause", self._tts_toggle),
                ("Stop", self._tts_stop),
                ("Play from Cursor", self._qt_play_from_cursor),
                ("Faster", lambda: self._rate_change(+20)),
                ("Slower", lambda: self._rate_change(-20)),
                ("Choose Voice…", self._voice_picker_qt),
                (
                    "Speech Cursor Mode",
                    lambda: (
                        self._qt_sc_exit() if self._qt_sc_mode else self._qt_sc_enter()
                    ),
                ),
                ("Next Sentence", self._qt_skip_next_sentence),
                ("Previous Sentence", self._qt_skip_prev_sentence),
                ("Replay Sentence", self._qt_replay_sentence),
                ("Next Paragraph", self._qt_skip_next_paragraph),
                ("Previous Paragraph", self._qt_skip_prev_paragraph),
                ("Replay Paragraph", self._qt_replay_paragraph),
                ("Next Heading", self._qt_read_next_heading),
                ("Previous Heading", self._qt_read_prev_heading),
                ("Next Table", self._qt_skip_next_table),
                ("Previous Table", self._qt_skip_prev_table),
                ("Add Note", self._qt_add_annotation),
                ("Toggle Notes Panel", self._qt_toggle_annotations),
                ("Export Notes…", self._qt_export_annotations),
                ("Import Citations…", self._qt_import_citations),
                ("Export Citations…", self._qt_export_citations),
                ("Add Citation…", self._qt_add_citation),
                ("Add Citation by DOI…", self._qt_add_citation_by_doi),
                ("Insert Citation at Cursor…", self._qt_insert_citation),
                ("Browse Citations…", self._qt_manage_citations),
                ("Toggle Contents Panel", self._qt_toggle_toc),
                ("Transcribe Audio File…", self._qt_transcribe_file),
                ("Dictate Note…", self._qt_dictate_note),
                ("Copy", self._qt_copy),
                ("Toggle Edit Mode", self._qt_edit_mode_toggle),
                ("Save", self._qt_save),
                ("Export as Markdown…", self._qt_export_markdown),
                ("Export as PDF…", self._qt_export_pdf),
                ("Export as Braille (BRF)…", self._qt_export_brf),
                ("Export as Audio…", self._qt_export_audio),
                ("Reading Level", self._qt_reading_level),
                ("Change Font…", self._qt_change_font_dialog),
                ("Next Theme", self._next_theme),
                ("Choose Theme…", self._qt_pick_theme),
                ("Customize Shortcuts…", self._qt_customize_shortcuts),
                ("Keyboard Shortcuts", self._qt_show_shortcuts),
                ("Open README (Help)", self._show_about),
            ]

        def _qt_command_palette(self) -> None:
            """A searchable command palette (F2) — type to filter, Enter to run."""
            cmds = self._qt_command_registry()
            by_label = {label: fn for label, fn in cmds}
            dlg = QDialog(self)
            dlg.setWindowTitle("Command Palette")
            dlg.resize(520, 440)
            lay = QVBoxLayout(dlg)
            box = QLineEdit()
            box.setPlaceholderText(
                "Type to search commands…  (Enter runs, Esc cancels)"
            )
            lst = QListWidget()
            lay.addWidget(box)
            lay.addWidget(lst)

            def _populate(query: str = "") -> None:
                lst.clear()
                terms = query.lower().split()
                for label, _fn in cmds:
                    if all(t in label.lower() for t in terms):
                        lst.addItem(QListWidgetItem(label))
                if lst.count():
                    lst.setCurrentRow(0)

            def _run() -> None:
                it = lst.currentItem() or (lst.item(0) if lst.count() else None)
                if it is None:
                    return
                fn = by_label.get(it.text())
                dlg.accept()
                if fn:
                    fn()

            _populate()
            box.textChanged.connect(_populate)
            box.returnPressed.connect(_run)
            lst.itemActivated.connect(lambda _it: _run())
            # Down-arrow from the search box moves into the results list.
            box.setFocus()
            dlg.exec() if _QT == "PyQt6" else dlg.exec_()

        def _qt_toggle_dyslexia_font(self) -> None:
            """Toggle the dyslexia-friendly display font preference."""
            new = not bool(self.settings.get("qt_dyslexia_font", False))
            self.settings["qt_dyslexia_font"] = new
            if hasattr(self, "_dyslexia_font_act"):
                self._dyslexia_font_act.setChecked(new)
            self.editor.setFont(self._make_editor_font())
            self._apply_qt_theme(str(self.settings.get("theme", "dark")))
            if new:
                fam = self._find_dyslexia_font()
                if fam:
                    self.statusBar().showMessage(f"Dyslexia-friendly font: {fam}")
                else:
                    self.statusBar().showMessage(
                        "No dyslexia-friendly font found — install OpenDyslexic, "
                        "Atkinson Hyperlegible, or Lexend"
                    )
            else:
                self.statusBar().showMessage("Dyslexia-friendly font: OFF")

        def _qt_toggle_bionic(self) -> None:
            """Toggle bionic-reading emphasis and re-render the document."""
            new = not bool(self.settings.get("qt_bionic_reading", False))
            self.settings["qt_bionic_reading"] = new
            if hasattr(self, "_bionic_act"):
                self._bionic_act.setChecked(new)
            self._apply_qt_theme(str(self.settings.get("theme", "dark")))
            self.statusBar().showMessage("Bionic reading: " + ("ON" if new else "OFF"))

        def _qt_toggle_current_line(self) -> None:
            """Toggle the current-line focus band shown while reading."""
            new = not bool(self.settings.get("qt_current_line_highlight", False))
            self.settings["qt_current_line_highlight"] = new
            if hasattr(self, "_current_line_act"):
                self._current_line_act.setChecked(new)
            self.statusBar().showMessage(
                "Current-line highlight: " + ("ON" if new else "OFF")
            )

        def _qt_text_spacing_dialog(self) -> None:
            """Adjust line height, letter spacing, and word spacing.

            Changes preview live; OK keeps them, Cancel reverts.  Spacing is
            a recognized accommodation that reduces crowding for dyslexic and
            low-vision readers (WCAG 1.4.12 Text Spacing).
            """
            keys = ("qt_line_height", "qt_letter_spacing", "qt_word_spacing")
            orig = {k: self.settings.get(k) for k in keys}

            dlg = QDialog(self)
            dlg.setWindowTitle("Text Spacing")
            form = QFormLayout(dlg)
            info = QLabel(
                "Generous spacing reduces crowding (WCAG 1.4.12).\n"
                "Changes preview live — OK to keep, Cancel to revert."
            )
            info.setWordWrap(True)
            form.addRow(info)

            lh = QDoubleSpinBox()
            lh.setRange(1.0, 3.0)
            lh.setSingleStep(0.1)
            lh.setDecimals(2)
            lh.setValue(float(orig["qt_line_height"] or 1.5))
            form.addRow("Line height (×):", lh)

            le = QDoubleSpinBox()
            le.setRange(0.0, 40.0)
            le.setSingleStep(1.0)
            le.setDecimals(1)
            le.setSuffix(" %")
            le.setValue(float(orig["qt_letter_spacing"] or 0.0))
            form.addRow("Extra letter spacing:", le)

            wd = QDoubleSpinBox()
            wd.setRange(0.0, 40.0)
            wd.setSingleStep(1.0)
            wd.setDecimals(1)
            wd.setSuffix(" px")
            wd.setValue(float(orig["qt_word_spacing"] or 0.0))
            form.addRow("Extra word spacing:", wd)

            def _preview() -> None:
                # Mutate in-memory only (no disk write) so Cancel can revert.
                self.settings._data["qt_line_height"] = lh.value()
                self.settings._data["qt_letter_spacing"] = le.value()
                self.settings._data["qt_word_spacing"] = wd.value()
                self._apply_text_spacing()

            lh.valueChanged.connect(lambda _v: _preview())
            le.valueChanged.connect(lambda _v: _preview())
            wd.valueChanged.connect(lambda _v: _preview())

            try:
                _ok = QDialogButtonBox.StandardButton.Ok  # PyQt6
                _cancel = QDialogButtonBox.StandardButton.Cancel
            except AttributeError:
                _ok = QDialogButtonBox.Ok  # type: ignore[attr-defined]
                _cancel = QDialogButtonBox.Cancel  # type: ignore[attr-defined]
            buttons = QDialogButtonBox(_ok | _cancel)
            buttons.accepted.connect(dlg.accept)
            buttons.rejected.connect(dlg.reject)
            form.addRow(buttons)

            result = dlg.exec() if _QT == "PyQt6" else dlg.exec_()
            if result:
                _preview()
                self.settings.save()
                self.statusBar().showMessage(
                    f"Spacing — line {lh.value():.2f}×, letter +{le.value():.0f}%, "
                    f"word +{wd.value():.0f}px"
                )
            else:
                for k, v in orig.items():
                    self.settings._data[k] = v
                self._apply_text_spacing()

        def _qt_karaoke_dialog(self) -> None:
            """Tune the per-word karaoke highlight: style, color, speed, lead.

            Different readers track best with different cues — a sharp filled
            highlight, a colored underline, or bold text — and at different
            lead/lag offsets relative to the audio.
            """
            keys = (
                "highlight_style",
                "highlight_color",
                "highlight_speed",
                "highlight_lead_words",
            )
            orig = {k: self.settings.get(k) for k in keys}

            dlg = QDialog(self)
            dlg.setWindowTitle("Karaoke Highlight")
            form = QFormLayout(dlg)
            info = QLabel(
                "Tune how the spoken word is marked as it is read.\n"
                "Changes preview live — OK to keep, Cancel to revert."
            )
            info.setWordWrap(True)
            form.addRow(info)

            _STYLES = ["background", "underline", "box", "bold", "color"]
            style_box = QComboBox()
            style_box.addItems(_STYLES)
            cur_style = str(orig["highlight_style"] or "background")
            style_box.setCurrentIndex(
                _STYLES.index(cur_style) if cur_style in _STYLES else 0
            )
            form.addRow("Style:", style_box)

            _COLORS = [
                "cyan",
                "yellow",
                "green",
                "magenta",
                "orange",
                "red",
                "blue",
                "white",
            ]
            color_box = QComboBox()
            color_box.setEditable(True)
            color_box.addItems(_COLORS)
            color_box.setEditText(str(orig["highlight_color"] or "cyan"))
            form.addRow("Color:", color_box)

            speed = QDoubleSpinBox()
            speed.setRange(0.5, 1.5)
            speed.setSingleStep(0.05)
            speed.setDecimals(2)
            speed.setValue(float(orig["highlight_speed"] or 1.0))
            speed.setToolTip(
                "Highlight pacing relative to speech (1.0 = match WPM).\n"
                "Applies on the next play."
            )
            form.addRow("Speed (× WPM):", speed)

            lead = QSpinBox()
            lead.setRange(-5, 5)
            lead.setValue(int(orig["highlight_lead_words"] or 0))
            lead.setToolTip("Words the highlight leads (+) or lags (−) the audio")
            form.addRow("Lead / lag (words):", lead)

            def _preview() -> None:
                self.settings._data["highlight_style"] = style_box.currentText()
                self.settings._data["highlight_color"] = color_box.currentText().strip()
                self.settings._data["highlight_speed"] = speed.value()
                self.settings._data["highlight_lead_words"] = lead.value()
                self._rebuild_hl_fmt()

            style_box.currentTextChanged.connect(lambda _v: _preview())
            color_box.editTextChanged.connect(lambda _v: _preview())
            speed.valueChanged.connect(lambda _v: _preview())
            lead.valueChanged.connect(lambda _v: _preview())

            try:
                _ok = QDialogButtonBox.StandardButton.Ok  # PyQt6
                _cancel = QDialogButtonBox.StandardButton.Cancel
            except AttributeError:
                _ok = QDialogButtonBox.Ok  # type: ignore[attr-defined]
                _cancel = QDialogButtonBox.Cancel  # type: ignore[attr-defined]
            buttons = QDialogButtonBox(_ok | _cancel)
            buttons.accepted.connect(dlg.accept)
            buttons.rejected.connect(dlg.reject)
            form.addRow(buttons)

            result = dlg.exec() if _QT == "PyQt6" else dlg.exec_()
            if result:
                _preview()
                self.settings.save()
                self.statusBar().showMessage(
                    f"Highlight — {style_box.currentText()}, "
                    f"{color_box.currentText().strip()}, "
                    f"{speed.value():.2f}×, lead {lead.value():+d}"
                )
            else:
                for k, v in orig.items():
                    self.settings._data[k] = v
                self._rebuild_hl_fmt()

        def _qt_change_font_dialog(self) -> None:
            """Open the system font picker to choose family, style, and size.

            Uses Qt's built-in QFontDialog so the user gets the full OS
            font browser with live preview.  Both the chosen family and the
            chosen point size are applied immediately and persisted.
            """
            fam = str(self.settings.get("qt_font_family", "Georgia"))
            fsz = int(self.settings.get("qt_font_size", 14))
            current_font = QFont(fam, fsz)
            font, ok = QFontDialog.getFont(current_font, self, "Choose Display Font")
            if ok:
                self._set_font(family=font.family(), size=max(6, font.pointSize()))

        def closeEvent(self, event: Any) -> None:
            # Persist position, then silence.
            self._qt_save_reading_position()
            if self._qt_sc_reader is not None:
                self._qt_sc_reader.close()
                self._qt_sc_reader = None
            self.tts_manager.stop()
            self.settings["gui_width"] = self.width()
            self.settings["gui_height"] = self.height()
            self.settings.save()
            event.accept()

    # ────────────────────────────────────────────────────────────────────────
    class _HelpWindow(QDialog):
        """Help window that mirrors the main window\'s TTS controls.

        Opened by StarWindow._show_about().  Uses the parent StarWindow\'s
        TTSManager so rate / volume changes propagate immediately.  On open
        it saves and takes over the manager\'s word-map and highlight
        callback; on close it restores them so the main window can resume
        normal highlighting.
        """

        # Thread-safe word-highlight delivery (same pattern as StarWindow).
        _word_signal = pyqtSignal(int)

        def __init__(self, parent: StarWindow) -> None:  # type: ignore[name-defined]
            super().__init__(parent)
            self._sw = parent
            self._doc_plain: str = ""
            self._word_map: List = []
            self._qt_word_map: List[int] = []
            # Saved parent TTS state — restored when this window closes.
            self._saved_word_map: List = []
            self._saved_hl_cb = None

            self._hl_fmt = QTextCharFormat()
            self._hl_fmt.setBackground(QColor("#06b6d4"))
            self._hl_fmt.setForeground(QColor("#000000"))
            self._hl_fmt.setFontWeight(700)

            self._setup_ui()
            self._word_signal.connect(self._apply_highlight, _QUEUED)
            # Kick off async word-map build so TTS can highlight words.
            self._load_async()

        # ── UI ────────────────────────────────────────────────────────────

        def _setup_ui(self) -> None:
            self.setWindowTitle(f"Help — {APP_TITLE}")
            self.resize(780, 600)

            root = QVBoxLayout(self)
            root.setContentsMargins(0, 0, 0, 0)
            root.setSpacing(0)

            # ─ Toolbar ─ identical actions to the main window ──────────────
            tb = QToolBar("TTS Controls")
            tb.setMovable(False)

            def _act(label: str, shortcut: str, fn: Any) -> None:
                a = QAction(label, self)
                if shortcut:
                    a.setShortcut(shortcut)
                a.triggered.connect(fn)
                tb.addAction(a)

            _act("Play/Pause ▶⏸", "Space", self._tts_toggle)
            _act("Stop ■", "Escape", self._tts_stop)
            _act("+ Speed", "Ctrl+=", lambda: self._rate_change(+20))
            _act("− Speed", "Ctrl+-", lambda: self._rate_change(-20))
            tb.addSeparator()
            _act("Close", "Ctrl+W", self.close)
            root.addWidget(tb)

            # ─ Editor ────────────────────────────────────────────
            pal = self._sw._PALETTES.get(
                self._sw.settings.get("theme", "dark"),
                self._sw._PALETTES["dark"],
            )
            font_size = int(self._sw.settings.get("font_size", 0)) or 14
            self.editor = QTextEdit()
            self.editor.setReadOnly(True)
            self.editor.setFont(
                QFont(self._sw.settings.get("qt_font_family", "Georgia"), font_size)
            )
            self.editor.setCursorWidth(0)
            self.editor.setStyleSheet(
                "QTextEdit {"
                f" background-color:{pal['bg']}; color:{pal['fg']};"
                f" font-size:{font_size}pt;"
                f" selection-background-color:{pal['sel']};"
                "}"
            )
            root.addWidget(self.editor)

            # ─ Status label ──────────────────────────────────────────
            self._status = QLabel(" ")
            self._status.setContentsMargins(4, 0, 0, 2)
            root.addWidget(self._status)

            # Show HTML immediately; TTS word map loads asynchronously.
            self.editor.setHtml(self._sw._md_to_html(_HELP_TEXT))

        # ── Content / word map ──────────────────────────────────────────

        def _load_async(self) -> None:
            """Build the TTS plain text and word map in a background thread
            so the window appears instantly."""
            qt_text = self.editor.document().toPlainText()

            def _work() -> None:
                plain = _strip_markdown_for_tts(_HELP_TEXT, self._sw.settings)
                flat = qt_text.splitlines()
                wm = _build_word_map(plain, flat)

                # Build Qt character-offset map (same algorithm as StarWindow).
                qt_lower = qt_text.lower()
                tok = re.compile(r"\b\w[\w'-]*")
                result: List[int] = []
                sfrom = 0
                for m in tok.finditer(plain):
                    w = m.group().lower()
                    p = qt_lower.find(w, sfrom)
                    if p >= 0:
                        result.append(p)
                        sfrom = p + 1
                    else:
                        p = qt_lower.find(w, 0)
                        result.append(p if p >= 0 else 0)

                self._doc_plain = plain
                self._word_map = wm
                self._qt_word_map = result

            threading.Thread(target=_work, daemon=True).start()

        # ── TTS ───────────────────────────────────────────────────────────

        def _tts_play(self) -> None:
            if not self._doc_plain:
                self._status.setText("Help text still loading … try again in a moment")
                return
            tm = self._sw.tts_manager
            wm = self._word_map
            text_offset = wm[0].tts_offset if wm else 0
            # Take over the manager\'s word map and highlight callback.
            self._saved_word_map = tm._word_map
            self._saved_hl_cb = tm._on_highlight
            tm.set_word_map(wm)
            tm.set_on_highlight(lambda idx: self._word_signal.emit(idx))
            tm.speak(
                self._doc_plain[text_offset:],
                start_word_idx=0,
                text_offset=text_offset,
            )
            self._status.setText(
                f"Reading at {self._sw.settings['tts_rate']} wpm"
                f"  —  via {tm.backend_name}"
            )

        def _tts_stop(self) -> None:
            self._sw.tts_manager.stop()
            self.editor.setExtraSelections([])
            self._status.setText("Stopped.")

        def _tts_toggle(self) -> None:
            if self._sw.tts_manager.speaking:
                self._tts_stop()
            else:
                self._tts_play()

        def _rate_change(self, delta: int) -> None:
            new_rate = max(50, min(600, int(self._sw.settings["tts_rate"]) + delta))
            self._sw.tts_manager.set_rate(new_rate)
            if self._sw._qt_sc_reader is not None:
                self._sw._qt_sc_reader.update_rate(new_rate)
            self._status.setText(f"Rate: {new_rate} wpm")

        # ── Word highlight ─────────────────────────────────────────────────

        def _apply_highlight(self, word_idx: int) -> None:
            self.editor.setExtraSelections([])
            if word_idx < 0 or not self._qt_word_map:
                return
            if word_idx >= len(self._qt_word_map):
                return
            char_pos = self._qt_word_map[word_idx]
            word_len = 1
            if word_idx < len(self._word_map):
                word_len = max(1, self._word_map[word_idx].tts_len)
            doc_obj = self.editor.document()
            doc_len = doc_obj.characterCount()
            if char_pos >= doc_len:
                return
            word_len = min(word_len, doc_len - char_pos - 1)
            if word_len <= 0:
                return
            cursor = QTextCursor(doc_obj)
            cursor.setPosition(char_pos)
            cursor.setPosition(char_pos + word_len, _KEEP_ANCHOR)
            sel = QTextEdit.ExtraSelection()
            sel.format = self._hl_fmt
            sel.cursor = cursor
            self.editor.setExtraSelections([sel])
            self.editor.setTextCursor(cursor)
            self.editor.ensureCursorVisible()
            if word_idx < len(self._word_map):
                word_text = self._word_map[word_idx].word
                pct = int(100 * word_idx / max(1, len(self._word_map)))
                self._status.setText(
                    f"▶  “{word_text}”  —  {pct}%"
                    f"  —  {self._sw.settings['tts_rate']} wpm"
                )

        # ── Close ──────────────────────────────────────────────────────────────

        def closeEvent(self, event: Any) -> None:
            """Stop speech and restore the main window\'s TTS context."""
            tm = self._sw.tts_manager
            tm.stop()
            self.editor.setExtraSelections([])
            # Restore the main window\'s word map.
            restore_wm = self._saved_word_map
            if not restore_wm and self._sw.doc:
                restore_wm = getattr(self._sw.doc, "word_map", []) or []
            tm.set_word_map(restore_wm)
            # Restore highlight callback — rewire to the main window\'s signal.
            if self._saved_hl_cb is not None:
                tm.set_on_highlight(self._saved_hl_cb)
            else:
                tm.set_on_highlight(lambda idx: self._sw._word_signal.emit(idx))
            event.accept()

    # ──────────────────────────────────────────────────────────────────────
    window = StarWindow()
    window.show()
    sys.exit(app.exec() if _QT == "PyQt6" else app.exec_())


# =============================================================================
# Command-line interface and entry point
# =============================================================================


def main() -> None:
    ap = argparse.ArgumentParser(
        prog=APP_NAME,
        description=f"{APP_TITLE} — keyboard-driven reading with built-in TTS.",
        epilog=(
            "Keyboard shortcuts:  Space=play/pause  Ctrl+O=open  "
            "F2=commands  F1=help  Ctrl+Q=quit"
        ),
    )
    ap.add_argument(
        "file", nargs="?", default="", help="Document to open (file path or URL)"
    )
    ap.add_argument("--version", action="version", version=f"%(prog)s {APP_VERSION}")
    ap.add_argument(
        "--gui",
        action="store_true",
        help="Launch the Qt GUI (requires PyQt6 or PyQt5; now the default when Qt is available)",
    )
    ap.add_argument(
        "--tui",
        action="store_true",
        help="Force terminal UI mode even when Qt is available",
    )
    ap.add_argument(
        "--theme", default="", help=f"Color theme: {', '.join(THEME_NAMES)}"
    )
    ap.add_argument(
        "--rate",
        type=int,
        default=0,
        help="TTS rate in words per minute (default: 265)",
    )
    ap.add_argument(
        "--backend",
        default="",
        help="TTS backend: auto|pyttsx3|espeak|festival|coqui|dectalk|none",
    )
    ap.add_argument(
        "--plain",
        action="store_true",
        help="Extract text to stdout and exit (no TUI — ideal for piping to other tools)",
    )
    ap.add_argument(
        "--list-themes", action="store_true", help="Print available themes and exit"
    )
    ap.add_argument(
        "--list-voices", action="store_true", help="Print available TTS voices and exit"
    )
    ap.add_argument(
        "--keytest", action="store_true", help="Run the key-code inspector (diagnostic)"
    )
    args = ap.parse_args()

    settings = Settings()

    if args.list_themes:
        for name in THEME_NAMES:
            print(name)
        return

    if args.theme and args.theme in THEMES:
        settings["theme"] = args.theme
    if args.rate > 0:
        settings["tts_rate"] = args.rate
    if args.backend:
        settings["tts_backend"] = args.backend

    if args.list_voices:
        mgr = TTSManager(settings)
        voices = mgr.list_voices()
        if voices:
            for v in voices:
                print(f"{v.get('id', '?')}\t{v.get('name', '?')}\t{v.get('lang', '?')}")
        else:
            print("No voices available or TTS not installed.")
        return

    if args.plain and args.file:
        doc = load_document(args.file, settings)
        sys.stdout.write(doc.plain_text)
        sys.stdout.write("\n")
        return

    if args.keytest:
        _run_keytest()
        return

    # GUI is the default mode when Qt is available; use --tui to force terminal.
    # --gui keeps working as an explicit opt-in (and errors if Qt missing).
    if args.gui:
        _run_qt_gui(settings, args.file)  # errors internally if _QT is None
        return

    if not args.tui and _QT:
        _run_qt_gui(settings, args.file)
        return

    os.environ.setdefault("ESCDELAY", "25")

    def _tui(stdscr: "curses.window") -> None:
        app = StarApp(stdscr, settings, initial_path=args.file)
        app.run()

    try:
        curses.wrapper(_tui)
    except KeyboardInterrupt:
        pass
    except Exception as e:
        # Surface any crash outside curses for debugging
        print(f"star crashed: {e}", file=sys.stderr)
        traceback.print_exc()
        sys.exit(1)


def _run_keytest() -> None:
    """Interactive key-code inspector — shows raw curses key values."""

    def _inner(scr: "curses.window") -> None:
        scr.keypad(True)
        scr.timeout(500)
        try:
            scr.addstr(
                0, 0, "star key tester — press any key to see its code.  q = quit."
            )
            scr.addstr(1, 0, "Try: Alt+x, Esc, F1-F12, Ctrl+letter, arrow keys …")
        except curses.error:
            pass
        row = 3
        while True:
            ch = scr.getch()
            if ch == -1:
                continue
            if ch == ord("q"):
                break
            extra = ""
            if ch == 27:
                scr.timeout(150)
                ch2 = scr.getch()
                scr.timeout(500)
                if ch2 != -1:
                    extra = f"  +  ch2={ch2} (0x{ch2:04x})"
            label = repr(chr(ch)) if 32 <= ch <= 126 else "?"
            line = f"ch={ch:6d}  0x{ch:04x}  {label}{extra}"
            try:
                scr.addstr(row, 0, line)
                scr.clrtoeol()
            except curses.error:
                pass
            row += 1
            if row > curses.LINES - 2:
                row = 3
                scr.clear()
            scr.refresh()

    curses.wrapper(_inner)


if __name__ == "__main__":
    main()
